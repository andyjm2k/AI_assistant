#!/usr/bin/env python3
"""
Python FastAPI replacement for the Node.js proxy server.
Provides the same functionality but in Python for better integration with the MCP ecosystem.
"""

import asyncio
import collections
import json
import html
import io
import logging
import os
import mimetypes
import re
import sys
import time
import base64
import binascii
import hmac
import hashlib
import secrets
import glob
import socket
import struct
import traceback
import shutil
from typing import Dict, List, Optional, Any, Set, Tuple
from pathlib import Path, PurePosixPath
from datetime import datetime, timedelta, timezone
from urllib.parse import urlencode, urljoin, urlparse, urlunparse, unquote_to_bytes
from dataclasses import dataclass
from contextlib import suppress

import httpx
from fastapi import Depends, FastAPI, File, Form, Header, HTTPException, Query, Request, UploadFile
from fastapi.exceptions import RequestValidationError
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse, HTMLResponse, RedirectResponse, Response
from starlette.middleware.base import BaseHTTPMiddleware
from pydantic import AliasChoices, BaseModel, ConfigDict, Field
import uvicorn

from src.utils.token_budget import (
    estimate_tokens_from_messages,
    format_messages_for_summary,
    get_max_token_limit,
    is_context_limit_error,
)
from src.utils.file_readers import (
    FILE_READERS_AVAILABLE,
    MISSING_FILE_READER_DEPENDENCIES,
    read_docx_file,
    read_pdf_file,
    read_png_file,
    read_supported_file_text,
    read_text_file,
    read_xlsx_file,
)
from src.utils.openai_compat import (
    coerce_message_text,
    is_minimax_chat_request,
    normalize_chat_completion_message,
    normalize_temperature_for_minimax,
    preferred_api_key_env_names,
    prepare_openai_compatible_chat_payload,
)
try:
    from src.skills.bootstrap import create_default_skill_manager
    from src.skills.skill_server import create_skill_router
    from src.skills.models import SkillContext
    SKILLS_FRAMEWORK_AVAILABLE = True
except Exception as e:
    SKILLS_FRAMEWORK_AVAILABLE = False
    create_default_skill_manager = None
    create_skill_router = None
    SkillContext = None
    print(f"[WARN] Skills framework not available: {e}")
try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BeautifulSoup = None
    BS4_AVAILABLE = False

try:
    from playwright.async_api import async_playwright
    PLAYWRIGHT_FETCH_AVAILABLE = True
except ImportError:
    async_playwright = None
    PLAYWRIGHT_FETCH_AVAILABLE = False

try:
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options as SeleniumChromeOptions
    from selenium.webdriver.common.by import By as SeleniumBy
    from selenium.webdriver.support.ui import WebDriverWait as SeleniumWebDriverWait
    from selenium.webdriver.support import expected_conditions as SeleniumExpectedConditions
    SELENIUM_FETCH_AVAILABLE = True
except ImportError:
    webdriver = None
    SeleniumChromeOptions = None
    SeleniumBy = None
    SeleniumWebDriverWait = None
    SeleniumExpectedConditions = None
    SELENIUM_FETCH_AVAILABLE = False

# Import dotenv to load .env file
try:
    from dotenv import load_dotenv
    DOTENV_AVAILABLE = True
except ImportError:
    DOTENV_AVAILABLE = False
    print("[WARN] python-dotenv not available. Install with: pip install python-dotenv")

# Import file operations libraries
try:
    from docx import Document  # python-docx for Word documents
    import openpyxl  # openpyxl for Excel files
    from openpyxl.styles import Font, Alignment  # For Excel formatting
    FILE_OPS_AVAILABLE = FILE_READERS_AVAILABLE
    if FILE_OPS_AVAILABLE:
        print("[OK] File operations libraries loaded successfully")
    else:
        missing = ", ".join(MISSING_FILE_READER_DEPENDENCIES)
        print(f"[WARN] File operations libraries not available: missing {missing}")
except ImportError as e:
    print(f"[WARN] File operations libraries not available: {e}")
    FILE_OPS_AVAILABLE = False

# Import AutoGen components for team-based chat
try:
    from autogen_agentchat.teams import SelectorGroupChat
    from autogen_core import Component, FunctionCall, ComponentLoader
    AUTOGEN_AVAILABLE = True
    print("[OK] AutoGen imports successful")
except ImportError as e:
    print(f"[WARN] AutoGen not available: {e}")
    AUTOGEN_AVAILABLE = False
    SelectorGroupChat = None
    Component = None
    ComponentLoader = None

# Optional: code execution tool and Docker executor for actor workbench (requires autogen-ext[docker])
AUTOGEN_CODE_EXEC_AVAILABLE = False
PythonCodeExecutionTool = None
DockerCommandLineCodeExecutor = None
if AUTOGEN_AVAILABLE:
    try:
        from autogen_ext.tools.code_execution import PythonCodeExecutionTool as _PythonCodeExecutionTool
        from autogen_ext.code_executors.docker import DockerCommandLineCodeExecutor as _DockerCommandLineCodeExecutor
        PythonCodeExecutionTool = _PythonCodeExecutionTool
        DockerCommandLineCodeExecutor = _DockerCommandLineCodeExecutor
        AUTOGEN_CODE_EXEC_AVAILABLE = True
        print("[OK] AutoGen code execution (Docker) available")
    except ImportError as e:
        print(f"[WARN] AutoGen code execution not available (install autogen-ext[docker]): {e}")

# Import MCP SDK (with error handling)
try:
    from mcp import ClientSession, stdio_client, StdioServerParameters
    MCP_AVAILABLE = True
except ImportError as e:
    print(f"MCP import error: {e}")
    MCP_AVAILABLE = False
    ClientSession = None
    stdio_client = None
    StdioServerParameters = None

# Browser-use HTTP server URL (must run: uv run mcp-server-browser-use server)
MCP_BROWSER_USE_HTTP_URL = os.environ.get("MCP_BROWSER_USE_HTTP_URL", "http://127.0.0.1:8383/mcp").strip()
BROWSER_USE_HTTP_UNAVAILABLE_MSG = (
    "Browser-use HTTP server not available. Start it with: uv run mcp-server-browser-use server (in mcp-browser-use directory)."
)
OPEN_METEO_FORECAST_BASE_URL = os.environ.get("OPEN_METEO_FORECAST_BASE_URL", "https://api.open-meteo.com/v1/forecast").strip().rstrip("/")
OPEN_METEO_GEOCODING_BASE_URL = os.environ.get("OPEN_METEO_GEOCODING_BASE_URL", "https://geocoding-api.open-meteo.com/v1/search").strip().rstrip("/")
_open_meteo_timeout_value = os.environ.get("OPEN_METEO_TIMEOUT_SECONDS", os.environ.get("BOM_API_TIMEOUT_SECONDS", "12"))
try:
    OPEN_METEO_TIMEOUT_SECONDS = float(_open_meteo_timeout_value)
except ValueError:
    OPEN_METEO_TIMEOUT_SECONDS = 12.0
_open_meteo_trust_env_value = os.environ.get("OPEN_METEO_TRUST_ENV", os.environ.get("BOM_API_TRUST_ENV", ""))
OPEN_METEO_TRUST_ENV = str(_open_meteo_trust_env_value).strip().lower() in {"1", "true", "yes", "y", "on"}
OPEN_METEO_USER_AGENT = (
    os.environ.get("OPEN_METEO_USER_AGENT")
    or os.environ.get("BOM_API_USER_AGENT")
    or "CATBot/1.0 (+open-meteo-weather-tool)"
).strip() or "CATBot/1.0"
OPEN_METEO_ALLOWED_HOST_SUFFIX = "open-meteo.com"

def _env_bool(name: str, default: bool = False) -> bool:
    value = os.environ.get(name, "")
    if value == "":
        return default
    return value.strip().lower() in {"1", "true", "yes", "y", "on"}


def _env_str(name: str) -> Optional[str]:
    value = os.getenv(name)
    if value is None:
        return None
    value = value.strip()
    return value or None

PROXY_RELOAD = _env_bool("PROXY_RELOAD", default=False)
PROXY_RESTART_ENABLED = _env_bool("PROXY_RESTART_ENABLED", default=True)
try:
    PROXY_RESTART_DELAY_SECONDS = max(0.2, float(os.environ.get("PROXY_RESTART_DELAY_SECONDS", "1.0")))
except ValueError:
    PROXY_RESTART_DELAY_SECONDS = 1.0

_proxy_restart_scheduled = False
PROXY_START_TIME = time.time()

# Embedded Kitten TTS settings (disabled by default; set EMBEDDED_KITTEN_TTS_ENABLED=true)
EMBEDDED_KITTEN_TTS_ENABLED = _env_bool("EMBEDDED_KITTEN_TTS_ENABLED", default=False)
EMBEDDED_KITTEN_MODEL = os.environ.get("EMBEDDED_KITTEN_MODEL", "KittenML/kitten-tts-nano-0.2").strip()
EMBEDDED_KITTEN_DEFAULT_VOICE = os.environ.get("EMBEDDED_KITTEN_DEFAULT_VOICE", "expr-voice-2-f").strip() or "expr-voice-2-f"
EMBEDDED_KITTEN_VOICES = [
    v.strip() for v in os.environ.get(
        "EMBEDDED_KITTEN_VOICES",
        "expr-voice-2-f,expr-voice-3-f,expr-voice-4-m,expr-voice-5-m,expr-voice-2-m,expr-voice-5-f",
    ).split(",") if v.strip()
]
try:
    EMBEDDED_KITTEN_SAMPLE_RATE = max(8000, int(os.environ.get("EMBEDDED_KITTEN_SAMPLE_RATE", "24000")))
except ValueError:
    EMBEDDED_KITTEN_SAMPLE_RATE = 24000
try:
    EMBEDDED_KITTEN_STREAM_CHUNK_BYTES = max(512, int(os.environ.get("EMBEDDED_KITTEN_STREAM_CHUNK_BYTES", "8192")))
except ValueError:
    EMBEDDED_KITTEN_STREAM_CHUNK_BYTES = 8192
try:
    EMBEDDED_KITTEN_MAX_INPUT_CHARS = max(120, int(os.environ.get("EMBEDDED_KITTEN_MAX_INPUT_CHARS", "320")))
except ValueError:
    EMBEDDED_KITTEN_MAX_INPUT_CHARS = 320
try:
    EMBEDDED_KITTEN_CHUNK_SILENCE_MS = max(0, int(os.environ.get("EMBEDDED_KITTEN_CHUNK_SILENCE_MS", "120")))
except ValueError:
    EMBEDDED_KITTEN_CHUNK_SILENCE_MS = 120

# Embedded Pocket TTS settings (disabled by default; set EMBEDDED_POCKET_TTS_ENABLED=true)
EMBEDDED_POCKET_TTS_ENABLED = _env_bool("EMBEDDED_POCKET_TTS_ENABLED", default=False)
EMBEDDED_POCKET_MODEL = os.environ.get("EMBEDDED_POCKET_MODEL", "pocket-tts-realtime").strip() or "pocket-tts-realtime"
EMBEDDED_POCKET_DEFAULT_VOICE = os.environ.get("EMBEDDED_POCKET_DEFAULT_VOICE", "alba").strip() or "alba"
EMBEDDED_POCKET_VOICES = [
    v.strip() for v in os.environ.get(
        "EMBEDDED_POCKET_VOICES",
        "alba,marius,javert,jean,fantine,cosette,eponine,azelma",
    ).split(",") if v.strip()
]
try:
    EMBEDDED_POCKET_STREAM_CHUNK_BYTES = max(512, int(os.environ.get("EMBEDDED_POCKET_STREAM_CHUNK_BYTES", "8192")))
except ValueError:
    EMBEDDED_POCKET_STREAM_CHUNK_BYTES = 8192
try:
    TTS_PROXY_TIMEOUT_SECONDS = max(30.0, float(os.environ.get("TTS_PROXY_TIMEOUT_SECONDS", "120")))
except ValueError:
    TTS_PROXY_TIMEOUT_SECONDS = 120.0

_embedded_kitten_model_instance = None
_embedded_kitten_model_lock = asyncio.Lock()
_embedded_kitten_model_repo_id: Optional[str] = None
_embedded_kitten_voice_aliases: Dict[str, str] = {}
_embedded_pocket_model_instance = None
_embedded_pocket_model_lock = asyncio.Lock()
_embedded_pocket_voice_states: Dict[str, Any] = {}
_embedded_pocket_voice_states_lock = asyncio.Lock()

EMBEDDED_KITTEN_COMPAT_VOICE_ALIASES: Dict[str, str] = {
    # OpenAI-style names
    "alloy": "expr-voice-5-m",
    "echo": "expr-voice-2-m",
    "fable": "expr-voice-3-f",
    "onyx": "expr-voice-4-m",
    "nova": "expr-voice-5-f",
    "shimmer": "expr-voice-2-f",
    # Kitten mini aliases and project-specific names seen in configs
    "bella": "expr-voice-2-f",
    "jasper": "expr-voice-2-m",
    "luna": "expr-voice-3-f",
    "bruno": "expr-voice-3-m",
    "rosie": "expr-voice-4-f",
    "hugo": "expr-voice-4-m",
    "kiki": "expr-voice-5-f",
    "leo": "expr-voice-5-m",
    "empress": "expr-voice-4-f",
}
EMBEDDED_POCKET_MODEL_ALIASES: Set[str] = {
    "pocket-tts",
    "pocket-tts-realtime",
    "kyutai/pocket-tts",
}
EMBEDDED_POCKET_COMPAT_VOICE_ALIASES: Dict[str, str] = {
    "alloy": "alba",
    "echo": "marius",
    "fable": "fantine",
    "onyx": "javert",
    "nova": "cosette",
    "shimmer": "eponine",
    "empress": "alba",
}


# Import memory system (with error handling)
MEMORY_AVAILABLE = False
MemoryManager = None
memory_import_error = None

try:
    # Check for required dependencies first
    try:
        import numpy
    except ImportError:
        raise ImportError("numpy is required for the memory system. Install with: pip install numpy")
    
    from src.memory import MemoryManager
    MEMORY_AVAILABLE = True
    print("[OK] Memory system imports successful")
except ImportError as e:
    memory_import_error = str(e)
    print(f"[WARN] Memory system not available: {e}")
    if "numpy" in str(e).lower():
        print("   Install numpy with: pip install numpy")
    MEMORY_AVAILABLE = False
    MemoryManager = None
except Exception as e:
    memory_import_error = str(e)
    print(f"[WARN] Memory system not available (unexpected error): {e}")
    MEMORY_AVAILABLE = False
    MemoryManager = None

# Import philosopher mode
try:
    from src.features.philosopher_mode import PhilosopherMode
    PHILOSOPHER_MODE_AVAILABLE = True
    print("[OK] Philosopher mode imports successful")
except ImportError as e:
    print(f"[WARN] Philosopher mode not available: {e}")
    PHILOSOPHER_MODE_AVAILABLE = False
    PhilosopherMode = None
except Exception as e:
    print(f"[WARN] Philosopher mode not available (unexpected error): {e}")
    PHILOSOPHER_MODE_AVAILABLE = False
    PhilosopherMode = None

# Optional embedded Kitten TTS (OpenAI-compatible /v1/audio/speech + /v1/audio/voices)
try:
    from kittentts import KittenTTS as EmbeddedKittenTTS
    EMBEDDED_KITTEN_IMPORT_AVAILABLE = True
except Exception as e:
    EmbeddedKittenTTS = None
    EMBEDDED_KITTEN_IMPORT_AVAILABLE = False
    print(f"[WARN] Embedded Kitten TTS import not available: {e}")

try:
    from huggingface_hub import hf_hub_download as _hf_hub_download
    from kittentts.onnx_model import KittenTTS_1_Onnx as _EmbeddedKittenOnnxModel
    EMBEDDED_KITTEN_FALLBACK_LOADER_AVAILABLE = True
except Exception as e:
    _hf_hub_download = None
    _EmbeddedKittenOnnxModel = None
    EMBEDDED_KITTEN_FALLBACK_LOADER_AVAILABLE = False
    print(f"[WARN] Embedded Kitten fallback loader unavailable: {e}")

try:
    from pocket_tts import TTSModel as EmbeddedPocketTTSModel
    EMBEDDED_POCKET_IMPORT_AVAILABLE = True
except Exception as e:
    EmbeddedPocketTTSModel = None
    EMBEDDED_POCKET_IMPORT_AVAILABLE = False
    print(f"[WARN] Embedded Pocket TTS import not available: {e}")

# Telegram tool parsing and execution (for Telegram tools parity with web client)
try:
    from src.servers import telegram_tools as _telegram_tools
    TELEGRAM_TOOLS_MODULE_AVAILABLE = True
except ImportError as e:
    print(f"[WARN] Telegram tools module not available: {e}")

# Persistent todo store (per-user file storage)
try:
    from src.servers import todo_store as _todo_store
    TODO_STORE_AVAILABLE = True
except ImportError as e:
    print(f"[WARN] Todo store not available: {e}")
    _todo_store = None
    TODO_STORE_AVAILABLE = False

# Load environment variables from .env file in project root
if DOTENV_AVAILABLE:
    # Load from project root (two levels up from src/servers/)
    env_path = Path(__file__).resolve().parent.parent.parent / '.env'
    if env_path.exists():
        load_dotenv(env_path)
        print(f"✅ Loaded environment variables from {env_path}")
    else:
        print(f"âš ï¸  No .env file found. Using system environment variables.")
        print(f"   Looked in: {env_path}")
else:
    print("âš ï¸  python-dotenv not available. Using system environment variables only.")

# Primary hostname for SSL cert discovery (same as https_server; configurable via HTTPS_CERT_HOSTNAME in .env)
_SSL_CERT_HOSTNAME = (os.environ.get("HTTPS_CERT_HOSTNAME") or "anton.local").strip()
# Sanitize for glob to avoid matching unintended files
_SSL_CERT_HOSTNAME_GLOB = _SSL_CERT_HOSTNAME.replace("*", "").replace("?", "") or "anton.local"

# Pydantic models for request/response validation
# Note: 'command' is intentionally not accepted from clients; only server-side presets are used.
class ServerConfig(BaseModel):
    model_config = ConfigDict(extra="forbid")  # Reject any extra fields (e.g. 'command') in request body

    id: Optional[str] = None
    name: Optional[str] = None
    preset_id: Optional[str] = None  # Required for add/update; must be in MCP_PRESETS
    apiKey: Optional[str] = None
    model: Optional[str] = None
    url: Optional[str] = None
    wsUrl: Optional[str] = None
    status: Optional[str] = None
    enabled: Optional[bool] = None
    action: Optional[str] = None

class ToolCallRequest(BaseModel):
    toolName: str
    parameters: Optional[Dict[str, Any]] = None

# Pydantic models for file operations
class ReadFileRequest(BaseModel):
    filename: str = Field(validation_alias=AliasChoices("filename", "path", "file"))
    max_chars: Optional[int] = Field(default=None, ge=1, le=50000)
    start_line: Optional[int] = Field(default=None, ge=1)
    end_line: Optional[int] = Field(default=None, ge=1)
    include_line_numbers: bool = False

class WriteFileRequest(BaseModel):
    filename: str = Field(validation_alias=AliasChoices("filename", "path", "file"))
    content: str = Field(validation_alias=AliasChoices("content", "text", "body"))
    format: Optional[str] = Field(
        default="txt",
        validation_alias=AliasChoices("format", "ext", "extension"),
    )
    append: bool = False

class FileResponse(BaseModel):
    success: bool  # Whether the operation was successful
    message: str  # Human-readable message
    data: Optional[Dict[str, Any]] = None  # Optional data payload

class AuthSignupRequest(BaseModel):
    username: str
    password: str


class AuthLoginRequest(BaseModel):
    username: str
    password: str


class AuthTokenResponse(BaseModel):
    access_token: str
    token_type: str = "bearer"
    expires_in: int
    username: str


class AuthUserResponse(BaseModel):
    username: str
    created_at: str


# Todo REST API request/response models (all endpoints require auth)
class TodoRecurrenceRequest(BaseModel):
    frequency: str
    interval: int = Field(default=1, ge=1, le=10000)


class TodoAddRequest(BaseModel):
    taskDescription: str
    scheduledFor: Optional[str] = None
    recurrence: Optional[TodoRecurrenceRequest] = None


class TodoUpdateRequest(BaseModel):
    taskDescription: Optional[str] = None
    scheduledFor: Optional[str] = None
    recurrence: Optional[TodoRecurrenceRequest] = None
    clearSchedule: bool = False
    clearRecurrence: bool = False


class TodoTaskItemResponse(BaseModel):
    taskId: int
    taskDescription: str
    scheduledFor: Optional[str] = None
    nextRunAt: Optional[str] = None
    recurrence: Optional[Dict[str, Any]] = None
    lastCompletedAt: Optional[str] = None
    createdAt: Optional[str] = None
    updatedAt: Optional[str] = None
    isDue: bool = False


class TodoCompletionMetaResponse(BaseModel):
    taskId: int
    rescheduled: bool = False
    nextRunAt: Optional[str] = None


class TodoListResponse(BaseModel):
    tasks: List[str]
    taskItems: List[TodoTaskItemResponse] = []
    updated_at: Optional[str] = None
    completion: Optional[TodoCompletionMetaResponse] = None


class TodoExecuteRequest(BaseModel):
    taskId: int
    promptOverride: Optional[str] = None


class TodoResumeRequest(BaseModel):
    userMessage: str
    taskId: Optional[int] = None


class TodoCancelRequest(BaseModel):
    taskId: Optional[int] = None


class CodexExecRequest(BaseModel):
    prompt: str


class TodoExecuteResponse(BaseModel):
    status: str
    message: str
    taskId: Optional[int] = None


class TodoExecutionStatusResponse(BaseModel):
    active: bool
    activeTaskIds: List[int] = []
    runs: List[Dict[str, Any]] = []
    message: Optional[str] = None
    task: Optional[Dict[str, Any]] = None


class TelegramChatMessage(BaseModel):
    role: str
    content: str


class ChatAttachment(BaseModel):
    filename: str
    content_base64: str = Field(validation_alias=AliasChoices("content_base64", "contentBase64", "data"))
    mime_type: Optional[str] = Field(
        default=None,
        validation_alias=AliasChoices("mime_type", "mimeType", "content_type", "contentType"),
    )
    size_bytes: Optional[int] = Field(
        default=None,
        validation_alias=AliasChoices("size_bytes", "sizeBytes", "size"),
    )


class TelegramChatRequest(BaseModel):
    message: str = ""
    conversation_id: Optional[str] = None
    user_id: Optional[str] = None
    request_id: Optional[str] = None
    history: Optional[List[TelegramChatMessage]] = None
    system_prompt: Optional[str] = None
    model: Optional[str] = None
    temperature: Optional[float] = None
    max_output_tokens: Optional[int] = None
    attachments: Optional[List[ChatAttachment]] = None


class TelegramChatResponse(BaseModel):
    reply: str
    conversation_id: str
    usage: Optional[Dict[str, Any]] = None


class EmbeddedTtsSpeechRequest(BaseModel):
    model: Optional[str] = None
    voice: Optional[str] = None
    input: str
    response_format: Optional[str] = "wav"
    stream: Optional[bool] = False
    sample_rate: Optional[int] = None
    channels: Optional[int] = 1
    speed: Optional[float] = None


class StatusStartRequest(BaseModel):
    conversation_id: str
    request_id: str
    channel: Optional[str] = "web"
    state: Optional[str] = None


class StatusUpdateRequest(BaseModel):
    conversation_id: str
    request_id: str
    state: str
    phase: Optional[str] = None


class StatusFinishRequest(BaseModel):
    conversation_id: str
    request_id: str
    final_state: Optional[str] = None
    phase: Optional[str] = None

# Pydantic models for memory operations
class MemoryStoreRequest(BaseModel):
    text: str
    category: Optional[str] = None
    source: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = None

class MemorySearchRequest(BaseModel):
    query: str
    limit: Optional[int] = None
    similarity_threshold: Optional[float] = None
    category: Optional[str] = None

class MemoryExtractRequest(BaseModel):
    messages: List[Dict[str, str]]
    max_memories: Optional[int] = 3

class MemoryLearningContextRequest(BaseModel):
    task_description: str
    limit: Optional[int] = None
    similarity_threshold: Optional[float] = None

class MemoryResponse(BaseModel):
    success: bool
    message: str
    data: Optional[Dict[str, Any]] = None

# Pydantic models for philosopher mode
class PhilosopherStartRequest(BaseModel):
    conversation_id: Optional[str] = None
    user_id: Optional[str] = None

class PhilosopherStopRequest(BaseModel):
    conversation_id: Optional[str] = None
    user_id: Optional[str] = None

class PhilosopherContemplateRequest(BaseModel):
    conversation_id: Optional[str] = None
    user_id: Optional[str] = None
    question: Optional[str] = None  # Optional: if not provided, generate one

class PhilosopherResponse(BaseModel):
    success: bool
    message: str
    data: Optional[Dict[str, Any]] = None


# Pydantic models for companion config API (CATBot tool settings saved as named companions)
class CompanionCreateRequest(BaseModel):
    """Request body for creating a new companion."""
    name: str
    settings: Dict[str, Any]  # Full tool settings snapshot (serializable)


class CompanionResponse(BaseModel):
    """Single companion record (id, name, optional full settings)."""
    id: str
    name: str
    settings: Optional[Dict[str, Any]] = None  # Omitted in list view, included in GET by id


class ProxyFetchRequest(BaseModel):
    """Request body for POST /v1/proxy/fetch (avoids URL length limits on iOS Safari)."""
    url: Optional[str] = None
    urls: Optional[List[str]] = None  # Optional list: try each until one succeeds (scrape-with-retry)
    crawl: bool = True
    max_pages: int = 3
    max_depth: int = 1
    render_js: bool = False
    render_engine: str = "auto"  # auto|playwright|selenium
    wait_for_selector: Optional[str] = None
    js_wait_ms: int = 2200


# Project root (two levels up from src/servers/proxy_server.py)
_PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent

# Global state (similar to the Node.js version)
mcp_clients = {}
mcp_servers = {}
SERVERS_FILE = _PROJECT_ROOT / "config" / "mcp_servers.json"

# Server-side allowlist: only these presets may be used for MCP server execution.
# Never execute user-supplied command strings; resolve command/args only from here.
MCP_PRESETS = {
    "browser-use": {"type": "inprocess"},  # No subprocess; handled in-process in call_tool/list_tools
    # Future: "stdio": {"type": "stdio", "command": sys.executable, "allowed_args": [["-m", "mcp_server_browser_use"], ...]},
}
TEAM_CONFIG_FILE = _PROJECT_ROOT / "config" / "team-config.json"
AUTOGEN_TEAM_BUILDER_FILE = _PROJECT_ROOT / "src" / "autogen" / "team_builder.py"
# Optional: same system prompt / rules as web UI; when present, used for Telegram (overrides TELEGRAM_SYSTEM_PROMPT env)
CATBOT_SYSTEM_PROMPT_FILE = _PROJECT_ROOT / "config" / "catbot_system_prompt.txt"
SOUL_PROMPT_FILE = _PROJECT_ROOT / "config" / "soul.md"
SCRATCH_DIR = _PROJECT_ROOT / "scratch"
# Companion configs stored as one JSON file per companion (server filesystem only)
COMPANIONS_DIR = _PROJECT_ROOT / "config" / "companions"

# Allowed file extensions for scratch file operations (path traversal mitigation)
TEXT_FILE_EXTENSIONS = {".txt", ".md", ".csv", ".py", ".js", ".html"}
READ_ALLOWED_EXTENSIONS = TEXT_FILE_EXTENSIONS | {".docx", ".xlsx", ".xls", ".pdf", ".pptx", ".png", ".jpg", ".jpeg"}
WRITE_ALLOWED_EXTENSIONS = TEXT_FILE_EXTENSIONS | {".docx", ".xlsx", ".xls", ".pdf", ".pptx"}
SEARCHABLE_TEXT_EXTENSIONS = TEXT_FILE_EXTENSIONS | {".docx", ".xlsx", ".xls", ".pdf"}
# Allowed extensions for Google Drive upload (scratch workspace only; path exfiltration mitigation)
DRIVE_UPLOAD_EXTENSIONS = {".txt", ".md", ".docx", ".xlsx", ".xls", ".pdf", ".pptx", ".png", ".jpg", ".jpeg"}
ATTACHMENT_ALLOWED_EXTENSIONS = READ_ALLOWED_EXTENSIONS
# Max file size for read/write in bytes (10MB default), configurable via env
FILE_OPS_MAX_SIZE_BYTES = int(os.getenv("FILE_OPS_MAX_SIZE", "10485760"))
SEARCH_FILE_MAX_SIZE_BYTES = int(os.getenv("SEARCH_FILE_MAX_SIZE", "1048576"))
ATTACHMENT_MAX_FILES_PER_REQUEST = max(1, min(12, int(os.getenv("ATTACHMENT_MAX_FILES_PER_REQUEST", "6"))))

# Telegram chat session storage (simple in-memory cache)
telegram_conversations: Dict[str, List[Dict[str, str]]] = {}
# Per-conversation todo list and memory cache for Telegram tools (same semantics as web client)
telegram_todo: Dict[str, List[str]] = {}
telegram_memory_cache: Dict[str, List[str]] = {}
# Guard long-running ad hoc deep-research calls to avoid overlapping runs per Telegram conversation.
telegram_deep_research_active: Dict[str, Dict[str, Any]] = {}
telegram_deep_research_lock = asyncio.Lock()
try:
    TELEGRAM_DEEP_RESEARCH_STALE_SECONDS = max(
        300,
        int(os.getenv("TELEGRAM_DEEP_RESEARCH_STALE_SECONDS", "21600")),
    )
except ValueError:
    TELEGRAM_DEEP_RESEARCH_STALE_SECONDS = 21600

# Optional: tool-capable system prompt for Telegram when TELEGRAM_TOOLS_ENABLED=true
CATBOT_SYSTEM_PROMPT_WITH_TOOLS_FILE = _PROJECT_ROOT / "config" / "catbot_system_prompt_with_tools.txt"
TELEGRAM_TOOLS_ENABLED = os.getenv("TELEGRAM_TOOLS_ENABLED", "false").lower() == "true"
TELEGRAM_TOOLS_MAX_ITERATIONS = max(1, min(10, int(os.getenv("TELEGRAM_TOOLS_MAX_ITERATIONS", "10"))))
# Automatic memory injection quality controls (Telegram/web auto-context paths)
MEMORY_AUTO_SEARCH_MIN_SIMILARITY = max(
    0.0, min(1.0, float(os.getenv("MEMORY_AUTO_SEARCH_MIN_SIMILARITY", "0.72")))
)
MEMORY_AUTO_SEARCH_CANDIDATE_THRESHOLD = max(
    0.0, min(1.0, float(os.getenv("MEMORY_AUTO_SEARCH_CANDIDATE_THRESHOLD", "0.55")))
)
MEMORY_AUTO_SEARCH_LIMIT = max(1, min(10, int(os.getenv("MEMORY_AUTO_SEARCH_LIMIT", "3"))))
MEMORY_AUTO_SEARCH_SCORE_WINDOW = max(
    0.0, min(1.0, float(os.getenv("MEMORY_AUTO_SEARCH_SCORE_WINDOW", "0.12")))
)
MEMORY_CONTEXT_BLOCKED_CATEGORIES = {"task_experience", "task_learning"}
MEMORY_CONTEXT_BLOCKED_SOURCES = {"task_execution", "task_scheduler", "status_system"}
MEMORY_CONTEXT_OPERATIONAL_PATTERN = re.compile(
    r"\b(todo|to-?do|task list|my tasks?|due tasks?|overdue tasks?|task execution|"
    r"execution status|status update|awaiting confirmation|paused awaiting feedback|"
    r"pending tasks?|completed tasks?|cancelled tasks?|task id|current state|"
    r"list state|working:|done:|failed:)\b",
    re.IGNORECASE,
)
# Optional: map Telegram user_id or conversation_id to app username for shared todo list
TELEGRAM_USER_LINKS_FILE = _PROJECT_ROOT / "config" / "telegram_user_links.json"
_TELEGRAM_CHAT_ID_RE = re.compile(r"^-?\d+$")

# ============================================================================
# STATUS EVENTS (persistent progress updates)
# ============================================================================

STATUS_UPDATE_INTERVAL_SECONDS = max(5, int(os.environ.get("STATUS_UPDATE_INTERVAL_SECONDS", "60")))
STATUS_DATA_DIR = _PROJECT_ROOT / "status_data"
STATUS_EVENTS_FILE = STATUS_DATA_DIR / "status_events.jsonl"


@dataclass
class StatusSession:
    conversation_id: str
    request_id: str
    channel: str
    state: str
    phase: str
    seq: int = 0
    started_at: float = 0.0
    heartbeat_task: Optional[asyncio.Task] = None
    done: bool = False


status_sessions: Dict[Tuple[str, str], StatusSession] = {}
status_latest_index: Dict[Tuple[str, str], Dict[str, Any]] = {}
status_write_lock = asyncio.Lock()
MONITOR_RUN_HISTORY_LIMIT = max(10, int(os.environ.get("MONITOR_RUN_HISTORY_LIMIT", "25")))
MONITOR_LOG_PREVIEW_LINES = max(10, int(os.environ.get("MONITOR_LOG_PREVIEW_LINES", "80")))
MONITOR_RUN_LOG_MAX_BYTES = max(32768, int(os.environ.get("MONITOR_RUN_LOG_MAX_BYTES", "1048576")))
MONITOR_BROWSER_HEALTH_MAX_AGE_SECONDS = max(
    5.0,
    float(os.environ.get("MONITOR_BROWSER_HEALTH_MAX_AGE_SECONDS", "15")),
)
BROWSER_USE_LOG_FILE = _env_str("BROWSER_USE_LOG_FILE") or _env_str("MCP_BROWSER_USE_LOG_FILE")
monitor_recent_runs: Dict[str, collections.deque] = {
    "autogen": collections.deque(maxlen=MONITOR_RUN_HISTORY_LIMIT),
    "browser_use": collections.deque(maxlen=MONITOR_RUN_HISTORY_LIMIT),
    "philosopher": collections.deque(maxlen=MONITOR_RUN_HISTORY_LIMIT),
    "task_execution": collections.deque(maxlen=MONITOR_RUN_HISTORY_LIMIT),
}
monitor_active_runs: Dict[str, Dict[str, Any]] = {}
monitor_browser_health_snapshot: Dict[str, Any] = {
    "checked_at": None,
    "ok": False,
    "message": "Browser-use health has not been checked yet.",
    "result": None,
}
PROXY_LOG_FILE = Path(_env_str("PROXY_LOG_FILE") or (_PROJECT_ROOT / "logs" / "proxy_server.log")).expanduser()
if not PROXY_LOG_FILE.is_absolute():
    PROXY_LOG_FILE = (_PROJECT_ROOT / PROXY_LOG_FILE).resolve()
os.environ.setdefault("PROXY_LOG_FILE", str(PROXY_LOG_FILE))
_proxy_log_file_handle: Optional[io.TextIOWrapper] = None


class _ProxyLogTeeStream:
    """Mirror stdout/stderr into the proxy log file while preserving the active console stream."""

    def __init__(self, stream: Any, log_handle: io.TextIOWrapper):
        self._stream = stream
        self._log_handle = log_handle
        self.encoding = getattr(stream, "encoding", "utf-8")
        self.errors = getattr(stream, "errors", "replace")
        self._proxy_log_tee = True
        self._proxy_log_path = str(PROXY_LOG_FILE)

    def write(self, data: Any) -> int:
        text = data if isinstance(data, str) else str(data)
        written = self._stream.write(text)
        self._log_handle.write(text)
        return written

    def flush(self) -> None:
        self._stream.flush()
        self._log_handle.flush()

    def isatty(self) -> bool:
        return bool(getattr(self._stream, "isatty", lambda: False)())

    def writable(self) -> bool:
        return True

    def fileno(self) -> int:
        return self._stream.fileno()

    @property
    def buffer(self) -> Any:
        return getattr(self._stream, "buffer", None)

    def __getattr__(self, name: str) -> Any:
        return getattr(self._stream, name)


def _install_proxy_log_capture() -> None:
    global _proxy_log_file_handle

    current_path = str(PROXY_LOG_FILE)
    if getattr(sys.stdout, "_proxy_log_path", None) == current_path and getattr(sys.stderr, "_proxy_log_path", None) == current_path:
        return

    try:
        PROXY_LOG_FILE.parent.mkdir(parents=True, exist_ok=True)
        _proxy_log_file_handle = open(PROXY_LOG_FILE, "a", encoding="utf-8", buffering=1)
    except OSError as exc:
        fallback_stream = getattr(sys, "__stderr__", None) or sys.stderr
        fallback_stream.write(f"[WARN] Failed to initialize proxy log file {PROXY_LOG_FILE}: {exc}\n")
        fallback_stream.flush()
        return

    if getattr(sys.stdout, "_proxy_log_path", None) != current_path:
        sys.stdout = _ProxyLogTeeStream(sys.stdout, _proxy_log_file_handle)
    if getattr(sys.stderr, "_proxy_log_path", None) != current_path:
        sys.stderr = _ProxyLogTeeStream(sys.stderr, _proxy_log_file_handle)

    logging.captureWarnings(True)
    print(f"[LOG] Proxy log capture initialized at {PROXY_LOG_FILE}", flush=True)


_install_proxy_log_capture()


def _truncate_monitor_text(value: Any, max_chars: int = 240) -> str:
    text = "" if value is None else str(value)
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) <= max_chars:
        return text
    return text[: max(0, max_chars - 3)].rstrip() + "..."


def _read_monitor_log_excerpt(path: Optional[Path], max_lines: int = MONITOR_LOG_PREVIEW_LINES) -> List[str]:
    if not path:
        return []
    text = _tail_text_file(path, max_lines=max_lines)
    return text.splitlines()[-max_lines:] if text else []


def _resolve_monitor_log_path(log_file: Optional[str]) -> Optional[Path]:
    if not log_file:
        return None
    path = Path(log_file)
    if not path.is_absolute():
        path = SCRATCH_DIR / log_file
    return path


def _find_monitor_run(run_id: str) -> Optional[Dict[str, Any]]:
    entry = monitor_active_runs.get(run_id)
    if entry is not None:
        return entry
    for runs in monitor_recent_runs.values():
        for run in runs:
            if run.get("id") == run_id:
                return run
    return None


def _read_monitor_run_log(path: Optional[Path], max_bytes: int = MONITOR_RUN_LOG_MAX_BYTES) -> Dict[str, Any]:
    if not path or not path.exists():
        return {"available": False, "content": "", "truncated": False, "path": str(path) if path else None}
    try:
        file_size = path.stat().st_size
        if file_size <= max_bytes:
            return {
                "available": True,
                "content": path.read_text(encoding="utf-8"),
                "truncated": False,
                "path": str(path),
            }
        head_bytes = max_bytes // 2
        tail_bytes = max_bytes - head_bytes
        with path.open("rb") as f:
            head = f.read(head_bytes)
            f.seek(max(0, file_size - tail_bytes))
            tail = f.read(tail_bytes)
        marker = b"\n\n...[log truncated for monitor response]...\n\n"
        content = (head + marker + tail).decode("utf-8", errors="replace")
        return {"available": True, "content": content, "truncated": True, "path": str(path)}
    except OSError:
        return {"available": False, "content": "", "truncated": False, "path": str(path)}


def _stringify_autogen_message_content(content: Any) -> str:
    if content is None:
        return ""
    if isinstance(content, str):
        return (
            normalize_chat_completion_message(
                {"role": "assistant", "content": content},
                preserve_reasoning_details=False,
            ).get("content")
            or ""
        )
    if isinstance(content, (dict, list)):
        try:
            normalized = normalize_chat_completion_message(
                {"role": "assistant", "content": content},
                preserve_reasoning_details=False,
            )
            rendered = normalized.get("content") or ""
            if rendered:
                return rendered
            return json.dumps(content, ensure_ascii=False, indent=2, default=str)
        except TypeError:
            return str(content)
    return str(content)


def _format_autogen_conversation_log(
    input_text: str,
    messages: List[Dict[str, str]],
    conversation_summary: str,
    *,
    timestamp_human: str,
    status: str = "completed",
    progress_notes: Optional[List[str]] = None,
    error_text: Optional[str] = None,
) -> str:
    lines = [
        "AutoGen team conversation log",
        f"Updated: {timestamp_human}",
        f"Status: {status}",
        "",
        "Input:",
        input_text or "(empty)",
        "",
        "--- Progress ---",
    ]
    if progress_notes:
        for note in progress_notes:
            lines.append(note if note else "(empty)")
    else:
        lines.append("(No progress updates recorded yet)")
    lines.extend(
        [
            "",
            "--- Messages ---",
        ]
    )
    if messages:
        for i, msg in enumerate(messages, 1):
            source = msg.get("source", "unknown")
            content = _stringify_autogen_message_content(msg.get("content", ""))
            lines.append(f"[{i}] {source}:")
            lines.append(content if content else "(empty)")
            lines.append("")
    else:
        lines.append("(No messages returned from AutoGen team)")
        lines.append("")
    lines.extend(
        [
            "--- Conversation Summary ---",
            "",
            conversation_summary or "(pending)",
        ]
    )
    if error_text:
        lines.extend(
            [
                "",
                "--- Error ---",
                "",
                error_text,
            ]
        )
    return "\n".join(lines)


def _monitor_browser_health_is_stale(snapshot: Optional[Dict[str, Any]]) -> bool:
    if not snapshot:
        return True
    checked_at = snapshot.get("checked_at")
    if not checked_at:
        return True
    try:
        checked_dt = datetime.fromisoformat(str(checked_at).replace("Z", "+00:00"))
    except ValueError:
        return True
    if checked_dt.tzinfo is None:
        checked_dt = checked_dt.replace(tzinfo=timezone.utc)
    age_seconds = (datetime.now(timezone.utc) - checked_dt).total_seconds()
    return age_seconds >= MONITOR_BROWSER_HEALTH_MAX_AGE_SECONDS


def _monitor_run_start(agent: str, kind: str, input_text: str = "", metadata: Optional[Dict[str, Any]] = None) -> str:
    run_id = secrets.token_hex(8)
    now_ts = time.time()
    entry = {
        "id": run_id,
        "agent": agent,
        "kind": kind,
        "status": "running",
        "started_at_ts": now_ts,
        "started_at": datetime.fromtimestamp(now_ts, tz=timezone.utc).isoformat(),
        "ended_at_ts": None,
        "ended_at": None,
        "duration_ms": None,
        "input_preview": _truncate_monitor_text(input_text, max_chars=320),
        "summary": "",
        "progress": [],
        "metadata": dict(metadata or {}),
        "log_file": None,
        "log_excerpt": [],
    }
    monitor_recent_runs.setdefault(agent, collections.deque(maxlen=MONITOR_RUN_HISTORY_LIMIT)).append(entry)
    monitor_active_runs[run_id] = entry
    return run_id


def _monitor_run_note(run_id: str, note: str) -> None:
    entry = monitor_active_runs.get(run_id)
    if not entry:
        return
    progress = entry.setdefault("progress", [])
    progress.append(
        {
            "ts": time.time(),
            "iso": datetime.now(timezone.utc).isoformat(),
            "message": _truncate_monitor_text(note, max_chars=300),
        }
    )
    if len(progress) > 25:
        del progress[:-25]


def _monitor_run_update(
    run_id: str,
    *,
    summary: Optional[str] = None,
    metadata: Optional[Dict[str, Any]] = None,
    log_file: Optional[str] = None,
    log_excerpt: Optional[List[str]] = None,
) -> None:
    entry = monitor_active_runs.get(run_id)
    if not entry:
        return
    if summary is not None:
        entry["summary"] = _truncate_monitor_text(summary, max_chars=500)
    if metadata:
        entry.setdefault("metadata", {}).update(metadata)
    if log_file:
        entry["log_file"] = log_file
    if log_excerpt is not None:
        entry["log_excerpt"] = log_excerpt[-MONITOR_LOG_PREVIEW_LINES:]


def _monitor_run_finish(
    run_id: str,
    *,
    status: str,
    summary: str = "",
    metadata: Optional[Dict[str, Any]] = None,
    log_file: Optional[str] = None,
    log_excerpt: Optional[List[str]] = None,
) -> None:
    entry = monitor_active_runs.pop(run_id, None)
    if not entry:
        return
    now_ts = time.time()
    entry["status"] = status
    entry["ended_at_ts"] = now_ts
    entry["ended_at"] = datetime.fromtimestamp(now_ts, tz=timezone.utc).isoformat()
    entry["duration_ms"] = int(max(0.0, now_ts - float(entry.get("started_at_ts") or now_ts)) * 1000)
    entry["summary"] = _truncate_monitor_text(summary, max_chars=500)
    if metadata:
        entry.setdefault("metadata", {}).update(metadata)
    if log_file:
        entry["log_file"] = log_file
    if log_excerpt is not None:
        entry["log_excerpt"] = log_excerpt[-MONITOR_LOG_PREVIEW_LINES:]


def _get_monitor_runs_payload(agent: str) -> Dict[str, Any]:
    recent = list(monitor_recent_runs.get(agent, []))
    active = [run for run in recent if run.get("status") == "running"]
    return {
        "active_count": len(active),
        "recent": list(reversed(recent)),
    }


def _ensure_status_storage() -> None:
    STATUS_DATA_DIR.mkdir(parents=True, exist_ok=True)
    if not STATUS_EVENTS_FILE.exists():
        STATUS_EVENTS_FILE.write_text("", encoding="utf-8")


def _tail_text_file(path: Path, max_lines: int = 200, max_bytes: int = 262144) -> str:
    """Return the last max_lines from a text file, reading up to max_bytes from the end."""
    if not path.exists():
        return ""
    try:
        with path.open("rb") as f:
            f.seek(0, os.SEEK_END)
            size = f.tell()
            read_size = min(size, max_bytes)
            f.seek(-read_size, os.SEEK_END)
            data = f.read(read_size)
        text = data.decode("utf-8", errors="replace")
        lines = text.splitlines()
        return "\n".join(lines[-max_lines:])
    except Exception as exc:
        print(f"[WARN] Failed to tail file {path}: {exc}")
        return ""


def _get_recent_status_events(limit: int = 50) -> List[Dict[str, Any]]:
    """Read recent status events from the JSONL store."""
    text = _tail_text_file(STATUS_EVENTS_FILE, max_lines=limit * 2)
    events: List[Dict[str, Any]] = []
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            events.append(json.loads(line))
        except json.JSONDecodeError:
            continue
    return events[-limit:]


def _load_status_index() -> None:
    status_latest_index.clear()
    if not STATUS_EVENTS_FILE.exists():
        return
    try:
        with STATUS_EVENTS_FILE.open("r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    event = json.loads(line)
                except json.JSONDecodeError:
                    continue
                conv = event.get("conversation_id")
                req = event.get("request_id")
                if not conv or not req:
                    continue
                status_latest_index[(conv, req)] = event
    except Exception as exc:
        print(f"[WARN] Failed to load status index: {exc}")


def _init_status_store() -> None:
    _ensure_status_storage()
    _load_status_index()


# Initialize status storage on module load
_init_status_store()


async def _record_status_event(event: Dict[str, Any]) -> None:
    async with status_write_lock:
        try:
            with STATUS_EVENTS_FILE.open("a", encoding="utf-8") as f:
                f.write(json.dumps(event, ensure_ascii=False) + "\n")
            conv = event.get("conversation_id")
            req = event.get("request_id")
            if conv and req:
                status_latest_index[(conv, req)] = event
        except Exception as exc:
            print(f"[WARN] Failed to write status event: {exc}")


async def _emit_status_event(
    session: StatusSession,
    state: Optional[str] = None,
    phase: Optional[str] = None,
    done: bool = False,
    event_type: str = "heartbeat",
) -> Dict[str, Any]:
    if state is not None:
        session.state = state
    if phase is not None:
        session.phase = phase
    session.seq += 1
    now = datetime.now(timezone.utc)
    event = {
        "ts": now.timestamp(),
        "iso": now.isoformat(),
        "conversation_id": session.conversation_id,
        "request_id": session.request_id,
        "channel": session.channel,
        "seq": session.seq,
        "state": session.state,
        "phase": session.phase,
        "done": done,
        "type": event_type,
    }
    await _record_status_event(event)
    return event


async def _status_heartbeat(session: StatusSession) -> None:
    while not session.done:
        await asyncio.sleep(STATUS_UPDATE_INTERVAL_SECONDS)
        if session.done:
            break
        await _emit_status_event(session, event_type="heartbeat")


async def _start_status_session(
    conversation_id: str,
    request_id: str,
    channel: str,
    initial_state: str,
    phase: str = "start",
) -> StatusSession:
    key = (conversation_id, request_id)
    existing = status_sessions.get(key)
    if existing and existing.heartbeat_task:
        existing.done = True
        with suppress(Exception):
            existing.heartbeat_task.cancel()
    session = StatusSession(
        conversation_id=conversation_id,
        request_id=request_id,
        channel=channel,
        state=initial_state,
        phase=phase,
        seq=0,
        started_at=time.time(),
        heartbeat_task=None,
        done=False,
    )
    status_sessions[key] = session
    await _emit_status_event(session, event_type="start")
    session.heartbeat_task = asyncio.create_task(_status_heartbeat(session))
    return session


async def _update_status_session(
    conversation_id: str,
    request_id: str,
    state: str,
    phase: Optional[str] = None,
) -> Optional[Dict[str, Any]]:
    key = (conversation_id, request_id)
    session = status_sessions.get(key)
    if not session:
        return None
    return await _emit_status_event(session, state=state, phase=phase, event_type="update")


async def _finish_status_session(
    conversation_id: str,
    request_id: str,
    final_state: Optional[str] = None,
    phase: str = "done",
) -> Optional[Dict[str, Any]]:
    key = (conversation_id, request_id)
    session = status_sessions.get(key)
    if not session:
        return None
    if final_state:
        session.state = final_state
    session.done = True
    event = await _emit_status_event(session, phase=phase, done=True, event_type="finish")
    if session.heartbeat_task:
        with suppress(Exception):
            session.heartbeat_task.cancel()
    status_sessions.pop(key, None)
    return event


def _format_telegram_tool_status(tool_name: Any) -> str:
    """Return a user-facing Telegram status update for the given tool."""
    raw_name = str(tool_name or "").strip()
    lowered = raw_name.lower()
    status_map = {
        "websearch": "On it. I'm looking for the best sources now.",
        "scrapewebsite": "On it. I'm reading through the page now.",
        "fetchnews": "On it. I'm pulling together the latest updates.",
        "weatherinfo": "On it. I'm checking the weather details now.",
        "runbrowseragent": "On it. I'm working through that in the browser now.",
        "rundeepresearch": "On it. I'm gathering sources and comparing them now.",
        "healthcheck": "On it. I'm checking the browser task status now.",
        "runworkflow": "On it. I'm running that workflow now.",
        "createslidespresentation": "On it. I'm building the presentation now.",
        "pdftopowerpoint": "On it. I'm converting the document into a PowerPoint now.",
    }
    if lowered in status_map:
        return status_map[lowered]
    if lowered.startswith("googleworkspace_cli."):
        return "On it. I'm checking your Google Workspace data now."
    if "." in raw_name:
        prefix = lowered.split(".", 1)[0]
        if prefix == "filesystem":
            return "On it. I'm checking the workspace files now."
        if prefix == "github":
            return "On it. I'm checking GitHub for that now."
    return f"On it. I'm using {raw_name or 'a tool'} for that now."


def _get_latest_status_event(
    conversation_id: str,
    request_id: Optional[str] = None,
) -> Optional[Dict[str, Any]]:
    if request_id:
        event = status_latest_index.get((conversation_id, request_id))
        if event:
            return event
    # Fallback: scan file for latest event for conversation
    if not STATUS_EVENTS_FILE.exists():
        return None
    latest = None
    try:
        with STATUS_EVENTS_FILE.open("r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    event = json.loads(line)
                except json.JSONDecodeError:
                    continue
                if event.get("conversation_id") != conversation_id:
                    continue
                if request_id and event.get("request_id") != request_id:
                    continue
                latest = event
    except Exception as exc:
        print(f"[WARN] Failed reading status events: {exc}")
    return latest


def _get_status_events_since(
    conversation_id: str,
    request_id: str,
    since_seq: int,
) -> List[Dict[str, Any]]:
    events: List[Dict[str, Any]] = []
    if not STATUS_EVENTS_FILE.exists():
        return events
    try:
        with STATUS_EVENTS_FILE.open("r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    event = json.loads(line)
                except json.JSONDecodeError:
                    continue
                if event.get("conversation_id") != conversation_id:
                    continue
                if event.get("request_id") != request_id:
                    continue
                seq = event.get("seq") or 0
                if seq > since_seq:
                    events.append(event)
    except Exception as exc:
        print(f"[WARN] Failed reading status events: {exc}")
    return events


def _load_telegram_user_links() -> Dict[str, str]:
    """Load Telegram user-link mappings from config/telegram_user_links.json."""
    mapping: Dict[str, str] = {}
    if not TELEGRAM_USER_LINKS_FILE.exists():
        return mapping
    try:
        raw = TELEGRAM_USER_LINKS_FILE.read_text(encoding="utf-8")
        loaded = json.loads(raw)
        if not isinstance(loaded, dict):
            return mapping
        for key, value in loaded.items():
            key_text = str(key or "").strip()
            value_text = str(value or "").strip()
            if key_text and value_text:
                mapping[key_text] = value_text
    except (json.JSONDecodeError, OSError):
        return {}
    return mapping


def _resolve_todo_user_for_telegram(conversation_id: str, user_id: Optional[str]) -> str:
    """
    Resolve the todo user key for Telegram: if telegram_user_links.json maps this conversation
    or user to an app username, return that; else return conversation_id (or user_id) for persistent per-chat todo.
    """
    mapping = _load_telegram_user_links()
    # Normalize to string (API may send user_id as number)
    uid = str(user_id or "").strip()
    cid = str(conversation_id or "").strip() or "default"
    if uid and uid in mapping:
        return str(mapping[uid]).strip() or cid
    if cid in mapping:
        return str(mapping[cid]).strip() or cid
    return cid


def _resolve_telegram_chat_ids_for_todo_user(user_key: str) -> List[str]:
    """
    Resolve Telegram chat IDs that should receive task-execution notifications for a todo user key.

    Supports:
    - direct numeric user_key values (Telegram chat/user IDs)
    - reverse lookup in telegram_user_links.json where value == user_key
    """
    key = str(user_key or "").strip()
    if not key:
        return []
    chat_ids: List[str] = []
    if _TELEGRAM_CHAT_ID_RE.match(key):
        chat_ids.append(key)
    mapping = _load_telegram_user_links()
    for tg_id, mapped_user in mapping.items():
        if mapped_user != key:
            continue
        tg_id_text = str(tg_id or "").strip()
        if tg_id_text and _TELEGRAM_CHAT_ID_RE.match(tg_id_text):
            chat_ids.append(tg_id_text)
    return list(dict.fromkeys(chat_ids))


def _task_execution_register_telegram_target(
    user_key: str,
    chat_id: Optional[Any],
    task_id: Optional[int] = None,
) -> bool:
    """Attach a Telegram chat ID to in-flight task execution(s) so completion can notify the user."""
    chat_id_text = str(chat_id or "").strip()
    if not chat_id_text or not _TELEGRAM_CHAT_ID_RE.match(chat_id_text):
        return False

    targets: List[Dict[str, Any]] = []
    normalized_tid = _coerce_task_id(task_id)
    if normalized_tid is not None:
        target = _get_task_run_state(user_key, normalized_tid)
        if isinstance(target, dict) and not _is_task_execution_terminal_status(_state_status_lower(target)):
            targets.append(target)
    else:
        targets.extend(state for _, state in _active_task_runs(user_key))

    # Legacy single-state fallback (pre-migration tests or stale in-memory state).
    if not targets:
        legacy = task_execution_state.get(user_key)
        if isinstance(legacy, dict) and "runs" not in legacy:
            targets.append(legacy)

    if not targets:
        return False

    for state in targets:
        existing = state.get("telegram_chat_ids")
        if not isinstance(existing, list):
            existing = []
            state["telegram_chat_ids"] = existing
        if chat_id_text not in existing:
            existing.append(chat_id_text)
    return True


def _build_task_execution_telegram_message(
    *,
    task_id: Optional[int],
    task_description: str,
    status: str,
    result_message: str,
) -> str:
    now_utc = datetime.now(timezone.utc).isoformat()
    body = (result_message or "").strip()
    if len(body) > 1800:
        body = body[:1800].rstrip() + "\n...[truncated]"
    lines = [
        "Scheduled task execution completed.",
        f"Task ID: {task_id if task_id is not None else 'unknown'}",
        f"Task: {task_description or '(no description)'}",
        f"Status: {status or 'unknown'}",
        f"Completed at (UTC): {now_utc}",
    ]
    if body:
        lines.append("")
        lines.append("Result:")
        lines.append(body)
    return "\n".join(lines)


async def _send_telegram_bot_message(chat_id: str, text: str) -> bool:
    """Send a direct Telegram message via Bot API. Returns True when sent successfully."""
    token = str(os.getenv("TELEGRAM_BOT_TOKEN") or "").strip()
    chat_id_text = str(chat_id or "").strip()
    message_text = (text or "").strip()
    if not token or not chat_id_text or not message_text:
        return False
    url = f"https://api.telegram.org/bot{token}/sendMessage"
    payload = {"chat_id": chat_id_text, "text": message_text}
    try:
        async with httpx.AsyncClient(timeout=20.0) as client:
            response = await client.post(url, json=payload)
        if response.status_code != 200:
            print(f"[WARN] Telegram sendMessage failed ({response.status_code}) for chat_id={chat_id_text}", flush=True)
            return False
        data = response.json()
        ok = bool(data.get("ok"))
        if not ok:
            print(f"[WARN] Telegram sendMessage returned ok=false for chat_id={chat_id_text}", flush=True)
        return ok
    except Exception as exc:
        print(f"[WARN] Telegram sendMessage error for chat_id={chat_id_text}: {exc}", flush=True)
        return False


async def _send_telegram_file_internal(
    chat_id: Any,
    filename: str,
    caption: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Send a file from scratch to a Telegram chat via sendDocument.
    Used by Telegram tools only.
    """
    token = str(os.getenv("TELEGRAM_BOT_TOKEN") or "").strip()
    chat_id_text = str(chat_id or "").strip()
    logical_name = str(filename or "").strip()
    if not token:
        return {"success": False, "message": "TELEGRAM_BOT_TOKEN is not configured."}
    if not chat_id_text or not _TELEGRAM_CHAT_ID_RE.match(chat_id_text):
        return {"success": False, "message": "Current Telegram chat id is unavailable."}
    if not logical_name:
        return {"success": False, "message": "filename is required."}

    try:
        filepath = resolve_scratch_path(logical_name)
    except HTTPException as e:
        return {"success": False, "message": e.detail or "Invalid filename."}

    if not filepath.exists() or not filepath.is_file():
        return {"success": False, "message": f"File not found: {logical_name}"}

    try:
        size_bytes = filepath.stat().st_size
    except OSError as exc:
        return {"success": False, "message": f"Failed to read file metadata: {exc}"}

    if size_bytes > FILE_OPS_MAX_SIZE_BYTES:
        return {
            "success": False,
            "message": f"File is too large ({size_bytes} bytes). Limit is {FILE_OPS_MAX_SIZE_BYTES} bytes.",
        }

    caption_text = (caption or "").strip()
    if len(caption_text) > 1024:
        caption_text = caption_text[:1024]

    mime_type, _ = mimetypes.guess_type(str(filepath))
    url = f"https://api.telegram.org/bot{token}/sendDocument"
    payload = {"chat_id": chat_id_text}
    if caption_text:
        payload["caption"] = caption_text

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            with filepath.open("rb") as f:
                response = await client.post(
                    url,
                    data=payload,
                    files={
                        "document": (
                            filepath.name,
                            f,
                            mime_type or "application/octet-stream",
                        )
                    },
                )
        if response.status_code != 200:
            detail = response.text
            try:
                parsed = response.json()
                detail = parsed.get("description") or parsed.get("error_code") or detail
            except ValueError:
                pass
            return {"success": False, "message": f"Telegram sendDocument failed ({response.status_code}): {detail}"}
        parsed = response.json()
        if not parsed.get("ok"):
            return {"success": False, "message": parsed.get("description", "Telegram sendDocument returned ok=false.")}
        result = parsed.get("result") or {}
        return {
            "success": True,
            "message": f"Sent {filepath.name} to Telegram.",
            "data": {
                "filename": filepath.name,
                "size": size_bytes,
                "message_id": result.get("message_id"),
                "chat_id": chat_id_text,
            },
        }
    except Exception as exc:
        return {"success": False, "message": f"Failed to send file to Telegram: {exc}"}


async def _maybe_notify_telegram_task_completion(
    user_key: str,
    state: Dict[str, Any],
    status: str,
    message: str,
) -> None:
    """
    Send Telegram notification when a scheduled task execution reaches completion state.
    No-op when task is not scheduled or no resolvable Telegram chat IDs are available.
    """
    if status != STATUS_AWAITING_CONFIRMATION:
        return
    if not bool(state.get("is_scheduled")):
        return
    chat_ids = state.get("telegram_chat_ids")
    if not isinstance(chat_ids, list) or not chat_ids:
        return
    resolved_task_id = _resolve_task_id_for_state(user_key, state)
    text = _build_task_execution_telegram_message(
        task_id=resolved_task_id or state.get("task_id"),
        task_description=str(state.get("task_description") or ""),
        status=status,
        result_message=message or "",
    )
    for chat_id in chat_ids:
        await _send_telegram_bot_message(str(chat_id), text)


def _auto_complete_scheduled_execution(user_key: str, state: Dict[str, Any], status: str) -> Optional[str]:
    """
    Automatically complete scheduled executions once the run loop reaches
    awaiting_confirmation so recurring tasks advance their lifecycle.
    """
    if status != STATUS_AWAITING_CONFIRMATION:
        return None
    if not bool(state.get("is_scheduled")):
        return None
    if not TODO_STORE_AVAILABLE or not _todo_store:
        return "Scheduled run finished, but automatic completion is unavailable (todo store not available)."

    task_id = _resolve_task_id_for_state(user_key, state)
    if task_id is None:
        return "Scheduled run finished, but automatic completion was skipped (invalid task id)."

    try:
        result = _todo_store.complete_task(user_key, task_id)
    except Exception as exc:
        return f"Scheduled run finished, but automatic completion failed: {exc}"

    if bool(result.get("rescheduled")):
        next_run = str(result.get("next_run_at") or "").strip() or "the next scheduled run"
        return f"Scheduled task {task_id} auto-completed and was rescheduled for {next_run}."
    return f"Scheduled task {task_id} auto-completed and was removed from the todo list."

# Philosopher mode state storage (per conversation)
philosopher_mode_active: Dict[str, bool] = {}
philosopher_mode_instances: Dict[str, Any] = {}

# Task execution: max iterations per run (configurable via .env)
TASK_EXECUTION_MAX_ITERATIONS = max(1, min(200, int(os.getenv("TASK_EXECUTION_MAX_ITERATIONS", "200"))))
# Per-user execution state:
# user_key -> {"runs": {task_id: run_state}}
task_execution_state: Dict[str, Dict[str, Any]] = {}

# Task execution module (bounded LLM+tools loop for todo tasks)
try:
    from src.features.task_execution import TodoTaskExecutor
    from src.features.task_execution import (
        STATUS_EXECUTING,
        STATUS_PAUSED_AWAITING_FEEDBACK,
        STATUS_AWAITING_CONFIRMATION,
        STATUS_CANCELLED,
    )
    TASK_EXECUTION_AVAILABLE = True
except ImportError as e:
    print(f"[WARN] Task execution not available: {e}")
    TodoTaskExecutor = None
    TASK_EXECUTION_AVAILABLE = False
    STATUS_EXECUTING = "executing"
    STATUS_PAUSED_AWAITING_FEEDBACK = "paused_awaiting_feedback"
    STATUS_AWAITING_CONFIRMATION = "awaiting_confirmation"
    STATUS_CANCELLED = "cancelled"


def _is_task_execution_terminal_status(status: Optional[str]) -> bool:
    """Statuses that should be treated as no-longer-active execution runs."""
    return str(status or "").strip().lower() in {
        STATUS_AWAITING_CONFIRMATION,
        STATUS_CANCELLED,
    }


def _coerce_task_id(value: Any) -> Optional[int]:
    try:
        task_id = int(value)
    except (TypeError, ValueError):
        return None
    return task_id if task_id >= 1 else None


def _state_status_lower(state: Dict[str, Any]) -> str:
    return str(state.get("status") or "").strip().lower()


def _get_user_task_runs(user_key: str, *, create: bool = False) -> Dict[int, Dict[str, Any]]:
    """
    Return per-task run map for a user.
    Backward-compatible with legacy single-state shape by migrating it in-memory.
    """
    entry = task_execution_state.get(user_key)
    if entry is None:
        if not create:
            return {}
        task_execution_state[user_key] = {"runs": {}}
        return task_execution_state[user_key]["runs"]
    if not isinstance(entry, dict):
        if not create:
            return {}
        task_execution_state[user_key] = {"runs": {}}
        return task_execution_state[user_key]["runs"]

    runs_raw = entry.get("runs")
    if isinstance(runs_raw, dict):
        normalized: Dict[int, Dict[str, Any]] = {}
        mutated = False
        for raw_key, raw_state in runs_raw.items():
            if not isinstance(raw_state, dict):
                mutated = True
                continue
            tid = _coerce_task_id(raw_key)
            if tid is None:
                tid = _coerce_task_id(raw_state.get("task_id"))
            if tid is None:
                mutated = True
                continue
            raw_state["task_id"] = tid
            normalized[tid] = raw_state
            if raw_key != tid:
                mutated = True
        if mutated:
            entry["runs"] = normalized
        return entry.get("runs", {})

    # Legacy single-run shape: {"task_id", "status", "executor", ...}
    legacy_tid = _coerce_task_id(entry.get("task_id"))
    if legacy_tid is None and any(k in entry for k in ("status", "executor", "telegram_chat_ids", "message")):
        legacy_tid = 1
    migrated_runs: Dict[int, Dict[str, Any]] = {}
    if legacy_tid is not None:
        entry["task_id"] = legacy_tid
        migrated_runs[legacy_tid] = entry
    task_execution_state[user_key] = {"runs": migrated_runs}
    return task_execution_state[user_key]["runs"]


def _set_task_run_state(user_key: str, task_id: int, state: Dict[str, Any]) -> None:
    runs = _get_user_task_runs(user_key, create=True)
    state["task_id"] = task_id
    runs[int(task_id)] = state


def _remove_task_run_state(user_key: str, task_id: int) -> None:
    runs = _get_user_task_runs(user_key, create=False)
    runs.pop(int(task_id), None)
    if not runs:
        task_execution_state.pop(user_key, None)


def _get_task_run_state(user_key: str, task_id: int) -> Optional[Dict[str, Any]]:
    runs = _get_user_task_runs(user_key, create=False)
    return runs.get(int(task_id))


def _cleanup_terminal_task_runs(user_key: str) -> None:
    runs = _get_user_task_runs(user_key, create=False)
    if not runs:
        return
    remove_ids = [tid for tid, state in runs.items() if _is_task_execution_terminal_status(_state_status_lower(state))]
    for tid in remove_ids:
        runs.pop(tid, None)
    if not runs:
        task_execution_state.pop(user_key, None)


def _active_task_runs(user_key: str) -> List[Tuple[int, Dict[str, Any]]]:
    _cleanup_terminal_task_runs(user_key)
    runs = _get_user_task_runs(user_key, create=False)
    out: List[Tuple[int, Dict[str, Any]]] = []
    for tid, state in runs.items():
        if _is_task_execution_terminal_status(_state_status_lower(state)):
            continue
        out.append((tid, state))
    out.sort(key=lambda item: item[0])
    return out


def _resolve_task_id_for_state(user_key: str, state: Dict[str, Any]) -> Optional[int]:
    """Resolve current stable task ID for a run, preferring task item id when available."""
    fallback = _coerce_task_id(state.get("task_id"))
    task_item_id = str(state.get("task_item_id") or "").strip()
    if not task_item_id or not TODO_STORE_AVAILABLE or not _todo_store:
        return fallback
    try:
        meta = _todo_store.load_tasks_with_meta(user_key)
        task_items = meta.get("task_items") if isinstance(meta, dict) else None
        if isinstance(task_items, list):
            for item in task_items:
                if isinstance(item, dict) and str(item.get("id") or "").strip() == task_item_id:
                    resolved = _coerce_task_id(item.get("task_id"))
                    if resolved is not None:
                        return resolved
    except Exception:
        return fallback
    return fallback


def _state_brief_for_response(user_key: str, state: Dict[str, Any]) -> Dict[str, Any]:
    resolved_task_id = _resolve_task_id_for_state(user_key, state)
    original_task_id = _coerce_task_id(state.get("task_id"))
    diagnostics = _safe_task_execution_diagnostics(state.get("executor"))
    out = {
        "status": state.get("status"),
        "task_id": resolved_task_id or original_task_id,
        "message": state.get("message"),
        "current_step": diagnostics.get("iterations"),
        "total_steps": diagnostics.get("max_iterations"),
        "elapsed_seconds": diagnostics.get("elapsed_seconds"),
    }
    if state.get("run_id"):
        out["run_id"] = state.get("run_id")
    if original_task_id is not None and resolved_task_id is not None and original_task_id != resolved_task_id:
        out["original_task_id"] = original_task_id
    return out

# Assistant context: current date, timezone, and knowledge-gap awareness (prepended to all chat system prompts)
def _get_assistant_context_block() -> str:
    """Returns context block with server date, timezone, and knowledge-awareness instructions."""
    now = datetime.now().astimezone()
    date_str = now.strftime("%A, %d %B %Y")  # e.g. Thursday, 13 February 2025
    tz_str = now.strftime("%Z (UTC%z)")
    return (
        f"Context: Today's date is {date_str}. You are running in timezone {tz_str}. "
        "Use this when interpreting dates, times, and relative references (e.g. 'today', 'this week') unless the user specifies otherwise.\n"
        "Knowledge awareness: Your training has a cutoff. Acknowledge your knowledge gap; do not assume the current year or recent events. "
        "When the user provides current facts, corrections, or information that differs from your training "
        '(e.g., "it\'s 2025 now", "that API changed"), accept them as authoritative and do not contradict them.\n\n'
    )


def _read_optional_prompt_file(path: Path) -> str:
    if not path.exists():
        return ""
    try:
        return path.read_text(encoding="utf-8").strip()
    except Exception as e:
        print(f"Warning: Could not read {path}: {e}")
        return ""


def _get_soul_prompt_text() -> str:
    """Return optional persona text loaded from config/soul.md."""
    return _read_optional_prompt_file(SOUL_PROMPT_FILE)


def _compose_system_prompt_with_context(system_prompt: Optional[str]) -> str:
    """Prepend runtime context and optional soul prompt to a base system prompt."""
    parts: List[str] = [_get_assistant_context_block().strip()]
    soul_prompt = _get_soul_prompt_text()
    if soul_prompt:
        parts.append(soul_prompt)
    base_prompt = (system_prompt or "").strip()
    if base_prompt:
        parts.append(base_prompt)
    return "\n\n".join(part for part in parts if part).strip()


def _normalize_chat_endpoint(endpoint: str) -> str:
    if not endpoint:
        return ""
    if endpoint.endswith("/chat/completions"):
        return endpoint
    return endpoint.rstrip("/") + "/chat/completions"


def _estimate_total_tokens(messages: List[Dict[str, Any]], max_tokens: int) -> int:
    return estimate_tokens_from_messages(messages) + max_tokens


def _get_max_tokens_from_payload(payload: Dict[str, Any]) -> int:
    raw = payload.get("max_tokens")
    try:
        if raw is None:
            return 0
        return max(0, int(raw))
    except (TypeError, ValueError):
        return 0


_shared_chat_http_client: Optional[httpx.AsyncClient] = None
_shared_chat_http_client_lock = asyncio.Lock()


async def _get_shared_chat_http_client() -> httpx.AsyncClient:
    global _shared_chat_http_client
    client = _shared_chat_http_client
    if client is not None and not client.is_closed:
        return client
    async with _shared_chat_http_client_lock:
        client = _shared_chat_http_client
        if client is None or client.is_closed:
            _shared_chat_http_client = httpx.AsyncClient()
        return _shared_chat_http_client


async def _close_shared_chat_http_client() -> None:
    global _shared_chat_http_client
    async with _shared_chat_http_client_lock:
        client = _shared_chat_http_client
        _shared_chat_http_client = None
    if client is not None and not client.is_closed:
        await client.aclose()


async def _call_chat_completion(
    endpoint: str,
    headers: Dict[str, str],
    payload: Dict[str, Any],
    timeout_seconds: float = 120.0,
) -> httpx.Response:
    prepared_payload = prepare_openai_compatible_chat_payload(
        payload,
        api_base=endpoint,
        model=payload.get("model"),
    )
    async with httpx.AsyncClient(timeout=timeout_seconds) as client:
        return await client.post(endpoint, json=prepared_payload, headers=headers)


_NO_KEY_LLM_PROVIDERS = frozenset({"ollama", "bedrock"})
_MCP_PROVIDER_API_KEY_ENV_CANDIDATES: Dict[str, List[str]] = {
    "openai": ["OPENAI_API_KEY", "MCP_LLM_OPENAI_API_KEY"],
    "minimax": [
        "MINIMAX_API_KEY",
        "MCP_LLM_MINIMAX_API_KEY",
        "MCP_LLM_OPENAI_API_KEY",
        "OPENAI_API_KEY",
    ],
    "anthropic": ["ANTHROPIC_API_KEY", "MCP_LLM_ANTHROPIC_API_KEY"],
    "google": ["GEMINI_API_KEY", "GOOGLE_API_KEY", "MCP_LLM_GOOGLE_API_KEY"],
    "azure_openai": ["AZURE_OPENAI_API_KEY", "MCP_LLM_AZURE_OPENAI_API_KEY"],
    "groq": ["GROQ_API_KEY", "MCP_LLM_GROQ_API_KEY"],
    "deepseek": ["DEEPSEEK_API_KEY", "MCP_LLM_DEEPSEEK_API_KEY"],
    "cerebras": ["CEREBRAS_API_KEY", "MCP_LLM_CEREBRAS_API_KEY"],
    "browser_use": ["BROWSER_USE_API_KEY", "MCP_LLM_BROWSER_USE_API_KEY"],
    "openrouter": [
        "OPENROUTER_API_KEY",
        "MCP_LLM_OPENROUTER_API_KEY",
        "MCP_LLM_OPENAI_API_KEY",
        "OPENAI_API_KEY",
    ],
    "vercel": ["VERCEL_API_KEY", "MCP_LLM_VERCEL_API_KEY"],
}


def _first_non_empty_env(var_names: List[str]) -> Optional[str]:
    for var_name in var_names:
        value = os.getenv(var_name)
        if value and value.strip():
            return value.strip()
    return None


def _get_mcp_llm_provider() -> str:
    return (os.getenv("MCP_LLM_PROVIDER") or "").strip().lower()


def _get_mcp_llm_chat_endpoint() -> Optional[str]:
    base = (os.getenv("MCP_LLM_BASE_URL") or "").strip()
    if not base:
        return None
    return _normalize_chat_endpoint(base)


def _get_mcp_llm_model_name() -> Optional[str]:
    model_name = (os.getenv("MCP_LLM_MODEL_NAME") or "").strip()
    return model_name or None


def _resolve_mcp_llm_api_key(provider: Optional[str] = None) -> Optional[str]:
    normalized_provider = (provider or _get_mcp_llm_provider() or "").strip().lower()
    candidates: List[str] = ["MCP_LLM_API_KEY"]
    candidates.extend(_MCP_PROVIDER_API_KEY_ENV_CANDIDATES.get(normalized_provider, []))

    if normalized_provider:
        candidates.append(f"MCP_LLM_{normalized_provider.upper()}_API_KEY")
    else:
        # OpenAI-compatible default fallback for generic/unspecified provider.
        candidates.extend(["MCP_LLM_OPENAI_API_KEY", "OPENAI_API_KEY"])

    api_key = _first_non_empty_env(candidates)
    if api_key:
        return api_key
    if normalized_provider in _NO_KEY_LLM_PROVIDERS:
        return None
    return None


def _build_mcp_fallback_headers(primary_headers: Dict[str, str]) -> Dict[str, str]:
    headers: Dict[str, str] = {"Content-Type": "application/json"}

    for header_name in ("OpenAI-Organization", "OpenAI-Project", "HTTP-Referer", "X-Title"):
        header_value = primary_headers.get(header_name)
        if header_value:
            headers[header_name] = header_value

    provider = _get_mcp_llm_provider()
    api_key = _resolve_mcp_llm_api_key(provider)
    inherited_auth = (primary_headers.get("Authorization") or "").strip()

    if provider == "azure_openai":
        if api_key:
            headers["api-key"] = api_key
        elif inherited_auth.lower().startswith("bearer "):
            inherited_token = inherited_auth[7:].strip()
            if inherited_token:
                headers["api-key"] = inherited_token
        return headers

    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"
    elif inherited_auth:
        headers["Authorization"] = inherited_auth
    return headers


def _build_mcp_fallback_payload(payload: Dict[str, Any]) -> Dict[str, Any]:
    fallback_payload = dict(payload)
    fallback_model = _get_mcp_llm_model_name()
    if fallback_model:
        fallback_payload["model"] = fallback_model
    return fallback_payload


async def _attempt_mcp_chat_fallback(
    primary_headers: Dict[str, str],
    payload: Dict[str, Any],
    timeout_seconds: float,
    source_label: str,
) -> Tuple[Optional[httpx.Response], Optional[str]]:
    fallback_endpoint = _get_mcp_llm_chat_endpoint()
    if not fallback_endpoint:
        return None, "MCP_LLM_BASE_URL is not configured"

    fallback_payload = _build_mcp_fallback_payload(payload)
    fallback_headers = _build_mcp_fallback_headers(primary_headers)
    fallback_model = fallback_payload.get("model", "")
    fallback_provider = _get_mcp_llm_provider() or "openai-compatible"
    print(
        f"[LLM_FALLBACK] {source_label}: trying provider={fallback_provider}, model={fallback_model}, endpoint={fallback_endpoint}",
        flush=True,
    )
    try:
        response = await _call_chat_completion(
            fallback_endpoint,
            fallback_headers,
            fallback_payload,
            timeout_seconds=timeout_seconds,
        )
        print(f"[LLM_FALLBACK] {source_label}: status={response.status_code}", flush=True)
        return response, None
    except httpx.RequestError as exc:
        err = str(exc)
        print(f"[LLM_FALLBACK] {source_label}: request error: {err}", flush=True)
        return None, err


async def _extract_memories_from_recent_messages_async(recent_messages: List[Dict[str, str]]) -> None:
    if not MEMORY_AVAILABLE or not memory_manager or not recent_messages:
        return
    try:
        await memory_manager.extract_memories_from_conversation(
            messages=recent_messages,
            max_memories=3,
        )
    except Exception as e:
        print(f"Warning: Failed to extract memories: {e}")


async def _summarize_messages_for_budget_proxy(
    messages: List[Dict[str, Any]],
    endpoint: str,
    headers: Dict[str, str],
    model_name: str,
    max_tokens: int,
    large_model: Optional[str] = None,
    large_endpoint: Optional[str] = None,
) -> List[Dict[str, Any]]:
    if not messages:
        return messages

    system_msg = messages[0] if messages and messages[0].get("role") == "system" else None
    remainder = messages[1:] if system_msg else messages[:]
    if len(remainder) <= 4:
        return messages

    keep_tail_options = [8, 6, 4, 2]
    for keep_tail in keep_tail_options:
        if len(remainder) <= keep_tail:
            continue
        head = remainder[:-keep_tail]
        tail = remainder[-keep_tail:]
        summary_source_text = format_messages_for_summary(head, max_chars=40000)
        if not summary_source_text.strip():
            break

        summary_prompt = (
            "Summarize the prior conversation so the assistant can continue accurately. "
            "Include key requirements, decisions, file paths, commands, tool outputs, errors, "
            "and remaining TODOs. Keep it concise and structured."
        )
        summary_messages = [
            {"role": "system", "content": "You are a summarization assistant. Summarize accurately and concisely."},
            {"role": "user", "content": f"{summary_prompt}\n\nConversation:\n{summary_source_text}"},
        ]

        summary_model = large_model or model_name
        summary_endpoint = _normalize_chat_endpoint(large_endpoint or endpoint)
        summary_payload = {
            "model": summary_model,
            "messages": summary_messages,
            "temperature": 0.2,
            "max_tokens": 800,
        }
        try:
            summary_response = await _call_chat_completion(summary_endpoint, headers, summary_payload, timeout_seconds=120.0)
            if summary_response.status_code == 200:
                data = summary_response.json()
                summary_text = (data.get("choices") or [{}])[0].get("message", {}).get("content", "").strip()
            else:
                summary_text = ""
        except Exception:
            summary_text = ""

        if not summary_text:
            summary_text = summary_source_text[-2000:] if summary_source_text else ""

        new_messages: List[Dict[str, Any]] = []
        if system_msg:
            new_messages.append(system_msg)
        new_messages.append({"role": "system", "content": f"Summary of previous context:\n{summary_text}"})
        new_messages.extend(tail)

        if _estimate_total_tokens(new_messages, max_tokens) <= get_max_token_limit():
            return new_messages

    return messages


# Telegram/OpenAI configuration
def _parse_telegram_history_limit() -> int:
    """Parse TELEGRAM_HISTORY_LIMIT with default 12 on invalid value."""
    try:
        return max(1, int(os.getenv("TELEGRAM_HISTORY_LIMIT", "12")))
    except ValueError:
        return 12


def _parse_telegram_chat_timeout() -> float:
    """Parse TELEGRAM_CHAT_TIMEOUT with default 30 on invalid value."""
    try:
        raw_timeout = float(os.getenv("TELEGRAM_CHAT_TIMEOUT", "30"))
    except ValueError:
        return 30.0
    hard_cap_raw = os.getenv("TELEGRAM_CHAT_TIMEOUT_HARD_CAP", "120")
    try:
        hard_cap = max(1.0, float(hard_cap_raw))
    except ValueError:
        hard_cap = 120.0
    if raw_timeout > hard_cap:
        print(
            f"[WARN] TELEGRAM_CHAT_TIMEOUT={raw_timeout} exceeds hard cap {hard_cap}; using {hard_cap}.",
            flush=True,
        )
    return max(1.0, min(raw_timeout, hard_cap))


def _parse_telegram_tool_followup_timeout() -> float:
    """Parse TELEGRAM_TOOL_FOLLOWUP_TIMEOUT with default 45 and cap to chat timeout."""
    try:
        raw = float(os.getenv("TELEGRAM_TOOL_FOLLOWUP_TIMEOUT", "45"))
    except ValueError:
        raw = 45.0
    return max(1.0, min(raw, TELEGRAM_CHAT_TIMEOUT))


TELEGRAM_DEFAULT_MODEL = os.getenv("TELEGRAM_MODEL") or os.getenv("OPENAI_MODEL") or os.getenv("MCP_LLM_MODEL_NAME", "gpt-4o-mini")
TELEGRAM_SYSTEM_PROMPT_ENV = os.getenv(
    "TELEGRAM_SYSTEM_PROMPT",
    "You are CATBot, a helpful AI assistant that responds concisely for Telegram users.",
)


def _get_telegram_system_prompt_base() -> str:
    """Return the base system prompt for Telegram: config file if present, else TELEGRAM_SYSTEM_PROMPT env."""
    content = _read_optional_prompt_file(CATBOT_SYSTEM_PROMPT_FILE)
    if content:
        return content
    return TELEGRAM_SYSTEM_PROMPT_ENV


def _sanitize_telegram_legacy_tool_prompt(content: str) -> str:
    """Rewrite legacy XML-only Telegram tool prompt text toward structured tool calling."""
    if not content:
        return content

    structured_intro = (
        "To use a tool, prefer structured tool calls with the exact tool name and JSON schema provided.\n"
        "Return only the tool call when a tool is required, and return plain user-facing text when answering directly.\n"
        "Use XML tool markup only as a legacy fallback when structured tool calls are unavailable."
    )
    content = re.sub(
        r"To use a tool, you MUST ALWAYS respond in this EXACT format:[\s\S]*?"
        r"IMPORTANT: Always use the XML-style format shown above\. Never return raw JSON or other formats\.",
        structured_intro,
        content,
        count=1,
        flags=re.IGNORECASE,
    )
    content = re.sub(
        r'\nUser: "Add a task to call John"[\s\S]*?\nIMPORTANT REMINDER:',
        (
            "\nStructured tool-call XML examples have been removed. "
            "Use the generated tool schemas below and prefer structured tool calls.\n\n"
            "IMPORTANT REMINDER:"
        ),
        content,
        count=1,
    )
    content = content.replace(
        "IMPORTANT REMINDER: Always wrap your tool calls in <tool> and <parameters> tags as shown in the examples above. Never return raw JSON.",
        "IMPORTANT REMINDER: Prefer structured tool calls. Use XML tool markup only as a legacy fallback.",
    )
    content = content.replace(
        "- If another tool is needed, output only the XML tool call and nothing else.",
        "- If another tool is needed, return only the structured tool call with no user-facing text.",
    )
    return content


def _get_telegram_system_prompt_with_tools(conversation_id: str, todo_user_key: Optional[str] = None) -> str:
    """Return the tool-capable system prompt for Telegram, with current todo and memory cache for this conversation."""
    content = _read_optional_prompt_file(CATBOT_SYSTEM_PROMPT_WITH_TOOLS_FILE)
    if not content:
        content = _get_telegram_system_prompt_base()
    content = _sanitize_telegram_legacy_tool_prompt(content)
    # Use persistent todo store when available and todo_user_key provided; else in-memory fallback
    if todo_user_key and TODO_STORE_AVAILABLE and _todo_store:
        todo_list = _todo_store.load_tasks(todo_user_key)
    else:
        todo_list = telegram_todo.get(conversation_id, [])
    mem_cache = telegram_memory_cache.get(conversation_id, [])
    todo_block = "\n".join([f"{i + 1}. {t}" for i, t in enumerate(todo_list)]) if todo_list else "(empty)"
    mem_block = "\n".join([f"{i + 1}. {m}" for i, m in enumerate(mem_cache)]) if mem_cache else "(empty)"
    content = content.replace("{{MEMORY_CACHE}}", mem_block).replace("{{TODO_LIST}}", todo_block)
    native_tool_block = _build_telegram_native_tools_prompt_block()
    if native_tool_block:
        content = f"{content.rstrip()}\n\n{native_tool_block}"
    dynamic_skill_block = _build_telegram_skill_tools_prompt_block()
    if dynamic_skill_block:
        content = f"{content.rstrip()}\n\n{dynamic_skill_block}"
    return content


TELEGRAM_HISTORY_LIMIT = _parse_telegram_history_limit()
TELEGRAM_CHAT_TIMEOUT = _parse_telegram_chat_timeout()
TELEGRAM_TOOL_FOLLOWUP_TIMEOUT = _parse_telegram_tool_followup_timeout()
TELEGRAM_OPENAI_BASE_URL = (
    os.getenv("TELEGRAM_OPENAI_BASE_URL")
    or os.getenv("OPENAI_API_BASE")
    or os.getenv("MCP_LLM_BASE_URL")
    or "https://api.openai.com/v1"
)
TELEGRAM_OPENAI_CHAT_PATH = os.getenv("TELEGRAM_OPENAI_CHAT_PATH", "/chat/completions")
OPENAI_ORG_ID = os.getenv("OPENAI_ORG_ID") or os.getenv("OPENAI_ORGANIZATION")
OPENAI_PROJECT_ID = os.getenv("OPENAI_PROJECT_ID")
# Optional shared secrets for internal service-to-proxy auth.
# TELEGRAM_SECRET protects Telegram-specific flows.
# AUTOGEN_TEAM_SECRET protects internal AutoGen team tool calls to selected proxy routes.
TELEGRAM_SECRET = os.getenv("TELEGRAM_SECRET")
AUTOGEN_TEAM_SECRET = (os.getenv("AUTOGEN_TEAM_SECRET") or os.getenv("CATBOT_AGENT_SECRET") or "").strip() or None

# Large payload model fallback (optional)
LARGE_PAYLOAD_MODEL = (os.getenv("LARGE_PAYLOAD_MODEL") or "").strip() or None
LARGE_PAYLOAD_ENDPOINT = (os.getenv("LARGE_PAYLOAD_ENDPOINT") or "").strip() or None

# Codex CLI configuration
CODEX_ENABLED = os.getenv("CODEX_ENABLED", "true").lower() == "true"
CODEX_CLI_PATH = (os.getenv("CODEX_CLI_PATH") or "codex").strip() or "codex"
CODEX_SANDBOX_MODE = (os.getenv("CODEX_SANDBOX_MODE") or "workspace-write").strip() or "workspace-write"
CODEX_APPROVAL_POLICY = (os.getenv("CODEX_APPROVAL_POLICY") or "never").strip() or "never"
CODEX_ENABLE_SEARCH = os.getenv("CODEX_ENABLE_SEARCH", "true").lower() == "true"
CODEX_TIMEOUT_SECONDS = int(os.getenv("CODEX_TIMEOUT_SECONDS", "1800"))
CODEX_JSON_EVENTS = os.getenv("CODEX_JSON_EVENTS", "true").lower() == "true"
CODEX_OUTPUT_LAST_MESSAGE = os.getenv("CODEX_OUTPUT_LAST_MESSAGE", "true").lower() == "true"
CODEX_AUTOGEN_WORKSPACES_DIRNAME = "autogen"

# Auth configuration
AUTH_USERS_FILE = _PROJECT_ROOT / "config" / "auth_users.json"
ENV_FILE = _PROJECT_ROOT / ".env"
JWT_SECRET = os.getenv("JWT_SECRET", "change-this-secret-in-production")
JWT_ALGORITHM = "HS256"
JWT_EXPIRATION_SECONDS = int(os.getenv("JWT_EXPIRATION_SECONDS", "3600"))
SPOTIFY_ACCOUNTS_BASE = "https://accounts.spotify.com"
SPOTIFY_AUTH_BASE = "https://accounts.spotify.com/api"
SPOTIFY_OAUTH_TIMEOUT_SECONDS = 20.0
SPOTIFY_OAUTH_STATE_TTL_SECONDS = 600
SPOTIFY_PLAYBACK_SCOPES = ("user-modify-playback-state", "user-read-playback-state")
spotify_oauth_pending_states: Dict[str, float] = {}

# Simple in-memory user storage with JSON persistence
users_db: Dict[str, Dict[str, str]] = {}


def _base64url_encode(data: bytes) -> str:
    return base64.urlsafe_b64encode(data).decode().rstrip("=")


def _base64url_decode(data: str) -> bytes:
    padding = "=" * (-len(data) % 4)
    return base64.urlsafe_b64decode(f"{data}{padding}")


def hash_password(password: str, salt: str) -> str:
    hashed = hashlib.pbkdf2_hmac(
        "sha256",
        password.encode("utf-8"),
        salt.encode("utf-8"),
        100_000,
    )
    return _base64url_encode(hashed)


def create_password_record(password: str) -> Dict[str, str]:
    salt = secrets.token_hex(16)
    return {
        "salt": salt,
        "password_hash": hash_password(password, salt),
    }


def verify_password(password: str, salt: str, expected_hash: str) -> bool:
    candidate_hash = hash_password(password, salt)
    return hmac.compare_digest(candidate_hash, expected_hash)


def save_users_db() -> None:
    AUTH_USERS_FILE.write_text(json.dumps(users_db, indent=2), encoding="utf-8")


def load_users_db() -> None:
    global users_db
    if not AUTH_USERS_FILE.exists():
        users_db = {}
        return

    try:
        users_db = json.loads(AUTH_USERS_FILE.read_text(encoding="utf-8"))
    except Exception as e:
        print(f"âš ï¸ Failed to load users database: {e}")
        users_db = {}


def create_jwt(payload: Dict[str, Any], expires_in: int = JWT_EXPIRATION_SECONDS) -> str:
    now = datetime.now(timezone.utc)
    body = payload.copy()
    body["iat"] = int(now.timestamp())
    body["exp"] = int((now + timedelta(seconds=expires_in)).timestamp())

    header = {"alg": JWT_ALGORITHM, "typ": "JWT"}
    header_b64 = _base64url_encode(json.dumps(header, separators=(",", ":")).encode("utf-8"))
    payload_b64 = _base64url_encode(json.dumps(body, separators=(",", ":")).encode("utf-8"))

    signing_input = f"{header_b64}.{payload_b64}".encode("utf-8")
    signature = hmac.new(JWT_SECRET.encode("utf-8"), signing_input, hashlib.sha256).digest()
    signature_b64 = _base64url_encode(signature)
    return f"{header_b64}.{payload_b64}.{signature_b64}"


def decode_and_validate_jwt(token: str) -> Dict[str, Any]:
    try:
        header_b64, payload_b64, signature_b64 = token.split(".")
    except ValueError as exc:
        raise HTTPException(status_code=401, detail="Invalid token format") from exc

    signing_input = f"{header_b64}.{payload_b64}".encode("utf-8")
    expected_sig = hmac.new(JWT_SECRET.encode("utf-8"), signing_input, hashlib.sha256).digest()
    actual_sig = _base64url_decode(signature_b64)

    if not hmac.compare_digest(expected_sig, actual_sig):
        raise HTTPException(status_code=401, detail="Invalid token signature")

    try:
        payload = json.loads(_base64url_decode(payload_b64).decode("utf-8"))
    except Exception as exc:
        raise HTTPException(status_code=401, detail="Invalid token payload") from exc

    exp = payload.get("exp")
    if not isinstance(exp, int):
        raise HTTPException(status_code=401, detail="Token missing expiration")
    if datetime.now(timezone.utc).timestamp() > exp:
        raise HTTPException(status_code=401, detail="Token expired")

    return payload


def get_current_user_from_headers(authorization: Optional[str], x_auth_token: Optional[str]) -> Dict[str, Any]:
    auth_value = authorization
    if not auth_value and x_auth_token:
        auth_value = f"Bearer {x_auth_token}"

    if not auth_value:
        raise HTTPException(status_code=401, detail="Missing Authorization header")

    # Strip whitespace and handle potential formatting issues
    auth_value = auth_value.strip()
    
    try:
        scheme, token = auth_value.split(" ", 1)
    except ValueError as exc:
        raise HTTPException(status_code=401, detail="Invalid Authorization header") from exc

    # Strip token whitespace as well
    token = token.strip()
    
    if scheme.lower() != "bearer":
        raise HTTPException(status_code=401, detail="Authorization scheme must be Bearer")

    # Debug logging for token issues
    if not token or len(token.split(".")) != 3:
        print(f"🔒 Token format issue - token length: {len(token) if token else 0}, parts: {len(token.split('.')) if token else 0}")
        print(f"   Token preview: {token[:50] if token else 'None'}...")
        raise HTTPException(status_code=401, detail="Invalid token format")

    payload = decode_and_validate_jwt(token)
    username = payload.get("sub")
    if not username or username not in users_db:
        raise HTTPException(status_code=401, detail="Invalid token subject")

    return {"username": username, **users_db[username]}


def get_current_user(
    authorization: Optional[str] = Header(default=None),
    x_auth_token: Optional[str] = Header(default=None, alias="X-Auth-Token"),
) -> Dict[str, Any]:
    return get_current_user_from_headers(authorization, x_auth_token)


def get_current_user_or_autogen_team(
    request: Request,
    authorization: Optional[str] = Header(default=None),
    x_auth_token: Optional[str] = Header(default=None, alias="X-Auth-Token"),
) -> Dict[str, Any]:
    if _autogen_team_secret_matches(request):
        return {"username": "autogen_team", "auth_type": "agent_secret"}
    return get_current_user_from_headers(authorization, x_auth_token)


def _spotify_redirect_uri() -> str:
    """Return the configured Spotify OAuth redirect URI."""

    redirect_uri = str(os.getenv("SPOTIFY_REDIRECT_URI") or "").strip()
    if not redirect_uri:
        raise HTTPException(
            status_code=500,
            detail=(
                "SPOTIFY_REDIRECT_URI is not configured. Register CATBot's callback URL in the Spotify "
                "Developer Dashboard and set SPOTIFY_REDIRECT_URI in .env."
            ),
        )

    parsed = urlparse(redirect_uri)
    if not parsed.scheme or not parsed.netloc:
        raise HTTPException(status_code=500, detail="SPOTIFY_REDIRECT_URI must be an absolute URL.")
    hostname = (parsed.hostname or "").strip().lower()
    if hostname == "localhost":
        raise HTTPException(
            status_code=500,
            detail=(
                "SPOTIFY_REDIRECT_URI cannot use localhost. Spotify's current redirect URI rules prohibit "
                "localhost aliases, and CATBot's HTTPS certificate will not validate for https://localhost:8002. "
                f"Use CATBot's trusted HTTPS hostname instead, for example "
                f"'https://{_SSL_CERT_HOSTNAME}:8002/spotify/callback', and register that exact URI in the Spotify app dashboard."
            ),
        )
    return redirect_uri


def _spotify_client_id() -> str:
    """Return the configured Spotify client ID."""

    client_id = str(os.getenv("SPOTIFY_CLIENT_ID") or "").strip()
    if not client_id:
        raise HTTPException(status_code=500, detail="SPOTIFY_CLIENT_ID is not configured.")
    return client_id


def _spotify_client_secret() -> str:
    """Return the configured Spotify client secret."""

    client_secret = str(os.getenv("SPOTIFY_CLIENT_SECRET") or "").strip()
    if not client_secret:
        raise HTTPException(status_code=500, detail="SPOTIFY_CLIENT_SECRET is not configured.")
    return client_secret


def _cleanup_spotify_oauth_states(now: Optional[float] = None) -> None:
    """Remove expired Spotify OAuth state values."""

    current_time = time.time() if now is None else now
    cutoff = current_time - SPOTIFY_OAUTH_STATE_TTL_SECONDS
    expired = [state for state, created_at in spotify_oauth_pending_states.items() if created_at < cutoff]
    for state in expired:
        spotify_oauth_pending_states.pop(state, None)


def _format_env_value(value: str) -> str:
    """Format a value for safe .env persistence."""

    if re.fullmatch(r"[A-Za-z0-9._:/@+-]+", value):
        return value
    escaped = value.replace("\\", "\\\\").replace('"', '\\"')
    return f'"{escaped}"'


def _set_key_in_env_content(content: str, key: str, value: str) -> str:
    """Replace or append a KEY=value pair in .env-style text."""

    line_match = re.compile(r"^(\s*)" + re.escape(key) + r"\s*=[^\n]*", re.MULTILINE)
    new_line = f"{key}={_format_env_value(value)}\n"
    if line_match.search(content):
        return line_match.sub(new_line, content, count=1)
    return content.rstrip() + ("\n" if content.strip() else "") + new_line


def _persist_env_key(key: str, value: str, *, env_file: Optional[Path] = None) -> None:
    """Persist a runtime-generated secret to .env and the current process environment."""

    target = Path(env_file or ENV_FILE)
    existing = ""
    if target.exists():
        existing = target.read_text(encoding="utf-8")
    updated = _set_key_in_env_content(existing, key, value)
    target.write_text(updated, encoding="utf-8")
    os.environ[key] = value


def _spotify_token_request_headers() -> Dict[str, str]:
    """Build Basic auth headers for Spotify Accounts token exchanges."""

    basic_token = base64.b64encode(
        f"{_spotify_client_id()}:{_spotify_client_secret()}".encode("utf-8")
    ).decode("utf-8")
    return {
        "Authorization": f"Basic {basic_token}",
        "Content-Type": "application/x-www-form-urlencoded",
    }


def _extract_spotify_error_message(payload: Any) -> str:
    """Extract a compact Spotify OAuth error message."""

    if isinstance(payload, dict):
        error_description = str(payload.get("error_description") or "").strip()
        if error_description:
            return error_description
        error = payload.get("error")
        if isinstance(error, dict):
            message = str(error.get("message") or "").strip()
            if message:
                return message
        if isinstance(error, str) and error.strip():
            return error.strip()
    return "Spotify OAuth request failed."


def security_log(action: str, user: str, server_id: Optional[str], detail: str) -> None:
    """Log MCP config/connect actions for audit; do not log secrets."""
    ts = datetime.now(timezone.utc).isoformat()
    sid = server_id or ""
    print(f"[SEC] {ts} action={action} user={user} server_id={sid} detail={detail}", flush=True)


# Helper utilities for Telegram integration
def build_openai_url(path: str) -> str:
    """Return absolute URL for OpenAI-compatible endpoints."""

    if path.startswith("http://") or path.startswith("https://"):
        return path

    base = TELEGRAM_OPENAI_BASE_URL.rstrip("/")
    if not path.startswith("/"):
        path = f"/{path}"
    return f"{base}{path}"


def trim_telegram_history(history: List[Dict[str, str]]) -> None:
    """Trim stored history to the configured limit (in-place)."""

    max_messages = max(TELEGRAM_HISTORY_LIMIT, 1) * 2
    if len(history) > max_messages:
        del history[: len(history) - max_messages]


def _telegram_secret_matches(request: Request) -> bool:
    """Return True when request carries TELEGRAM_SECRET via X-Telegram-Secret or Bearer token."""
    if not TELEGRAM_SECRET:
        return False
    secret_header = request.headers.get("X-Telegram-Secret")
    if secret_header is not None and secret_header.strip() == TELEGRAM_SECRET:
        return True
    auth_header = request.headers.get("Authorization")
    if auth_header and auth_header.strip().startswith("Bearer "):
        token = auth_header.strip()[7:].strip()
        if token == TELEGRAM_SECRET:
            return True
    return False


def _validate_telegram_secret(request: Request) -> None:
    """If TELEGRAM_SECRET is set, require X-Telegram-Secret or Authorization Bearer to match; else raise 401."""
    if not TELEGRAM_SECRET:
        return
    if not _telegram_secret_matches(request):
        raise HTTPException(status_code=401, detail="Telegram secret required or invalid")


def _autogen_team_secret_matches(request: Request) -> bool:
    """Return True when request carries AUTOGEN_TEAM_SECRET via X-Agent-Secret or Bearer token."""
    if not AUTOGEN_TEAM_SECRET:
        return False
    secret_header = request.headers.get("X-Agent-Secret")
    if secret_header is not None and secret_header.strip() == AUTOGEN_TEAM_SECRET:
        return True
    auth_header = request.headers.get("Authorization")
    if auth_header and auth_header.strip().startswith("Bearer "):
        token = auth_header.strip()[7:].strip()
        if token == AUTOGEN_TEAM_SECRET:
            return True
    return False


# Create scratch and companions directories if they don't exist
SCRATCH_DIR.mkdir(parents=True, exist_ok=True)
COMPANIONS_DIR.mkdir(parents=True, exist_ok=True)
load_users_db()

# Global AutoGen team instance
autogen_team = None

# Global memory manager instance
memory_manager = None

# Initialize memory manager if available
if MEMORY_AVAILABLE:
    try:
        memory_enabled = os.getenv("MEMORY_ENABLED", "true").lower() == "true"
        if memory_enabled:
            memory_manager = MemoryManager()
            print(f"✅ Memory system initialized with {memory_manager.count()} existing memories")
        else:
            print("âš ï¸  Memory system disabled via MEMORY_ENABLED=false")
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"âš ï¸  Failed to initialize memory system: {e}")
        print(f"   Full traceback:\n{error_trace}")
        memory_manager = None

# Initialize skills manager if framework is available
skill_manager = None
if SKILLS_FRAMEWORK_AVAILABLE and create_default_skill_manager is not None:
    try:
        _skills_manifest_dir = _env_str("SKILLS_MANIFEST_DIR")
        skill_manager = create_default_skill_manager(manifest_dir=_skills_manifest_dir)
        loaded_skill_names = [spec.name for spec in skill_manager.list_skills()]
        print(f"[OK] Skills framework initialized with {len(loaded_skill_names)} skills: {loaded_skill_names}")
    except Exception as e:
        skill_manager = None
        print(f"[WARN] Failed to initialize skills framework: {e}")


def _resolve_skill_tool_qualified_name(tool_name: str) -> Optional[str]:
    """Resolve a skill tool name to its qualified form if available."""
    if not tool_name or skill_manager is None:
        return None
    alias_map = {
        "google_slides.slides_create_presentation_from_markdown": (
            "googleworkspace_cli.slides_create_presentation_from_markdown"
        ),
        "google_slides.slides_batch_update_presentation": (
            "googleworkspace_cli.slides_batch_update_presentation"
        ),
        "slides_create_presentation_from_markdown": (
            "googleworkspace_cli.slides_create_presentation_from_markdown"
        ),
        "markdown_to_slides": "googleworkspace_cli.slides_create_presentation_from_markdown",
        "markdownToSlides": "googleworkspace_cli.slides_create_presentation_from_markdown",
    }
    candidate_name = alias_map.get(str(tool_name).strip(), str(tool_name).strip())
    try:
        resolved = skill_manager.registry.resolve_tool(candidate_name)
        return resolved.qualified_name
    except Exception:
        return None


def _filter_overlapping_file_skill_tools(
    tools: List[Dict[str, Any]],
    *,
    openai_schema: bool,
) -> List[Dict[str, Any]]:
    """Return skill tools unchanged; filesystem tools are now skill-backed."""
    return [tool for tool in tools if isinstance(tool, dict)]


def _get_skill_tools_openai_schema() -> List[Dict[str, Any]]:
    """Return skill tools formatted for OpenAI-style tool-calling."""
    if skill_manager is None:
        return []
    try:
        tools = skill_manager.openai_tools(qualified_names=True)
        return _filter_overlapping_file_skill_tools(tools, openai_schema=True)
    except Exception as exc:
        print(f"[WARN] Failed to list skill tools (openai schema): {exc}")
        return []


def _get_skill_tools_mcp_schema() -> List[Dict[str, Any]]:
    """Return skill tools formatted as MCP-style entries for server-managed agents."""
    if skill_manager is None:
        return []
    try:
        tools = skill_manager.mcp_tools(qualified_names=True)
        return _filter_overlapping_file_skill_tools(tools, openai_schema=False)
    except Exception as exc:
        print(f"[WARN] Failed to list skill tools (mcp schema): {exc}")
        return []


def _get_telegram_excluded_skill_tool_names() -> Set[str]:
    """Return skill tools that should not be exposed directly to Telegram."""
    return {"googleworkspace_cli.slides_create_presentation_from_markdown"}


def _get_telegram_skill_tools_openai_schema() -> List[Dict[str, Any]]:
    """Return Telegram-safe skill tools for OpenAI-style tool calling."""
    excluded = _get_telegram_excluded_skill_tool_names()
    filtered: List[Dict[str, Any]] = []
    for item in _get_skill_tools_openai_schema():
        function = item.get("function") if isinstance(item, dict) else None
        name = str(function.get("name") or "").strip() if isinstance(function, dict) else ""
        if name and name in excluded:
            continue
        filtered.append(item)
    return filtered


def _get_telegram_skill_tools_mcp_schema() -> List[Dict[str, Any]]:
    """Return Telegram-safe skill tools for MCP-style prompt rendering."""
    excluded = _get_telegram_excluded_skill_tool_names()
    return [
        item
        for item in _get_skill_tools_mcp_schema()
        if str(item.get("name") or "").strip() not in excluded
    ]


async def _execute_skill_framework_tool(
    tool_name: str,
    arguments: Optional[Dict[str, Any]] = None,
    *,
    conversation_id: Optional[str] = None,
    user_id: Optional[str] = None,
    metadata: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """Execute a skill framework tool and return a normalized payload."""
    if skill_manager is None or SkillContext is None:
        return {
            "success": False,
            "message": "Skill framework is not available.",
            "error_code": "framework_unavailable",
            "tool_name": tool_name,
        }

    qualified_name = _resolve_skill_tool_qualified_name(tool_name)
    if not qualified_name:
        return {
            "success": False,
            "message": f"Skill tool '{tool_name}' is not registered.",
            "error_code": "tool_not_found",
            "tool_name": tool_name,
        }

    context_metadata = dict(metadata or {})
    context = SkillContext(
        conversation_id=conversation_id,
        user_id=user_id,
        scratch_dir=SCRATCH_DIR,
        metadata=context_metadata,
    )
    context.set_service("telegram_send_message", _send_telegram_bot_message)
    context.set_service(
        "telegram_send_file",
        lambda chat_id, filename: _send_telegram_file_internal(chat_id, filename),
    )
    context.set_service(
        "telegram_admin_chat_ids",
        lambda: [item.strip() for item in str(os.getenv("TELEGRAM_ADMIN_IDS") or "").split(",") if item.strip()],
    )
    result = await skill_manager.execute_tool(
        tool_name=qualified_name,
        arguments=arguments or {},
        context=context,
        raise_errors=False,
    )
    payload = result.to_dict() if hasattr(result, "to_dict") else {
        "success": False,
        "message": "Skill execution returned an invalid response.",
        "error_code": "invalid_result",
        "tool_name": qualified_name,
    }
    if not payload.get("tool_name"):
        payload["tool_name"] = qualified_name
    return payload


def _format_filesystem_skill_tool_output(
    qualified_name: str,
    result: Dict[str, Any],
) -> Optional[str]:
    """Render filesystem skill results with the same concise text style as legacy file tools."""
    if not result.get("success", False):
        return None
    data = result.get("data")
    if not isinstance(data, dict):
        return None

    if qualified_name == "filesystem.read_text":
        path = str(data.get("path") or "").strip() or "(unknown)"
        content = data.get("content", "")
        content = content if isinstance(content, str) else str(content)
        header_parts = [f"File: {path}"]
        start_meta = data.get("excerpt_start_line")
        end_meta = data.get("excerpt_end_line")
        total_lines = data.get("total_lines")
        if isinstance(start_meta, int) and start_meta > 0 and isinstance(end_meta, int) and end_meta > 0:
            if start_meta != 1 or (isinstance(total_lines, int) and end_meta != total_lines):
                header_parts.append(f"lines {start_meta}-{end_meta}")
        if isinstance(total_lines, int) and total_lines > 0:
            header_parts.append(f"total_lines={total_lines}")
        if data.get("truncated"):
            header_parts.append("truncated")
        return " | ".join(header_parts) + "\n\n" + (content or "(empty file)")

    if qualified_name == "filesystem.write_text":
        path = str(data.get("path") or "").strip() or "(unknown)"
        bytes_written = data.get("bytes_written")
        action = "Appended to" if bool(data.get("appended")) else "Wrote"
        if isinstance(bytes_written, int):
            return f"{action} {path} ({bytes_written} bytes)."
        return f"{action} {path}."

    if qualified_name == "filesystem.list_files":
        items = data.get("items", [])
        converted_items = []
        for item in items if isinstance(items, list) else []:
            if not isinstance(item, dict):
                continue
            converted_items.append(
                {
                    "name": str(item.get("relative_path") or item.get("name") or "").strip(),
                    "type": str(item.get("type") or ""),
                    "size": item.get("size_bytes"),
                }
            )
        skipped_count = int(data.get("skipped_count", 0) or 0)
        if not converted_items:
            path_label = data.get("path") or "."
            total_count = int(data.get("total_count", 0) or 0)
            recursive = bool(data.get("recursive", False))
            message = (
                "Scratch workspace is empty for this scope. "
                f"(Directory: {data.get('root', 'scratch')}, Path: {path_label}, Recursive: {recursive})"
            )
            if total_count > 0 and int(data.get("offset", 0) or 0) > 0:
                message = (
                    "No files returned for this page. "
                    f"(Directory: {data.get('root', 'scratch')}, Path: {path_label}, "
                    f"Recursive: {recursive}, Offset: {data.get('offset', 0)}, Total: {total_count})"
                )
            if skipped_count > 0:
                message += f" Skipped {skipped_count} inaccessible or unsafe entries."
            return message
        rendered = _format_list_files_for_tool_output(
            converted_items,
            include_sizes=True,
            total_count=data.get("total_count"),
            offset=data.get("offset"),
            has_more=data.get("has_more"),
            next_offset=data.get("next_offset"),
            limit=data.get("max_entries"),
        )
        if skipped_count > 0:
            rendered += f"\n... skipped {skipped_count} inaccessible or unsafe entries."
        return rendered

    if qualified_name == "filesystem.search_files":
        items = data.get("items", [])
        matches = [item for item in items if isinstance(item, dict)] if isinstance(items, list) else []
        rendered = _format_search_files_for_tool_output(
            matches,
            query=str(data.get("query") or ""),
            total_matches=data.get("total_matches"),
            offset=data.get("offset"),
            has_more=data.get("has_more"),
            next_offset=data.get("next_offset"),
            read_tool_name="filesystem.read_text",
        )
        skipped_count = int(data.get("skipped_count", 0) or 0)
        if skipped_count > 0:
            rendered += f"\nSkipped {skipped_count} inaccessible, oversized, or unsupported files."
        return rendered

    return None


def _format_generic_skill_tool_output(
    qualified_name: str,
    result: Dict[str, Any],
) -> str:
    """Prefer descriptive skill messages over raw payload dumps for user-facing tool output."""
    message = str(result.get("message") or "").strip()
    if message and message.lower() not in {"ok", "success", "done"}:
        return message

    data = result.get("data")
    if data is None:
        return message or f"Skill tool '{qualified_name}' executed successfully."
    if isinstance(data, (dict, list)):
        return json.dumps(data, ensure_ascii=False, default=str)
    return str(data)


def _build_telegram_skill_tools_prompt_block() -> str:
    """Render dynamic skill-tool instructions for Telegram XML tool-calling."""
    skill_tools = _get_telegram_skill_tools_mcp_schema()
    if not skill_tools:
        return ""
    excluded_tool_names = {"googleworkspace_cli.run_readonly_command"}
    tool_names = {str(item.get("name") or "").strip() for item in skill_tools if isinstance(item, dict)}

    lines: List[str] = [
        "Additional Skill Framework tools (dynamically loaded):",
        "Prefer structured tool calls with the exact tool name and JSON schema shown below.",
        "Use XML tool markup only as a legacy fallback when structured tool calls are unavailable.",
    ]
    filesystem_tool_names = [
        "filesystem.read_text",
        "filesystem.write_text",
        "filesystem.list_files",
        "filesystem.search_files",
    ]
    if any(name in tool_names for name in filesystem_tool_names):
        lines.extend(
            [
                "Filesystem skill tools are available in this conversation.",
                "Use these exact qualified names and prefer them over legacy readFile/writeFile/listFiles/searchFiles:",
                *[f"- {name}" for name in filesystem_tool_names if name in tool_names],
            ]
        )
    for item in skill_tools:
        name = str(item.get("name") or "").strip()
        if not name:
            continue
        if name in excluded_tool_names:
            continue
        description = str(item.get("description") or "No description.").strip()
        input_schema = item.get("inputSchema")
        if not isinstance(input_schema, dict):
            input_schema = {"type": "object", "properties": {}}
        schema_text = json.dumps(input_schema, ensure_ascii=False, separators=(",", ":"), default=str)
        lines.append(f"- {name}: {description}")
        lines.append(f"  Parameters JSON schema: {schema_text}")

    if "googleworkspace_cli.gmail_list_unread" in tool_names:
        lines.extend(
            [
                "",
                "Gmail tool examples (prefer these over generic run_readonly_command):",
                "<tool>googleworkspace_cli.gmail_list_unread</tool>",
                "<parameters>{\"max_results\": 5}</parameters>",
                "<tool>googleworkspace_cli.gmail_list_all</tool>",
                "<parameters>{\"max_results\": 10, \"page_token\": \"NEXT_PAGE_TOKEN\"}</parameters>",
                "<tool>googleworkspace_cli.gmail_get_message</tool>",
                "<parameters>{\"message_id\": \"18c...\", \"format\": \"full\"}</parameters>",
                "<tool>googleworkspace_cli.gmail_compose_draft</tool>",
                "<parameters>{\"to\": \"user@example.com\", \"subject\": \"Status\", \"body_text\": \"Hi\"}</parameters>",
                "<tool>googleworkspace_cli.gmail_send_message</tool>",
                "<parameters>{\"to\": \"user@example.com\", \"subject\": \"Status\", \"body_text\": \"Hi\"}</parameters>",
                "<tool>googleworkspace_cli.gmail_mark_read</tool>",
                "<parameters>{\"message_id\": \"18c...\"}</parameters>",
            ]
        )
    if "googleworkspace_cli.slides_create_presentation_from_markdown" in tool_names:
        lines.extend(
            [
                "",
                "Google Slides markdown-file example:",
                "<tool>googleworkspace_cli.slides_create_presentation_from_markdown</tool>",
                "<parameters>{\"markdown_path\":\"<relative_path_to_scratch>/file.md\"}</parameters>",
                "Use markdown_path for markdown files that already exist under scratch. Do not call any standalone slides skill.",
                "For slide-creation tasks, reading a markdown file with filesystem.read_text is only preparation, not completion.",
                "A slide task is only complete after googleworkspace_cli.slides_create_presentation_from_markdown succeeds and returns a presentation URL or saved result file path.",
            ]
        )

    if len(lines) <= 3:
        return ""
    return "\n".join(lines)


def _get_telegram_native_tools_mcp_schema() -> List[Dict[str, Any]]:
    """Return native Telegram tools using the same MCP-like schema shape as skills."""
    return [
        {
            "name": "manageTodoList",
            "description": "Persistent todo list operations.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "action": {"type": "string"},
                    "taskId": {"type": "integer"},
                    "taskDescription": {"type": "string"},
                    "scheduledFor": {"type": "string"},
                    "recurrence": {"type": "object"},
                    "repeatFrequency": {"type": "string"},
                    "repeatInterval": {"type": "integer"},
                    "clearSchedule": {"type": "boolean"},
                    "clearRecurrence": {"type": "boolean"},
                },
                "required": ["action"],
            },
        },
        {
            "name": "executeTodoTask",
            "description": "Execute an existing todo task.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "taskId": {"type": "integer"},
                    "promptOverride": {"type": "string"},
                },
                "required": ["taskId"],
            },
        },
        {
            "name": "getTodoExecutionStatus",
            "description": "Check active todo execution status.",
            "inputSchema": {"type": "object", "properties": {}},
        },
        {
            "name": "cancelTodoExecution",
            "description": "Cancel the current todo execution.",
            "inputSchema": {"type": "object", "properties": {}},
        },
        {
            "name": "navigateToUrl",
            "description": "Return a URL the Telegram user should open in a browser.",
            "inputSchema": {
                "type": "object",
                "properties": {"url": {"type": "string"}},
                "required": ["url"],
            },
        },
        {
            "name": "openChatToUser",
            "description": "Return a Teams chat URL for the Telegram user to open.",
            "inputSchema": {
                "type": "object",
                "properties": {"url": {"type": "string"}},
                "required": ["url"],
            },
        },
        {
            "name": "calculate",
            "description": "Evaluate a simple arithmetic expression.",
            "inputSchema": {
                "type": "object",
                "properties": {"expression": {"type": "string"}},
                "required": ["expression"],
            },
        },
        {
            "name": "runWorkflow",
            "description": "Run an AutoGen workflow.",
            "inputSchema": {
                "type": "object",
                "properties": {"contentPrompt": {"type": "string"}},
                "required": ["contentPrompt"],
            },
        },
        {
            "name": "runCodexCli",
            "description": "Run Codex CLI to make CATBot code or tool changes.",
            "inputSchema": {
                "type": "object",
                "properties": {"prompt": {"type": "string"}},
                "required": ["prompt"],
            },
        },
        {
            "name": "restartProxyServer",
            "description": "Restart the proxy server after explicit confirmation.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "confirm": {"type": "boolean"},
                    "reason": {"type": "string"},
                },
                "required": ["confirm"],
            },
        },
        {
            "name": "scrapeWebsite",
            "description": "Fetch one URL or try multiple URLs in order.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "url": {"type": "string"},
                    "urls": {"type": "array", "items": {"type": "string"}},
                    "render_js": {"type": "boolean"},
                    "render_engine": {"type": "string"},
                    "wait_for_selector": {"type": "string"},
                    "js_wait_ms": {"type": "number"},
                },
            },
        },
        {
            "name": "webSearch",
            "description": "Web search query.",
            "inputSchema": {
                "type": "object",
                "properties": {"query": {"type": "string"}},
                "required": ["query"],
            },
        },
        {
            "name": "fetchNews",
            "description": "Fetch news and write a CSV file.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "searchTerm": {"type": "string"},
                    "filename": {"type": "string"},
                },
                "required": ["searchTerm"],
            },
        },
        {
            "name": "pdfToPowerPoint",
            "description": "Convert a PDF or Markdown document into a PowerPoint (.pptx). Supports uploaded attachments, scratch-relative paths, URLs, inline Markdown, and base64 content.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "source": {
                        "type": "object",
                        "properties": {
                            "type": {"type": "string"},
                            "value": {"type": "string"},
                            "url": {"type": "string"},
                            "path": {"type": "string"},
                            "relativePath": {"type": "string"},
                            "content": {"type": "string"},
                            "contentBase64": {"type": "string"},
                            "mimeType": {"type": "string"},
                            "filename": {"type": "string"},
                        },
                    },
                    "sourceUrl": {"type": "string"},
                    "sourceType": {"type": "string"},
                    "pdfUrl": {"type": "string"},
                    "title": {"type": "string"},
                    "filename": {"type": "string"},
                },
            },
        },
        {
            "name": "createSlidesPresentation",
            "description": (
                "Create a Google Slides presentation for Telegram from a prompt, inline markdown, "
                "or a markdown file under scratch. Prefer this over direct Google Slides skill calls."
            ),
            "inputSchema": {
                "type": "object",
                "properties": {
                    "prompt": {"type": "string"},
                    "title": {"type": "string"},
                    "markdown": {"type": "string"},
                    "markdown_path": {"type": "string"},
                    "audience": {"type": "string"},
                    "purpose": {"type": "string"},
                    "tone": {"type": "string"},
                    "max_slides": {"type": "integer"},
                    "attach_scratch_images": {"type": "boolean"},
                    "image_dir": {"type": "string"},
                    "max_images_per_slide": {"type": "integer"},
                    "include_image_requests": {"type": "boolean"},
                    "timeout_seconds": {"type": "number"},
                },
            },
        },
        {
            "name": "uploadToGoogleDrive",
            "description": "Upload a scratch file to Google Drive.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "filePath": {"type": "string"},
                    "filename": {"type": "string"},
                    "fileName": {"type": "string"},
                },
                "required": ["filePath"],
            },
        },
        {
            "name": "readFile",
            "description": "Legacy file read tool. Prefer filesystem.read_text when available.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "path": {"type": "string"},
                    "file": {"type": "string"},
                    "start_line": {"type": "integer"},
                    "end_line": {"type": "integer"},
                    "max_chars": {"type": "integer"},
                    "include_line_numbers": {"type": "boolean"},
                },
                "required": ["filename"],
            },
        },
        {
            "name": "listFiles",
            "description": "Legacy file listing tool. Prefer filesystem.list_files when available.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "recursive": {"type": "boolean"},
                    "offset": {"type": "integer"},
                    "max_entries": {"type": "integer"},
                },
            },
        },
        {
            "name": "searchFiles",
            "description": "Legacy file search tool. Prefer filesystem.search_files when available.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "query": {"type": "string"},
                    "path": {"type": "string"},
                    "recursive": {"type": "boolean"},
                    "filename_only": {"type": "boolean"},
                    "case_sensitive": {"type": "boolean"},
                    "offset": {"type": "integer"},
                    "max_results": {"type": "integer"},
                },
                "required": ["query"],
            },
        },
        {
            "name": "sendTelegramFile",
            "description": "Send a scratch file back to the current Telegram chat.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "filePath": {"type": "string"},
                    "caption": {"type": "string"},
                },
                "required": ["filename"],
            },
        },
        {
            "name": "writeFile",
            "description": "Legacy file write tool. Prefer filesystem.write_text when available.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "path": {"type": "string"},
                    "file": {"type": "string"},
                    "content": {"type": "string"},
                    "text": {"type": "string"},
                    "body": {"type": "string"},
                    "format": {"type": "string"},
                    "append": {"type": "boolean"},
                },
            },
        },
        {
            "name": "storeMemory",
            "description": "Store a persistent memory.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "text": {"type": "string"},
                    "content": {"type": "string"},
                    "category": {"type": "string"},
                },
                "required": ["text"],
            },
        },
        {
            "name": "searchMemories",
            "description": "Search persistent memories.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "query": {"type": "string"},
                    "limit": {"type": "integer"},
                },
                "required": ["query"],
            },
        },
        {
            "name": "listMemories",
            "description": "List recent persistent memories.",
            "inputSchema": {
                "type": "object",
                "properties": {"limit": {"type": "integer"}},
            },
        },
        {
            "name": "deleteMemory",
            "description": "Delete a persistent memory by ID.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "memory_id": {"type": "string"},
                    "id": {"type": "string"},
                },
                "required": ["memory_id"],
            },
        },
        {
            "name": "manageMemoryCache",
            "description": "Inspect or edit the lightweight in-session Telegram memory cache.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "action": {"type": "string"},
                    "memoryId": {"type": "integer"},
                    "memId": {"type": "integer"},
                    "memoryDescription": {"type": "string"},
                    "memDescription": {"type": "string"},
                },
                "required": ["action"],
            },
        },
        {
            "name": "runBrowserAgent",
            "description": "Browser automation task.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "task": {"type": "string"},
                    "instruction": {"type": "string"},
                    "url": {"type": "string"},
                },
            },
        },
        {
            "name": "runDeepResearch",
            "description": "Deep browser-based research task.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "research_task": {"type": "string"},
                    "researchTask": {"type": "string"},
                    "max_parallel_browsers": {"type": "integer"},
                },
            },
        },
        {
            "name": "healthCheck",
            "description": "Browser-use health and running jobs.",
            "inputSchema": {"type": "object", "properties": {}},
        },
        {
            "name": "llmQuery",
            "description": "Send a direct query to the language model without other tools.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "contentPrompt": {"type": "string"},
                    "query": {"type": "string"},
                    "message": {"type": "string"},
                },
            },
        },
        {
            "name": "weatherInfo",
            "description": "Weather lookup.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "location": {"type": "string"},
                    "requestType": {"type": "string"},
                    "detail": {"type": "string"},
                },
            },
        },
    ]


def _mcp_tool_entries_to_openai_tools(tools: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Convert MCP-style tool entries into OpenAI-style tool definitions."""
    converted: List[Dict[str, Any]] = []
    for item in tools:
        if not isinstance(item, dict):
            continue
        name = str(item.get("name") or "").strip()
        if not name:
            continue
        input_schema = item.get("inputSchema")
        if not isinstance(input_schema, dict):
            input_schema = {"type": "object", "properties": {}}
        converted.append(
            {
                "type": "function",
                "function": {
                    "name": name,
                    "description": str(item.get("description") or "").strip(),
                    "parameters": input_schema,
                },
            }
        )
    return converted


def _merge_openai_tool_lists(*tool_lists: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Merge OpenAI-style tool lists while keeping the first definition for each name."""
    merged: List[Dict[str, Any]] = []
    seen_names: Set[str] = set()
    for tool_list in tool_lists:
        for item in tool_list:
            if not isinstance(item, dict):
                continue
            function = item.get("function")
            if not isinstance(function, dict):
                continue
            name = str(function.get("name") or "").strip()
            if not name or name in seen_names:
                continue
            seen_names.add(name)
            merged.append(item)
    return merged


def _get_telegram_native_tools_openai_schema() -> List[Dict[str, Any]]:
    """Return native Telegram tools formatted for OpenAI-compatible tool calling."""
    return _mcp_tool_entries_to_openai_tools(_get_telegram_native_tools_mcp_schema())


def _get_telegram_combined_openai_tools() -> List[Dict[str, Any]]:
    """Return the merged native+skill tool list for Telegram chat payloads."""
    return _merge_openai_tool_lists(
        _get_telegram_native_tools_openai_schema(),
        _get_telegram_skill_tools_openai_schema(),
    )


def _build_telegram_native_tools_prompt_block() -> str:
    """Render schema-driven instructions for native Telegram tools."""
    native_tools = _get_telegram_native_tools_mcp_schema()

    lines: List[str] = [
        "Native Telegram tools:",
        "Prefer structured tool calls using the provided tool schema and exact tool name.",
        "Use XML tool markup only as a legacy fallback when structured tool calls are unavailable.",
        "If you must use XML, inside <parameters> output exactly one JSON object and nothing else.",
        "Do not use nested XML child tags, bare keys, comments, trailing commas, or prose inside <parameters>.",
    ]

    skill_tool_names = {
        str(item.get("name") or "").strip()
        for item in (_get_telegram_skill_tools_mcp_schema() or [])
        if isinstance(item, dict)
    }
    if {
        "filesystem.read_text",
        "filesystem.write_text",
        "filesystem.list_files",
        "filesystem.search_files",
    } & skill_tool_names:
        lines.append(
            "When filesystem skill tools are available, prefer filesystem.read_text, "
            "filesystem.write_text, filesystem.list_files, and filesystem.search_files over "
            "readFile, writeFile, listFiles, and searchFiles."
        )
    lines.append(
        "For presentation requests in Telegram, use createSlidesPresentation. "
        "Do not call googleworkspace_cli.slides_create_presentation_from_markdown directly."
    )
    lines.append(
        "Use pdfToPowerPoint when the user explicitly wants a .pptx file from a PDF or Markdown source. "
        "Use createSlidesPresentation when the user wants a Google Slides deck."
    )

    for item in native_tools:
        schema_text = json.dumps(item["inputSchema"], ensure_ascii=False, separators=(",", ":"), default=str)
        lines.append(f"- {item['name']}: {item['description']}")
        lines.append(f"  Parameters JSON schema: {schema_text}")
        if item["name"] == "createSlidesPresentation":
            lines.append('  Example: {"prompt":"Create investor slides for PermitFlow AI","title":"PermitFlow AI"}')
            lines.append('  Use markdown_path when a slide markdown file already exists under scratch.')

    return "\n".join(lines)


def _slugify_telegram_slides_component(value: str, *, fallback: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", str(value or "").strip().lower()).strip("-")
    return slug or fallback


def _coerce_telegram_slides_max_slides(value: Any, *, default: int = 10) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = default
    return max(3, min(parsed, 20))


def _strip_markdown_fence_block(value: str) -> str:
    text = str(value or "").strip()
    fenced_match = re.fullmatch(r"```(?:markdown)?\s*([\s\S]*?)\s*```", text, flags=re.IGNORECASE)
    if fenced_match:
        return fenced_match.group(1).strip()
    return text


def _validate_telegram_slides_markdown(
    markdown: str,
    *,
    title: str,
    max_slides: int,
) -> Tuple[bool, Dict[str, int]]:
    from src.skills.builtin.googleworkspace_slides_support import _parse_markdown_to_slides

    slides = _parse_markdown_to_slides(
        markdown,
        fallback_title=title or "Presentation",
        scratch_root=SCRATCH_DIR,
        source_dir=SCRATCH_DIR / "presentations",
    )[:max_slides]
    slide_count = len(slides)
    bullet_count = sum(len(slide.get("bullets") or []) for slide in slides)
    text_length = len(str(markdown or "").strip())
    valid = bool(text_length >= 120 and slide_count >= 2 and bullet_count >= max(3, min(slide_count, 6)))
    return valid, {
        "slide_count": slide_count,
        "bullet_count": bullet_count,
        "text_length": text_length,
    }


def _write_telegram_slides_markdown(markdown: str, *, title: str) -> str:
    slug = _slugify_telegram_slides_component(title, fallback="presentation")
    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d-%H%M%S")
    relative_path = f"presentations/{slug}-{timestamp}.md"
    output_path = SCRATCH_DIR / relative_path
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(markdown, encoding="utf-8")
    return relative_path.replace("\\", "/")


async def _generate_telegram_slides_markdown(
    *,
    prompt: str,
    title: str,
    max_slides: int,
    model_name: Optional[str] = None,
    audience: str = "",
    purpose: str = "",
    tone: str = "",
    retry_reason: str = "",
) -> str:
    resolved_model = (model_name or TELEGRAM_DEFAULT_MODEL or "").strip() or "gpt-4o-mini"
    endpoint = _normalize_chat_endpoint(build_openai_url(TELEGRAM_OPENAI_CHAT_PATH))
    api_key = _first_non_empty_env(preferred_api_key_env_names(TELEGRAM_OPENAI_BASE_URL, resolved_model))
    if not api_key:
        raise RuntimeError("No compatible API key is configured for Telegram slide generation.")

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
    }
    if OPENAI_ORG_ID:
        headers["OpenAI-Organization"] = OPENAI_ORG_ID
    if OPENAI_PROJECT_ID:
        headers["OpenAI-Project"] = OPENAI_PROJECT_ID

    guidance_parts = [
        f"Topic or brief: {prompt}",
        f"Presentation title: {title or 'Presentation'}",
        f"Target slide count: {max_slides}",
    ]
    if audience.strip():
        guidance_parts.append(f"Audience: {audience.strip()}")
    if purpose.strip():
        guidance_parts.append(f"Purpose: {purpose.strip()}")
    if tone.strip():
        guidance_parts.append(f"Tone: {tone.strip()}")
    if retry_reason.strip():
        guidance_parts.append(f"Revision requirement: {retry_reason.strip()}")

    payload = {
        "model": resolved_model,
        "messages": [
            {
                "role": "system",
                "content": (
                    "You are writing slide-ready markdown for Google Slides creation.\n"
                    "Return markdown only. Do not include code fences or commentary.\n"
                    "Format rules:\n"
                    "- First line must be a single H1 title.\n"
                    "- Then create 5 to 12 H2 sections depending on requested slide count.\n"
                    "- Each H2 section must contain 2 to 5 bullet lines.\n"
                    "- Keep each bullet concise and presentation-ready.\n"
                    "- Do not output blank sections.\n"
                    "- Do not mention these instructions."
                ),
            },
            {"role": "user", "content": "\n".join(guidance_parts)},
        ],
        "temperature": 0.4,
    }
    response = await _call_chat_completion(
        endpoint,
        headers,
        payload,
        timeout_seconds=max(20.0, min(120.0, TELEGRAM_CHAT_TIMEOUT)),
    )
    if response.status_code != 200:
        detail = response.text
        try:
            error_json = response.json()
            detail = error_json.get("error", {}).get("message") or error_json.get("message") or detail
        except ValueError:
            pass
        raise RuntimeError(f"Slide markdown generation failed: {detail}")

    data = response.json()
    choices = data.get("choices") or []
    if not choices:
        raise RuntimeError("Slide markdown generation returned no choices.")
    normalized_message = normalize_chat_completion_message(choices[0].get("message") or {})
    content = _strip_markdown_fence_block(coerce_message_text(normalized_message.get("content") or ""))
    if not content:
        raise RuntimeError("Slide markdown generation returned empty content.")
    return content


def _coerce_pdf_to_powerpoint_max_slides(value: Any, *, default: int = 10) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = default
    return max(3, min(parsed, 20))


def _first_non_empty_presentation_string(*values: Any) -> str:
    for value in values:
        if isinstance(value, str) and value.strip():
            return value.strip()
    return ""


def _looks_like_http_url(value: str) -> bool:
    return bool(re.match(r"^https?://", str(value or "").strip(), flags=re.IGNORECASE))


def _decode_data_url_bytes(value: str) -> Tuple[str, bytes]:
    raw = str(value or "").strip()
    match = re.match(r"^data:(?P<mime>[^;,]+)?(?P<params>(?:;[^,]+)*?),(?P<data>.*)$", raw, flags=re.IGNORECASE | re.DOTALL)
    if not match:
        raise RuntimeError("Invalid data URL source.")
    mime_type = str(match.group("mime") or "application/octet-stream").strip() or "application/octet-stream"
    params = str(match.group("params") or "").lower()
    payload = str(match.group("data") or "")
    if ";base64" in params:
        try:
            return mime_type, base64.b64decode(payload, validate=True)
        except (ValueError, binascii.Error) as exc:
            raise RuntimeError("Invalid base64 data URL source.") from exc
    return mime_type, unquote_to_bytes(payload)


def _infer_presentation_source_type(source_input: Any, explicit_type: str = "") -> str:
    normalized_type = str(explicit_type or "").strip().lower()
    if normalized_type == "md":
        return "markdown"
    if normalized_type in {"pdf", "markdown"}:
        return normalized_type
    if normalized_type:
        raise RuntimeError(f"Unsupported source type: {explicit_type}")

    descriptor_type = ""
    file_name = ""
    mime_type = ""
    source_value = ""
    inline_text = ""
    if isinstance(source_input, dict):
        raw_descriptor_type = source_input.get("type")
        if isinstance(raw_descriptor_type, str) and "/" not in raw_descriptor_type:
            descriptor_type = raw_descriptor_type.strip().lower()
        file_name = _first_non_empty_presentation_string(
            source_input.get("filename"),
            source_input.get("fileName"),
            source_input.get("name"),
            source_input.get("original_filename"),
            source_input.get("originalFilename"),
        ).lower()
        mime_type = _first_non_empty_presentation_string(
            source_input.get("mimeType"),
            source_input.get("mime_type"),
            source_input.get("contentType"),
            source_input.get("content_type"),
            raw_descriptor_type if isinstance(raw_descriptor_type, str) and "/" in raw_descriptor_type else "",
        ).lower()
        source_value = _first_non_empty_presentation_string(
            source_input.get("sourceUrl"),
            source_input.get("url"),
            source_input.get("href"),
            source_input.get("uri"),
            source_input.get("src"),
            source_input.get("pdfUrl"),
            source_input.get("path"),
            source_input.get("filePath"),
            source_input.get("relative_path"),
            source_input.get("relativePath"),
            source_input.get("value") if descriptor_type in {"url", "path", "attachment", "file"} else "",
        ).lower()
        inline_text = _first_non_empty_presentation_string(
            source_input.get("markdown"),
            source_input.get("text"),
            source_input.get("value") if descriptor_type in {"inline", "markdown", "text"} else "",
            source_input.get("content"),
        )
    elif isinstance(source_input, str):
        source_value = source_input.strip().lower()

    if mime_type == "application/pdf" or file_name.endswith(".pdf") or re.match(r"^data:application/pdf[;,]", source_value, flags=re.IGNORECASE):
        return "pdf"
    if (
        ((descriptor_type in {"inline", "markdown", "text"}) and inline_text)
        or mime_type in {"text/markdown", "text/x-markdown"}
        or (mime_type == "text/plain" and (file_name.endswith(".md") or file_name.endswith(".markdown")))
        or file_name.endswith(".md")
        or file_name.endswith(".markdown")
        or re.search(r"\.md(?:[?#].*)?$", source_value, flags=re.IGNORECASE)
        or re.search(r"\.markdown(?:[?#].*)?$", source_value, flags=re.IGNORECASE)
        or re.match(r"^data:text/markdown[;,]", source_value, flags=re.IGNORECASE)
    ):
        return "markdown"
    if isinstance(source_input, str) and source_input.strip():
        return "pdf"
    raise RuntimeError('Unable to determine source type. Use sourceType with "pdf" or "markdown".')


def _normalize_presentation_source_input(source_input: Any, explicit_type: str = "") -> Dict[str, Any]:
    source_type = _infer_presentation_source_type(source_input, explicit_type)
    if isinstance(source_input, str):
        locator = source_input.strip()
        if not locator:
            raise RuntimeError("Missing source document.")
        return {
            "source_type": source_type,
            "locator": locator,
            "inline_text": None,
            "content_bytes": None,
            "mime_type": "",
            "filename": Path(locator).name,
        }

    if not isinstance(source_input, dict):
        raise RuntimeError("Unsupported source input.")

    descriptor_type = str(source_input.get("type") or source_input.get("kind") or source_input.get("sourceKind") or "").strip().lower()
    mime_type = _first_non_empty_presentation_string(
        source_input.get("mimeType"),
        source_input.get("mime_type"),
        source_input.get("contentType"),
        source_input.get("content_type"),
    )
    file_name = _first_non_empty_presentation_string(
        source_input.get("filename"),
        source_input.get("fileName"),
        source_input.get("name"),
        source_input.get("original_filename"),
        source_input.get("originalFilename"),
    )
    locator = _first_non_empty_presentation_string(
        source_input.get("sourceUrl"),
        source_input.get("url"),
        source_input.get("href"),
        source_input.get("uri"),
        source_input.get("src"),
        source_input.get("pdfUrl"),
        source_input.get("path"),
        source_input.get("filePath"),
        source_input.get("relative_path"),
        source_input.get("relativePath"),
        source_input.get("value") if descriptor_type in {"url", "path", "attachment", "file"} else "",
        source_input.get("name") if descriptor_type == "attachment" else "",
    )
    inline_text = _first_non_empty_presentation_string(
        source_input.get("markdown"),
        source_input.get("text"),
        source_input.get("value") if descriptor_type in {"inline", "markdown", "text"} else "",
        source_input.get("content"),
    )
    encoded_content = _first_non_empty_presentation_string(
        source_input.get("contentBase64"),
        source_input.get("content_base64"),
        source_input.get("base64"),
    )

    if source_type == "markdown" and inline_text and (descriptor_type in {"inline", "markdown", "text"} or not locator):
        return {
            "source_type": source_type,
            "locator": "",
            "inline_text": inline_text,
            "content_bytes": None,
            "mime_type": mime_type,
            "filename": file_name,
        }

    if encoded_content:
        if re.match(r"^data:", encoded_content, flags=re.IGNORECASE):
            decoded_mime, content_bytes = _decode_data_url_bytes(encoded_content)
        else:
            try:
                content_bytes = base64.b64decode(encoded_content, validate=True)
            except (ValueError, binascii.Error) as exc:
                raise RuntimeError("Invalid base64 source content.") from exc
            decoded_mime = mime_type or ("text/markdown" if source_type == "markdown" else "application/pdf")
        return {
            "source_type": source_type,
            "locator": "",
            "inline_text": None,
            "content_bytes": content_bytes,
            "mime_type": decoded_mime,
            "filename": file_name,
        }

    if locator:
        return {
            "source_type": _infer_presentation_source_type(
                {"name": file_name, "type": mime_type, "sourceUrl": locator},
                explicit_type or source_type,
            ),
            "locator": locator,
            "inline_text": None,
            "content_bytes": None,
            "mime_type": mime_type,
            "filename": file_name or Path(locator).name,
        }

    raise RuntimeError("Unsupported source input. Use a URL, scratch-relative path, attachment, inline Markdown, or base64 content.")


async def _fetch_presentation_source_url(source_url: str) -> Tuple[str, bytes]:
    timeout = httpx.Timeout(connect=15.0, read=60.0, write=15.0, pool=15.0)
    async with httpx.AsyncClient(follow_redirects=True, timeout=timeout) as client:
        response = await client.get(source_url)
    if response.status_code != 200:
        raise RuntimeError(f"Failed to fetch source URL: HTTP {response.status_code}")
    content = response.content
    if len(content) > FILE_OPS_MAX_SIZE_BYTES:
        raise RuntimeError(f"Source file is too large. Limit is {FILE_OPS_MAX_SIZE_BYTES} bytes.")
    mime_type = str(response.headers.get("content-type") or "").split(";", 1)[0].strip().lower()
    return mime_type, content


def _decode_text_bytes(content_bytes: bytes) -> str:
    try:
        return content_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return content_bytes.decode("latin-1")


def _extract_pdf_text_from_bytes(content_bytes: bytes) -> str:
    try:
        from PyPDF2 import PdfReader
    except ImportError as exc:
        raise RuntimeError("PyPDF2 is required for PDF to PowerPoint conversion.") from exc
    reader = PdfReader(io.BytesIO(content_bytes))
    parts: List[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        page_text = (page.extract_text() or "").strip()
        if page_text:
            parts.append(f"=== Page {page_index} ===\n{page_text}")
    return "\n\n".join(parts).strip()


def _clean_markdown_for_presentation_text(markdown_text: str) -> str:
    return (
        str(markdown_text or "")
        .replace("\r\n", "\n")
        .replace("\r", "\n")
        .replace("<!--", "\n<!--")
    )


def _build_markdown_from_slide_sections(
    slides: List[Dict[str, Any]],
    *,
    title: str,
    max_slides: int,
) -> str:
    normalized_slides = list(slides or [])
    if normalized_slides:
        first_slide = normalized_slides[0]
        first_bullets = [str(item).strip() for item in (first_slide.get("bullets") or []) if str(item).strip()]
        first_title = str(first_slide.get("title") or "").strip().lower()
        if not first_bullets and first_title and first_title == str(title or "").strip().lower():
            normalized_slides = normalized_slides[1:]
    lines = [f"# {title or 'Presentation'}"]
    for slide in normalized_slides[:max_slides]:
        slide_title = str(slide.get("title") or "").strip() or "Slide"
        bullets = [str(item).strip() for item in (slide.get("bullets") or []) if str(item).strip()]
        lines.append("")
        lines.append(f"## {slide_title}")
        if bullets:
            lines.extend(f"- {bullet}" for bullet in bullets[:6])
        else:
            lines.append("- Key point")
    return "\n".join(lines).strip()


def _build_fallback_presentation_markdown(
    document_text: str,
    *,
    title: str,
    max_slides: int,
) -> str:
    text = re.sub(r"\s+", " ", str(document_text or "")).strip()
    if not text:
        raise RuntimeError("Source document is empty.")
    sentences = [
        sentence.strip(" -•\t\r\n")
        for sentence in re.split(r"(?<=[.!?])\s+", text)
        if sentence and sentence.strip()
    ]
    if not sentences:
        sentences = [text]
    section_count = max(3, min(max_slides, max(3, min(8, (len(sentences) + 2) // 3))))
    title_bank = [
        "Overview",
        "Key Points",
        "Details",
        "Evidence",
        "Implications",
        "Recommendations",
        "Summary",
        "Next Steps",
    ]
    chunk_size = max(1, (len(sentences) + section_count - 1) // section_count)
    lines = [f"# {title or 'Presentation'}"]
    cursor = 0
    for index in range(section_count):
        chunk = sentences[cursor : cursor + chunk_size]
        cursor += chunk_size
        if not chunk:
            break
        section_title = title_bank[index] if index < len(title_bank) else f"Section {index + 1}"
        lines.append("")
        lines.append(f"## {section_title}")
        for bullet in chunk[:4]:
            lines.append(f"- {bullet[:220].strip()}")
    return "\n".join(lines).strip()


async def _generate_presentation_markdown_from_document(
    *,
    document_text: str,
    source_type: str,
    title: str,
    max_slides: int,
    model_name: Optional[str] = None,
) -> str:
    normalized_text = re.sub(r"\s+", " ", str(document_text or "")).strip()
    if not normalized_text:
        raise RuntimeError("Source document is empty.")
    resolved_model = (model_name or TELEGRAM_DEFAULT_MODEL or "").strip() or "gpt-4o-mini"
    endpoint = _normalize_chat_endpoint(build_openai_url(TELEGRAM_OPENAI_CHAT_PATH))
    api_key = _first_non_empty_env(preferred_api_key_env_names(TELEGRAM_OPENAI_BASE_URL, resolved_model))
    if not api_key:
        return _build_fallback_presentation_markdown(normalized_text, title=title, max_slides=max_slides)

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
    }
    if OPENAI_ORG_ID:
        headers["OpenAI-Organization"] = OPENAI_ORG_ID
    if OPENAI_PROJECT_ID:
        headers["OpenAI-Project"] = OPENAI_PROJECT_ID

    payload = {
        "model": resolved_model,
        "messages": [
            {
                "role": "system",
                "content": (
                    "You are creating slide-ready markdown for a PowerPoint deck from a source document.\n"
                    "Return markdown only. Do not include code fences or commentary.\n"
                    "Format rules:\n"
                    "- First line must be a single H1 title.\n"
                    "- Then create 4 to 12 H2 sections depending on the requested slide count.\n"
                    "- Each H2 section must contain 2 to 5 bullet lines.\n"
                    "- Keep wording concise, concrete, and presentation-ready.\n"
                    "- Preserve important facts, numbers, dates, and names from the source.\n"
                    "- Do not mention these instructions."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Source type: {source_type}\n"
                    f"Requested title: {title or 'Presentation'}\n"
                    f"Target slides: {max_slides}\n\n"
                    "Document text:\n"
                    f"{normalized_text[:12000]}"
                ),
            },
        ],
        "temperature": 0.3,
    }
    try:
        response = await _call_chat_completion(
            endpoint,
            headers,
            payload,
            timeout_seconds=max(20.0, min(120.0, TELEGRAM_CHAT_TIMEOUT)),
        )
        if response.status_code == 200:
            data = response.json()
            choices = data.get("choices") or []
            if choices:
                normalized_message = normalize_chat_completion_message(choices[0].get("message") or {})
                content = _strip_markdown_fence_block(coerce_message_text(normalized_message.get("content") or ""))
                if content.strip():
                    return content.strip()
    except Exception:
        pass
    return _build_fallback_presentation_markdown(normalized_text, title=title, max_slides=max_slides)


async def _load_presentation_source_document(normalized_source: Dict[str, Any]) -> Dict[str, Any]:
    source_type = str(normalized_source.get("source_type") or "").strip().lower()
    locator = str(normalized_source.get("locator") or "").strip()
    inline_text = normalized_source.get("inline_text")
    content_bytes = normalized_source.get("content_bytes")
    file_name = str(normalized_source.get("filename") or "").strip()
    source_descriptor = locator or file_name or source_type

    if source_type == "markdown":
        if isinstance(inline_text, str) and inline_text.strip():
            return {"source_type": source_type, "text": inline_text, "source_descriptor": source_descriptor}
        if isinstance(content_bytes, (bytes, bytearray)):
            return {
                "source_type": source_type,
                "text": _decode_text_bytes(bytes(content_bytes)),
                "source_descriptor": source_descriptor,
            }
        if locator:
            if re.match(r"^data:", locator, flags=re.IGNORECASE):
                _, raw_bytes = _decode_data_url_bytes(locator)
                return {"source_type": source_type, "text": _decode_text_bytes(raw_bytes), "source_descriptor": source_descriptor}
            if _looks_like_http_url(locator):
                _, raw_bytes = await _fetch_presentation_source_url(locator)
                return {"source_type": source_type, "text": _decode_text_bytes(raw_bytes), "source_descriptor": locator}
            filepath = resolve_scratch_path(locator, READ_ALLOWED_EXTENSIONS)
            if not filepath.exists() or not filepath.is_file():
                raise RuntimeError(f"Markdown source not found: {locator}")
            return {"source_type": source_type, "text": read_text_file(filepath), "source_descriptor": locator}
        raise RuntimeError("Missing Markdown source.")

    if source_type != "pdf":
        raise RuntimeError(f"Unsupported source type: {source_type}")
    if isinstance(content_bytes, (bytes, bytearray)):
        return {
            "source_type": source_type,
            "text": _extract_pdf_text_from_bytes(bytes(content_bytes)),
            "source_descriptor": source_descriptor,
        }
    if not locator:
        raise RuntimeError("Missing PDF source.")
    if re.match(r"^data:", locator, flags=re.IGNORECASE):
        _, raw_bytes = _decode_data_url_bytes(locator)
        return {"source_type": source_type, "text": _extract_pdf_text_from_bytes(raw_bytes), "source_descriptor": source_descriptor}
    if _looks_like_http_url(locator):
        _, raw_bytes = await _fetch_presentation_source_url(locator)
        return {"source_type": source_type, "text": _extract_pdf_text_from_bytes(raw_bytes), "source_descriptor": locator}
    filepath = resolve_scratch_path(locator, READ_ALLOWED_EXTENSIONS)
    if not filepath.exists() or not filepath.is_file():
        raise RuntimeError(f"PDF source not found: {locator}")
    return {"source_type": source_type, "text": read_pdf_file(filepath), "source_descriptor": locator}


def _choose_pdf_to_powerpoint_output_path(filename: str, *, title: str) -> Tuple[str, Path]:
    logical_name = str(filename or "").strip()
    if not logical_name:
        slug = _slugify_telegram_slides_component(title, fallback="presentation")
        timestamp = datetime.now(timezone.utc).strftime("%Y%m%d-%H%M%S")
        logical_name = f"presentations/{slug}-{timestamp}.pptx"
    logical_name = _normalize_scratch_relative_input(logical_name)
    if not logical_name.lower().endswith(".pptx"):
        logical_name = f"{logical_name}.pptx"
    filepath = resolve_scratch_path(logical_name, WRITE_ALLOWED_EXTENSIONS)
    filepath.parent.mkdir(parents=True, exist_ok=True)
    return logical_name.replace("\\", "/"), filepath


def _pick_latest_presentation_attachment(attachment_records: List[Dict[str, Any]]) -> str:
    for item in reversed(list(attachment_records or [])):
        relative_path = str(item.get("relative_path") or "").strip()
        suffix = Path(relative_path).suffix.lower()
        if suffix in {".pdf", ".md", ".markdown"}:
            return relative_path
    return ""


def _render_markdown_to_powerpoint(
    markdown: str,
    *,
    title: str,
    output_path: Path,
    max_slides: int,
    source_label: str,
) -> Dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PowerPoint generation.") from exc

    from src.skills.builtin.googleworkspace_slides_support import _parse_markdown_to_slides

    slides = _parse_markdown_to_slides(
        markdown,
        fallback_title=title or "Presentation",
        scratch_root=SCRATCH_DIR,
        source_dir=SCRATCH_DIR / "presentations",
    )[:max_slides]
    if not slides:
        raise RuntimeError("No slide content could be generated from the source document.")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    title_slide_hint = slides[0]
    title_hint_bullets = [str(item).strip() for item in (title_slide_hint.get("bullets") or []) if str(item).strip()]
    use_first_slide_as_title = not title_hint_bullets

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = str(title or title_slide_hint.get("title") or "Presentation").strip() or "Presentation"
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Generated from {source_label}"

    content_slides = slides[1:] if use_first_slide_as_title else slides
    if not content_slides:
        content_slides = slides[:1]

    rendered_slide_count = 0
    for slide_data in content_slides[:max_slides]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(slide_data.get("title") or "Slide").strip() or "Slide"
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        bullets = [str(item).strip() for item in (slide_data.get("bullets") or []) if str(item).strip()]
        if not bullets:
            bullets = ["Key point"]
        for index, bullet in enumerate(bullets[:6]):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = bullet[:240]
            paragraph.level = 0
            if paragraph.runs:
                paragraph.runs[0].font.size = Pt(20)
        rendered_slide_count += 1

    presentation.save(str(output_path))
    return {
        "rendered_slide_count": rendered_slide_count,
        "parsed_slide_count": len(slides),
    }


async def _handle_pdf_to_powerpoint_internal(
    arguments: Dict[str, Any],
    *,
    conversation_id: str,
    user_id: str,
    attachment_records: Optional[List[Dict[str, Any]]] = None,
    model_name: Optional[str] = None,
) -> Dict[str, Any]:
    source = arguments.get("source")
    source_url = _first_non_empty_presentation_string(
        arguments.get("sourceUrl"),
        arguments.get("pdfUrl"),
    )
    title = str(arguments.get("title") or "").strip() or "Presentation"
    filename = str(arguments.get("filename") or "").strip()
    max_slides = _coerce_pdf_to_powerpoint_max_slides(arguments.get("maxSlides") or arguments.get("max_slides"))

    try:
        resolved_source_input: Any = source if source is not None else source_url
        if resolved_source_input in (None, ""):
            latest_attachment = _pick_latest_presentation_attachment(attachment_records or [])
            if latest_attachment:
                resolved_source_input = latest_attachment
        if resolved_source_input in (None, ""):
            return {
                "success": False,
                "message": "Provide a PDF or Markdown source via upload, scratch path, or URL.",
            }

        normalized_source = _normalize_presentation_source_input(
            resolved_source_input,
            str(arguments.get("sourceType") or "").strip(),
        )
        source_document = await _load_presentation_source_document(normalized_source)
        source_type = str(source_document.get("source_type") or "").strip().lower()
        document_text = str(source_document.get("text") or "").strip()
        if len(document_text) < 40:
            raise RuntimeError(f"Could not extract sufficient text content from the {source_type or 'source'} document.")

        markdown_output_path: Optional[str] = None
        markdown_to_render = ""
        if source_type == "markdown":
            from src.skills.builtin.googleworkspace_slides_support import _parse_markdown_to_slides

            parsed_slides = _parse_markdown_to_slides(
                document_text,
                fallback_title=title or "Presentation",
                scratch_root=SCRATCH_DIR,
                source_dir=SCRATCH_DIR,
            )[:max_slides]
            if len(parsed_slides) >= 2:
                markdown_to_render = _build_markdown_from_slide_sections(parsed_slides, title=title, max_slides=max_slides)
            else:
                cleaned_markdown_text = _clean_markdown_for_presentation_text(document_text)
                markdown_to_render = await _generate_presentation_markdown_from_document(
                    document_text=cleaned_markdown_text,
                    source_type=source_type,
                    title=title,
                    max_slides=max_slides,
                    model_name=model_name,
                )
                markdown_output_path = _write_telegram_slides_markdown(markdown_to_render, title=title)
        else:
            markdown_to_render = await _generate_presentation_markdown_from_document(
                document_text=document_text,
                source_type=source_type,
                title=title,
                max_slides=max_slides,
                model_name=model_name,
            )
            markdown_output_path = _write_telegram_slides_markdown(markdown_to_render, title=title)

        relative_output_path, output_path = _choose_pdf_to_powerpoint_output_path(filename, title=title)
        render_stats = _render_markdown_to_powerpoint(
            markdown_to_render,
            title=title,
            output_path=output_path,
            max_slides=max_slides,
            source_label="Markdown" if source_type == "markdown" else "PDF",
        )
        data: Dict[str, Any] = {
            "file_path": relative_output_path,
            "source_type": source_type,
            "source_descriptor": str(source_document.get("source_descriptor") or ""),
            "requested_max_slides": max_slides,
            **render_stats,
        }
        if markdown_output_path:
            data["generated_markdown_path"] = markdown_output_path
        return {
            "success": True,
            "message": (
                f"Created {relative_output_path} from {source_type or 'document'} with "
                f"{render_stats.get('rendered_slide_count', 0)} content slide(s)."
            ),
            "data": data,
        }
    except HTTPException as exc:
        return {"success": False, "message": str(exc.detail or "Invalid presentation source.")}
    except Exception as exc:
        return {"success": False, "message": f"PowerPoint conversion failed: {exc}"}


async def _create_telegram_slides_presentation_internal(
    arguments: Dict[str, Any],
    *,
    conversation_id: str,
    user_id: str,
    model_name: Optional[str] = None,
) -> Dict[str, Any]:
    prompt = str(
        arguments.get("prompt")
        or arguments.get("topic")
        or arguments.get("brief")
        or arguments.get("description")
        or ""
    ).strip()
    title = str(arguments.get("title") or "").strip() or "Presentation"
    markdown = str(arguments.get("markdown") or "").strip()
    markdown_path = str(arguments.get("markdown_path") or arguments.get("markdownPath") or "").strip()
    max_slides = _coerce_telegram_slides_max_slides(arguments.get("max_slides"))
    audience = str(arguments.get("audience") or "").strip()
    purpose = str(arguments.get("purpose") or "").strip()
    tone = str(arguments.get("tone") or "").strip()

    if not any([prompt, markdown, markdown_path]):
        return {"success": False, "message": "Provide at least one of prompt, markdown, or markdown_path."}

    generated_markdown_path: Optional[str] = None
    generated_stats: Optional[Dict[str, int]] = None

    if not markdown_path:
        if not markdown:
            retry_reason = ""
            for _ in range(2):
                markdown = await _generate_telegram_slides_markdown(
                    prompt=prompt,
                    title=title,
                    max_slides=max_slides,
                    model_name=model_name,
                    audience=audience,
                    purpose=purpose,
                    tone=tone,
                    retry_reason=retry_reason,
                )
                is_valid, stats = _validate_telegram_slides_markdown(
                    markdown,
                    title=title,
                    max_slides=max_slides,
                )
                generated_stats = stats
                if is_valid:
                    break
                retry_reason = (
                    "The draft was too thin. Produce a fuller presentation with multiple H2 slides "
                    "and several bullet points per slide."
                )
            else:
                return {
                    "success": False,
                    "message": (
                        "I couldn't generate a strong enough slide draft for Google Slides. "
                        "Please provide a more specific prompt or pass markdown directly."
                    ),
                    "data": {"validation": generated_stats or {}},
                }
        else:
            is_valid, stats = _validate_telegram_slides_markdown(
                markdown,
                title=title,
                max_slides=max_slides,
            )
            generated_stats = stats
            if not is_valid:
                return {
                    "success": False,
                    "message": (
                        "The provided markdown is too thin to create a useful presentation. "
                        "Please add more slide sections and bullet points."
                    ),
                    "data": {"validation": stats},
                }

        generated_markdown_path = _write_telegram_slides_markdown(markdown, title=title)
        markdown_path = generated_markdown_path

    max_images_per_slide_raw = arguments.get("max_images_per_slide", 1)
    try:
        max_images_per_slide = int(max_images_per_slide_raw)
    except (TypeError, ValueError):
        max_images_per_slide = 1

    tool_args: Dict[str, Any] = {
        "title": title,
        "markdown_path": markdown_path,
        "max_slides": max_slides,
        "attach_scratch_images": _coerce_bool(arguments.get("attach_scratch_images", True), default=True),
        "image_dir": str(arguments.get("image_dir", "images") or "images").strip() or "images",
        "max_images_per_slide": max(0, min(max_images_per_slide, 4)),
        "include_image_requests": _coerce_bool(arguments.get("include_image_requests", True), default=True),
    }
    timeout_seconds = arguments.get("timeout_seconds")
    if timeout_seconds is not None:
        tool_args["timeout_seconds"] = timeout_seconds

    skill_result = await _execute_skill_framework_tool(
        tool_name="googleworkspace_cli.slides_create_presentation_from_markdown",
        arguments=tool_args,
        conversation_id=conversation_id,
        user_id=user_id,
        metadata={"channel": "telegram", "source_tool": "createSlidesPresentation"},
    )

    data = dict(skill_result.get("data") or {}) if isinstance(skill_result.get("data"), dict) else {}
    if generated_markdown_path:
        data["generated_markdown_path"] = generated_markdown_path
    if generated_stats:
        data["generated_markdown_validation"] = generated_stats
    if data:
        skill_result["data"] = data

    if skill_result.get("success") and generated_markdown_path:
        base_message = str(skill_result.get("message") or "").strip()
        suffix = f" Source markdown saved to {generated_markdown_path}."
        if suffix.strip() not in base_message:
            skill_result["message"] = f"{base_message}{suffix}".strip()
    return skill_result

# MCP Client Manager class to handle transport lifecycle
class MCPClientManager:
    """Manages MCP client and transport lifecycle."""

    def __init__(self, server_config: Dict[str, Any]):
        self.server_config = server_config
        self.client = None
        self.transport = None
        self.read = None
        self.write = None

    async def connect(self):
        """Connect to MCP server using server-side allowlisted preset only; never execute user-supplied command."""
        if not MCP_AVAILABLE:
            raise Exception("MCP SDK not available")

        preset_id = self.server_config.get("preset_id")
        if not preset_id or preset_id not in MCP_PRESETS:
            raise ValueError(
                "Server has no valid preset_id. Reconfigure the server with a preset (e.g. preset_id: 'browser-use')."
            )
        preset = MCP_PRESETS[preset_id]
        if preset.get("type") == "inprocess":
            # browser-use: connect is handled in connect_server without calling this code path
            raise ValueError(
                f"Preset '{preset_id}' is inprocess; connect is handled separately. Do not call MCPClientManager.connect for this preset."
            )
        if preset.get("type") != "stdio":
            raise ValueError(
                f"Preset '{preset_id}' has no executable command. Reconfigure the server with a valid stdio preset."
            )

        # Resolve command and args only from allowlist; never use server_config['command']
        command = preset.get("command")
        allowed_args = preset.get("allowed_args", [])
        if not command or not allowed_args:
            raise ValueError(f"Preset '{preset_id}' has no command/args defined in MCP_PRESETS.")
        # Use first allowed_args entry for this preset (e.g. ["-m", "mcp_server_browser_use"])
        args = list(allowed_args[0]) if isinstance(allowed_args[0], (list, tuple)) else []

        try:
            print(f"🔧 Creating MCP client with command: {command} and args: {args}")

            # Prepare environment variables from server config (apiKey, model only)
            env = os.environ.copy()
            if self.server_config.get("apiKey"):
                model = self.server_config.get("model", "").lower()
                if "gemini" in model:
                    env["GOOGLE_API_KEY"] = self.server_config["apiKey"]
                    env["MCP_MODEL_PROVIDER"] = "google"
                elif "claude" in model:
                    env["ANTHROPIC_API_KEY"] = self.server_config["apiKey"]
                    env["MCP_MODEL_PROVIDER"] = "anthropic"
                else:
                    env["OPENAI_API_KEY"] = self.server_config["apiKey"]
                    env["MCP_MODEL_PROVIDER"] = "openai"
            if self.server_config.get("model"):
                env["MCP_MODEL_NAME"] = self.server_config["model"]
            env.setdefault("BROWSER_USE_HEADLESS", "true")
            env.setdefault("BROWSER_USE_DISABLE_SECURITY", "false")

            print(f"ðŸ” Environment variables for MCP server: {list(env.keys())}")

            server_params = StdioServerParameters(command=command, args=args, env=env)

            try:
                transport_cm = stdio_client(server_params)
                async with transport_cm as (self.read, self.write):
                    self.client = ClientSession(self.read, self.write)
                    await self.client.initialize()
                    print("✅ MCP client setup complete")
                    return self.client
            except Exception as e:
                print(f"Failed to create stdio client: {e}")
                raise

        except Exception as e:
            print(f"MCP connection error: {e}")
            raise Exception(f"Failed to connect to MCP server: {str(e)}")

    async def disconnect(self):
        """Disconnect from MCP server."""
        if self.client:
            await self.client.close()
        if self.transport:
            # The transport context manager will handle cleanup
            pass

# FastAPI app
app = FastAPI(title="CATBot Proxy Server", version="2.0.0")

if skill_manager is not None and create_skill_router is not None:
    app.include_router(create_skill_router(skill_manager, auth_dependency=get_current_user))

# Startup event to verify app initialization
@app.on_event("startup")
async def startup_event():
    """Log that the application has started successfully."""
    import sys
    await _get_shared_chat_http_client()
    print("🚀 FastAPI application startup event fired", flush=True)
    sys.stdout.flush()
    print(f"🚀 App routes registered: {len(app.routes)} routes", flush=True)
    sys.stdout.flush()
    # List all registered routes
    for route in app.routes:
        if hasattr(route, 'path') and hasattr(route, 'methods'):
            print(f"   Route: {list(route.methods)} {route.path}", flush=True)
            sys.stdout.flush()


@app.on_event("shutdown")
async def shutdown_event():
    """Stop AutoGen code executors (e.g. Docker containers) on app shutdown."""
    global autogen_team
    await _close_shared_chat_http_client()
    if autogen_team is not None:
        await _stop_code_executors(autogen_team)
        autogen_team = None

# Request logging middleware to debug CORS issues
class RequestLoggingMiddleware(BaseHTTPMiddleware):
    """Middleware to log all incoming requests for debugging."""
    async def dispatch(self, request: Request, call_next):
        import sys
        # Safely log the incoming request with error handling
        try:
            method = getattr(request, 'method', 'UNKNOWN')
            path = getattr(request.url, 'path', 'unknown') if hasattr(request, 'url') else 'unknown'
            query = getattr(request.url, 'query', '') if hasattr(request, 'url') and hasattr(request.url, 'query') else ''
            origin = request.headers.get('origin', 'none') if hasattr(request, 'headers') else 'none'
            print(f"ðŸŒ [{method}] {path}?{query}", flush=True)
            sys.stdout.flush()
            print(f"   Origin: {origin}", flush=True)
            sys.stdout.flush()
            if hasattr(request, 'headers'):
                print(f"   Headers: {dict(request.headers)}", flush=True)
                sys.stdout.flush()
        except Exception as log_error:
            # If logging fails, continue anyway - don't break the request
            print(f"âš ï¸ Error logging request: {log_error}", flush=True)
            sys.stdout.flush()
            import traceback
            print(traceback.format_exc(), flush=True)
            sys.stdout.flush()
        
        try:
            # Process the request
            response = await call_next(request)
            # Log successful response
            try:
                method = getattr(request, 'method', 'UNKNOWN')
                path = getattr(request.url, 'path', 'unknown') if hasattr(request, 'url') else 'unknown'
                status = getattr(response, 'status_code', 'unknown')
                print(f"✅ [{method}] {path} -> {status}", flush=True)
                sys.stdout.flush()
            except Exception as log_err:
                print(f"âš ï¸ Error logging response: {log_err}", flush=True)
                sys.stdout.flush()
            return response
        except Exception as e:
            # Log the exception
            try:
                method = getattr(request, 'method', 'UNKNOWN')
                path = getattr(request.url, 'path', 'unknown') if hasattr(request, 'url') else 'unknown'
                print(f"âŒ [{method}] {path} -> Exception: {e}", flush=True)
                sys.stdout.flush()
                import traceback
                print(traceback.format_exc(), flush=True)
                sys.stdout.flush()
            except Exception:
                print("âŒ Error in exception logging", flush=True)
                sys.stdout.flush()
            # Re-raise the exception so it can be handled by exception handlers
            raise

# Add request logging middleware first (before CORS)
app.add_middleware(RequestLoggingMiddleware)

# Add CORS middleware
# Allow specific origins for development and production
# Note: Using allow_credentials=False allows more flexible origin matching
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins for development (be more permissive to debug)
    allow_credentials=False,  # Set to False to allow more flexible CORS handling
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS", "PATCH"],
    allow_headers=["*", "Authorization", "X-Auth-Token", "Content-Type", "Accept"],  # Explicitly allow Authorization header
    expose_headers=["*"],
)


@app.middleware("http")
async def require_auth_for_v1_routes(request: Request, call_next):
    path = request.url.path

    if request.method == "OPTIONS":
        return await call_next(request)

    # Exempt certain routes from authentication (public endpoints)
    exempt_paths = {
        "/v1/auth/signup",
        "/v1/auth/login",
        "/v1/tools/log",  # Tool invocation log sink for HTML UI tool calls
        "/v1/audio/transcriptions",  # Whisper endpoint - public for audio transcription
        "/v1/audio/speech",  # Embedded OpenAI-compatible TTS speech endpoint
        "/v1/audio/voices",  # Embedded OpenAI-compatible TTS voices endpoint
        "/v1/proxy/chat/completions",  # Chat completions proxy - public to avoid mixed content
        "/v1/proxy/models",  # Models list proxy - public to avoid mixed content
        "/v1/proxy/autogen",  # AutoGen workflow proxy - public to avoid mixed content
        "/v1/proxy/tts/voices",  # TTS voices endpoint - public
        "/v1/proxy/tts/speech",  # TTS speech endpoint - public
        "/v1/proxy/search",  # Search proxy - public
        "/v1/proxy/news",  # News proxy - public
        "/v1/status/start",
        "/v1/status/update",
        "/v1/status/finish",
        "/v1/status/latest",
        "/v1/status/events",
    }
    autogen_team_secret_paths = {
        "/v1/proxy/browser-agent",
        "/v1/proxy/deep-research",
        "/v1/proxy/browser-health",
        "/v1/proxy/fetch",
        "/v1/proxy/codex",
    }
    # Telegram bot endpoints are unauthenticated (bot uses TELEGRAM_SECRET when set)
    require_auth = (
        path.startswith("/v1/")
        and path not in exempt_paths
        and not path.startswith("/v1/telegram/chat")
    )
    if require_auth:
        try:
            if path in autogen_team_secret_paths and _autogen_team_secret_matches(request):
                return await call_next(request)
            # Get authorization header - FastAPI headers are case-insensitive, but check both for robustness
            # Use get() with case-insensitive lookup
            auth_header = None
            x_auth_token = None
            
            # Try to get authorization header (case-insensitive)
            for key, value in request.headers.items():
                if key.lower() == "authorization":
                    auth_header = value
                    break
                elif key.lower() == "x-auth-token":
                    x_auth_token = value
                    break
            
            # Fallback to direct get if not found in loop
            if not auth_header:
                auth_header = request.headers.get("authorization") or request.headers.get("Authorization")
            if not x_auth_token:
                x_auth_token = request.headers.get("x-auth-token") or request.headers.get("X-Auth-Token")
            
            # Debug logging for auth issues
            if not auth_header and not x_auth_token:
                print(f"🔒 Auth check failed for {path}: No authorization header found")
                print(f"   Available headers: {list(request.headers.keys())}")
            elif auth_header:
                # Log token preview for debugging (first 50 chars)
                token_preview = auth_header[:50] + "..." if len(auth_header) > 50 else auth_header
                print(f"🔒 Auth check for {path}: Found auth header (length: {len(auth_header)}, preview: {token_preview})")
            
            get_current_user_from_headers(
                auth_header,
                x_auth_token,
            )
        except HTTPException as exc:
            print(f"🔒 Auth check failed for {path}: {exc.detail}")
            # Log the actual header value for debugging (truncated)
            auth_debug = request.headers.get("authorization") or request.headers.get("Authorization") or "None"
            if auth_debug != "None":
                print(f"   Auth header value (first 100 chars): {auth_debug[:100]}")
            # Include CORS headers in error response
            cors_headers = build_cors_headers(request)
            return JSONResponse(
                status_code=exc.status_code,
                content={"detail": exc.detail},
                headers=cors_headers
            )

    return await call_next(request)

def build_cors_headers(request: Request) -> Dict[str, str]:
    """Build CORS headers for the request origin. Supports both localhost and remote access."""
    try:
        # Safely get origin from request headers, with fallback
        if hasattr(request, 'headers'):
            origin = request.headers.get("origin")
            # If no origin header (e.g., same-origin request), allow all origins for network access
            if not origin:
                origin = "*"
        else:
            origin = "*"
    except Exception:
        # If we can't access headers, allow all origins for network access
        origin = "*"
    
    return {
        "Access-Control-Allow-Origin": origin,
        "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS, PATCH",
        "Access-Control-Allow-Headers": "*",
    }

# Global exception handler to ensure CORS headers are always included
@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    """Handle HTTP exceptions and ensure CORS headers are included."""
    # Safely build CORS headers - if this fails, use minimal headers
    try:
        cors_headers = build_cors_headers(request)
    except Exception as header_error:
        print(f"âš ï¸ Error building CORS headers in HTTPException handler: {header_error}")
        # Use minimal safe headers if build_cors_headers fails
        cors_headers = {
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS, PATCH",
            "Access-Control-Allow-Headers": "*",
        }
    
    return JSONResponse(
        status_code=exc.status_code,
        content={"detail": exc.detail},
        headers=cors_headers,
    )

@app.exception_handler(RequestValidationError)
async def validation_exception_handler(request: Request, exc: RequestValidationError):
    """Handle validation errors and ensure CORS headers are included."""
    # Safely build CORS headers - if this fails, use minimal headers
    try:
        cors_headers = build_cors_headers(request)
    except Exception as header_error:
        print(f"âš ï¸ Error building CORS headers in ValidationException handler: {header_error}")
        # Use minimal safe headers if build_cors_headers fails
        cors_headers = {
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS, PATCH",
            "Access-Control-Allow-Headers": "*",
        }
    
    return JSONResponse(
        status_code=422,
        content={"detail": exc.errors()},
        headers=cors_headers,
    )

@app.exception_handler(Exception)
async def general_exception_handler(request: Request, exc: Exception):
    """Handle all other exceptions and ensure CORS headers are included."""
    import sys
    import traceback
    error_trace = traceback.format_exc()
    print(f"âŒ Unhandled exception in general_exception_handler: {exc}", flush=True)
    sys.stdout.flush()
    print(error_trace, flush=True)
    sys.stdout.flush()
    
    # Safely build CORS headers - if this fails, use minimal headers
    try:
        cors_headers = build_cors_headers(request)
    except Exception as header_error:
        print(f"âš ï¸ Error building CORS headers: {header_error}", flush=True)
        sys.stdout.flush()
        import traceback
        print(traceback.format_exc(), flush=True)
        sys.stdout.flush()
        # Use minimal safe headers if build_cors_headers fails
        cors_headers = {
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS, PATCH",
            "Access-Control-Allow-Headers": "*",
        }
    
    try:
        response = JSONResponse(
            status_code=500,
            content={"detail": f"Internal server error: {str(exc)}"},
            headers=cors_headers,
        )
        print(f"✅ Created error response with status 500", flush=True)
        sys.stdout.flush()
        return response
    except Exception as response_error:
        print(f"âŒ Error creating error response: {response_error}", flush=True)
        sys.stdout.flush()
        import traceback
        print(traceback.format_exc(), flush=True)
        sys.stdout.flush()
        # Last resort - return a simple response
        from fastapi.responses import PlainTextResponse
        return PlainTextResponse(
            content=f"Internal server error: {str(exc)}",
            status_code=500,
            headers=cors_headers,
        )

# Helper function to clean HTML text (similar to Node.js version)
def clean_text(text: str) -> str:
    """Clean HTML text by removing tags and decoding entities."""
    if not text:
        return ""

    # Remove HTML tags
    text = re.sub(r'</?[^>]+(>|$)', '', text)

    # Decode HTML entities
    html_entities = {
        '&amp;': '&', '&lt;': '<', '&gt;': '>', '&quot;': '"', '&#039;': "'",
        '&rsquo;': "'", '&lsquo;': "'", '&rdquo;': '"', '&ldquo;': '"',
        '&ndash;': '-', '&mdash;': '—'
    }

    for entity, replacement in html_entities.items():
        text = text.replace(entity, replacement)

    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

# Helper function to parse dates (similar to Node.js version)
def parse_date(date_str: Optional[str]) -> Optional[float]:
    """Parse date string to timestamp."""
    if not date_str:
        return None
    try:
        return datetime.fromisoformat(date_str.replace('Z', '+00:00')).timestamp()
    except (ValueError, AttributeError):
        return None

# Load servers from disk; migrate legacy 'command' to preset_id and never retain command
def load_servers():
    """Load MCP servers from JSON file. Migrate legacy config: set preset_id from name, drop command."""
    global mcp_servers
    try:
        if SERVERS_FILE.exists():
            with open(SERVERS_FILE, "r", encoding="utf-8") as f:
                servers = json.load(f)
            result = {}
            for server in servers:
                sid = server.get("id")
                if not sid:
                    continue
                # Migrate: infer preset_id from name when missing (browser-use detection)
                if not server.get("preset_id"):
                    name = (server.get("name") or "").lower()
                    if "mcp-browser-use" in name or "browser-use" in name or ("browser" in name and "use" in name):
                        server["preset_id"] = "browser-use"
                # Never retain command in memory
                server.pop("command", None)
                result[sid] = server
            mcp_servers = result
            print(f"Loaded {len(mcp_servers)} MCP servers from disk")
    except Exception as e:
        print(f"No existing servers file found, starting with empty state: {e}")

# Save servers to disk; never persist 'command'
def save_servers():
    """Save MCP servers to disk. Only persist safe keys; never write command."""
    global mcp_servers
    try:
        servers = []
        for s in mcp_servers.values():
            safe = {k: s[k] for k in MCP_SERVER_SAFE_KEYS if k in s}
            servers.append(safe)
        with open(SERVERS_FILE, "w", encoding="utf-8") as f:
            json.dump(servers, f, indent=2, ensure_ascii=False)
        print(f"Saved {len(servers)} MCP servers to disk")
    except Exception as e:
        print(f"Error saving servers to disk: {e}")

# Load AutoGen team from config
def load_autogen_team():
    """Load AutoGen team from team-config.json."""
    global autogen_team
    
    if not AUTOGEN_AVAILABLE:
        print("âš ï¸  AutoGen not available, skipping team load")
        return None
    
    try:
        if not TEAM_CONFIG_FILE.exists():
            print(f"âš ï¸  Team config file not found: {TEAM_CONFIG_FILE}")
            return None
            
        print(f"📂 Loading AutoGen team from {TEAM_CONFIG_FILE}...")
        
        with open(TEAM_CONFIG_FILE, 'r', encoding='utf-8') as f:
            team_config = json.load(f)

        patched_agents = _patch_autogen_team_config_for_model_compatibility(team_config)
        if patched_agents:
            print(
                "[AUTOGEN] Disabled reflect_on_tool_use for OpenRouter-backed agents: "
                + ", ".join(patched_agents),
                flush=True,
            )
        
        # Load the team from the configuration using ComponentLoader
        loader = ComponentLoader()
        team = loader.load_component(team_config)

        # Inject PythonCodeExecutionTool (Docker) into the lead engineer when available
        if AUTOGEN_CODE_EXEC_AVAILABLE and PythonCodeExecutionTool and DockerCommandLineCodeExecutor:
            coding_dir = _PROJECT_ROOT / "coding"
            try:
                coding_dir.mkdir(exist_ok=True)
            except OSError:
                pass
            work_dir = str(coding_dir)
            executor = DockerCommandLineCodeExecutor(
                work_dir=work_dir,
                image="python:3.12-slim",
                timeout=120,
            )
            code_tool = PythonCodeExecutionTool(executor)
            participants = getattr(team, "participants", None) or getattr(team, "_participants", [])
            if participants:
                target_agent = None
                for participant in participants:
                    candidate_name = getattr(participant, "name", None) or getattr(participant, "_name", None)
                    if candidate_name == "lead_engineer_agent":
                        target_agent = participant
                        break
                if target_agent is None:
                    target_agent = participants[0]

                wb = getattr(target_agent, "workbench", getattr(target_agent, "_workbench", None))
                if wb is not None:
                    wb_list = wb if isinstance(wb, list) else [wb]
                    for wb_item in wb_list:
                        tools = getattr(wb_item, "tools", getattr(wb_item, "_tools", None))
                        if tools is not None and isinstance(tools, list):
                            tools.insert(0, code_tool)
                            target_name = getattr(target_agent, "name", None) or getattr(target_agent, "_name", "unknown_agent")
                            print(f"✅ Injected PythonCodeExecutionTool (Docker) into {target_name} workbench")
                            break

        print(f"✅ AutoGen team loaded successfully: {team_config.get('label', 'Unknown')}")
        return team
        
    except Exception as e:
        import traceback
        print(f"âŒ Error loading AutoGen team: {e}")
        print(traceback.format_exc())
        error_text = str(e)
        if "No user query found in messages" in error_text:
            autogen_base = (
                os.getenv("AUTOGEN_OPENROUTER_BASE_URL")
                or os.getenv("OPENROUTER_API_BASE")
                or os.getenv("MCP_LLM_BASE_URL")
                or os.getenv("OPENAI_API_BASE")
                or ""
            ).strip()
            autogen_model = (
                os.getenv("AUTOGEN_TEAM_MODEL")
                or os.getenv("OPENROUTER_AUTOGEN_MODEL")
                or os.getenv("MCP_LLM_MODEL_NAME")
                or os.getenv("OPENAI_MODEL")
                or ""
            ).strip()
            error_text = (
                "AutoGen team execution failed because the configured AutoGen model endpoint rejected an internal "
                "multi-agent prompt with 'No user query found in messages'. This usually means the current chat "
                "template requires a trailing user message and is not compatible with AutoGen's internal "
                f"assistant/tool turns. Current settings: AUTOGEN_OPENROUTER_BASE_URL={autogen_base or '(unset)'}, "
                f"AUTOGEN_TEAM_MODEL={autogen_model or '(unset)'}. Fix this by either pointing AutoGen to a "
                "compatible hosted endpoint such as https://openrouter.ai/api/v1 with a stable chat model, or by "
                "switching the local model/server to one whose chat template supports assistant continuation and tool "
                "calling. Restart the proxy server after changing .env."
            )
        return None


def _patch_autogen_team_config_for_model_compatibility(team_config: Dict[str, Any]) -> List[str]:
    """Patch known provider incompatibilities in exported AutoGen team configs."""
    patched_agents: List[str] = []
    config = team_config.get("config") if isinstance(team_config, dict) else None
    if not isinstance(config, dict):
        return patched_agents

    def _patch_model_client_name_compat(model_client: Any) -> bool:
        if not isinstance(model_client, dict):
            return False
        provider = str(model_client.get("provider") or "")
        if not provider.endswith("OpenAIChatCompletionClient"):
            return False
        model_client_config = model_client.get("config")
        if not isinstance(model_client_config, dict):
            return False
        base_url = str(model_client_config.get("base_url") or "")
        model_name = str(model_client_config.get("model") or "")
        if not is_minimax_chat_request(base_url, model_name):
            return False
        changed = False
        if model_client_config.get("include_name_in_message") is not False:
            model_client_config["include_name_in_message"] = False
            changed = True
        if model_client_config.get("add_name_prefixes") is not True:
            model_client_config["add_name_prefixes"] = True
            changed = True
        return changed

    selector_model_client = config.get("model_client")
    if _patch_model_client_name_compat(selector_model_client):
        patched_agents.append("selector_model_client")

    participants = config.get("participants")
    if not isinstance(participants, list):
        return patched_agents

    for participant in participants:
        if not isinstance(participant, dict):
            continue
        participant_config = participant.get("config")
        if not isinstance(participant_config, dict):
            continue
        participant_name = str(participant_config.get("name") or participant.get("label") or "unknown_agent")

        model_client = participant_config.get("model_client")
        if _patch_model_client_name_compat(model_client):
            patched_agents.append(f"{participant_name}:minimax_name_compat")

        if participant_config.get("reflect_on_tool_use") is not True:
            continue

        workbench = participant_config.get("workbench")
        if not isinstance(workbench, list) or not workbench:
            continue

        if not isinstance(model_client, dict):
            continue
        provider = str(model_client.get("provider") or "")
        model_client_config = model_client.get("config")
        if not isinstance(model_client_config, dict):
            continue
        base_url = str(model_client_config.get("base_url") or "").lower()

        # OpenRouter + AssistantAgent tool reflection is currently unstable for this team:
        # after a tool call, the reflection turn may come back without text, which AutoGen
        # raises as "Reflect on tool use produced no valid text response."
        if provider.endswith("OpenAIChatCompletionClient") and "openrouter.ai" in base_url:
            participant_config["reflect_on_tool_use"] = False
            patched_agents.append(participant_name)

    return patched_agents


def _autogen_team_definition_mtime() -> float:
    """Return the latest mtime for the active AutoGen team definition sources."""
    mtimes: List[float] = []
    if AUTOGEN_TEAM_BUILDER_FILE.exists():
        mtimes.append(AUTOGEN_TEAM_BUILDER_FILE.stat().st_mtime)
    if TEAM_CONFIG_FILE.exists():
        mtimes.append(TEAM_CONFIG_FILE.stat().st_mtime)
    return max(mtimes) if mtimes else 0.0


def load_autogen_team_runtime():
    """Load the AutoGen team from the Python builder, with JSON fallback."""
    if not AUTOGEN_AVAILABLE:
        print("AutoGen not available, skipping runtime team load")
        return None

    try:
        team = None
        team_label = "VirtualProductCompany"

        if AUTOGEN_TEAM_BUILDER_FILE.exists():
            print(f"Loading AutoGen team from Python builder {AUTOGEN_TEAM_BUILDER_FILE}...")
            from src.autogen.team_builder import (
                build_virtual_product_company_team,
                export_virtual_product_company_team_config,
            )

            team = build_virtual_product_company_team()
            export_virtual_product_company_team_config()
            team_label = getattr(team, "name", None) or getattr(team, "_name", None) or team_label
        elif TEAM_CONFIG_FILE.exists():
            team = load_autogen_team()
            team_label = getattr(team, "name", None) or getattr(team, "_name", None) or team_label
        else:
            print(
                "AutoGen team definition not found: "
                f"{AUTOGEN_TEAM_BUILDER_FILE} or {TEAM_CONFIG_FILE}"
            )
            return None

        if team is None:
            return None

        if AUTOGEN_CODE_EXEC_AVAILABLE and PythonCodeExecutionTool and DockerCommandLineCodeExecutor:
            coding_dir = _PROJECT_ROOT / "coding"
            try:
                coding_dir.mkdir(exist_ok=True)
            except OSError:
                pass
            work_dir = str(coding_dir)
            executor = DockerCommandLineCodeExecutor(
                work_dir=work_dir,
                image="python:3.12-slim",
                timeout=120,
            )
            code_tool = PythonCodeExecutionTool(executor)
            participants = getattr(team, "participants", None) or getattr(team, "_participants", [])
            if participants:
                target_agent = None
                for participant in participants:
                    candidate_name = getattr(participant, "name", None) or getattr(participant, "_name", None)
                    if candidate_name == "lead_engineer_agent":
                        target_agent = participant
                        break
                if target_agent is None:
                    target_agent = participants[0]

                wb = getattr(target_agent, "workbench", getattr(target_agent, "_workbench", None))
                if wb is not None:
                    wb_list = wb if isinstance(wb, list) else [wb]
                    for wb_item in wb_list:
                        tools = getattr(wb_item, "tools", getattr(wb_item, "_tools", None))
                        if tools is not None and isinstance(tools, list):
                            tool_names = [getattr(t, "name", "") for t in tools]
                            if "python_code_execution" not in tool_names:
                                tools.insert(0, code_tool)
                                print("Injected PythonCodeExecutionTool (Docker) into lead_engineer_agent workbench")
                            break

        try:
            team._config_mtime = _autogen_team_definition_mtime()
        except Exception:
            pass

        print(f"AutoGen team loaded successfully: {team_label}")
        return team
    except Exception as e:
        import traceback

        print(f"Error loading runtime AutoGen team: {e}")
        print(traceback.format_exc())
        return None


async def _start_code_executors(team: Any) -> None:
    """Start any code executors (Docker/local) attached to team participants or their tools."""
    if not AUTOGEN_AVAILABLE or team is None:
        return
    participants = getattr(team, "participants", None) or getattr(team, "_participants", [])
    for agent in participants or []:
        # CodeExecutorAgent has _code_executor
        executor = getattr(agent, "_code_executor", None)
        if executor is not None and hasattr(executor, "start"):
            try:
                await executor.start()
            except Exception as e:
                print(f"âš ï¸ Code executor start warning: {e}")
        # AssistantAgent workbench tools may have .executor (e.g. PythonCodeExecutionTool)
        wb = getattr(agent, "workbench", getattr(agent, "_workbench", None))
        if wb is not None:
            wb_list = wb if isinstance(wb, list) else [wb]
            for wb_item in wb_list:
                tools = getattr(wb_item, "tools", getattr(wb_item, "_tools", None))
                if tools is not None:
                    for t in tools:
                        ex = getattr(t, "executor", None) or getattr(t, "_executor", None)
                        if ex is not None and hasattr(ex, "start"):
                            try:
                                await ex.start()
                            except Exception as e:
                                print(f"âš ï¸ Tool executor start warning: {e}")


async def _stop_code_executors(team: Any) -> None:
    """Stop any code executors attached to team participants or their tools."""
    if not AUTOGEN_AVAILABLE or team is None:
        return
    participants = getattr(team, "participants", None) or getattr(team, "_participants", [])
    for agent in participants or []:
        executor = getattr(agent, "_code_executor", None)
        if executor is not None and hasattr(executor, "stop"):
            try:
                await executor.stop()
            except Exception as e:
                print(f"âš ï¸ Code executor stop warning: {e}")
        wb = getattr(agent, "workbench", getattr(agent, "_workbench", None))
        if wb is not None:
            wb_list = wb if isinstance(wb, list) else [wb]
            for wb_item in wb_list:
                tools = getattr(wb_item, "tools", getattr(wb_item, "_tools", None))
                if tools is not None:
                    for t in tools:
                        ex = getattr(t, "executor", None) or getattr(t, "_executor", None)
                        if ex is not None and hasattr(ex, "stop"):
                            try:
                                await ex.stop()
                            except Exception as e:
                                print(f"âš ï¸ Tool executor stop warning: {e}")


# Load servers on startup (with error handling to prevent startup failures)
try:
    load_servers()
    print(f"✅ Loaded {len(mcp_servers)} MCP servers from disk")
except Exception as e:
    import traceback
    print(f"âš ï¸ Warning: Could not load servers on startup: {e}")
    print(traceback.format_exc())
    # Continue anyway - server should still work without pre-loaded servers

# Load AutoGen team on startup (with error handling to prevent startup failures)
try:
    autogen_team = load_autogen_team_runtime()
    if autogen_team is not None:
        print("✅ AutoGen team loaded successfully on startup")
except Exception as e:
    import traceback
    print(f"âš ï¸ Warning: Could not load AutoGen team on startup: {e}")
    print(traceback.format_exc())
    # Continue anyway - server should still work without AutoGen team
    autogen_team = None

# Web proxy endpoint for fetching content (GET for backward compat, POST for iOS Safari / long URLs)
def _is_dns_or_network_error(exc: BaseException) -> bool:
    """True if the exception is DNS (getaddrinfo) or network unreachable."""
    if isinstance(exc, socket.gaierror):
        return True
    # Windows: OSError can have winerror 11002 (WSAHOST_NOT_FOUND)
    if isinstance(exc, OSError) and getattr(exc, "winerror", None) == 11002:
        return True
    # Check wrapped cause (e.g. httpx wraps socket errors)
    cause = getattr(exc, "__cause__", None)
    if cause and _is_dns_or_network_error(cause):
        return True
    errstr = str(exc).lower()
    if "getaddrinfo failed" in errstr or "name or service not known" in errstr or "nodename nor servname" in errstr or "errno 11002" in errstr:
        return True
    return False


def _coerce_bool(value: Any, default: bool = False) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    if isinstance(value, str):
        lowered = value.strip().lower()
        if lowered in {"1", "true", "yes", "y", "on"}:
            return True
        if lowered in {"0", "false", "no", "n", "off"}:
            return False
    return default


def _coerce_bounded_int(
    value: Any,
    *,
    default: int,
    minimum: int = 0,
    maximum: Optional[int] = None,
) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = default
    if maximum is not None:
        parsed = min(parsed, maximum)
    return max(minimum, parsed)


def _normalize_render_engine(value: Optional[str]) -> str:
    engine = str(value).strip().lower() if value is not None else "auto"
    if engine not in {"auto", "playwright", "selenium"}:
        raise HTTPException(
            status_code=400,
            detail="render_engine must be one of: auto, playwright, selenium",
        )
    return engine


def _normalize_js_wait_ms(value: Any) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = 2200
    return max(0, min(parsed, 20000))


async def _render_page_playwright(
    *,
    url: str,
    headers: Dict[str, str],
    timeout_seconds: float,
    wait_for_selector: Optional[str],
    js_wait_ms: int,
) -> Dict[str, Any]:
    if not PLAYWRIGHT_FETCH_AVAILABLE or async_playwright is None:
        raise RuntimeError(
            "Playwright is not installed. Install with: pip install playwright && playwright install chromium"
        )

    browser = None
    context = None
    timeout_ms = max(1000, int(timeout_seconds * 1000))
    selector = (wait_for_selector or "").strip() or None
    try:
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            context = await browser.new_context(
                user_agent=headers.get("User-Agent", "CATBot/1.0"),
                extra_http_headers={k: v for k, v in headers.items() if k.lower() != "user-agent"},
            )
            page = await context.new_page()
            response = await page.goto(url, wait_until="domcontentloaded", timeout=timeout_ms)
            if selector:
                await page.wait_for_selector(selector, timeout=timeout_ms)
            else:
                with suppress(Exception):
                    await page.wait_for_load_state("networkidle", timeout=min(timeout_ms, 9000))
            if js_wait_ms > 0:
                await page.wait_for_timeout(js_wait_ms)

            return {
                "raw_html": await page.content(),
                "final_url": page.url or url,
                "status_code": response.status if response else None,
                "render_engine": "playwright",
            }
    except Exception as exc:
        raise RuntimeError(f"Playwright rendering failed: {exc}") from exc
    finally:
        if context is not None:
            with suppress(Exception):
                await context.close()
        if browser is not None:
            with suppress(Exception):
                await browser.close()


def _render_page_selenium_sync(
    *,
    url: str,
    timeout_seconds: float,
    wait_for_selector: Optional[str],
    js_wait_ms: int,
) -> Dict[str, Any]:
    if not SELENIUM_FETCH_AVAILABLE or webdriver is None or SeleniumChromeOptions is None:
        raise RuntimeError(
            "Selenium is not installed. Install with: pip install selenium and ensure ChromeDriver is available."
        )

    options = SeleniumChromeOptions()
    options.add_argument("--headless=new")
    options.add_argument("--disable-gpu")
    options.add_argument("--window-size=1366,768")
    options.add_argument("--no-sandbox")
    options.add_argument("--disable-dev-shm-usage")

    selector = (wait_for_selector or "").strip() or None
    driver = webdriver.Chrome(options=options)
    try:
        driver.set_page_load_timeout(max(5, int(timeout_seconds)))
        driver.get(url)
        if selector:
            SeleniumWebDriverWait(driver, max(1, int(timeout_seconds))).until(
                SeleniumExpectedConditions.presence_of_element_located((SeleniumBy.CSS_SELECTOR, selector))
            )
        if js_wait_ms > 0:
            time.sleep(js_wait_ms / 1000.0)
        return {
            "raw_html": driver.page_source or "",
            "final_url": driver.current_url or url,
            "status_code": None,
            "render_engine": "selenium",
        }
    except Exception as exc:
        raise RuntimeError(f"Selenium rendering failed: {exc}") from exc
    finally:
        with suppress(Exception):
            driver.quit()


async def _render_page_selenium(
    *,
    url: str,
    timeout_seconds: float,
    wait_for_selector: Optional[str],
    js_wait_ms: int,
) -> Dict[str, Any]:
    return await asyncio.to_thread(
        _render_page_selenium_sync,
        url=url,
        timeout_seconds=timeout_seconds,
        wait_for_selector=wait_for_selector,
        js_wait_ms=js_wait_ms,
    )


async def _render_page_dynamic(
    *,
    url: str,
    headers: Dict[str, str],
    timeout_seconds: float,
    render_engine: str,
    wait_for_selector: Optional[str],
    js_wait_ms: int,
) -> Dict[str, Any]:
    engine = _normalize_render_engine(render_engine)
    candidates = [engine] if engine != "auto" else ["playwright", "selenium"]
    errors: List[str] = []

    for candidate in candidates:
        try:
            if candidate == "playwright":
                return await _render_page_playwright(
                    url=url,
                    headers=headers,
                    timeout_seconds=timeout_seconds,
                    wait_for_selector=wait_for_selector,
                    js_wait_ms=js_wait_ms,
                )
            return await _render_page_selenium(
                url=url,
                timeout_seconds=timeout_seconds,
                wait_for_selector=wait_for_selector,
                js_wait_ms=js_wait_ms,
            )
        except Exception as exc:
            errors.append(f"{candidate}: {exc}")

    raise HTTPException(
        status_code=503,
        detail=(
            "JavaScript rendering failed. "
            "Install Playwright or Selenium on the proxy host. "
            + " | ".join(errors[:2])
        ),
    )


async def _do_proxy_fetch(
    url: str,
    crawl: bool = True,
    max_pages: int = 3,
    max_depth: int = 1,
    render_js: bool = False,
    render_engine: str = "auto",
    wait_for_selector: Optional[str] = None,
    js_wait_ms: int = 2200,
) -> Dict[str, Any]:
    """Shared fetch logic: fetch URL(s), extract readable content, and optionally crawl."""
    if not url or not url.strip():
        raise HTTPException(status_code=400, detail="URL parameter is required")
    return await _fetch_and_extract_content(
        url=url,
        crawl=bool(crawl),
        max_pages=max_pages,
        max_depth=max_depth,
        render_js=_coerce_bool(render_js),
        render_engine=_normalize_render_engine(render_engine),
        wait_for_selector=(wait_for_selector or "").strip() or None,
        js_wait_ms=_normalize_js_wait_ms(js_wait_ms),
    )


def _normalize_url(raw_url: str) -> str:
    url = (raw_url or "").strip()
    if not url.startswith(("http://", "https://")):
        url = "https://" + url
    return url


def _is_same_domain(base_url: str, candidate_url: str) -> bool:
    try:
        base_host = (urlparse(base_url).hostname or "").lower()
        cand_host = (urlparse(candidate_url).hostname or "").lower()
        return bool(base_host) and base_host == cand_host
    except Exception:
        return False


def _extract_links_from_html(raw_html: str, base_url: str) -> List[str]:
    if not raw_html:
        return []
    links: List[str] = []
    seen: Set[str] = set()
    for match in re.finditer(r"""<a\s+[^>]*href=["']([^"']+)["']""", raw_html, re.IGNORECASE):
        href = (match.group(1) or "").strip()
        if not href or href.startswith(("#", "javascript:", "mailto:", "tel:")):
            continue
        absolute = urljoin(base_url, href)
        parsed = urlparse(absolute)
        if parsed.scheme not in {"http", "https"}:
            continue
        # Normalize by dropping fragments only.
        normalized = urlunparse((parsed.scheme, parsed.netloc, parsed.path, parsed.params, parsed.query, ""))
        if normalized in seen:
            continue
        seen.add(normalized)
        links.append(normalized)
    return links


def _extract_text_bs4(raw_html: str) -> Tuple[str, str]:
    soup = BeautifulSoup(raw_html, "html.parser")

    for tag in soup(["script", "style", "noscript", "svg", "canvas", "iframe", "form"]):
        tag.decompose()

    main = (
        soup.find("article")
        or soup.find("main")
        or soup.find(attrs={"role": "main"})
        or soup.find("body")
        or soup
    )

    for tag in main.find_all(["nav", "header", "footer", "aside"]):
        tag.decompose()

    title = (soup.title.get_text(" ", strip=True) if soup.title else "").strip()
    text = main.get_text(" ", strip=True)
    text = re.sub(r"\s+", " ", text).strip()
    return title, text


def _extract_text_fallback(raw_html: str) -> Tuple[str, str]:
    title_match = re.search(r"<title[^>]*>(.*?)</title>", raw_html, re.IGNORECASE | re.DOTALL)
    title = html.unescape(re.sub(r"<[^>]+>", " ", title_match.group(1))).strip() if title_match else ""

    cleaned = re.sub(r"<!--.*?-->", " ", raw_html, flags=re.DOTALL)
    cleaned = re.sub(r"<(script|style|noscript|svg|canvas|iframe|form)\b[^>]*>.*?</\1>", " ", cleaned, flags=re.IGNORECASE | re.DOTALL)

    # Prefer content-heavy semantic containers before full body fallback.
    for pattern in [
        r"<article\b[^>]*>(.*?)</article>",
        r"<main\b[^>]*>(.*?)</main>",
        r"<body\b[^>]*>(.*?)</body>",
    ]:
        m = re.search(pattern, cleaned, flags=re.IGNORECASE | re.DOTALL)
        if m:
            cleaned = m.group(1)
            break

    cleaned = re.sub(r"<(nav|header|footer|aside)\b[^>]*>.*?</\1>", " ", cleaned, flags=re.IGNORECASE | re.DOTALL)
    cleaned = re.sub(r"<[^>]+>", " ", cleaned)
    text = html.unescape(cleaned)
    text = re.sub(r"\s+", " ", text).strip()
    return title, text


def _extract_readable_content(raw_html: str) -> Tuple[str, str]:
    if BS4_AVAILABLE and BeautifulSoup is not None:
        try:
            return _extract_text_bs4(raw_html)
        except Exception:
            pass
    return _extract_text_fallback(raw_html)


async def _fetch_and_extract_content(
    url: str,
    crawl: bool = True,
    max_pages: int = 3,
    max_depth: int = 1,
    render_js: bool = False,
    render_engine: str = "auto",
    wait_for_selector: Optional[str] = None,
    js_wait_ms: int = 2200,
) -> Dict[str, Any]:
    normalized_start = _normalize_url(url)
    max_pages = max(1, min(int(max_pages or 1), 10))
    max_depth = max(0, min(int(max_depth or 0), 2))

    headers = {
        "User-Agent": "Mozilla/5.0 (iPhone; CPU iPhone OS 17_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.0 Mobile/15E148 Safari/604.1",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
        "Connection": "keep-alive",
    }

    queue = collections.deque([(normalized_start, 0)])
    visited: Set[str] = set()
    pages: List[Dict[str, Any]] = []
    last_raw_html = ""
    last_error: Optional[BaseException] = None
    rendered_engine_used: Optional[str] = None

    try:
        async with httpx.AsyncClient(timeout=15.0, follow_redirects=True) as client:
            while queue and len(pages) < max_pages:
                current_url, depth = queue.popleft()
                if current_url in visited:
                    continue
                visited.add(current_url)

                try:
                    should_render_this_page = bool(render_js) and depth == 0
                    if should_render_this_page:
                        rendered = await _render_page_dynamic(
                            url=current_url,
                            headers=headers,
                            timeout_seconds=15.0,
                            render_engine=render_engine,
                            wait_for_selector=wait_for_selector,
                            js_wait_ms=js_wait_ms,
                        )
                        raw_html = rendered.get("raw_html") or ""
                        final_url = rendered.get("final_url") or current_url
                        content_type = "text/html"
                        rendered_engine_used = rendered.get("render_engine") or rendered_engine_used
                    else:
                        response = await client.get(current_url, headers=headers)
                        response.raise_for_status()
                        raw_html = response.text
                        final_url = str(response.url)
                        content_type = (response.headers.get("content-type") or "").lower()

                    if "text/html" not in content_type and "<html" not in raw_html[:3000].lower():
                        # Non-HTML content is still captured as plain text.
                        extracted_title = ""
                        extracted_text = raw_html.strip()
                    else:
                        extracted_title, extracted_text = _extract_readable_content(raw_html)

                    if extracted_text:
                        pages.append(
                            {
                                "url": final_url,
                                "title": extracted_title,
                                "content": extracted_text,
                                "rendered": should_render_this_page,
                                "render_engine": rendered_engine_used if should_render_this_page else None,
                            }
                        )
                        last_raw_html = raw_html

                    if crawl and depth < max_depth and len(pages) < max_pages:
                        for link in _extract_links_from_html(raw_html, final_url):
                            if link in visited:
                                continue
                            if not _is_same_domain(normalized_start, link):
                                continue
                            queue.append((link, depth + 1))
                except HTTPException as page_error:
                    # Surface rendering/config errors for the root page immediately.
                    if depth == 0:
                        raise
                    last_error = page_error
                    continue
                except Exception as page_error:
                    last_error = page_error
                    continue
    except HTTPException:
        raise
    except Exception as e:
        if _is_dns_or_network_error(e):
            raise HTTPException(
                status_code=502,
                detail=(
                    "The proxy server could not resolve the website's hostname (DNS lookup failed). "
                    "This usually means the machine running the proxy has no internet or restricted DNS. "
                    "Ensure the proxy runs on a machine with working internet and DNS (e.g. try pinging the host from that machine)."
                ),
            )
        raise HTTPException(status_code=500, detail=f"Failed to fetch content: {str(e)}")

    if not pages:
        if isinstance(last_error, HTTPException):
            raise last_error
        if last_error and _is_dns_or_network_error(last_error):
            raise HTTPException(
                status_code=502,
                detail=(
                    "The proxy server could not resolve the website's hostname (DNS lookup failed). "
                    "This usually means the machine running the proxy has no internet or restricted DNS. "
                    "Ensure the proxy runs on a machine with working internet and DNS (e.g. try pinging the host from that machine)."
                ),
            )
        if last_error:
            raise HTTPException(status_code=500, detail=f"Failed to fetch content: {str(last_error)}")
        raise HTTPException(status_code=500, detail="Failed to fetch content: no pages returned")

    combined_parts: List[str] = []
    for i, page in enumerate(pages, start=1):
        page_title = (page.get("title") or "").strip()
        title_line = f"Title: {page_title}\n" if page_title else ""
        combined_parts.append(f"[Page {i}] {page.get('url')}\n{title_line}{page.get('content', '')}")
    combined_content = "\n\n".join(combined_parts).strip()

    return {
        "url": pages[0].get("url", normalized_start),
        "content": combined_content,
        "title": pages[0].get("title", ""),
        "pages": pages,
        "crawled": bool(crawl),
        "page_count": len(pages),
        "raw_html": last_raw_html,
        "render_js": bool(render_js),
        "render_engine": rendered_engine_used if render_js else None,
    }


@app.get("/v1/proxy/fetch")
async def proxy_fetch_get(
    url: str,
    request: Request,
    crawl: bool = True,
    max_pages: int = 3,
    max_depth: int = 1,
    render_js: bool = False,
    render_engine: str = "auto",
    wait_for_selector: Optional[str] = None,
    js_wait_ms: int = 2200,
):
    """Fetch web content via GET (query param). Use POST for long URLs (e.g. iOS Safari)."""
    try:
        result = await _do_proxy_fetch(
            url,
            crawl=crawl,
            max_pages=max_pages,
            max_depth=max_depth,
            render_js=render_js,
            render_engine=render_engine,
            wait_for_selector=wait_for_selector,
            js_wait_ms=js_wait_ms,
        )
        cors = build_cors_headers(request)
        return JSONResponse(content=result, headers=cors)
    except HTTPException:
        raise
    except Exception as e:
        print(f"Proxy fetch error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch content: {str(e)}")


@app.post("/v1/proxy/fetch")
async def proxy_fetch_post(body: ProxyFetchRequest, request: Request):
    """Fetch web content via POST body. Use url (single) or urls (list); if urls, try each until one succeeds."""
    # Build list: prefer urls if non-empty, else single url
    to_try: List[str] = []
    if body.urls:
        to_try = [u.strip() for u in body.urls if u and isinstance(u, str) and u.strip()]
    if not to_try and body.url:
        to_try = [body.url.strip()]
    if not to_try:
        raise HTTPException(status_code=400, detail="Either 'url' or 'urls' is required")
    last_error: Optional[Exception] = None
    for one_url in to_try:
        try:
            result = await _do_proxy_fetch(
                one_url,
                crawl=body.crawl,
                max_pages=body.max_pages,
                max_depth=body.max_depth,
                render_js=body.render_js,
                render_engine=body.render_engine,
                wait_for_selector=body.wait_for_selector,
                js_wait_ms=body.js_wait_ms,
            )
            cors = build_cors_headers(request)
            return JSONResponse(content=result, headers=cors)
        except HTTPException as e:
            # Retry next candidate URL for server-side fetch failures.
            if len(to_try) == 1 or e.status_code == 400:
                raise
            last_error = e
            print(f"Proxy fetch failed for {one_url[:60]}...: {e.detail}")
            continue
        except Exception as e:
            last_error = e
            print(f"Proxy fetch failed for {one_url[:60]}...: {e}")
            continue
    if last_error:
        if isinstance(last_error, HTTPException):
            raise HTTPException(status_code=last_error.status_code, detail=last_error.detail)
        raise HTTPException(status_code=500, detail=f"Failed to fetch content: {str(last_error)}")
    raise HTTPException(status_code=400, detail="No URLs to try")

# Shared search logic for route and Telegram tool runner
async def _do_proxy_search(query: str) -> Dict[str, Any]:
    """Search the web using Brave Search API or DuckDuckGo fallback. Raises HTTPException on failure."""
    if not query:
        raise HTTPException(status_code=400, detail="Search query is required")
    brave_api_key = os.getenv('BRAVE_API_KEY')
    if not brave_api_key:
        print("âš ï¸  BRAVE_API_KEY not configured. Falling back to DuckDuckGo.")
    else:
        try:
            print(f"ðŸ” Using Brave Search API for query: {query[:50]}...")
            async with httpx.AsyncClient(timeout=15.0) as client:
                response = await client.get(
                    'https://api.search.brave.com/res/v1/web/search',
                    headers={
                        'Accept': 'application/json',
                        'Accept-Encoding': 'gzip',
                        'X-Subscription-Token': brave_api_key
                    },
                    params={
                        'q': query,
                        'count': 10,
                        'search_lang': 'en',
                        'safesearch': 'moderate',
                        'freshness': 'past_month'
                    }
                )

            if response.status_code == 200:
                data = response.json()
                if data.get('web', {}).get('results'):
                    results = []
                    for result in data['web']['results']:
                        date_str = result.get('age') or result.get('published')
                        parsed_date = parse_date(date_str) if date_str else None
                        results.append({
                            'url': result['url'],
                            'title': clean_text(result.get('title', '')),
                            'snippet': clean_text(result.get('description', '')),
                            'date': parsed_date
                        })
                    results = [r for r in results if r['title'] and r['snippet']]
                    results.sort(key=lambda x: parse_date(x.get('date')) or 0, reverse=True)
                    print(f"✅ Brave Search returned {len(results)} results")
                    return {"results": results[:5], "source": "brave"}
                else:
                    print(f"âš ï¸  Brave Search returned no results in response")
            elif response.status_code == 401:
                print(f"âŒ Brave Search API authentication failed (401). Check your BRAVE_API_KEY.")
                raise HTTPException(
                    status_code=500,
                    detail="Brave Search API authentication failed. Please check your BRAVE_API_KEY configuration."
                )
            elif response.status_code == 429:
                print(f"âš ï¸  Brave Search API rate limit exceeded (429). Falling back to DuckDuckGo.")
            else:
                print(f"âš ï¸  Brave Search API returned status {response.status_code}. Falling back to DuckDuckGo.")
                try:
                    error_data = response.json()
                    print(f"   Error details: {error_data}")
                except Exception:
                    print(f"   Error text: {response.text[:200]}")

        except httpx.RequestError as e:
            error_msg = str(e) if str(e) else f"Network error: {type(e).__name__}"
            print(f"âŒ Brave Search network error: {error_msg}")
            print("   Falling back to DuckDuckGo...")
        except httpx.HTTPStatusError as e:
            print(f"âŒ Brave Search HTTP error: {e.response.status_code if e.response else 'Unknown'}")
            print("   Falling back to DuckDuckGo...")
        except HTTPException:
            raise
        except Exception as e:
            error_msg = str(e) if str(e) else f"Unknown error: {type(e).__name__}"
            print(f"âŒ Brave Search failed: {error_msg}")
            print("   Falling back to DuckDuckGo...")
            import traceback
            print(traceback.format_exc())

    print("🦆 Falling back to DuckDuckGo search...")
    try:
        search_url = f"https://html.duckduckgo.com/html/?q={query}"
        async with httpx.AsyncClient(timeout=15.0) as client:
            response = await client.get(
                search_url,
                headers={
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                    'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
                    'Accept-Language': 'en-US,en;q=0.5',
                    'Referer': 'https://duckduckgo.com/'
                },
                follow_redirects=True
            )
            if response.status_code != 200:
                raise HTTPException(
                    status_code=500,
                    detail=f"DuckDuckGo search returned HTTP {response.status_code}. The search service may be temporarily unavailable."
                )
        results = []
        html = response.text
        patterns = [
            r'<div class="links_main links_deep result__body">.*?<a class="result__a" href="([^"]+)".*?>(.*?)</a>.*?<a class="result__snippet".*?>(.*?)</a>',
            r'<div class="result__body">.*?<a class="result__url" href="([^"]+)".*?>(.*?)</a>.*?<div class="result__snippet">(.*?)</div>',
            r'<div class="result__body">.*?<a class="result__a" href="([^"]+)".*?>(.*?)</a>.*?<div class="result__snippet">(.*?)</div>',
            r'<a[^>]*class="[^"]*result[^"]*"[^>]*href="([^"]+)"[^>]*>(.*?)</a>.*?<div[^>]*class="[^"]*snippet[^"]*"[^>]*>(.*?)</div>',
            r'<a[^>]*href="([^"]+)"[^>]*class="[^"]*result__a[^"]*"[^>]*>(.*?)</a>.*?<span[^>]*class="[^"]*result__snippet[^"]*"[^>]*>(.*?)</span>'
        ]
        for pattern in patterns:
            matches = re.finditer(pattern, html, re.DOTALL)
            for match in matches:
                if len(results) >= 5:
                    break
                try:
                    url, title, snippet = match.groups()
                    url = url.replace('&amp;', '&')
                    title = clean_text(title)
                    snippet = clean_text(snippet)
                    if url and 'duckduckgo.com' not in url and title and snippet:
                        results.append({'url': url, 'title': title, 'snippet': snippet})
                except (ValueError, IndexError):
                    continue
            if len(results) >= 5:
                break
        if len(results) == 0:
            print(f"âš ï¸  DuckDuckGo search: No results found with regex patterns. HTML preview: {html[:1000]}")
            return {"results": [], "source": "duckduckgo", "message": "No results found. DuckDuckGo HTML structure may have changed."}
        print(f"✅ DuckDuckGo returned {len(results)} results")
        return {"results": results, "source": "duckduckgo"}
    except httpx.RequestError as e:
        error_msg = str(e) if str(e) else f"Network error: {type(e).__name__}"
        print(f"Search error (network): {error_msg}")
        raise HTTPException(status_code=500, detail=f"Failed to perform search: Network error - {error_msg}")
    except httpx.HTTPStatusError as e:
        error_msg = f"HTTP {e.response.status_code}: {e.response.text[:200]}" if e.response else str(e)
        print(f"Search error (HTTP): {error_msg}")
        raise HTTPException(status_code=500, detail=f"Failed to perform search: {error_msg}")
    except Exception as e:
        error_msg = str(e) if str(e) else f"Unknown error: {type(e).__name__}"
        print(f"Search error: {error_msg}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to perform search: {error_msg}")


@app.get("/v1/proxy/search")
async def proxy_search(query: str):
    """Search the web using Brave Search API or DuckDuckGo fallback."""
    return await _do_proxy_search(query)


# Shared news logic for route and Telegram tool runner
async def _do_proxy_news(query: str) -> Dict[str, Any]:
    """Fetch news articles for query. Raises HTTPException on failure."""
    if not query:
        raise HTTPException(status_code=400, detail="Search query is required")
    news_api_key = os.getenv('NEWS_API_KEY')
    if not news_api_key:
        raise HTTPException(
            status_code=503,
            detail="NEWS_API_KEY is not configured. Please set it in your .env file."
        )
    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            response = await client.get(
                'https://newsapi.org/v2/everything',
                headers={'Accept': 'application/json'},
                params={
                    'q': query,
                    'apiKey': news_api_key,
                    'sortBy': 'publishedAt',
                    'language': 'en',
                    'pageSize': 100
                }
            )
        if response.status_code == 200:
            data = response.json()
            articles = data.get('articles', [])
            if not articles:
                return {"success": False, "message": f"No articles found for \"{query}\"", "articles": []}
            formatted_articles = []
            for article in articles:
                formatted_articles.append({
                    'title': clean_text(article.get('title', '')),
                    'url': article.get('url', ''),
                    'description': clean_text(article.get('description', '')),
                    'publishedAt': article.get('publishedAt', ''),
                    'source': article.get('source', {}).get('name', 'Unknown')
                })
            return {
                "success": True,
                "message": f"Found {len(formatted_articles)} articles",
                "articles": formatted_articles,
                "totalResults": data.get('totalResults', len(formatted_articles))
            }
        error_data = response.json() if response.headers.get('content-type', '').startswith('application/json') else {}
        error_message = error_data.get('message', f"News API returned status {response.status_code}")
        print(f"News API error: {error_message}")
        raise HTTPException(status_code=response.status_code, detail=f"News API error: {error_message}")
    except httpx.HTTPStatusError as e:
        print(f"News API HTTP error: {e}")
        raise HTTPException(status_code=e.response.status_code, detail=f"News API request failed: {str(e)}")
    except Exception as e:
        print(f"News API error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch news: {str(e)}")


@app.get("/v1/proxy/news")
async def proxy_news_search(query: str):
    """Search news articles using News API."""
    return await _do_proxy_news(query)


def _sanitize_weather_location(location: str) -> str:
    if location is None:
        return ""
    if not isinstance(location, str):
        location = str(location)
    return re.sub(r"\s+", " ", location.strip())


def _normalize_weather_detail(detail: Optional[str]) -> str:
    """Normalize weather detail/request type to supported values."""
    value = (detail or "summary")
    if not isinstance(value, str):
        value = str(value)
    key = value.strip().lower()
    aliases = {
        "overview": "summary",
        "all": "summary",
        "rain": "forecast",
        "wind": "current",
        "alerts": "summary",
    }
    return aliases.get(key, key if key in {"summary", "current", "forecast"} else "summary")


def _safe_float(value: Any) -> Optional[float]:
    try:
        if value is None or value == "":
            return None
        return float(value)
    except (TypeError, ValueError):
        return None


def _open_meteo_weather_code_to_text(code: Any) -> str:
    code_int: Optional[int]
    try:
        code_int = int(float(code))
    except (TypeError, ValueError):
        code_int = None
    mapping = {
        0: "clear sky",
        1: "mainly clear",
        2: "partly cloudy",
        3: "overcast",
        45: "fog",
        48: "depositing rime fog",
        51: "light drizzle",
        53: "moderate drizzle",
        55: "dense drizzle",
        56: "light freezing drizzle",
        57: "dense freezing drizzle",
        61: "slight rain",
        63: "moderate rain",
        65: "heavy rain",
        66: "light freezing rain",
        67: "heavy freezing rain",
        71: "slight snow fall",
        73: "moderate snow fall",
        75: "heavy snow fall",
        77: "snow grains",
        80: "slight rain showers",
        81: "moderate rain showers",
        82: "violent rain showers",
        85: "slight snow showers",
        86: "heavy snow showers",
        95: "thunderstorm",
        96: "thunderstorm with slight hail",
        99: "thunderstorm with heavy hail",
    }
    if code_int is None:
        return "unknown"
    return mapping.get(code_int, f"weather code {code_int}")


def _extract_memory_location(memories: List[Dict[str, Any]]) -> Optional[str]:
    if not memories:
        return None
    patterns = [
        r"(?:i(?:\s|'| a)?m|i live|i am based|my location is|i'm in)\s+(?:in\s+)?([A-Za-z\s,'-]{2,80})",
        r"\b([A-Za-z\s'-]{2,60},\s*(?:nsw|vic|qld|sa|wa|tas|nt|act))\b",
        r"\b(\d{4})\b",
    ]
    for m in memories:
        txt = (m.get("text") or "").strip()
        if not txt:
            continue
        for pat in patterns:
            mo = re.search(pat, txt, re.IGNORECASE)
            if mo:
                return _sanitize_weather_location(mo.group(1).strip(" .,"))
    return None


async def _resolve_weather_location(location: Optional[str], user_id: Optional[str], memory_manager: Any) -> Tuple[str, str]:
    requested = _sanitize_weather_location(location or "")
    if requested:
        return requested, "request"
    if memory_manager is not None and user_id:
        try:
            memories = await memory_manager.search_memories(
                query=f"{user_id} location city suburb postcode", limit=8, similarity_threshold=0.3
            )
            inferred = _extract_memory_location(memories)
            if inferred:
                return inferred, "memory"
        except Exception as e:
            print(f"[WEATHER] Memory lookup failed: {e}")
    raise HTTPException(status_code=400, detail="Location is required. Please provide a city/suburb/postcode.")


async def _do_proxy_weather(location: Optional[str], detail: str = "summary", user_id: Optional[str] = None, memory_manager: Any = None) -> Dict[str, Any]:
    resolved_location, source = await _resolve_weather_location(location, user_id, memory_manager)
    detail = _normalize_weather_detail(detail)
    parsed_geocoding_base = urlparse(OPEN_METEO_GEOCODING_BASE_URL)
    if not parsed_geocoding_base.scheme.startswith("http"):
        raise HTTPException(status_code=500, detail="Invalid OPEN_METEO_GEOCODING_BASE_URL configuration")
    if not parsed_geocoding_base.hostname or not parsed_geocoding_base.hostname.endswith(OPEN_METEO_ALLOWED_HOST_SUFFIX):
        raise HTTPException(status_code=500, detail="OPEN_METEO_GEOCODING_BASE_URL host is not allowlisted")
    parsed_forecast_base = urlparse(OPEN_METEO_FORECAST_BASE_URL)
    if not parsed_forecast_base.scheme.startswith("http"):
        raise HTTPException(status_code=500, detail="Invalid OPEN_METEO_FORECAST_BASE_URL configuration")
    if not parsed_forecast_base.hostname or not parsed_forecast_base.hostname.endswith(OPEN_METEO_ALLOWED_HOST_SUFFIX):
        raise HTTPException(status_code=500, detail="OPEN_METEO_FORECAST_BASE_URL host is not allowlisted")

    weather_headers = {"Accept": "application/json", "User-Agent": OPEN_METEO_USER_AGENT}
    try:
        async with httpx.AsyncClient(timeout=OPEN_METEO_TIMEOUT_SECONDS, trust_env=OPEN_METEO_TRUST_ENV, follow_redirects=True) as client:
            loc_resp = await client.get(
                OPEN_METEO_GEOCODING_BASE_URL,
                params={"name": resolved_location, "count": 10, "language": "en", "format": "json"},
                headers=weather_headers,
            )
            loc_resp.raise_for_status()
            try:
                loc_json = loc_resp.json()
            except ValueError as e:
                raise HTTPException(status_code=502, detail=f"Invalid Open-Meteo geocoding payload: {str(e)}")
            loc_items = loc_json.get("results") if isinstance(loc_json, dict) else []
            if not isinstance(loc_items, list) or not loc_items:
                raise HTTPException(status_code=404, detail=f"No Open-Meteo location found for '{resolved_location}'.")
            requested_lc = resolved_location.lower()
            loc = next(
                (
                    item for item in loc_items
                    if isinstance(item, dict) and str(item.get("name") or "").strip().lower() == requested_lc
                ),
                None,
            )
            if not loc:
                loc = next(
                    (
                        item for item in loc_items
                        if isinstance(item, dict) and str(item.get("name") or "").strip().lower().startswith(requested_lc)
                ),
                None,
            )
            if not loc:
                loc = next(
                    (
                        item for item in loc_items
                        if isinstance(item, dict) and requested_lc in " ".join(
                            str(item.get(part) or "").strip().lower()
                            for part in ("name", "admin1", "country", "country_code")
                        )
                    ),
                    None,
                )
            if not loc:
                loc = loc_items[0] if isinstance(loc_items[0], dict) else {}

            latitude = _safe_float(loc.get("latitude"))
            longitude = _safe_float(loc.get("longitude"))
            if latitude is None or longitude is None:
                raise HTTPException(status_code=502, detail="Open-Meteo geocoding response missing coordinates")

            forecast_resp = await client.get(
                OPEN_METEO_FORECAST_BASE_URL,
                params={
                    "latitude": latitude,
                    "longitude": longitude,
                    "current": "temperature_2m,apparent_temperature,relative_humidity_2m,wind_speed_10m,weather_code",
                    "daily": "weather_code,temperature_2m_max,temperature_2m_min,precipitation_probability_max",
                    "forecast_days": 7,
                    "timezone": "auto",
                },
                headers=weather_headers,
            )
            forecast_resp.raise_for_status()
            try:
                forecast_json = forecast_resp.json()
            except ValueError as e:
                raise HTTPException(status_code=502, detail=f"Invalid Open-Meteo weather payload: {str(e)}")
    except HTTPException:
        raise
    except httpx.ConnectError as e:
        hint = ""
        if not OPEN_METEO_TRUST_ENV:
            hint = " If you require an outbound proxy, set OPEN_METEO_TRUST_ENV=true and configure HTTP(S)_PROXY."
        raise HTTPException(status_code=502, detail=f"Could not connect to Open-Meteo service: {str(e)}.{hint}")
    except httpx.TimeoutException:
        raise HTTPException(status_code=504, detail="Timed out contacting Open-Meteo weather service")
    except httpx.HTTPStatusError as e:
        raise HTTPException(status_code=502, detail=f"Open-Meteo weather service returned status {e.response.status_code}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve weather data: {str(e)}")

    current_data = (forecast_json or {}).get("current") if isinstance(forecast_json, dict) else {}
    if not isinstance(current_data, dict):
        current_data = {}
    daily_data = (forecast_json or {}).get("daily") if isinstance(forecast_json, dict) else {}
    if not isinstance(daily_data, dict):
        daily_data = {}

    daily_dates = daily_data.get("time") if isinstance(daily_data.get("time"), list) else []
    daily_mins = daily_data.get("temperature_2m_min") if isinstance(daily_data.get("temperature_2m_min"), list) else []
    daily_maxs = daily_data.get("temperature_2m_max") if isinstance(daily_data.get("temperature_2m_max"), list) else []
    daily_rain = daily_data.get("precipitation_probability_max") if isinstance(daily_data.get("precipitation_probability_max"), list) else []
    daily_codes = daily_data.get("weather_code") if isinstance(daily_data.get("weather_code"), list) else []

    def _daily_value(values: Any, index: int) -> Any:
        if isinstance(values, list) and index < len(values):
            return values[index]
        return None

    current = {
        "temperature_c": _safe_float(current_data.get("temperature_2m")),
        "feels_like_c": _safe_float(current_data.get("apparent_temperature")),
        "humidity_pct": _safe_float(current_data.get("relative_humidity_2m")),
        "wind_kph": _safe_float(current_data.get("wind_speed_10m")),
        "condition": _open_meteo_weather_code_to_text(current_data.get("weather_code")),
        "observation_time": current_data.get("time") or "",
    }

    forecast = []
    for idx, day_date in enumerate(daily_dates[:7]):
        day_code = _daily_value(daily_codes, idx)
        forecast.append({
            "date": day_date,
            "min_c": _safe_float(_daily_value(daily_mins, idx)),
            "max_c": _safe_float(_daily_value(daily_maxs, idx)),
            "rain_chance_pct": _safe_float(_daily_value(daily_rain, idx)),
            "condition": _open_meteo_weather_code_to_text(day_code),
        })

    loc_name = loc.get("name") or resolved_location
    current_temp = current.get("temperature_c")
    current_temp_text = f"{current_temp}C" if current_temp is not None else "N/A"
    summary = (
        f"Weather for {loc_name}: {current_temp_text}, {current.get('condition', 'unknown')}"
        f". Forecast entries: {len(forecast)}."
    )

    payload = {
        "success": True,
        "summary": summary,
        "resolved_location": loc_name,
        "location_source": source,
        "current": current,
        "forecast": forecast,
        "detail": detail,
        "source": "open-meteo.com",
    }

    if detail == "current":
        payload["forecast"] = []
    elif detail == "forecast":
        payload["current"] = {}
    return payload


@app.get("/v1/proxy/weather")
async def proxy_weather(
    location: Optional[str] = None,
    detail: str = "summary",
    requestType: Optional[str] = None,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Fetch parsed weather information from Open-Meteo. Auth required."""
    user_id = current_user.get("username") if isinstance(current_user, dict) else None
    final_detail = _normalize_weather_detail(requestType or detail)
    return await _do_proxy_weather(location=location, detail=final_detail, user_id=user_id, memory_manager=memory_manager if MEMORY_AVAILABLE else None)


def _is_retryable_autogen_provider_error(exc: BaseException) -> bool:
    """True when the upstream OpenRouter/OpenAI-compatible client returned a malformed response."""
    text = str(exc)
    if "NoneType" in text and "subscriptable" in text:
        return True
    if "result.choices" in text:
        return True
    traceback_text = traceback.format_exc()
    return "_openai_client.py" in traceback_text and "choices[0]" in traceback_text


# Shared AutoGen logic for route and Telegram tool runner
async def _do_autogen(input_text: str) -> Dict[str, Any]:
    """Run AutoGen team with input_text. Returns dict with output/response/messages. Raises HTTPException on failure."""
    global autogen_team
    if not input_text:
        raise HTTPException(status_code=400, detail="Input parameter is required")
    monitor_run_id = _monitor_run_start("autogen", "team-run", input_text=input_text)
    progress_notes: List[str] = []
    log_filename: Optional[str] = None

    def _persist_autogen_log(
        *,
        status: str,
        messages: Optional[List[Dict[str, str]]] = None,
        conversation_summary: str = "",
        error_text: Optional[str] = None,
        suppress_monitor_note: bool = False,
    ) -> None:
        nonlocal log_filename
        try:
            log_filename = _write_autogen_conversation_to_scratch(
                input_text,
                messages or [],
                conversation_summary,
                filename=log_filename,
                status=status,
                progress_notes=progress_notes,
                error_text=error_text,
            )
            log_path = _resolve_monitor_log_path(log_filename)
            _monitor_run_update(
                monitor_run_id,
                log_file=log_filename,
                log_excerpt=_read_monitor_log_excerpt(log_path),
            )
        except Exception as log_err:
            if not suppress_monitor_note:
                _monitor_run_note(monitor_run_id, f"Failed to update AutoGen scratch log: {log_err}")
            print(f"[AUTOGEN] Failed to update conversation scratch log: {log_err}", flush=True)

    def _note_autogen_progress(note: str, *, conversation_summary: str = "AutoGen run in progress.") -> None:
        progress_notes.append(note)
        _monitor_run_note(monitor_run_id, note)
        _persist_autogen_log(status="running", conversation_summary=conversation_summary)

    _persist_autogen_log(
        status="running",
        conversation_summary="AutoGen run created.",
        suppress_monitor_note=True,
    )
    if not AUTOGEN_AVAILABLE:
        _persist_autogen_log(
            status="error",
            conversation_summary="AutoGen run failed before execution started.",
            error_text="AutoGen is not available on this server.",
        )
        _monitor_run_finish(
            monitor_run_id,
            status="error",
            summary="AutoGen is not available on this server.",
            metadata={"autogen_available": False},
            log_file=log_filename,
            log_excerpt=_read_monitor_log_excerpt((SCRATCH_DIR / log_filename) if log_filename else None),
        )
        raise HTTPException(
            status_code=503,
            detail="AutoGen not available. Please install: pip install autogen-agentchat autogen-ext"
        )
    if autogen_team is None:
        _note_autogen_progress("Loading AutoGen team.")
        print("🔄 Loading AutoGen team for the first time...")
        autogen_team = load_autogen_team_runtime()
        if autogen_team is None:
            _persist_autogen_log(
                status="error",
                conversation_summary="AutoGen run failed during team initialization.",
                error_text="AutoGen team could not be loaded.",
            )
            _monitor_run_finish(
                monitor_run_id,
                status="error",
                summary="AutoGen team could not be loaded.",
                log_file=log_filename,
                log_excerpt=_read_monitor_log_excerpt((SCRATCH_DIR / log_filename) if log_filename else None),
            )
            raise HTTPException(
                status_code=503,
                detail="AutoGen team not loaded. Check the Python builder or fallback JSON config."
            )
    try:
        config_mtime = _autogen_team_definition_mtime()
        if not hasattr(autogen_team, '_config_mtime') or autogen_team._config_mtime != config_mtime:
            _note_autogen_progress("Reloading AutoGen team after config change.")
            print("Reloading AutoGen team after definition change...")
            await _stop_code_executors(autogen_team)
            new_team = load_autogen_team_runtime()
            if new_team is not None:
                autogen_team = new_team
                autogen_team._config_mtime = config_mtime
                if hasattr(autogen_team, '_executors_started'):
                    delattr(autogen_team, '_executors_started')
    except Exception as e:
        print(f"âš ï¸  Error checking team config modification time: {e}")
        _note_autogen_progress(f"Config reload check warning: {e}")
    if not getattr(autogen_team, '_executors_started', False):
        _note_autogen_progress("Starting AutoGen code executors.")
        await _start_code_executors(autogen_team)
        try:
            autogen_team._executors_started = True
        except Exception:
            pass
    if hasattr(autogen_team, "reset"):
        _note_autogen_progress("Resetting AutoGen team state.")
        await autogen_team.reset()
    try:
        print(f"🚀 Running AutoGen team with input: {input_text[:100]}...")
        _note_autogen_progress("Running AutoGen team.")
        try:
            result = await autogen_team.run(task=input_text)
        except Exception as first_exc:
            if _is_retryable_autogen_provider_error(first_exc):
                _note_autogen_progress("Retrying AutoGen run after malformed provider response.")
                await _stop_code_executors(autogen_team)
                fresh_team = load_autogen_team_runtime()
                if fresh_team is None:
                    raise first_exc
                autogen_team = fresh_team
                if not getattr(autogen_team, '_executors_started', False):
                    _note_autogen_progress("Restarting AutoGen code executors for retry.")
                    await _start_code_executors(autogen_team)
                    try:
                        autogen_team._executors_started = True
                    except Exception:
                        pass
                if hasattr(autogen_team, "reset"):
                    _note_autogen_progress("Resetting AutoGen team state for retry.")
                    await autogen_team.reset()
                result = await autogen_team.run(task=input_text)
            else:
                raise
        messages = []
        if hasattr(result, 'messages'):
            messages = [
                {
                    "source": msg.source if hasattr(msg, 'source') else 'unknown',
                    "content": _stringify_autogen_message_content(msg.content) if hasattr(msg, 'content') else str(msg)
                }
                for msg in result.messages
            ]
        _note_autogen_progress(
            f"AutoGen returned {len(messages)} messages.",
            conversation_summary="AutoGen returned messages and is writing the full transcript.",
        )
        final_message = _stringify_autogen_message_content(messages[-1].get("content", "")) if messages else ""
        if final_message:
            conversation_summary = (
                f"Completed with {len(messages)} messages. "
                f"Final message from {messages[-1].get('source', 'unknown')}:\n{final_message}"
            )
        else:
            conversation_summary = (
                f"Completed with {len(messages)} messages."
                if messages
                else "No messages returned from AutoGen team."
            )
        print(f"✅ AutoGen team completed with {len(messages)} messages")
        _persist_autogen_log(
            status="completed",
            messages=messages,
            conversation_summary=conversation_summary,
        )
        log_path = _resolve_monitor_log_path(log_filename)
        log_payload = _read_monitor_run_log(log_path)
        transcript_text = log_payload.get("content") if isinstance(log_payload, dict) else ""
        if not isinstance(transcript_text, str) or not transcript_text:
            transcript_text = _format_autogen_conversation_log(
                input_text,
                messages,
                conversation_summary,
                timestamp_human=datetime.now().astimezone().strftime("%Y-%m-%d %H:%M:%S %Z"),
                status="completed",
                progress_notes=progress_notes,
            )
        _monitor_run_finish(
            monitor_run_id,
            status="completed",
            summary=f"Completed with {len(messages)} messages.",
            metadata={
                "message_count": len(messages),
                "sources": [msg.get("source", "unknown") for msg in messages[:12]],
            },
            log_file=log_filename,
            log_excerpt=_read_monitor_log_excerpt((SCRATCH_DIR / log_filename) if log_filename else None),
        )
        return {
            "output": conversation_summary,
            "response": conversation_summary,
            "messages": messages,
            "message_count": len(messages),
            "log_file": log_filename,
            "log_content": transcript_text,
            "summary": conversation_summary,
        }
    except Exception as e:
        import traceback
        print(f"âŒ AutoGen team execution error: {e}")
        print(traceback.format_exc())
        error_text = str(e)
        if "No user query found in messages" in error_text:
            autogen_base = (
                os.getenv("AUTOGEN_OPENROUTER_BASE_URL")
                or os.getenv("OPENROUTER_API_BASE")
                or os.getenv("MCP_LLM_BASE_URL")
                or os.getenv("OPENAI_API_BASE")
                or ""
            ).strip()
            autogen_model = (
                os.getenv("AUTOGEN_TEAM_MODEL")
                or os.getenv("OPENROUTER_AUTOGEN_MODEL")
                or os.getenv("MCP_LLM_MODEL_NAME")
                or os.getenv("OPENAI_MODEL")
                or ""
            ).strip()
            error_text = (
                "AutoGen team execution failed because the configured AutoGen model endpoint rejected an internal "
                "multi-agent prompt with 'No user query found in messages'. This usually means the current chat "
                "template requires a trailing user message and is not compatible with AutoGen's internal "
                f"assistant/tool turns. Current settings: AUTOGEN_OPENROUTER_BASE_URL={autogen_base or '(unset)'}, "
                f"AUTOGEN_TEAM_MODEL={autogen_model or '(unset)'}. Fix this by either pointing AutoGen to a "
                "compatible hosted endpoint such as https://openrouter.ai/api/v1 with a stable chat model, or by "
                "switching the local model/server to one whose chat template supports assistant continuation and tool "
                "calling. Restart the proxy server after changing .env."
            )
        elif "user name must be consistent (2013)" in error_text:
            autogen_base = (
                os.getenv("AUTOGEN_MINIMAX_BASE_URL")
                or os.getenv("AUTOGEN_BASE_URL")
                or os.getenv("MCP_LLM_BASE_URL")
                or os.getenv("OPENAI_API_BASE")
                or ""
            ).strip()
            autogen_model = (
                os.getenv("AUTOGEN_TEAM_MODEL")
                or os.getenv("AUTOGEN_MINIMAX_MODEL")
                or os.getenv("MCP_LLM_MODEL_NAME")
                or os.getenv("OPENAI_MODEL")
                or ""
            ).strip()
            error_text = (
                "AutoGen team execution failed because MiniMax rejected the multi-agent message `name` fields with "
                "'user name must be consistent (2013)'. The proxy now applies a MiniMax compatibility mode that "
                "removes per-message names and prefixes speaker identity into message text instead. If you still see "
                f"this after the fix, restart the proxy so the AutoGen team is rebuilt. Current settings: "
                f"AUTOGEN_MINIMAX_BASE_URL={autogen_base or '(unset)'}, AUTOGEN_TEAM_MODEL={autogen_model or '(unset)' }."
            )
        _persist_autogen_log(
            status="error",
            conversation_summary="AutoGen execution failed.",
            error_text=error_text,
        )
        _monitor_run_finish(
            monitor_run_id,
            status="error",
            summary=f"AutoGen execution failed: {error_text}",
            log_file=log_filename,
            log_excerpt=_read_monitor_log_excerpt((SCRATCH_DIR / log_filename) if log_filename else None),
        )
        raise HTTPException(status_code=500, detail=error_text)


def _write_autogen_conversation_to_scratch(
    input_text: str,
    messages: List[Dict[str, str]],
    conversation_summary: str,
    *,
    filename: Optional[str] = None,
    status: str = "completed",
    progress_notes: Optional[List[str]] = None,
    error_text: Optional[str] = None,
) -> str:
    """Write or update an AutoGen conversation log in scratch and return filename."""
    SCRATCH_DIR.mkdir(parents=True, exist_ok=True)
    now = datetime.now().astimezone()
    timestamp_human = now.strftime("%Y-%m-%d %H:%M:%S %Z")
    if not filename:
        timestamp_file = now.strftime("%Y-%m-%d_%H-%M-%S")
        suffix = secrets.token_hex(4)
        filename = f"autogen_run_{timestamp_file}_{suffix}.txt"
    filepath = SCRATCH_DIR / filename
    filepath.write_text(
        _format_autogen_conversation_log(
            input_text,
            messages,
            conversation_summary,
            timestamp_human=timestamp_human,
            status=status,
            progress_notes=progress_notes,
            error_text=error_text,
        ),
        encoding="utf-8",
    )
    print(f"[AUTOGEN] Wrote conversation to {filepath}", flush=True)
    return filename


def _write_codex_summary_to_scratch(
    prompt: str,
    command: List[str],
    exit_code: Optional[int],
    stdout: str,
    stderr: str,
    duration_ms: int,
    timed_out: bool,
    events_file: Optional[str] = None,
    last_message_file: Optional[str] = None,
    workspace_dir: Optional[str] = None,
    workspace_mode: Optional[str] = None,
) -> str:
    """Write Codex execution summary to scratch and return filename."""
    SCRATCH_DIR.mkdir(parents=True, exist_ok=True)
    now = datetime.now().astimezone()
    timestamp_file = now.strftime("%Y-%m-%d_%H-%M-%S")
    timestamp_human = now.strftime("%Y-%m-%d %H:%M:%S %Z")
    suffix = secrets.token_hex(4)
    filename = f"codex_run_{timestamp_file}_{suffix}.txt"
    filepath = SCRATCH_DIR / filename
    lines = [
        "Codex CLI execution summary",
        f"Date-time: {timestamp_human}",
        f"Duration (ms): {duration_ms}",
        f"Timed out: {timed_out}",
        f"Exit code: {exit_code if exit_code is not None else 'N/A'}",
        f"Events file: {events_file or 'N/A'}",
        f"Last message file: {last_message_file or 'N/A'}",
        f"Workspace mode: {workspace_mode or 'project_root'}",
        f"Workspace dir: {workspace_dir or str(_PROJECT_ROOT)}",
        "",
        "Command:",
        " ".join(command),
        "",
        "Prompt:",
        prompt or "(empty)",
        "",
        "--- STDOUT ---",
        stdout or "(empty)",
        "",
        "--- STDERR ---",
        stderr or "(empty)",
    ]
    filepath.write_text("\n".join(lines), encoding="utf-8")
    return filename


def _prepare_autogen_codex_workspace() -> Path:
    """Create an empty isolated workspace for AutoGen Codex runs under scratch/autogen."""
    SCRATCH_DIR.mkdir(parents=True, exist_ok=True)
    autogen_root = SCRATCH_DIR / CODEX_AUTOGEN_WORKSPACES_DIRNAME
    autogen_root.mkdir(parents=True, exist_ok=True)
    now = datetime.now().astimezone()
    timestamp_file = now.strftime("%Y-%m-%d_%H-%M-%S")
    suffix = secrets.token_hex(4)
    workspace_dir = autogen_root / f"codex_run_{timestamp_file}_{suffix}"
    workspace_dir.mkdir(parents=True, exist_ok=False)
    marker_path = workspace_dir / "AUTOGEN_WORKSPACE_README.txt"
    marker_path.write_text(
        "\n".join(
            [
                "This workspace is an empty isolated directory for an AutoGen Codex CLI run.",
                "It was created under scratch/autogen so Codex can build a new project without copying CATBot.",
                "Changes here do not modify the live CATBot core repository.",
            ]
        ),
        encoding="utf-8",
    )
    return workspace_dir


def _write_codex_error_to_scratch(
    prompt: str,
    command: List[str],
    error_message: str,
) -> str:
    """Write Codex execution error to scratch and return filename."""
    SCRATCH_DIR.mkdir(parents=True, exist_ok=True)
    now = datetime.now().astimezone()
    timestamp_file = now.strftime("%Y-%m-%d_%H-%M-%S")
    timestamp_human = now.strftime("%Y-%m-%d %H:%M:%S %Z")
    suffix = secrets.token_hex(4)
    filename = f"codex_run_error_{timestamp_file}_{suffix}.txt"
    filepath = SCRATCH_DIR / filename
    lines = [
        "Codex CLI execution error",
        f"Date-time: {timestamp_human}",
        "",
        "Command:",
        " ".join(command) if command else "(none)",
        "",
        "Prompt:",
        prompt or "(empty)",
        "",
        "Error:",
        error_message or "(empty)",
    ]
    filepath.write_text("\n".join(lines), encoding="utf-8")
    return filename


async def _run_codex_cli(prompt: str, *, isolated_workspace: bool = False) -> Dict[str, Any]:
    if not CODEX_ENABLED:
        raise HTTPException(status_code=503, detail="Codex CLI tool is disabled.")
    prompt = (prompt or "").strip()
    if not prompt:
        raise HTTPException(status_code=400, detail="Prompt is required.")
    timeout = CODEX_TIMEOUT_SECONDS
    if timeout <= 0:
        raise HTTPException(status_code=500, detail="CODEX_TIMEOUT_SECONDS must be a positive integer.")
    timeout = min(max(timeout, 60), 7200)

    sandbox_mode = CODEX_SANDBOX_MODE if CODEX_SANDBOX_MODE in ("read-only", "workspace-write") else "workspace-write"
    approval_policy = CODEX_APPROVAL_POLICY if CODEX_APPROVAL_POLICY in ("untrusted", "on-request", "on-failure", "never") else "never"
    workspace_dir = _PROJECT_ROOT
    workspace_mode = "project_root"
    if isolated_workspace:
        try:
            workspace_dir = _prepare_autogen_codex_workspace()
            workspace_mode = "scratch_autogen_empty"
        except Exception as e:
            error_file = _write_codex_error_to_scratch(prompt, [], f"Failed to prepare AutoGen Codex workspace: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Failed to prepare isolated AutoGen Codex workspace. See {error_file} in scratch for details.",
            )

    cmd: List[str] = [CODEX_CLI_PATH]
    if CODEX_ENABLE_SEARCH:
        # --search is a top-level codex flag and must appear before the subcommand.
        cmd.append("--search")
    cmd.extend([
        "exec",
        "--sandbox",
        sandbox_mode,
        "-C",
        str(workspace_dir),
    ])
    if approval_policy == "never":
        # codex exec no longer accepts -a; --full-auto is the supported non-interactive mode
        cmd.append("--full-auto")
    events_file = None
    last_message_file = None
    if CODEX_JSON_EVENTS:
        cmd.append("--json")
    if CODEX_OUTPUT_LAST_MESSAGE:
        SCRATCH_DIR.mkdir(parents=True, exist_ok=True)
        now = datetime.now().astimezone()
        timestamp_file = now.strftime("%Y-%m-%d_%H-%M-%S")
        suffix = secrets.token_hex(4)
        last_message_file = f"codex_last_message_{timestamp_file}_{suffix}.txt"
        cmd.extend(["-o", str(SCRATCH_DIR / last_message_file)])
    cmd.append(prompt)

    start = time.time()
    stdout_text = ""
    stderr_text = ""
    exit_code: Optional[int] = None
    timed_out = False
    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            cwd=str(workspace_dir),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            out_bytes, err_bytes = await asyncio.wait_for(proc.communicate(), timeout=timeout)
        except asyncio.TimeoutError:
            timed_out = True
            proc.kill()
            out_bytes, err_bytes = await proc.communicate()
        stdout_text = (out_bytes or b"").decode("utf-8", errors="replace")
        stderr_text = (err_bytes or b"").decode("utf-8", errors="replace")
        exit_code = proc.returncode
    except FileNotFoundError:
        error_file = _write_codex_error_to_scratch(prompt, cmd, "Codex CLI not found. Ensure it is installed and on PATH.")
        raise HTTPException(status_code=500, detail=f"Codex CLI not found. See {error_file} in scratch for details.")
    except Exception as e:
        error_file = _write_codex_error_to_scratch(prompt, cmd, f"Failed to execute Codex CLI: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to execute Codex CLI. See {error_file} in scratch for details.")
    finally:
        duration_ms = int((time.time() - start) * 1000)

    if CODEX_JSON_EVENTS:
        SCRATCH_DIR.mkdir(parents=True, exist_ok=True)
        now = datetime.now().astimezone()
        timestamp_file = now.strftime("%Y-%m-%d_%H-%M-%S")
        suffix = secrets.token_hex(4)
        events_file = f"codex_events_{timestamp_file}_{suffix}.jsonl"
        (SCRATCH_DIR / events_file).write_text(stdout_text or "", encoding="utf-8")

    summary_file = _write_codex_summary_to_scratch(
        prompt=prompt,
        command=cmd,
        exit_code=exit_code,
        stdout=stdout_text,
        stderr=stderr_text,
        duration_ms=duration_ms,
        timed_out=timed_out,
        events_file=events_file,
        last_message_file=last_message_file,
        workspace_dir=str(workspace_dir),
        workspace_mode=workspace_mode,
    )
    return {
        "success": exit_code == 0 and not timed_out,
        "summaryFile": summary_file,
        "eventsFile": events_file,
        "lastMessageFile": last_message_file,
        "exitCode": exit_code,
        "timedOut": timed_out,
        "durationMs": duration_ms,
        "stdout": stdout_text,
        "stderr": stderr_text,
        "workspaceDir": str(workspace_dir),
        "workspaceMode": workspace_mode,
    }


# AutoGen team chat endpoint (integrated directly)
@app.post("/v1/proxy/autogen")
async def autogen_chat(request: Request):
    """Run AutoGen team conversation directly (no separate service needed)."""
    try:
        body = await request.json()
        input_text = body.get('input')
        if not input_text:
            raise HTTPException(status_code=400, detail="Input parameter is required")
        print(f"🤖 AutoGen team request: {input_text[:100]}...")
        return await _do_autogen(input_text)
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"âŒ AutoGen endpoint error: {e}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to process AutoGen request: {str(e)}")


@app.post("/v1/proxy/codex")
async def proxy_codex_exec(
    request: CodexExecRequest,
    current_user: Dict[str, Any] = Depends(get_current_user_or_autogen_team),
):
    """Run Codex CLI non-interactively in sandboxed mode. User JWT or AutoGen team secret required."""
    return await _run_codex_cli(
        request.prompt,
        isolated_workspace=current_user.get("auth_type") == "agent_secret",
    )


@app.post("/v1/proxy/restart")
async def proxy_restart_server(current_user: Dict[str, Any] = Depends(get_current_user)):
    """Schedule a proxy process restart. Auth required."""
    username = current_user.get("username", "unknown")
    return await _request_proxy_restart(trigger="web_api", requested_by=f"user:{username}")

# Allowed keys when persisting MCP server config (never persist 'command')
MCP_SERVER_SAFE_KEYS = {"id", "name", "preset_id", "apiKey", "model", "url", "wsUrl", "status", "enabled"}


# MCP server management endpoints (all require authentication)
@app.post("/v1/mcp/servers")
async def manage_servers(
    server_config: ServerConfig,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Manage MCP servers (create, update, clear). Never persist or execute client-supplied command."""
    # Log without secrets
    safe_dict = {k: v for k, v in server_config.model_dump().items() if k != "apiKey"}
    print(f"Received server config: {safe_dict}")

    global mcp_clients, mcp_servers

    try:
        # Handle clear action
        if server_config.action == "clear":
            for server_id, client in list(mcp_clients.items()):
                try:
                    await client.close()
                except Exception as e:
                    print(f"Error closing client {server_id}: {e}")
            mcp_servers.clear()
            mcp_clients.clear()
            print("Cleared all MCP servers and clients")
            save_servers()
            security_log("mcp_clear", current_user.get("username", ""), None, "all servers cleared")
            return {"message": "All MCP servers cleared successfully"}

        # Validate required fields
        if not server_config.id or not server_config.name:
            raise HTTPException(status_code=400, detail="Missing required fields: id, name")

        # Resolve preset_id: require or infer from name (browser-use)
        preset_id = server_config.preset_id
        if not preset_id and server_config.name and "mcp-browser-use" in server_config.name.lower():
            preset_id = "browser-use"
        if not preset_id:
            raise HTTPException(
                status_code=400,
                detail="Missing preset_id. Provide preset_id (e.g. 'browser-use') or use a name that implies it.",
            )
        if preset_id not in MCP_PRESETS:
            raise HTTPException(
                status_code=400,
                detail=f"Invalid preset_id '{preset_id}'. Allowed: {list(MCP_PRESETS.keys())}",
            )

        # Build stored server dict from allowed fields only; never store 'command'
        raw = server_config.model_dump()
        stored = {k: raw[k] for k in MCP_SERVER_SAFE_KEYS if k in raw and raw[k] is not None}
        stored["status"] = "disconnected"
        # Preserve connection status when updating
        existing_server = mcp_servers.get(server_config.id)
        if existing_server:
            stored["status"] = existing_server.get("status", "disconnected")
        stored["preset_id"] = preset_id

        if existing_server:
            mcp_servers[server_config.id] = {**existing_server, **stored}
            print(f"Updated MCP server: {server_config.name} ({server_config.id})")
            security_log("mcp_update", current_user.get("username", ""), server_config.id, f"preset_id={preset_id}")
        else:
            mcp_servers[server_config.id] = stored
            print(f"Added MCP server: {server_config.name} ({server_config.id})")
            security_log("mcp_add", current_user.get("username", ""), server_config.id, f"preset_id={preset_id}")

        save_servers()
        return {"message": "Server saved successfully"}

    except HTTPException:
        raise
    except Exception as e:
        print(f"Error saving MCP server: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to save MCP server: {str(e)}")

@app.get("/v1/mcp/servers")
async def get_servers(
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Get all MCP servers."""
    try:
        servers = list(mcp_servers.values())
        return {"servers": servers}
    except Exception as e:
        print(f"Error getting MCP servers: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to get MCP servers: {str(e)}")

def setup_browser_llm():
    """Set up the language model for browser automation."""
    model_provider = os.getenv("MCP_MODEL_PROVIDER", "google").lower()
    model_name = os.getenv("MCP_MODEL_NAME", "gemini-flash-latest")
    temperature = float(os.getenv("MCP_TEMPERATURE", "0.1"))
    model_base_url = (
        os.getenv("MCP_MODEL_BASE_URL")
        or os.getenv("MCP_LLM_BASE_URL")
        or os.getenv("OPENAI_API_BASE")
        or ""
    ).strip()

    if model_provider in {"openai", "openrouter", "minimax"}:
        from langchain_openai import ChatOpenAI

        api_key_candidates = (
            ["MINIMAX_API_KEY", "MCP_LLM_MINIMAX_API_KEY", "OPENAI_API_KEY", "MCP_LLM_OPENAI_API_KEY"]
            if model_provider == "minimax" or is_minimax_chat_request(model_base_url, model_name)
            else (
                ["OPENROUTER_API_KEY", "MCP_LLM_OPENROUTER_API_KEY", "OPENAI_API_KEY", "MCP_LLM_OPENAI_API_KEY"]
                if model_provider == "openrouter" or "openrouter.ai" in model_base_url.lower()
                else ["OPENAI_API_KEY", "MCP_LLM_OPENAI_API_KEY"]
            )
        )
        chat_kwargs: Dict[str, Any] = {
            "model": model_name,
            "temperature": (
                normalize_temperature_for_minimax(temperature)
                if is_minimax_chat_request(model_base_url, model_name)
                else temperature
            ),
            "api_key": _first_non_empty_env(api_key_candidates),
        }
        if model_base_url:
            chat_kwargs["base_url"] = model_base_url
        if model_provider == "openrouter" or "openrouter.ai" in model_base_url.lower():
            default_headers: Dict[str, str] = {}
            referer = (os.getenv("OPENROUTER_HTTP_REFERER") or os.getenv("OPENROUTER_REFERER") or "").strip()
            title = (os.getenv("OPENROUTER_X_TITLE") or "CATBot").strip()
            if referer:
                default_headers["HTTP-Referer"] = referer
            if title:
                default_headers["X-Title"] = title
            if default_headers:
                chat_kwargs["default_headers"] = default_headers
        return ChatOpenAI(**chat_kwargs)

    if model_provider == "anthropic":
        from langchain_anthropic import ChatAnthropic
        return ChatAnthropic(
            model=model_name,
            temperature=temperature,
            api_key=os.getenv("ANTHROPIC_API_KEY")
        )
    elif model_provider == "google":
        from langchain_google_genai import ChatGoogleGenerativeAI
        return ChatGoogleGenerativeAI(
            model=model_name,
            temperature=temperature,
            api_key=os.getenv("GOOGLE_API_KEY")
        )
    else:
        # Default to Google/Gemini
        from langchain_google_genai import ChatGoogleGenerativeAI
        return ChatGoogleGenerativeAI(
            model=model_name,
            temperature=temperature,
            api_key=os.getenv("GOOGLE_API_KEY", "dummy-key")
        )

async def create_mcp_client(server_config: Dict[str, Any]):
    """Create MCP client connection (similar to Node.js version)."""
    manager = MCPClientManager(server_config)
    client = await manager.connect()
    return client

@app.post("/v1/mcp/servers/{server_id}/connect")
async def connect_server(
    server_id: str,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Connect to an MCP server. Uses only server-side presets; never executes user-supplied command."""
    try:
        print(f"Attempting to connect to server: {server_id}")

        server = mcp_servers.get(server_id)
        if not server:
            print(f"Server not found: {server_id}")
            raise HTTPException(status_code=404, detail="Server not found")

        print(f"Found server: {server.get('name')} ({server_id})")

        # Inprocess preset (e.g. browser-use): no subprocess, mark connected
        preset_id = server.get("preset_id") or (
            "browser-use" if _is_browser_use_server(server) else None
        )
        if preset_id == "browser-use" or _is_browser_use_server(server):
            server["preset_id"] = preset_id or "browser-use"
            server["status"] = "connected"
            mcp_servers[server_id] = server
            print(f"Successfully connected to MCP Browser Use server: {server.get('name')}")
            security_log("mcp_connect", current_user.get("username", ""), server_id, "preset_id=browser-use")
            return {"message": "Server connected successfully"}

        if server_id in mcp_clients:
            print(f"Server already connected: {server_id}")
            raise HTTPException(status_code=409, detail="Server is already connected")

        if not preset_id or preset_id not in MCP_PRESETS:
            raise HTTPException(
                status_code=400,
                detail="Server has no valid preset_id. Reconfigure the server with a preset (e.g. preset_id: 'browser-use').",
            )
        if MCP_PRESETS[preset_id].get("type") != "stdio":
            raise HTTPException(
                status_code=400,
                detail=f"Preset '{preset_id}' does not support stdio connect. Use a stdio preset or browser-use.",
            )

        if not MCP_AVAILABLE:
            raise HTTPException(status_code=503, detail="MCP SDK not available")

        print(f"Creating MCP client for server: {server.get('name')}")
        client = await create_mcp_client(server)
        mcp_clients[server_id] = client

        server["status"] = "connected"
        mcp_servers[server_id] = server

        print(f"Successfully connected to MCP server: {server.get('name')}")
        security_log("mcp_connect", current_user.get("username", ""), server_id, f"preset_id={preset_id}")
        return {"message": "Server connected successfully"}

    except HTTPException:
        raise
    except ValueError as e:
        print(f"Error connecting to MCP server: {e}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        print(f"Error connecting to MCP server: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to connect to MCP server: {str(e)}")

@app.post("/v1/mcp/servers/{server_id}/disconnect")
async def disconnect_server(
    server_id: str,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Disconnect from an MCP server."""
    try:
        global mcp_clients, mcp_servers

        # Check if this is the MCP Browser Use server
        server = mcp_servers.get(server_id)
        if server and _is_browser_use_server(server):
            # Just mark as disconnected since we don't have a real connection
            server["status"] = "disconnected"
            mcp_servers[server_id] = server
            security_log("mcp_disconnect", current_user.get("username", ""), server_id, "browser-use")
            return {"message": "Server disconnected successfully"}

        if not MCP_AVAILABLE:
            raise HTTPException(status_code=503, detail="MCP SDK not available")

        client = mcp_clients.get(server_id)
        if not client:
            raise HTTPException(status_code=404, detail="Server is not connected")

        # Create a temporary manager to handle disconnect
        if server:
            manager = MCPClientManager(server)
            await manager.disconnect()

        del mcp_clients[server_id]

        server = mcp_servers.get(server_id)
        if server:
            server["status"] = "disconnected"
            mcp_servers[server_id] = server

        security_log("mcp_disconnect", current_user.get("username", ""), server_id, "stdio")
        return {"message": "Server disconnected successfully"}

    except HTTPException:
        raise
    except Exception as e:
        print(f"Error disconnecting MCP server: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to disconnect MCP server: {str(e)}")


async def _browser_use_http_list_tools() -> Dict[str, Any]:
    """List tools from the browser-use HTTP MCP server. Raises on connection failure."""
    from fastmcp import Client

    async with Client(MCP_BROWSER_USE_HTTP_URL) as client:
        tools = await client.list_tools()
    # Convert to the shape expected by the proxy API (name, description, inputSchema)
    tools_list = []
    for t in tools:
        entry = {
            "name": getattr(t, "name", "unknown"),
            "description": getattr(t, "description", None) or "",
            "inputSchema": getattr(t, "inputSchema", None) or {"type": "object", "properties": {}},
        }
        tools_list.append(entry)
    return {"tools": tools_list}


async def _browser_use_http_call_tool(tool_name: str, parameters: Dict[str, Any]) -> Dict[str, Any]:
    """Call a tool on the browser-use HTTP MCP server. Returns result with content list; raises on connection failure."""
    from fastmcp import Client

    _log_tool_invocation("browser_use_http", tool_name, parameters)

    # Map proxy parameter names to server names (e.g. instruction -> task for run_browser_agent)
    args = dict(parameters)
    if tool_name == "run_browser_agent" and "instruction" in args and "task" not in args:
        args["task"] = args.pop("instruction")

    async with Client(MCP_BROWSER_USE_HTTP_URL) as client:
        result = await client.call_tool(tool_name, args)

    # Build content list from result.content (list of items with .text or str)
    content = []
    for item in getattr(result, "content", []) or []:
        text = getattr(item, "text", None) or str(item)
        content.append({"type": "text", "text": text})
    if getattr(result, "is_error", False) and not content:
        content.append({"type": "text", "text": "Tool returned an error."})
    return {"content": content}


def _is_browser_use_server(server: Optional[Dict[str, Any]]) -> bool:
    """Return True if server is the browser-use preset (by preset_id or name). Used for routing tools and philosopher detection."""
    if not server:
        return False
    name = (server.get("name") or "").lower()
    preset_id = (server.get("preset_id") or "").lower()
    return (
        preset_id == "browser-use"
        or "mcp-browser-use" in name
        or "browser-use" in name
        or ("browser" in name and "use" in name)
    )


def _log_tool_invocation(source: str, tool_name: str, parameters: Optional[Dict[str, Any]] = None) -> None:
    """Unified console logging for all tool invocation entry points."""
    try:
        payload = parameters if isinstance(parameters, dict) else {}
        args_text = json.dumps(payload, ensure_ascii=False, default=str)
    except Exception:
        args_text = str(parameters)
    print(f"[TOOL][{source}] name={tool_name} args={args_text}", flush=True)


@app.post("/v1/tools/log")
async def log_tool_invocation(request: Request):
    """Record a tool invocation from clients that execute tool routing locally (e.g. HTML UI)."""
    body: Dict[str, Any] = {}
    try:
        parsed = await request.json()
        if isinstance(parsed, dict):
            body = parsed
    except Exception:
        body = {}

    tool_name = str(body.get("name") or "").strip()
    if not tool_name:
        raise HTTPException(status_code=400, detail="name is required")

    source = str(body.get("source") or "client")
    arguments = body.get("arguments")
    _log_tool_invocation(source, tool_name, arguments if isinstance(arguments, dict) else {})
    return {"success": True}


# Browser automation tool endpoint (browser-use via HTTP; other servers via MCP client)
@app.post("/v1/mcp/servers/{server_id}/tools/call")
async def call_tool(server_id: str, request: ToolCallRequest):
    """Call a tool on an MCP server. Browser-use preset uses HTTP client; others use connected MCP client."""
    if not MCP_AVAILABLE:
        raise HTTPException(status_code=503, detail="MCP SDK not available")

    try:
        print(f"🔧 [TOOLS/CALL] Server: {server_id}")

        server = mcp_servers.get(server_id)
        tool_name = request.toolName
        parameters = request.parameters or {}
        _log_tool_invocation(f"mcp_call:{server_id}", tool_name, parameters)
        print(f"ðŸ” [TOOLS/CALL] Tool name: {tool_name}")
        print(f"ðŸ” [TOOLS/CALL] Parameters: {parameters}")

        # Browser-use preset: use HTTP client (no mcp_clients entry)
        if server and _is_browser_use_server(server):
            try:
                result = await _browser_use_http_call_tool(tool_name, parameters)
                return {"result": result}
            except Exception as e:
                print(f"âŒ [TOOLS/CALL] Browser-use HTTP error: {e}")
                raise HTTPException(
                    status_code=503,
                    detail=BROWSER_USE_HTTP_UNAVAILABLE_MSG + " " + str(e),
                )

        client = mcp_clients.get(server_id)
        if not client:
            print(f"âŒ [TOOLS/CALL] Server {server_id} not found or not connected")
            raise HTTPException(status_code=404, detail="Server is not connected")

        if not tool_name:
            print("âŒ [TOOLS/CALL] toolName is required but missing")
            raise HTTPException(status_code=400, detail="toolName is required")

        result = await client.request(
            method="tools/call",
            params={"name": tool_name, "arguments": parameters},
        )
        return {"result": result}

    except HTTPException:
        raise
    except Exception as e:
        print(f"💥 [TOOLS/CALL] Error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to call tool on MCP server: {str(e)}")

@app.post("/v1/mcp/servers/{server_id}/tools/list")
async def list_tools(server_id: str):
    """List tools available on an MCP server. Browser-use preset uses HTTP client."""
    if not MCP_AVAILABLE:
        raise HTTPException(status_code=503, detail="MCP SDK not available")

    try:
        print(f"ðŸ” [TOOLS/LIST] Server: {server_id}")

        server = mcp_servers.get(server_id)
        # Browser-use preset: list tools from HTTP server
        if server and _is_browser_use_server(server):
            try:
                result = await _browser_use_http_list_tools()
                return {"result": result}
            except Exception as e:
                print(f"âŒ [TOOLS/LIST] Browser-use HTTP error: {e}")
                raise HTTPException(
                    status_code=503,
                    detail=BROWSER_USE_HTTP_UNAVAILABLE_MSG + " " + str(e),
                )

        global mcp_clients
        client = mcp_clients.get(server_id)
        if not client:
            print(f"âŒ [TOOLS/LIST] Server {server_id} not found or not connected")
            raise HTTPException(status_code=404, detail="Server is not connected")

        result = await client.request(
            method="tools/list",
            params={},
        )

        print(f"📨 [TOOLS/LIST] Raw response from MCP server: {result}")

        # Validate response structure
        if not result:
            print("âŒ [TOOLS/LIST] No result returned from MCP server")
        elif 'tools' not in result:
            print(f"âŒ [TOOLS/LIST] Missing 'tools' field in response: {list(result.keys())}")
        elif not isinstance(result['tools'], list):
            print(f"âŒ [TOOLS/LIST] 'tools' field is not an array: {type(result['tools'])}")
        else:
            print(f"✅ [TOOLS/LIST] Found {len(result['tools'])} tools in response")
            for i, tool in enumerate(result['tools']):
                print(f"  Tool {i}: {tool.get('name', 'unnamed')}")
                if 'name' not in tool:
                    print(f"    âŒ Missing name for tool {i}")
                if 'description' not in tool:
                    print(f"    âš ï¸  Missing description for tool {i}")
                if 'inputSchema' not in tool:
                    print(f"    âš ï¸  Missing inputSchema for tool {i}")
                else:
                    schema = tool['inputSchema']
                    print(f"    ✅ inputSchema type: {schema.get('type')}")
                    if 'properties' in schema:
                        print(f"    ✅ Has {len(schema['properties'])} properties")
                    else:
                        print("    âš ï¸  No properties in inputSchema")

        return {"result": result}

    except HTTPException:
        raise
    except Exception as e:
        print(f"💥 [TOOLS/LIST] Error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to list tools on MCP server: {str(e)}")

# Root endpoint
@app.get("/")
async def root():
    """Root endpoint."""
    return {"message": "CATBot Proxy Server", "version": "2.0.0"}


@app.get("/spotify/authorize")
async def spotify_authorize() -> RedirectResponse:
    """Redirect the browser to Spotify's authorization page for CATBot playback access."""

    _cleanup_spotify_oauth_states()
    state = secrets.token_urlsafe(24)
    spotify_oauth_pending_states[state] = time.time()
    query = urlencode(
        {
            "client_id": _spotify_client_id(),
            "response_type": "code",
            "redirect_uri": _spotify_redirect_uri(),
            "scope": " ".join(SPOTIFY_PLAYBACK_SCOPES),
            "state": state,
            "show_dialog": "true",
        }
    )
    return RedirectResponse(url=f"{SPOTIFY_ACCOUNTS_BASE}/authorize?{query}", status_code=307)


@app.get("/spotify/callback")
async def spotify_callback(code: Optional[str] = None, state: Optional[str] = None, error: Optional[str] = None):
    """Handle the Spotify authorization code callback and persist playback tokens to .env."""

    if error:
        raise HTTPException(status_code=400, detail=f"Spotify authorization failed: {error}")
    if not code:
        raise HTTPException(status_code=400, detail="Missing Spotify authorization code.")
    if not state:
        raise HTTPException(status_code=400, detail="Missing Spotify OAuth state.")

    _cleanup_spotify_oauth_states()
    created_at = spotify_oauth_pending_states.pop(state, None)
    if created_at is None:
        raise HTTPException(status_code=400, detail="Invalid or expired Spotify OAuth state.")
    if time.time() - created_at > SPOTIFY_OAUTH_STATE_TTL_SECONDS:
        raise HTTPException(status_code=400, detail="Expired Spotify OAuth state.")

    redirect_uri = _spotify_redirect_uri()
    async with httpx.AsyncClient(timeout=SPOTIFY_OAUTH_TIMEOUT_SECONDS) as client:
        response = await client.post(
            f"{SPOTIFY_AUTH_BASE}/token",
            headers=_spotify_token_request_headers(),
            data={
                "grant_type": "authorization_code",
                "code": code,
                "redirect_uri": redirect_uri,
            },
        )

    try:
        payload = response.json()
    except Exception:
        payload = None

    if response.status_code != 200:
        raise HTTPException(
            status_code=400,
            detail=f"Spotify token exchange failed: {_extract_spotify_error_message(payload)}",
        )
    if not isinstance(payload, dict):
        raise HTTPException(status_code=500, detail="Spotify token exchange returned an invalid payload.")

    access_token = str(payload.get("access_token") or "").strip()
    refresh_token = str(payload.get("refresh_token") or "").strip()
    if not access_token:
        raise HTTPException(status_code=500, detail="Spotify token exchange did not return an access token.")

    _persist_env_key("SPOTIFY_ACCESS_TOKEN", access_token)
    persisted_keys = ["SPOTIFY_ACCESS_TOKEN"]
    if refresh_token:
        _persist_env_key("SPOTIFY_REFRESH_TOKEN", refresh_token)
        persisted_keys.append("SPOTIFY_REFRESH_TOKEN")

    body_lines = [
        "<h1>Spotify authorization complete</h1>",
        "<p>CATBot saved the latest Spotify playback credentials to <code>.env</code>.</p>",
        f"<p>Updated keys: <code>{html.escape(', '.join(persisted_keys))}</code></p>",
    ]
    if refresh_token:
        body_lines.append("<p>You can use Spotify playback immediately. The current CATBot process was updated in memory too.</p>")
    else:
        body_lines.append(
            "<p>Spotify did not return a refresh token in this response. If playback renewal fails later, "
            "repeat the authorization flow after revoking the app in your Spotify account settings.</p>"
        )
    return HTMLResponse(
        "".join(
            [
                "<!doctype html><html><head><meta charset='utf-8'><title>Spotify Auth Complete</title></head><body>",
                *body_lines,
                "</body></html>",
            ]
        )
    )

@app.post("/v1/auth/signup", response_model=AuthTokenResponse)
async def auth_signup(request: AuthSignupRequest):
    """Create a new user account and return a signed JWT."""
    username = request.username.strip().lower()
    password = request.password

    if len(username) < 3:
        raise HTTPException(status_code=400, detail="username must be at least 3 characters")
    if len(password) < 8:
        raise HTTPException(status_code=400, detail="password must be at least 8 characters")
    if username in users_db:
        raise HTTPException(status_code=409, detail="username already exists")

    password_record = create_password_record(password)
    users_db[username] = {
        **password_record,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    save_users_db()

    token = create_jwt({"sub": username})
    return AuthTokenResponse(
        access_token=token,
        expires_in=JWT_EXPIRATION_SECONDS,
        username=username,
    )


@app.post("/v1/auth/login", response_model=AuthTokenResponse)
async def auth_login(request: AuthLoginRequest):
    """Authenticate a user and return a signed JWT."""
    username = request.username.strip().lower()
    user = users_db.get(username)
    if not user:
        raise HTTPException(status_code=401, detail="Invalid credentials")

    if not verify_password(request.password, user["salt"], user["password_hash"]):
        raise HTTPException(status_code=401, detail="Invalid credentials")

    token = create_jwt({"sub": username})
    return AuthTokenResponse(
        access_token=token,
        expires_in=JWT_EXPIRATION_SECONDS,
        username=username,
    )


@app.get("/v1/auth/me", response_model=AuthUserResponse)
async def auth_me(current_user: Dict[str, Any] = Depends(get_current_user)):
    """Return the profile of the authenticated user based on JWT bearer token."""
    return AuthUserResponse(
        username=current_user["username"],
        created_at=current_user["created_at"],
    )


# ============================================================================
# TODO REST API (all endpoints require authentication)
# ============================================================================

def _todo_recurrence_to_dict(recurrence: Optional[TodoRecurrenceRequest]) -> Optional[Dict[str, Any]]:
    if not recurrence:
        return None
    return {
        "frequency": str(recurrence.frequency).strip().lower(),
        "interval": int(recurrence.interval),
    }


def _parse_todo_datetime(value: Optional[str]) -> Optional[datetime]:
    if not value or not isinstance(value, str):
        return None
    text = value.strip()
    if not text:
        return None
    if text.endswith("Z"):
        text = text[:-1] + "+00:00"
    try:
        parsed = datetime.fromisoformat(text)
    except ValueError:
        return None
    if parsed.tzinfo is None:
        parsed = parsed.replace(tzinfo=timezone.utc)
    return parsed.astimezone(timezone.utc)


def _build_todo_list_response(
    meta: Dict[str, Any],
    due_only: bool = False,
    completion: Optional[Dict[str, Any]] = None,
) -> TodoListResponse:
    now = datetime.now(timezone.utc)
    raw_items = meta.get("task_items") if isinstance(meta, dict) else []
    if not isinstance(raw_items, list):
        raw_items = []

    task_items: List[TodoTaskItemResponse] = []
    tasks: List[str] = []
    for task_index, item in enumerate(raw_items, start=1):
        if not isinstance(item, dict):
            continue
        description = str(item.get("description") or "").strip()
        if not description:
            continue
        task_public_id = _coerce_task_id(item.get("task_id")) or task_index
        next_run = _parse_todo_datetime(item.get("next_run_at"))
        is_due = bool(next_run and next_run <= now)
        if due_only and not is_due:
            continue
        tasks.append(description)
        task_items.append(
            TodoTaskItemResponse(
                taskId=task_public_id,
                taskDescription=description,
                scheduledFor=item.get("scheduled_for"),
                nextRunAt=item.get("next_run_at"),
                recurrence=item.get("recurrence"),
                lastCompletedAt=item.get("last_completed_at"),
                createdAt=item.get("created_at"),
                updatedAt=item.get("updated_at"),
                isDue=is_due,
            )
        )

    # Backward compatibility with old store schema that only had "tasks": ["..."].
    if not task_items and not due_only:
        raw_legacy_tasks = meta.get("tasks") if isinstance(meta, dict) else []
        if isinstance(raw_legacy_tasks, list):
            tasks = [str(t) for t in raw_legacy_tasks]
            task_items = [
                TodoTaskItemResponse(taskId=i + 1, taskDescription=desc)
                for i, desc in enumerate(tasks)
            ]

    completion_payload: Optional[TodoCompletionMetaResponse] = None
    if isinstance(completion, dict):
        try:
            completion_payload = TodoCompletionMetaResponse(
                taskId=int(completion.get("taskId")),
                rescheduled=bool(completion.get("rescheduled", False)),
                nextRunAt=completion.get("nextRunAt"),
            )
        except (TypeError, ValueError):
            completion_payload = None

    return TodoListResponse(
        tasks=tasks,
        taskItems=task_items,
        updated_at=meta.get("updated_at"),
        completion=completion_payload,
    )


@app.get("/v1/todo", response_model=TodoListResponse)
async def todo_list(
    due_only: bool = False,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Return the authenticated user's todo list. Requires JWT."""
    if not TODO_STORE_AVAILABLE or not _todo_store:
        raise HTTPException(status_code=503, detail="Todo store is not available.")
    user_key = current_user["username"]
    meta = _todo_store.load_tasks_with_meta(user_key)
    return _build_todo_list_response(meta, due_only=due_only)


@app.get("/v1/todo/due", response_model=TodoListResponse)
async def todo_due_list(current_user: Dict[str, Any] = Depends(get_current_user)):
    """Return due scheduled tasks (next_run_at <= now) for the authenticated user."""
    if not TODO_STORE_AVAILABLE or not _todo_store:
        raise HTTPException(status_code=503, detail="Todo store is not available.")
    user_key = current_user["username"]
    meta = _todo_store.load_tasks_with_meta(user_key)
    return _build_todo_list_response(meta, due_only=True)


@app.post("/v1/todo", response_model=TodoListResponse)
async def todo_add(
    request: TodoAddRequest,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Add a task to the authenticated user's todo list. Requires JWT."""
    if not TODO_STORE_AVAILABLE or not _todo_store:
        raise HTTPException(status_code=503, detail="Todo store is not available.")
    user_key = current_user["username"]
    desc = (request.taskDescription or "").strip()
    if not desc:
        raise HTTPException(status_code=400, detail="taskDescription is required.")
    recurrence = _todo_recurrence_to_dict(request.recurrence)
    try:
        _todo_store.add_task(
            user_key,
            desc,
            scheduled_for=request.scheduledFor,
            recurrence=recurrence,
        )
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    meta = _todo_store.load_tasks_with_meta(user_key)
    return _build_todo_list_response(meta)


@app.patch("/v1/todo/{task_id}", response_model=TodoListResponse)
async def todo_update(
    task_id: int,
    request: TodoUpdateRequest,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Update a task by stable task ID. Requires JWT."""
    if not TODO_STORE_AVAILABLE or not _todo_store:
        raise HTTPException(status_code=503, detail="Todo store is not available.")
    user_key = current_user["username"]
    desc = request.taskDescription
    if desc is not None:
        desc = desc.strip()
        if not desc:
            raise HTTPException(status_code=400, detail="taskDescription is required.")
    if (
        desc is None
        and request.scheduledFor is None
        and request.recurrence is None
        and not request.clearSchedule
        and not request.clearRecurrence
    ):
        raise HTTPException(status_code=400, detail="No update fields provided.")
    try:
        _todo_store.update_task(
            user_key,
            task_id,
            desc,
            scheduled_for=request.scheduledFor,
            recurrence=_todo_recurrence_to_dict(request.recurrence),
            clear_schedule=bool(request.clearSchedule),
            clear_recurrence=bool(request.clearRecurrence),
        )
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    meta = _todo_store.load_tasks_with_meta(user_key)
    return _build_todo_list_response(meta)


@app.delete("/v1/todo/{task_id}", response_model=TodoListResponse)
async def todo_delete(
    task_id: int,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Delete a task by stable task ID. Requires JWT."""
    if not TODO_STORE_AVAILABLE or not _todo_store:
        raise HTTPException(status_code=503, detail="Todo store is not available.")
    user_key = current_user["username"]
    try:
        _todo_store.delete_task(user_key, task_id)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    meta = _todo_store.load_tasks_with_meta(user_key)
    return _build_todo_list_response(meta)


@app.delete("/v1/todo", response_model=TodoListResponse)
async def todo_clear(current_user: Dict[str, Any] = Depends(get_current_user)):
    """Clear all tasks for the authenticated user. Requires JWT."""
    if not TODO_STORE_AVAILABLE or not _todo_store:
        raise HTTPException(status_code=503, detail="Todo store is not available.")
    user_key = current_user["username"]
    _todo_store.clear_tasks(user_key)
    return TodoListResponse(tasks=[], taskItems=[], updated_at=None)


# ============================================================================
# TASK EXECUTION (internal helpers + REST; all REST require authentication)
# ============================================================================

def _write_task_exec_response_to_scratch(
    user_key: str, task_id: Optional[int], status: str, message: str
) -> Optional[str]:
    """
    Write the task execution agent response to a timestamped text file in scratch.
    Called whenever a run ends (completed, paused, awaiting confirmation, or cancelled).
    Includes a clear SUMMARY and "What you can do" so the file conclusively shows outcome and next steps.
    """
    try:
        SCRATCH_DIR.mkdir(parents=True, exist_ok=True)
        now = datetime.now().astimezone()
        timestamp_file = now.strftime("%Y-%m-%d_%H-%M-%S")
        timestamp_human = now.strftime("%Y-%m-%d %H:%M:%S %Z")
        safe_user = re.sub(r"[^\w\-]", "_", (user_key or "unknown")[:64])
        task_suffix = f"_task{task_id}" if task_id is not None else ""
        filename = f"task_exec_{timestamp_file}_{safe_user}{task_suffix}.txt"
        filepath = SCRATCH_DIR / filename
        # Human-readable summary so the file conclusively shows whether task is done, paused, or cancelled
        if status == STATUS_AWAITING_CONFIRMATION:
            summary = "Task run finished. Awaiting your confirmation to mark the task complete."
            what_to_do = "Confirm completion (e.g. 'mark task complete' or use todo complete), or ask for status."
        elif status == STATUS_PAUSED_AWAITING_FEEDBACK:
            summary = "Paused for your feedback or input."
            what_to_do = "Send your input to resume (e.g. via chat or execute/resume)."
        elif status == STATUS_CANCELLED:
            summary = "Execution was cancelled."
            what_to_do = "You can start a new execution for a task if needed."
        else:
            summary = f"Status: {status}"
            what_to_do = "Check status or resume/cancel as appropriate."
        lines = [
            "Task execution response",
            "SUMMARY: " + summary,
            "",
            f"Date-time: {timestamp_human}",
            f"Status: {status}",
            f"Task ID: {task_id}" if task_id is not None else "Task ID: (none)",
            "",
            "What you can do: " + what_to_do,
            "",
            "--- Agent response ---",
            "",
            message or "(No response text)",
        ]
        filepath.write_text("\n".join(lines), encoding="utf-8")
        print(f"[TASK_EXEC] Wrote response to {filepath}", flush=True)
        return filename
    except Exception as e:
        print(f"[TASK_EXEC] Failed to write response to scratch: {e}", flush=True)
        return None


def _safe_task_execution_diagnostics(executor: Any) -> Dict[str, Any]:
    """Best-effort fetch of executor diagnostics for learning capture."""
    if not executor or not hasattr(executor, "get_run_diagnostics"):
        return {}
    try:
        diagnostics = executor.get_run_diagnostics()
        return diagnostics if isinstance(diagnostics, dict) else {}
    except Exception as e:
        print(f"[TASK_EXEC] Failed to read executor diagnostics: {e}", flush=True)
        return {}


def _task_execution_metadata(
    user_key: str,
    state: Optional[Dict[str, Any]],
    diagnostics: Optional[Dict[str, Any]] = None,
    *,
    phase: Optional[str] = None,
) -> Dict[str, Any]:
    state = state if isinstance(state, dict) else {}
    diagnostics = diagnostics if isinstance(diagnostics, dict) else {}
    task_description = _truncate_monitor_text(state.get("task_description") or "", max_chars=220)
    status = str(state.get("status") or "").strip() or None
    metadata: Dict[str, Any] = {
        "user_key": user_key,
        "task_id": _coerce_task_id(state.get("task_id")),
        "workflow_name": task_description or f"Task {state.get('task_id') or 'unknown'}",
        "task_description": task_description,
        "status": status,
        "phase": phase or status or "running",
        "current_step": diagnostics.get("iterations"),
        "total_steps": diagnostics.get("max_iterations"),
        "elapsed_seconds": diagnostics.get("elapsed_seconds"),
        "tool_success_count": diagnostics.get("tool_success_count"),
        "tool_failure_count": diagnostics.get("tool_failure_count"),
    }
    run_id = str(state.get("run_id") or "").strip()
    if run_id:
        metadata["run_id"] = run_id
    return metadata


def _task_execution_summary(
    state: Optional[Dict[str, Any]],
    diagnostics: Optional[Dict[str, Any]] = None,
    fallback_message: str = "",
) -> str:
    diagnostics = diagnostics if isinstance(diagnostics, dict) else {}
    task_id = _coerce_task_id((state or {}).get("task_id"))
    iterations = diagnostics.get("iterations")
    max_iterations = diagnostics.get("max_iterations")
    message = _truncate_monitor_text(fallback_message or (state or {}).get("message") or "", max_chars=260)
    base = f"Task {task_id}" if task_id is not None else "Task execution"
    if isinstance(iterations, int) and isinstance(max_iterations, int) and max_iterations > 0:
        base += f" step {iterations}/{max_iterations}"
    if message:
        return f"{base}. {message}"
    return base


async def _record_task_execution_learning(
    user_key: str,
    state: Optional[Dict[str, Any]],
    status: str,
    message: str,
    source_phase: str,
) -> None:
    """Persist task outcome as experience memory for future task guidance."""
    if not MEMORY_AVAILABLE or not memory_manager:
        return
    if not hasattr(memory_manager, "record_task_outcome"):
        return
    if not isinstance(state, dict):
        return

    task_description = str(state.get("task_description") or "").strip()
    if not task_description:
        return

    diagnostics = _safe_task_execution_diagnostics(state.get("executor"))
    tool_usage = diagnostics.get("tool_usage_counts") or {}
    if isinstance(tool_usage, dict):
        tool_names = [str(name) for name in tool_usage.keys() if name]
    else:
        tool_names = []
    if not tool_names:
        tool_names = [str(name) for name in (diagnostics.get("tools_used") or []) if name]

    summary = str(message or "").strip()
    error_hint = str(diagnostics.get("last_error") or "").strip()
    if not error_hint and str(status or "").lower() in {"failed", "failure", "error"}:
        error_hint = summary

    metadata = {
        "user_key": user_key,
        "task_id": state.get("task_id"),
        "source_phase": source_phase,
        "iterations": diagnostics.get("iterations"),
        "max_iterations": diagnostics.get("max_iterations"),
        "elapsed_seconds": diagnostics.get("elapsed_seconds"),
        "tool_usage_counts": tool_usage if isinstance(tool_usage, dict) else {},
        "tool_success_count": diagnostics.get("tool_success_count"),
        "tool_failure_count": diagnostics.get("tool_failure_count"),
        "tool_error_messages": diagnostics.get("tool_error_messages") or [],
    }
    try:
        await memory_manager.record_task_outcome(
            task_description=task_description,
            status=status,
            message=message,
            summary=summary,
            error=error_hint,
            tool_names=tool_names,
            metadata=metadata,
            source="task_execution",
        )
    except Exception as e:
        print(f"[TASK_EXEC] Failed to record learning memory: {e}", flush=True)


async def _run_task_loop_background(user_key: str, task_id: int, executor: Any) -> None:
    """
    Run executor.run_loop() in background; update or clear task_execution_state in finally
    so we never leave state stuck as 'executing' on exception or cancel.
    Scheduled runs are auto-completed when they reach awaiting_confirmation.
    """
    status = STATUS_AWAITING_CONFIRMATION
    message = "Execution stopped unexpectedly."
    try:
        status, message = await executor.run_loop()
    except Exception as e:
        status = STATUS_AWAITING_CONFIRMATION
        message = str(e)
        print(f"[TASK_EXEC] run_loop error for user {user_key}: {e}", flush=True)
    finally:
        # Only update if this run still belongs to this executor instance.
        # Status may have changed (e.g. cancellation requested) while the loop was in-flight.
        state = _get_task_run_state(user_key, task_id)
        if state and state.get("executor") is executor:
            diagnostics = _safe_task_execution_diagnostics(executor)
            auto_completion_note = _auto_complete_scheduled_execution(user_key, state, status)
            if auto_completion_note:
                base_message = (message or "").strip()
                message = f"{base_message}\n\n{auto_completion_note}".strip() if base_message else auto_completion_note
            state["status"] = status
            state["message"] = message or ""
            await _record_task_execution_learning(
                user_key=user_key,
                state=state,
                status=status,
                message=message or "",
                source_phase="background_run",
            )
            # Always capture agent response to scratch with timestamp (paused, awaiting, cancelled, done)
            task_log_filename = _write_task_exec_response_to_scratch(user_key, state.get("task_id"), status, message or "")
            await _maybe_notify_telegram_task_completion(user_key, state, status, message or "")
            monitor_run_id = str(state.get("monitor_run_id") or "").strip()
            if monitor_run_id:
                _monitor_run_finish(
                    monitor_run_id,
                    status=status,
                    summary=_task_execution_summary(state, diagnostics, message or ""),
                    metadata=_task_execution_metadata(user_key, state, diagnostics, phase=status),
                    log_file=task_log_filename,
                    log_excerpt=_read_monitor_log_excerpt((SCRATCH_DIR / task_log_filename) if task_log_filename else None),
                )
                if not _is_task_execution_terminal_status(status):
                    state["monitor_run_id"] = None
            # Terminal runs should clear execution state so status returns to idle.
            if _is_task_execution_terminal_status(status):
                _remove_task_run_state(user_key, task_id)


async def _task_execute_start(user_key: str, task_id: int, prompt_override: Optional[str]) -> tuple:
    """Start task execution for user_key/task_id. Runs loop in background; returns immediately."""
    if not TASK_EXECUTION_AVAILABLE or not TodoTaskExecutor:
        raise HTTPException(status_code=503, detail="Task execution is not available.")
    if not TODO_STORE_AVAILABLE or not _todo_store:
        raise HTTPException(status_code=503, detail="Todo store is not available.")
    normalized_task_id = _coerce_task_id(task_id)
    if normalized_task_id is None:
        raise HTTPException(status_code=400, detail="Invalid task ID.")

    _cleanup_terminal_task_runs(user_key)
    active_same_task = _get_task_run_state(user_key, normalized_task_id)
    if active_same_task and not _is_task_execution_terminal_status(_state_status_lower(active_same_task)):
        raise HTTPException(status_code=409, detail=f"Task {normalized_task_id} is already executing.")

    task_description = ""
    task_is_scheduled = False
    task_item_id: Optional[str] = None
    try:
        meta = _todo_store.load_tasks_with_meta(user_key)
        task_items = meta.get("task_items") if isinstance(meta, dict) else None
        task_item: Optional[Dict[str, Any]] = None
        if isinstance(task_items, list):
            for item in task_items:
                if not isinstance(item, dict):
                    continue
                if _coerce_task_id(item.get("task_id")) == normalized_task_id:
                    task_item = item
                    break
            # Backward compatibility for legacy datasets with positional task IDs only.
            if task_item is None and 1 <= normalized_task_id <= len(task_items):
                candidate = task_items[normalized_task_id - 1]
                if isinstance(candidate, dict):
                    task_item = candidate
        if not isinstance(task_item, dict):
            raise HTTPException(status_code=400, detail="Invalid task ID.")

        item_desc = str(task_item.get("description") or "").strip()
        if item_desc:
            task_description = item_desc
        item_id = str(task_item.get("id") or "").strip()
        if item_id:
            task_item_id = item_id
        task_is_scheduled = bool(
            task_item.get("scheduled_for")
            or task_item.get("next_run_at")
            or task_item.get("recurrence")
        )
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Failed to load task metadata: {exc}")
    if not task_description:
        raise HTTPException(status_code=400, detail="Invalid task ID.")
    telegram_chat_ids = _resolve_telegram_chat_ids_for_todo_user(user_key)
    api_key = _first_non_empty_env(
        preferred_api_key_env_names(
            os.getenv("OPENAI_API_BASE", "https://api.openai.com/v1"),
            os.getenv("OPENAI_MODEL", "gpt-4o-mini"),
        )
    )
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail=(
                "No compatible API key is configured for task execution. "
                "Set OPENAI_API_KEY / MCP_LLM_OPENAI_API_KEY, or MINIMAX_API_KEY / "
                "MCP_LLM_MINIMAX_API_KEY when using Minimax."
            ),
        )
    monitor_run_id = _monitor_run_start(
        "task_execution",
        "task-run",
        input_text=prompt_override or task_description,
        metadata={
            "user_key": user_key,
            "task_id": normalized_task_id,
            "workflow_name": _truncate_monitor_text(task_description, max_chars=220),
            "task_description": _truncate_monitor_text(task_description, max_chars=220),
            "current_step": 0,
            "total_steps": TASK_EXECUTION_MAX_ITERATIONS,
            "phase": STATUS_EXECUTING,
        },
    )
    try:
        experience_guidance = ""
        if MEMORY_AVAILABLE and memory_manager and hasattr(memory_manager, "build_task_execution_guidance"):
            try:
                experience_guidance = await memory_manager.build_task_execution_guidance(task_description=task_description)
                if experience_guidance:
                    print(
                        f"[TASK_EXEC] Loaded experience guidance for user {user_key} task {normalized_task_id}",
                        flush=True,
                    )
            except Exception as e:
                print(f"[TASK_EXEC] Failed to load experience guidance: {e}", flush=True)
        executor = TodoTaskExecutor(
            api_key=api_key,
            task_id=normalized_task_id,
            task_description=task_description,
            prompt_override=prompt_override,
            max_iterations=TASK_EXECUTION_MAX_ITERATIONS,
            tool_executor=execute_tool_for_philosopher,
            get_tools_func=get_all_available_tools,
            experience_guidance=experience_guidance,
            progress_callback=None,
        )
        # Set state to 'executing' before starting so cancel/status work and we never leave state stuck.
        state = {
            "task_id": normalized_task_id,
            "task_item_id": task_item_id,
            "run_id": f"task-{normalized_task_id}-{secrets.token_hex(4)}",
            "monitor_run_id": monitor_run_id,
            "status": STATUS_EXECUTING,
            "executor": executor,
            "message": None,
            "task_description": task_description,
            "is_scheduled": task_is_scheduled,
            "telegram_chat_ids": telegram_chat_ids,
        }
    except Exception as exc:
        _monitor_run_finish(
            monitor_run_id,
            status="error",
            summary=f"Failed to start task execution: {exc}",
            metadata={
                "user_key": user_key,
                "task_id": normalized_task_id,
                "workflow_name": _truncate_monitor_text(task_description, max_chars=220),
                "phase": "error",
            },
        )
        raise HTTPException(status_code=500, detail=f"Failed to start task execution: {exc}")

    async def _task_monitor(event: str, payload: Dict[str, Any]) -> None:
        diagnostics = _safe_task_execution_diagnostics(executor)
        metadata = _task_execution_metadata(
            user_key,
            state,
            diagnostics,
            phase=payload.get("phase") or STATUS_EXECUTING,
        )
        for key in ("current_step", "total_steps", "tool_call_count"):
            value = payload.get(key)
            if value is not None:
                metadata[key] = value
        summary = payload.get("message") or _task_execution_summary(state, diagnostics)
        _monitor_run_update(monitor_run_id, summary=summary, metadata=metadata)
        if payload.get("message"):
            _monitor_run_note(monitor_run_id, payload["message"])

    executor.progress_callback = _task_monitor
    _set_task_run_state(user_key, normalized_task_id, state)
    asyncio.create_task(_run_task_loop_background(user_key, normalized_task_id, executor))

    active_ids = [tid for tid, _ in _active_task_runs(user_key)]
    if len(active_ids) <= 1:
        return (STATUS_EXECUTING, f"Task {normalized_task_id} execution started.")
    return (
        STATUS_EXECUTING,
        f"Task {normalized_task_id} execution started. Active task executions: {', '.join(str(t) for t in active_ids)}.",
    )


async def _task_execute_resume(user_key: str, user_message: str, task_id: Optional[int] = None) -> tuple:
    """Resume paused execution. Returns (status, message, task_id)."""
    if not TASK_EXECUTION_AVAILABLE:
        raise HTTPException(status_code=503, detail="Task execution is not available.")
    _cleanup_terminal_task_runs(user_key)
    runs = _get_user_task_runs(user_key, create=False)
    if not runs:
        raise HTTPException(status_code=400, detail="No paused execution to resume.")

    normalized_tid = _coerce_task_id(task_id)
    target_tid: Optional[int] = None
    state: Optional[Dict[str, Any]] = None
    if normalized_tid is not None:
        candidate = runs.get(normalized_tid)
        if not candidate or _state_status_lower(candidate) != STATUS_PAUSED_AWAITING_FEEDBACK:
            raise HTTPException(status_code=400, detail=f"Task {normalized_tid} is not paused awaiting feedback.")
        target_tid = normalized_tid
        state = candidate
    else:
        paused_items = [(tid, s) for tid, s in runs.items() if _state_status_lower(s) == STATUS_PAUSED_AWAITING_FEEDBACK]
        if not paused_items:
            raise HTTPException(status_code=400, detail="No paused execution to resume.")
        if len(paused_items) > 1:
            paused_ids = ", ".join(str(tid) for tid, _ in sorted(paused_items, key=lambda item: item[0]))
            raise HTTPException(
                status_code=400,
                detail=f"Multiple paused tasks found ({paused_ids}). Provide taskId to resume a specific task.",
            )
        target_tid, state = paused_items[0]

    if not state or target_tid is None:
        raise HTTPException(status_code=400, detail="Execution state lost. Start a new execution.")

    executor = state.get("executor")
    if not executor:
        _remove_task_run_state(user_key, target_tid)
        raise HTTPException(status_code=400, detail="Execution state lost. Start a new execution.")
    resume_monitor_run_id = _monitor_run_start(
        "task_execution",
        "resume-run",
        input_text=user_message or str(state.get("task_description") or ""),
        metadata={
            "user_key": user_key,
            "task_id": target_tid,
            "workflow_name": _truncate_monitor_text(state.get("task_description") or f"Task {target_tid}", max_chars=220),
            "task_description": _truncate_monitor_text(state.get("task_description") or "", max_chars=220),
            "current_step": _safe_task_execution_diagnostics(executor).get("iterations"),
            "total_steps": _safe_task_execution_diagnostics(executor).get("max_iterations"),
            "phase": STATUS_EXECUTING,
        },
    )
    state["monitor_run_id"] = resume_monitor_run_id
    state["status"] = STATUS_EXECUTING
    state["message"] = "Task execution resumed."

    async def _task_resume_monitor(event: str, payload: Dict[str, Any]) -> None:
        diagnostics = _safe_task_execution_diagnostics(executor)
        metadata = _task_execution_metadata(
            user_key,
            state,
            diagnostics,
            phase=payload.get("phase") or STATUS_EXECUTING,
        )
        for key in ("current_step", "total_steps", "tool_call_count"):
            value = payload.get(key)
            if value is not None:
                metadata[key] = value
        summary = payload.get("message") or _task_execution_summary(state, diagnostics)
        _monitor_run_update(resume_monitor_run_id, summary=summary, metadata=metadata)
        if payload.get("message"):
            _monitor_run_note(resume_monitor_run_id, payload["message"])

    previous_progress_callback = getattr(executor, "progress_callback", None)
    executor.progress_callback = _task_resume_monitor
    executor.add_user_message(user_message or "")
    try:
        status, message = await executor.run_loop()
        state["status"] = status
        state["message"] = message or ""
        await _record_task_execution_learning(
            user_key=user_key,
            state=state,
            status=status,
            message=message or "",
            source_phase="resume_run",
        )
        # Capture agent response to scratch (paused or awaiting confirmation after resume)
        task_log_filename = _write_task_exec_response_to_scratch(user_key, state.get("task_id"), status, message or "")
        await _maybe_notify_telegram_task_completion(user_key, state, status, message or "")
        _monitor_run_finish(
            resume_monitor_run_id,
            status=status,
            summary=_task_execution_summary(state, _safe_task_execution_diagnostics(executor), message or ""),
            metadata=_task_execution_metadata(user_key, state, _safe_task_execution_diagnostics(executor), phase=status),
            log_file=task_log_filename,
            log_excerpt=_read_monitor_log_excerpt((SCRATCH_DIR / task_log_filename) if task_log_filename else None),
        )
        if _is_task_execution_terminal_status(status):
            _remove_task_run_state(user_key, target_tid)
        return (status, message or "Resumed.", _coerce_task_id(state.get("task_id")) or target_tid)
    except Exception as e:
        state["status"] = STATUS_AWAITING_CONFIRMATION
        state["message"] = str(e)
        await _record_task_execution_learning(
            user_key=user_key,
            state=state,
            status="failed",
            message=str(e),
            source_phase="resume_error",
        )
        task_log_filename = _write_task_exec_response_to_scratch(user_key, state.get("task_id"), STATUS_AWAITING_CONFIRMATION, str(e))
        _monitor_run_finish(
            resume_monitor_run_id,
            status="error",
            summary=_task_execution_summary(state, _safe_task_execution_diagnostics(executor), str(e)),
            metadata=_task_execution_metadata(user_key, state, _safe_task_execution_diagnostics(executor), phase="error"),
            log_file=task_log_filename,
            log_excerpt=_read_monitor_log_excerpt((SCRATCH_DIR / task_log_filename) if task_log_filename else None),
        )
        return (STATUS_AWAITING_CONFIRMATION, str(e), _coerce_task_id(state.get("task_id")) or target_tid)
    finally:
        executor.progress_callback = previous_progress_callback
        if not _is_task_execution_terminal_status(_state_status_lower(state)):
            state["monitor_run_id"] = None


@app.post("/v1/todo/execute", response_model=TodoExecuteResponse)
async def todo_execute(
    request: TodoExecuteRequest,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Start task execution for the given todo task. Requires JWT. Supports parallel runs per task."""
    user_key = current_user["username"]
    status, message = await _task_execute_start(user_key, request.taskId, request.promptOverride)
    return TodoExecuteResponse(status=status, message=message, taskId=request.taskId)


@app.post("/v1/todo/execute/resume", response_model=TodoExecuteResponse)
async def todo_execute_resume(
    resume_request: TodoResumeRequest,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Resume a paused task execution with user feedback. Requires JWT."""
    user_key = current_user["username"]
    status, message, resumed_task_id = await _task_execute_resume(
        user_key,
        resume_request.userMessage or "",
        resume_request.taskId,
    )
    return TodoExecuteResponse(status=status, message=message, taskId=resumed_task_id)


@app.post("/v1/todo/{task_id}/complete", response_model=TodoListResponse)
async def todo_task_complete(
    task_id: int,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Human verification: complete a task. Repeating tasks are rescheduled; one-time tasks are removed."""
    if not TODO_STORE_AVAILABLE or not _todo_store:
        raise HTTPException(status_code=503, detail="Todo store is not available.")
    user_key = current_user["username"]
    _cleanup_terminal_task_runs(user_key)
    active_state: Optional[Dict[str, Any]] = None
    task_description_for_learning = ""
    completion_task_item_id: Optional[str] = None
    try:
        meta_before = _todo_store.load_tasks_with_meta(user_key)
        task_items_before = meta_before.get("task_items") if isinstance(meta_before, dict) else None
        if isinstance(task_items_before, list):
            matched_item: Optional[Dict[str, Any]] = None
            for item in task_items_before:
                if not isinstance(item, dict):
                    continue
                if _coerce_task_id(item.get("task_id")) == task_id:
                    matched_item = item
                    break
            if matched_item is None and 1 <= task_id <= len(task_items_before):
                # Legacy fallback for datasets without stable task_id values.
                candidate = task_items_before[task_id - 1]
                if isinstance(candidate, dict):
                    matched_item = candidate
            if isinstance(matched_item, dict):
                task_description_for_learning = str(matched_item.get("description") or "").strip()
                completion_task_item_id = str(matched_item.get("id") or "").strip() or None
    except Exception:
        task_description_for_learning = ""
    runs = _get_user_task_runs(user_key, create=False)
    if completion_task_item_id:
        for _, run_state in runs.items():
            if str(run_state.get("task_item_id") or "").strip() == completion_task_item_id:
                active_state = run_state
                break
    if active_state is None:
        active_state = runs.get(task_id)
    if isinstance(active_state, dict):
        task_description_for_learning = (
            str(active_state.get("task_description") or "").strip() or task_description_for_learning
        )
    if isinstance(active_state, dict):
        running_tid = _coerce_task_id(active_state.get("task_id"))
        if running_tid is not None:
            _remove_task_run_state(user_key, running_tid)
    try:
        result = _todo_store.complete_task(user_key, task_id)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    if task_description_for_learning and MEMORY_AVAILABLE and memory_manager and hasattr(memory_manager, "record_task_outcome"):
        completion_state = active_state if isinstance(active_state, dict) else {
            "task_id": task_id,
            "task_description": task_description_for_learning,
            "executor": None,
        }
        completion_note = (
            "Task marked complete by user confirmation."
            + (" Repeating task was rescheduled." if bool(result.get("rescheduled", False)) else "")
        )
        await _record_task_execution_learning(
            user_key=user_key,
            state=completion_state,
            status="confirmed_complete",
            message=completion_note,
            source_phase="user_confirmed_completion",
        )
    meta = {
        "tasks": result.get("tasks", []),
        "task_items": result.get("task_items", []),
        "updated_at": result.get("updated_at"),
    }
    completion = {
        "taskId": task_id,
        "rescheduled": bool(result.get("rescheduled", False)),
        "nextRunAt": result.get("next_run_at"),
    }
    return _build_todo_list_response(meta, completion=completion)


def _task_execute_cancel(user_key: str, task_id: Optional[int] = None) -> tuple:
    """Request soft cancel for an active run. Returns (ok, message, task_id)."""
    active_runs = _active_task_runs(user_key)
    if not active_runs:
        return (False, "No active execution to cancel.", None)

    target_tid: Optional[int] = _coerce_task_id(task_id)
    target_state: Optional[Dict[str, Any]] = None
    if target_tid is not None:
        candidate = _get_task_run_state(user_key, target_tid)
        if not candidate or _is_task_execution_terminal_status(_state_status_lower(candidate)):
            return (False, f"Task {target_tid} is not actively executing.", None)
        target_state = candidate
    elif len(active_runs) == 1:
        target_tid, target_state = active_runs[0]
    else:
        active_ids = ", ".join(str(tid) for tid, _ in active_runs)
        return (
            False,
            f"Multiple active tasks ({active_ids}). Provide taskId to cancel a specific task.",
            None,
        )

    if target_tid is None or target_state is None:
        return (False, "No active execution to cancel.", None)

    executor = target_state.get("executor")
    if executor and hasattr(executor, "request_cancel"):
        executor.request_cancel()
        # Update visible status immediately so callers don't remain stuck on 'executing'
        # while the loop finishes the current step.
        target_state["status"] = STATUS_CANCELLED
        target_state["message"] = "Cancellation requested. The task will stop after the current step."
        monitor_run_id = str(target_state.get("monitor_run_id") or "").strip()
        if monitor_run_id:
            diagnostics = _safe_task_execution_diagnostics(executor)
            _monitor_run_update(
                monitor_run_id,
                summary=_task_execution_summary(target_state, diagnostics, target_state["message"]),
                metadata=_task_execution_metadata(user_key, target_state, diagnostics, phase=STATUS_CANCELLED),
            )
            _monitor_run_note(monitor_run_id, target_state["message"])
        return (True, "Cancellation requested. The task will stop after the current step.", target_tid)
    _remove_task_run_state(user_key, target_tid)
    return (False, "No active execution to cancel.", None)


def _task_execution_status(user_key: str, task_id: Optional[int] = None) -> Optional[Dict[str, Any]]:
    """Return execution status for one task (if task_id provided) or user-level summary."""
    active_runs = _active_task_runs(user_key)
    if not active_runs:
        return None

    normalized_tid = _coerce_task_id(task_id)
    if normalized_tid is not None:
        state = _get_task_run_state(user_key, normalized_tid)
        if not state or _is_task_execution_terminal_status(_state_status_lower(state)):
            return None
        return _state_brief_for_response(user_key, state)

    if len(active_runs) == 1:
        return _state_brief_for_response(user_key, active_runs[0][1])

    run_summaries = [_state_brief_for_response(user_key, state) for _, state in active_runs]
    active_ids = [summary.get("task_id") for summary in run_summaries if _coerce_task_id(summary.get("task_id")) is not None]
    active_ids = sorted(set(int(tid) for tid in active_ids))
    return {
        "status": "multiple",
        "task_id": None,
        "task_ids": active_ids,
        "message": f"Multiple task executions are active: {', '.join(str(tid) for tid in active_ids)}.",
        "runs": run_summaries,
    }


@app.post("/v1/todo/execute/cancel", response_model=TodoExecuteResponse)
async def todo_execute_cancel(
    cancel_request: Optional[TodoCancelRequest] = None,
    taskId: Optional[int] = Query(default=None),
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Cancel an active task execution; task remains in list. Requires JWT."""
    user_key = current_user["username"]
    requested_task_id = taskId
    if requested_task_id is None and isinstance(cancel_request, TodoCancelRequest):
        requested_task_id = cancel_request.taskId
    ok, msg, cancelled_task_id = _task_execute_cancel(user_key, requested_task_id)
    response_status = STATUS_CANCELLED if ok else "idle"
    return TodoExecuteResponse(status=response_status, message=msg, taskId=cancelled_task_id)


@app.get("/v1/todo/execute/status", response_model=TodoExecutionStatusResponse)
async def todo_execute_status(
    taskId: Optional[int] = Query(default=None),
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Return active task execution status with all currently running task IDs."""
    user_key = current_user["username"]
    active_runs = _active_task_runs(user_key)
    run_summaries = [_state_brief_for_response(user_key, state) for _, state in active_runs]
    active_ids = [summary.get("task_id") for summary in run_summaries if _coerce_task_id(summary.get("task_id")) is not None]
    active_ids = sorted(set(int(tid) for tid in active_ids))

    task_summary = _task_execution_status(user_key, taskId) if taskId is not None else None
    summary = _task_execution_status(user_key) if taskId is None else task_summary

    message: Optional[str]
    if taskId is not None:
        if task_summary:
            message = task_summary.get("message") or f"Task {task_summary.get('task_id')} is {task_summary.get('status')}."
        else:
            message = f"No active run for task {taskId}."
    else:
        if summary:
            message = summary.get("message") or f"Active task IDs: {', '.join(str(t) for t in active_ids)}."
        else:
            message = "No task is currently running or paused."

    return TodoExecutionStatusResponse(
        active=bool(run_summaries),
        activeTaskIds=active_ids,
        runs=run_summaries,
        message=message,
        task=task_summary if isinstance(task_summary, dict) else None,
    )


def _is_todo_list_query(message_text: str) -> bool:
    """
    Return True if the message is clearly asking for the current todo/task list.
    Used to skip memory injection so the model uses manageTodoList tool instead of answering from memory.
    """
    if not message_text or len(message_text) > 500:
        return False
    lower = message_text.lower().strip()
    # Phrases that indicate user wants current list (use tool), not memory-based answer
    phrases = (
        "what's on my todo",
        "whats on my todo",
        "what is on my todo",
        "show my todo",
        "list my todo",
        "my todo list",
        "show my tasks",
        "list my tasks",
        "what are my tasks",
        "what's on my list",
        "whats on my list",
        "my tasks",
        "my todos",
        "what is due",
        "what's due",
        "whats due",
        "show due tasks",
        "list due tasks",
        "show overdue tasks",
        "list overdue tasks",
    )
    return any(p in lower for p in phrases)


def _is_memory_context_question(message_text: str) -> bool:
    """
    Return True when a message is likely asking for opinions/knowledge where memory context helps.
    Action/tool requests should return False to avoid noisy memory injection.
    """
    if not message_text:
        return False

    lower = message_text.lower().strip()
    if len(lower) > 800:
        return False

    action_patterns = (
        r"(search|find|look\s+up|get|fetch|retrieve)\s+(for|information\s+about|details\s+about)",
        r"how\s+much\s+(does|do|is|are|cost)",
        r"what'?s\s+the\s+(weather|price|cost|temperature|time)",
        r"(show|display|list|give\s+me)\s+(information|data|details|results)",
        r"(navigate|go\s+to|visit|open|browse)",
        r"(read|write|save|load|upload|download)\s+(the\s+)?(file|document|data)",
        r"(create|make|build|generate|produce)\s+(a|an|the)",
        r"(run|execute|perform|do)\s+(a|an|the)",
    )
    if any(re.search(p, lower, re.IGNORECASE) for p in action_patterns):
        return False

    opinion_patterns = (
        r"what\s+(do\s+you\s+)?think\s+(about|of)",
        r"what'?s\s+(your\s+)?(opinion|view|perspective|take|thoughts?)\s+(on|about|of|regarding)",
        r"what\s+(do\s+you\s+)?know\s+(about|of)",
        r"how\s+do\s+you\s+(feel|see|view|perceive)\s+(about|on|regarding)",
        r"tell\s+me\s+(what\s+you\s+)?(think|know|believe|feel)\s+(about|of|on)",
        r"share\s+(your\s+)?(thoughts?|views?|opinions?|perspective)",
        r"what\s+(do\s+you\s+)?(believe|understand|consider)\s+(about|of|regarding)",
    )
    if any(re.search(p, lower, re.IGNORECASE) for p in opinion_patterns):
        return True

    return bool(re.search(r"^what\s+(do\s+you|'?s\s+your)", lower, re.IGNORECASE))


def _extract_memory_search_query(message_text: str) -> str:
    """Extract likely topic phrase for memory search; falls back to full message."""
    if not message_text:
        return ""
    match = re.search(
        r"(?:think|opinion|view|know|thoughts?|beliefs?|feel|see|perceive|understand|consider|"
        r"insights?|reflections?|contemplations?)\s+(?:about|of|on|regarding)\s+(.+?)(?:\?|$)",
        message_text,
        re.IGNORECASE,
    )
    topic = (match.group(1).strip() if match and match.group(1) else message_text.strip())
    return topic[:500]


def _filter_high_relevance_memories(
    memories: List[Dict[str, Any]],
    memory_manager: Optional[Any] = None,
) -> List[Dict[str, Any]]:
    """
    Keep only high-confidence memory matches:
    - top hit must pass minimum similarity
    - keep matches in a score window close to top score
    - cap final count
    """
    if not memories:
        return []

    sorted_memories = sorted(memories, key=lambda m: float(m.get("similarity", 0.0)), reverse=True)
    top_score = float(sorted_memories[0].get("similarity", 0.0))
    if top_score < MEMORY_AUTO_SEARCH_MIN_SIMILARITY:
        return []

    floor = max(MEMORY_AUTO_SEARCH_MIN_SIMILARITY, top_score - MEMORY_AUTO_SEARCH_SCORE_WINDOW)
    filtered = [m for m in sorted_memories if float(m.get("similarity", 0.0)) >= floor]
    context_safe: List[Dict[str, Any]] = []
    for mem in filtered:
        category = (mem.get("category") or "").strip().lower()
        source = (mem.get("source") or "").strip().lower()
        memory_type = str(mem.get("memory_type") or "").strip().lower()
        text = str(mem.get("text") or "")
        if category in MEMORY_CONTEXT_BLOCKED_CATEGORIES:
            continue
        if memory_type in MEMORY_CONTEXT_BLOCKED_CATEGORIES:
            continue
        if source in MEMORY_CONTEXT_BLOCKED_SOURCES:
            continue
        if MEMORY_CONTEXT_OPERATIONAL_PATTERN.search(text):
            continue
        context_safe.append(mem)

    if memory_manager:
        try:
            filter_fn = getattr(memory_manager, "filter_memories_for_conversation_context", None)
            if callable(filter_fn):
                maybe_filtered = filter_fn(context_safe)
                if isinstance(maybe_filtered, list):
                    context_safe = maybe_filtered
        except Exception as e:
            print(f"Warning: MemoryManager context filter failed: {e}")
    filtered = context_safe
    return filtered[:MEMORY_AUTO_SEARCH_LIMIT]


def _message_requests_proxy_restart(message_text: str) -> bool:
    """
    Return True for explicit restart command-style phrases.
    This keeps restart deterministic for chat without relying on LLM tool selection.
    """
    if not message_text:
        return False
    normalized = re.sub(r"\s+", " ", message_text.strip().lower())
    normalized = normalized.rstrip(" .!?,;:")
    explicit_commands = {
        "/restartproxy",
        "/restart_proxy",
        "restart proxy",
        "restart proxy server",
        "restart the proxy",
        "restart the proxy server",
    }
    return normalized in explicit_commands


_TELEGRAM_SMALL_TALK_PATTERNS = [
    re.compile(r"^(?:hi|hello|hey|hiya|yo|howdy)(?:\s+(?:cat|catbot|there))?[!.?]*$", re.IGNORECASE),
    re.compile(
        r"^(?:hi|hello|hey)[,!\s]*(?:how are you|how'?s it going|how are things|what'?s up)[?.!\s]*$",
        re.IGNORECASE,
    ),
    re.compile(r"^(?:how are you|how'?s it going|how are things|what'?s up)[?.!\s]*$", re.IGNORECASE),
    re.compile(
        r"^(?:thanks|thank you|cheers|cool|awesome|great|nice|sounds good|got it|ok|okay|alright)[!.?\s]*$",
        re.IGNORECASE,
    ),
    re.compile(r"^(?:bye|goodbye|see ya|see you|cya|ttyl|talk to you later|good night)[!.?\s]*$", re.IGNORECASE),
]


def _is_telegram_small_talk_message(message_text: str) -> bool:
    """Return True for low-intent conversational turns that should not trigger progress chatter."""
    normalized = re.sub(r"\s+", " ", (message_text or "").strip())
    if not normalized:
        return False
    if len(normalized) > 120:
        return False
    return any(pattern.fullmatch(normalized) for pattern in _TELEGRAM_SMALL_TALK_PATTERNS)


def _should_emit_telegram_status_updates(message_text: str) -> bool:
    """
    Decide whether Telegram should emit progress/status messages for this turn.
    Suppress them for low-intent small talk so the chat feels conversational.
    """
    if _is_telegram_small_talk_message(message_text):
        return False
    return True


def _coerce_telegram_response_text(content: Any) -> str:
    """Normalize LLM message content (string or structured parts) into plain text."""
    return coerce_message_text(content)


async def _restart_proxy_after_delay(trigger: str, requested_by: str) -> None:
    """Restart this proxy process after a short delay so current HTTP response can complete."""
    global _proxy_restart_scheduled
    await asyncio.sleep(PROXY_RESTART_DELAY_SECONDS)
    print(
        f"[RESTART] Restarting proxy server now (trigger={trigger}, requested_by={requested_by})",
        flush=True,
    )
    # Restart using the canonical module entrypoint so startup behavior remains consistent.
    try:
        os.execv(sys.executable, [sys.executable, "-m", "src.servers.proxy_server"])
    except Exception as exc:
        _proxy_restart_scheduled = False
        print(f"[RESTART] Failed to restart proxy server: {exc}", flush=True)


async def _request_proxy_restart(trigger: str, requested_by: str) -> Dict[str, Any]:
    """Schedule a process restart once; subsequent requests before restart are acknowledged idempotently."""
    global _proxy_restart_scheduled
    if not PROXY_RESTART_ENABLED:
        return {
            "success": False,
            "scheduled": False,
            "message": "Proxy restart is disabled by server configuration.",
        }
    if _proxy_restart_scheduled:
        return {
            "success": True,
            "scheduled": True,
            "alreadyScheduled": True,
            "message": "Proxy restart is already scheduled and will start shortly.",
        }
    _proxy_restart_scheduled = True
    asyncio.create_task(_restart_proxy_after_delay(trigger=trigger, requested_by=requested_by))
    return {
        "success": True,
        "scheduled": True,
        "message": "Proxy restart scheduled. Reconnect in a few seconds to load updated tools.",
    }


@app.post("/v1/telegram/chat", response_model=TelegramChatResponse)
async def telegram_chat_endpoint(raw_request: Request, request: TelegramChatRequest):
    """Process a Telegram chat message via OpenAI-compatible API."""
    _validate_telegram_secret(raw_request)

    conversation_id = request.conversation_id or request.user_id or "default"
    attachment_records = _store_json_attachments(
        request.attachments,
        conversation_id=conversation_id,
        source="telegram",
    )
    message_text = _augment_message_with_attachments(request.message or "", attachment_records)
    if not message_text:
        raise HTTPException(status_code=400, detail="message or attachments are required")

    request_id = request.request_id or f"telegram-{conversation_id}-{int(time.time() * 1000)}"
    emit_status_updates = _should_emit_telegram_status_updates(message_text)

    if emit_status_updates:
        await _start_status_session(
            conversation_id=conversation_id,
            request_id=request_id,
            channel="telegram",
            initial_state="On it. I'm getting started now.",
            phase="llm_request",
        )

    # Explicit command path for reliable Telegram-triggered proxy restarts (no LLM/tool-call required).
    if _message_requests_proxy_restart(message_text):
        restart_result = await _request_proxy_restart(
            trigger="telegram_command",
            requested_by=f"telegram:{request.user_id or conversation_id}",
        )
        reply_text = restart_result.get("message", "Proxy restart requested.")
        await _finish_status_session(
            conversation_id=conversation_id,
            request_id=request_id,
            final_state="Done: proxy restart scheduled",
            phase="done",
        )
        history = telegram_conversations.setdefault(conversation_id, [])
        history.append({"role": "user", "content": message_text})
        history.append({"role": "assistant", "content": reply_text})
        trim_telegram_history(history)
        return TelegramChatResponse(reply=reply_text, conversation_id=conversation_id, usage=None)

    # Primary Telegram call uses the endpoint-appropriate OpenAI-compatible key.
    primary_api_key = _first_non_empty_env(
        preferred_api_key_env_names(
            TELEGRAM_OPENAI_BASE_URL,
            request.model or TELEGRAM_DEFAULT_MODEL,
        )
    )
    if not primary_api_key:
        await _finish_status_session(
            conversation_id=conversation_id,
            request_id=request_id,
            final_state="Failed: server missing API key",
            phase="error",
        )
        raise HTTPException(
            status_code=503,
            detail=(
                "No compatible API key is configured for the Telegram chat provider. "
                "Set OPENAI_API_KEY / MCP_LLM_OPENAI_API_KEY, or MINIMAX_API_KEY / "
                "MCP_LLM_MINIMAX_API_KEY when using Minimax."
            ),
        )

    todo_user_key = _resolve_todo_user_for_telegram(conversation_id, request.user_id)
    history = telegram_conversations.setdefault(conversation_id, [])

    if request.history is not None:
        history = [
            {"role": msg.role, "content": msg.content}
            for msg in request.history
            if msg.content
        ]
        telegram_conversations[conversation_id] = history
        trim_telegram_history(history)

    history.append({"role": "user", "content": message_text})
    trim_telegram_history(history)

    system_prompt = request.system_prompt
    if system_prompt is None:
        if TELEGRAM_TOOLS_ENABLED:
            system_prompt = _get_telegram_system_prompt_with_tools(conversation_id, todo_user_key=todo_user_key)
        else:
            system_prompt = _get_telegram_system_prompt_base()

    system_prompt = _compose_system_prompt_with_context(system_prompt)

    # Retrieve relevant memories only for opinion/knowledge-style prompts.
    memory_context = ""
    if (
        MEMORY_AVAILABLE
        and memory_manager
        and not _is_todo_list_query(message_text)
        and _is_memory_context_question(message_text)
    ):
        try:
            search_query = _extract_memory_search_query(message_text)
            candidate_memories = await memory_manager.search_memories(
                query=search_query,
                limit=max(5, MEMORY_AUTO_SEARCH_LIMIT * 2),
                similarity_threshold=MEMORY_AUTO_SEARCH_CANDIDATE_THRESHOLD,
            )

            relevant_memories = _filter_high_relevance_memories(
                candidate_memories,
                memory_manager=memory_manager,
            )
            if relevant_memories:
                print(
                    f"Memory context: kept {len(relevant_memories)}/{len(candidate_memories)} "
                    f"for query '{search_query[:50]}...'"
                )
                memory_context = "\n\nRelevant context from previous conversations:\n"
                for i, mem in enumerate(relevant_memories, 1):
                    memory_context += f"{i}. {mem.get('text', '')}\n"
                memory_context += "\nUse this context to provide more personalized and relevant responses."
            else:
                print(f"Memory context: no high-relevance matches for query '{search_query[:50]}...'")
        except Exception as e:
            print(f"Warning: Failed to retrieve memories: {e}")
            import traceback
            print(traceback.format_exc())

    # Add memory context to system prompt
    if memory_context:
        system_prompt = system_prompt + memory_context

    messages: List[Dict[str, str]] = []
    if system_prompt:
        messages.append({"role": "system", "content": system_prompt})
    messages.extend(history)
    messages = _attach_vision_parts_to_latest_user_message(messages, request.attachments)

    model_name = request.model or TELEGRAM_DEFAULT_MODEL
    if not model_name:
        raise HTTPException(status_code=400, detail="No model configured for Telegram chat")

    payload: Dict[str, Any] = {
        "model": model_name,
        "messages": messages,
    }
    if TELEGRAM_TOOLS_ENABLED and _telegram_tools is not None:
        telegram_tools_payload = _get_telegram_combined_openai_tools()
        if telegram_tools_payload:
            payload["tools"] = telegram_tools_payload
            payload["tool_choice"] = "auto"

    if request.temperature is not None:
        payload["temperature"] = request.temperature

    if request.max_output_tokens is not None:
        payload["max_tokens"] = request.max_output_tokens

    headers = {"Content-Type": "application/json"}
    if primary_api_key:
        headers["Authorization"] = f"Bearer {primary_api_key}"

    if OPENAI_ORG_ID:
        headers["OpenAI-Organization"] = OPENAI_ORG_ID

    if OPENAI_PROJECT_ID:
        headers["OpenAI-Project"] = OPENAI_PROJECT_ID

    url = build_openai_url(TELEGRAM_OPENAI_CHAT_PATH)
    url = _normalize_chat_endpoint(url)

    max_tokens = _get_max_tokens_from_payload(payload)
    summarized = False
    if isinstance(payload.get("messages"), list):
        estimated = _estimate_total_tokens(payload["messages"], max_tokens)
        if estimated > get_max_token_limit():
            payload["messages"] = await _summarize_messages_for_budget_proxy(
                payload["messages"],
                endpoint=url,
                headers=headers,
                model_name=model_name,
                max_tokens=max_tokens,
                large_model=LARGE_PAYLOAD_MODEL,
                large_endpoint=LARGE_PAYLOAD_ENDPOINT,
            )
            summarized = True

    used_mcp_fallback = False
    try:
        response = await _call_chat_completion(url, headers, payload, timeout_seconds=TELEGRAM_CHAT_TIMEOUT)
    except httpx.RequestError as exc:
        print(f"Telegram chat request error: {exc}")
        fallback_response, fallback_error = await _attempt_mcp_chat_fallback(
            primary_headers=headers,
            payload=payload,
            timeout_seconds=TELEGRAM_CHAT_TIMEOUT,
            source_label="telegram_chat_request_error",
        )
        if fallback_response is None:
            await _finish_status_session(
                conversation_id=conversation_id,
                request_id=request_id,
                final_state="Failed: could not contact model service",
                phase="error",
            )
            raise HTTPException(
                status_code=502,
                detail=f"Failed to contact language model service. Primary error: {exc}. Fallback error: {fallback_error or 'not configured'}",
            ) from exc
        response = fallback_response
        used_mcp_fallback = True

    if response.status_code != 200:
        print(f"Telegram chat API error {response.status_code}: {response.text}")
        detail = response.text
        try:
            error_json = response.json()
            detail = (
                error_json.get("error", {}).get("message")
                or error_json.get("message")
                or detail
            )
        except ValueError:
            pass
        if is_context_limit_error(response.status_code, detail):
            if not summarized and isinstance(payload.get("messages"), list):
                payload["messages"] = await _summarize_messages_for_budget_proxy(
                    payload["messages"],
                    endpoint=url,
                    headers=headers,
                    model_name=model_name,
                    max_tokens=max_tokens,
                    large_model=LARGE_PAYLOAD_MODEL,
                    large_endpoint=LARGE_PAYLOAD_ENDPOINT,
                )
                summarized = True
                try:
                    response = await _call_chat_completion(url, headers, payload, timeout_seconds=TELEGRAM_CHAT_TIMEOUT)
                except httpx.RequestError:
                    fallback_response, _ = await _attempt_mcp_chat_fallback(
                        primary_headers=headers,
                        payload=payload,
                        timeout_seconds=TELEGRAM_CHAT_TIMEOUT,
                        source_label="telegram_chat_context_retry_error",
                    )
                    if fallback_response is None:
                        raise
                    response = fallback_response
            if response.status_code != 200 and LARGE_PAYLOAD_MODEL:
                payload["model"] = LARGE_PAYLOAD_MODEL
                large_url = _normalize_chat_endpoint(LARGE_PAYLOAD_ENDPOINT or url)
                try:
                    response = await _call_chat_completion(large_url, headers, payload, timeout_seconds=TELEGRAM_CHAT_TIMEOUT)
                except httpx.RequestError:
                    fallback_response, _ = await _attempt_mcp_chat_fallback(
                        primary_headers=headers,
                        payload=payload,
                        timeout_seconds=TELEGRAM_CHAT_TIMEOUT,
                        source_label="telegram_chat_large_retry_error",
                    )
                    if fallback_response is None:
                        raise
                    response = fallback_response
            if response.status_code == 200:
                data = response.json()
            else:
                if not used_mcp_fallback:
                    fallback_response, _ = await _attempt_mcp_chat_fallback(
                        primary_headers=headers,
                        payload=payload,
                        timeout_seconds=TELEGRAM_CHAT_TIMEOUT,
                        source_label="telegram_chat_non_200_after_context_retry",
                    )
                    if fallback_response is not None:
                        response = fallback_response
                        used_mcp_fallback = True
                        if response.status_code == 200:
                            data = response.json()
                            # Continue with normal response parsing below.
                            pass
                        else:
                            detail = response.text or detail
                if response.status_code != 200:
                    await _finish_status_session(
                        conversation_id=conversation_id,
                        request_id=request_id,
                        final_state="Failed: model returned error",
                        phase="error",
                    )
                    raise HTTPException(status_code=response.status_code, detail=detail)
        else:
            if not used_mcp_fallback:
                fallback_response, _ = await _attempt_mcp_chat_fallback(
                    primary_headers=headers,
                    payload=payload,
                    timeout_seconds=TELEGRAM_CHAT_TIMEOUT,
                    source_label="telegram_chat_non_200",
                )
                if fallback_response is not None:
                    response = fallback_response
                    used_mcp_fallback = True
                    if response.status_code == 200:
                        data = response.json()
                    else:
                        detail = response.text or detail
            if response.status_code != 200:
                await _finish_status_session(
                    conversation_id=conversation_id,
                    request_id=request_id,
                    final_state="Failed: model returned error",
                    phase="error",
                )
                raise HTTPException(status_code=response.status_code, detail=detail)

    data = response.json()
    reply = None
    preserve_reasoning_details = is_minimax_chat_request(url, model_name)
    pending_native_tool_calls: List[Dict[str, Any]] = []
    pending_native_tool_message: Optional[Dict[str, Any]] = None
    choices = data.get("choices") or []
    if choices:
        normalized_message = normalize_chat_completion_message(
            choices[0].get("message") or {},
            preserve_reasoning_details=preserve_reasoning_details,
        )
        tool_calls = normalized_message.get("tool_calls")
        if isinstance(tool_calls, list):
            pending_native_tool_calls = tool_calls
            pending_native_tool_message = normalized_message.get("message")
        reply = normalized_message.get("content") or ""

    if not reply and not pending_native_tool_calls:
        reply = "I couldn't generate a response right now. Please try again shortly."

    # Tool loop: when tools enabled, parse for tool calls and execute up to TELEGRAM_TOOLS_MAX_ITERATIONS
    # Track last tool result so we can show it to the user if the LLM never returns a final text reply
    last_tool_result_message: Optional[str] = None
    last_tool_success: Optional[bool] = None
    preferred_tool_result_message: Optional[str] = None
    preferred_tool_success: Optional[bool] = None
    _tool_discovery_or_setup_names = {
        "list_available_commands",
        "googleworkspace_cli.list_available_commands",
        "check_auth",
        "googleworkspace_cli.check_auth",
        "check_cli",
        "googleworkspace_cli.check_cli",
        "listFiles",
        "list_files",
        "searchFiles",
        "search_files",
        "filesystem.list_files",
        "filesystem.search_files",
    }
    # Friendly message when tool failed or returned an error (avoid showing raw 404/500 to user)
    _telegram_tool_error_reply = "I wasn't able to get that information just now. Please try again or rephrase your question."
    if TELEGRAM_TOOLS_ENABLED and _telegram_tools is not None:
        working_messages: List[Dict[str, Any]] = []
        if system_prompt:
            working_messages.append({"role": "system", "content": system_prompt})
        working_messages.extend(history)

        def _build_telegram_tool_loop_controller_prompt() -> str:
            return (
                "You are continuing a Telegram tool-assisted reply.\n"
                "Return exactly one of these two outputs:\n"
                "1. Structured tool calls only, with no user-facing text, when another tool is required.\n"
                "2. A direct final answer to the user, with no XML or tool-call markup, when you already have enough information.\n"
                "Use XML tool markup only as a fallback when structured tool calls are unavailable.\n"
                "Do not narrate plans, next steps, or intentions.\n"
                "Do not ask the user to repeat URLs, filenames, search results, or prior tool output already in context.\n"
                "Reuse exact values from the latest tool result whenever possible.\n"
                "If the latest tool result already answers the request, summarize it directly for the user.\n"
                "For file tasks: if a likely filename is already known, use filesystem.read_text when available; otherwise use readFile. "
                "Use filesystem.search_files when available to find which file contains something; otherwise use searchFiles. "
                "Use filesystem.list_files when available only for discovery when no likely file or folder is known; otherwise use listFiles. "
                "Do not repeat broad filesystem.list_files or listFiles calls after a candidate file has already been identified."
            )

        def _build_telegram_followup_messages(
            extra_messages: Optional[List[Dict[str, Any]]] = None,
        ) -> List[Dict[str, Any]]:
            messages_for_payload = list(working_messages)
            controller_message = {"role": "system", "content": _build_telegram_tool_loop_controller_prompt()}
            if messages_for_payload and messages_for_payload[0].get("role") == "system":
                messages_for_payload.insert(1, controller_message)
            else:
                messages_for_payload.insert(0, controller_message)
            if extra_messages:
                messages_for_payload.extend(extra_messages)
            return messages_for_payload

        async def _request_telegram_tool_followup(
            extra_messages: Optional[List[Dict[str, Any]]] = None,
        ) -> Tuple[str, List[Dict[str, Any]], Optional[Dict[str, Any]]]:
            payload_tool = {"model": model_name, "messages": _build_telegram_followup_messages(extra_messages)}
            if TELEGRAM_TOOLS_ENABLED and _telegram_tools is not None:
                telegram_tools_payload = _get_telegram_combined_openai_tools()
                if telegram_tools_payload:
                    payload_tool["tools"] = telegram_tools_payload
                    payload_tool["tool_choice"] = "auto"
            if request.temperature is not None:
                payload_tool["temperature"] = request.temperature
            if request.max_output_tokens is not None:
                payload_tool["max_tokens"] = request.max_output_tokens
            max_tokens_tool = _get_max_tokens_from_payload(payload_tool)
            summarized_tool = False
            if isinstance(payload_tool.get("messages"), list):
                estimated_tool = _estimate_total_tokens(payload_tool["messages"], max_tokens_tool)
                if estimated_tool > get_max_token_limit():
                    payload_tool["messages"] = await _summarize_messages_for_budget_proxy(
                        payload_tool["messages"],
                        endpoint=url,
                        headers=headers,
                        model_name=payload_tool.get("model", ""),
                        max_tokens=max_tokens_tool,
                        large_model=LARGE_PAYLOAD_MODEL,
                        large_endpoint=LARGE_PAYLOAD_ENDPOINT,
                    )
                    summarized_tool = True
            await _update_status_session(
                conversation_id=conversation_id,
                request_id=request_id,
                state="Nearly there. I'm pulling the result together for you now.",
                phase="llm_followup",
            )
            try:
                response_tool = await _call_chat_completion(
                    url,
                    headers,
                    payload_tool,
                    timeout_seconds=TELEGRAM_TOOL_FOLLOWUP_TIMEOUT,
                )
            except httpx.RequestError as exc:
                print(f"Telegram tool-loop request error: {exc}")
                fallback_response_tool, _ = await _attempt_mcp_chat_fallback(
                    primary_headers=headers,
                    payload=payload_tool,
                    timeout_seconds=TELEGRAM_TOOL_FOLLOWUP_TIMEOUT,
                    source_label="telegram_tool_followup_request_error",
                )
                if fallback_response_tool is None:
                    fallback_reply = (
                        _telegram_tool_error_reply
                        if (not last_tool_success or _telegram_tools.tool_result_looks_like_error(result_message))
                        else f"Here's what I found:\n\n{result_message}"
                    )
                    return fallback_reply, [], None
                response_tool = fallback_response_tool
            if response_tool.status_code != 200:
                response_text = ""
                try:
                    response_text = str(response_tool.text or "")
                except Exception:
                    response_text = ""
                response_preview = response_text[:2000] if response_text else "<empty body>"
                print(
                    "Telegram tool follow-up returned status "
                    f"{response_tool.status_code}, body={response_preview}"
                )
                if is_context_limit_error(response_tool.status_code, response_tool.text or ""):
                    if not summarized_tool and isinstance(payload_tool.get("messages"), list):
                        payload_tool["messages"] = await _summarize_messages_for_budget_proxy(
                            payload_tool["messages"],
                            endpoint=url,
                            headers=headers,
                            model_name=payload_tool.get("model", ""),
                            max_tokens=max_tokens_tool,
                            large_model=LARGE_PAYLOAD_MODEL,
                            large_endpoint=LARGE_PAYLOAD_ENDPOINT,
                        )
                        summarized_tool = True
                        try:
                            response_tool = await _call_chat_completion(
                                url,
                                headers,
                                payload_tool,
                                timeout_seconds=TELEGRAM_TOOL_FOLLOWUP_TIMEOUT,
                            )
                        except httpx.RequestError:
                            fallback_response_tool, _ = await _attempt_mcp_chat_fallback(
                                primary_headers=headers,
                                payload=payload_tool,
                                timeout_seconds=TELEGRAM_TOOL_FOLLOWUP_TIMEOUT,
                                source_label="telegram_tool_followup_context_retry_error",
                            )
                            if fallback_response_tool is None:
                                raise
                            response_tool = fallback_response_tool
                    if response_tool.status_code != 200 and LARGE_PAYLOAD_MODEL:
                        payload_tool["model"] = LARGE_PAYLOAD_MODEL
                        large_url = _normalize_chat_endpoint(LARGE_PAYLOAD_ENDPOINT or url)
                        try:
                            response_tool = await _call_chat_completion(
                                large_url,
                                headers,
                                payload_tool,
                                timeout_seconds=TELEGRAM_TOOL_FOLLOWUP_TIMEOUT,
                            )
                        except httpx.RequestError:
                            fallback_response_tool, _ = await _attempt_mcp_chat_fallback(
                                primary_headers=headers,
                                payload=payload_tool,
                                timeout_seconds=TELEGRAM_TOOL_FOLLOWUP_TIMEOUT,
                                source_label="telegram_tool_followup_large_retry_error",
                            )
                            if fallback_response_tool is None:
                                raise
                            response_tool = fallback_response_tool
                if response_tool.status_code != 200:
                    fallback_response_tool, _ = await _attempt_mcp_chat_fallback(
                        primary_headers=headers,
                        payload=payload_tool,
                        timeout_seconds=TELEGRAM_TOOL_FOLLOWUP_TIMEOUT,
                        source_label="telegram_tool_followup_non_200",
                    )
                    if fallback_response_tool is not None:
                        response_tool = fallback_response_tool
                if response_tool.status_code != 200:
                    fallback_reply = (
                        _telegram_tool_error_reply
                        if (not last_tool_success or _telegram_tools.tool_result_looks_like_error(result_message))
                        else f"Here's what I found:\n\n{result_message}"
                    )
                    return fallback_reply, [], None
            data_tool = response_tool.json()
            choices_tool = data_tool.get("choices") or []
            if not choices_tool:
                print("Telegram tool follow-up returned no choices, using tool result as reply")
                fallback_reply = (
                    _telegram_tool_error_reply
                    if (not last_tool_success or _telegram_tools.tool_result_looks_like_error(result_message))
                    else f"Here's what I found:\n\n{result_message}"
                )
                return fallback_reply, [], None
            normalized_followup = normalize_chat_completion_message(
                choices_tool[0].get("message") or {},
                preserve_reasoning_details=preserve_reasoning_details,
            )
            new_pending_native_tool_calls = normalized_followup.get("tool_calls") or []
            new_content = normalized_followup.get("content") or ""
            if not new_content.strip() and not new_pending_native_tool_calls:
                print("Telegram tool follow-up returned empty content, using tool result as reply")
                fallback_reply = (
                    _telegram_tool_error_reply
                    if (not last_tool_success or _telegram_tools.tool_result_looks_like_error(result_message))
                    else f"Here's what I found:\n\n{result_message}"
                )
                return fallback_reply, [], None
            followup_history_message = (
                normalized_followup.get("message")
                if new_pending_native_tool_calls
                else None
            )
            return new_content, new_pending_native_tool_calls, followup_history_message

        def _build_telegram_result_reply_from_last_tool() -> str:
            result_for_user = preferred_tool_result_message or last_tool_result_message
            success_for_user = (
                preferred_tool_success
                if preferred_tool_result_message is not None
                else last_tool_success
            )
            if bool(success_for_user):
                return f"Here's what I found:\n\n{result_for_user}"
            return (
                _telegram_tool_error_reply
                if _telegram_tools.tool_result_looks_like_error(str(result_for_user or ""))
                else f"I ran into an issue while calling the tool:\n\n{result_for_user}"
            )

        planning_chatter_checker = getattr(_telegram_tools, "reply_looks_like_tool_planning", None)
        iterations = 0
        while iterations < TELEGRAM_TOOLS_MAX_ITERATIONS:
            native_tool_calls = pending_native_tool_calls if isinstance(pending_native_tool_calls, list) else []
            parsed = _telegram_tools.parse_telegram_tool_response(reply)
            result_message = last_tool_result_message or ""
            if not native_tool_calls and not parsed:
                if (
                    iterations > 0
                    and result_message
                    and callable(planning_chatter_checker)
                    and planning_chatter_checker(reply)
                ):
                    reply = _build_telegram_result_reply_from_last_tool()
                break
            # After at least one tool execution, prefer mixed natural-language replies
            # over re-entering the tool loop when the model includes incidental XML.
            if (
                not native_tool_calls
                and iterations > 0
                and parsed
                and not _telegram_tools.reply_looks_like_tool_call(reply)
            ):
                cleaned_reply = _telegram_tools.strip_tool_call_markup(reply)
                if cleaned_reply and not _telegram_tools.reply_looks_like_tool_planning(cleaned_reply):
                    reply = cleaned_reply
                    break
            # Build context for tool execution
            tool_ctx = {
                "conversation_id": conversation_id,
                "user_id": request.user_id,
                "todo_user_key": todo_user_key,
                "task_execute_start": _task_execute_start,
                "task_execute_resume": _task_execute_resume,
                "task_execute_cancel": _task_execute_cancel,
                "task_execution_status": _task_execution_status,
                "task_execution_register_telegram_target": _task_execution_register_telegram_target,
                "todo_store": telegram_todo,
                "memory_cache_store": telegram_memory_cache,
                "do_search": _do_proxy_search,
                "do_fetch": _do_proxy_fetch,
                "do_news": _do_proxy_news,
                "do_weather": _do_proxy_weather,
                "do_autogen": _do_autogen,
                "do_codex": _run_codex_cli,
                "do_restart_proxy": lambda reason=None: _request_proxy_restart(
                    trigger=f"telegram_tool:{(reason or '').strip() or 'requested'}",
                    requested_by=f"telegram:{request.user_id or conversation_id}",
                ),
                "do_browser_agent": _do_browser_agent,
                "do_deep_research": lambda args: _do_deep_research_for_telegram(
                    conversation_id=conversation_id,
                    request_id=request_id,
                    body=args if isinstance(args, dict) else {},
                ),
                "do_browser_health_check": _do_browser_health_check,
                "read_file_internal": _read_file_internal,
                "write_file_internal": _write_file_internal,
                "list_files_internal": _list_files_internal,
                "search_files_internal": _search_files_internal,
                "send_telegram_file_internal": lambda filename, caption=None: _send_telegram_file_internal(
                    request.user_id or conversation_id,
                    filename,
                    caption=caption,
                ),
                "upload_drive_internal": _upload_drive_internal,
                "create_telegram_slides_internal": lambda tool_args: _create_telegram_slides_presentation_internal(
                    tool_args if isinstance(tool_args, dict) else {},
                    conversation_id=conversation_id,
                    user_id=request.user_id or "",
                    model_name=model_name,
                ),
                "pdf_to_powerpoint_internal": lambda tool_args: _handle_pdf_to_powerpoint_internal(
                    tool_args if isinstance(tool_args, dict) else {},
                    conversation_id=conversation_id,
                    user_id=request.user_id or "",
                    attachment_records=attachment_records,
                    model_name=model_name,
                ),
                "attachment_records": attachment_records,
                "execute_skill_tool": lambda tool_name, tool_args: _execute_skill_framework_tool(
                    tool_name=tool_name,
                    arguments=tool_args,
                    conversation_id=conversation_id,
                    user_id=request.user_id,
                    metadata={"channel": "telegram"},
                ),
                "memory_manager": memory_manager if MEMORY_AVAILABLE else None,
            }

            if native_tool_calls:
                assistant_tool_message = pending_native_tool_message or {
                    "role": "assistant",
                    "tool_calls": native_tool_calls,
                }
                working_messages.append(assistant_tool_message)
                for native_call in native_tool_calls:
                    native_function = native_call.get("function") if isinstance(native_call, dict) else {}
                    tool_name = (
                        (native_function.get("name") if isinstance(native_function, dict) else None)
                        or (native_call.get("name") if isinstance(native_call, dict) else None)
                    )
                    if not tool_name:
                        continue
                    await _update_status_session(
                        conversation_id=conversation_id,
                        request_id=request_id,
                        state=_format_telegram_tool_status(tool_name),
                        phase=f"tool:{tool_name}",
                    )
                    raw_args = (
                        (native_function.get("arguments") if isinstance(native_function, dict) else None)
                        or (native_call.get("arguments") if isinstance(native_call, dict) else "{}")
                    )
                    try:
                        tool_args = json.loads(raw_args) if isinstance(raw_args, str) else raw_args
                    except (TypeError, json.JSONDecodeError):
                        tool_args = {}
                    if not isinstance(tool_args, dict):
                        tool_args = {}
                    try:
                        tool_result = await _telegram_tools.execute_telegram_tool(tool_name, tool_args, tool_ctx)
                    except Exception as e:
                        tool_result = {"success": False, "message": str(e)}
                    result_message = tool_result.get("message", str(tool_result))
                    last_tool_result_message = result_message
                    last_tool_success = tool_result.get("success", True)
                    normalized_tool_name = str(tool_name or "").strip()
                    if (
                        bool(last_tool_success)
                        and normalized_tool_name
                        and normalized_tool_name not in _tool_discovery_or_setup_names
                    ):
                        preferred_tool_result_message = result_message
                        preferred_tool_success = True
                    working_messages.append(
                        {
                            "role": "tool",
                            "tool_call_id": native_call.get("id") if isinstance(native_call, dict) else None,
                            "content": result_message,
                        }
                    )
            else:
                tool_name = parsed.get("name")
                if tool_name:
                    await _update_status_session(
                        conversation_id=conversation_id,
                        request_id=request_id,
                        state=_format_telegram_tool_status(tool_name),
                        phase=f"tool:{tool_name}",
                    )
                args_str = parsed.get("arguments", "{}")
                try:
                    tool_args = json.loads(args_str) if isinstance(args_str, str) else args_str
                except (TypeError, json.JSONDecodeError):
                    tool_args = {}
                if not isinstance(tool_args, dict):
                    tool_args = {}
                try:
                    tool_result = await _telegram_tools.execute_telegram_tool(tool_name, tool_args, tool_ctx)
                except Exception as e:
                    tool_result = {"success": False, "message": str(e)}
                result_message = tool_result.get("message", str(tool_result))
                last_tool_result_message = result_message
                last_tool_success = tool_result.get("success", True)
                normalized_tool_name = str(tool_name or "").strip()
                if (
                    bool(last_tool_success)
                    and normalized_tool_name
                    and normalized_tool_name not in _tool_discovery_or_setup_names
                ):
                    preferred_tool_result_message = result_message
                    preferred_tool_success = True
                canonical_tool_reply = _telegram_tools.format_telegram_tool_call(tool_name, tool_args)
                working_messages.append({"role": "assistant", "content": canonical_tool_reply})
                working_messages.append({"role": "user", "content": f"Tool result: {result_message}"})
            reply, pending_native_tool_calls, pending_native_tool_message = await _request_telegram_tool_followup()
            iterations += 1

        if (
            callable(planning_chatter_checker)
            and planning_chatter_checker(reply)
            and last_tool_result_message
        ):
            reply = _build_telegram_result_reply_from_last_tool()

    # Never send raw tool-call XML to the user: if reply still looks like a tool call, show last tool result instead
    if _telegram_tools is not None:
        strip_think_fn = getattr(_telegram_tools, "strip_think_markup", None)
        if callable(strip_think_fn):
            cleaned_reply = strip_think_fn(reply)
            if cleaned_reply:
                reply = cleaned_reply
            elif preferred_tool_result_message or last_tool_result_message:
                result_for_user = preferred_tool_result_message or last_tool_result_message
                success_for_user = (
                    preferred_tool_success
                    if preferred_tool_result_message is not None
                    else last_tool_success
                )
                if not success_for_user or _telegram_tools.tool_result_looks_like_error(str(result_for_user or "")):
                    reply = _telegram_tool_error_reply
                else:
                    reply = f"Here's what I found:\n\n{result_for_user}"
            else:
                reply = "I couldn't generate a response right now. Please try again shortly."

    # Never send raw tool-call XML to the user: if reply still looks like a tool call, show last tool result instead
    if _telegram_tools is not None and _telegram_tools.reply_looks_like_tool_call(reply):
        result_for_user = preferred_tool_result_message or last_tool_result_message
        success_for_user = (
            preferred_tool_success
            if preferred_tool_result_message is not None
            else last_tool_success
        )
        if result_for_user is not None:
            reply_preview = re.sub(r"\s+", " ", str(reply or "")).strip()[:300]
            print(
                "Telegram: reply was raw tool call, using last tool result for user. "
                f"preview={reply_preview}"
            )
            if not success_for_user or _telegram_tools.tool_result_looks_like_error(result_for_user):
                reply = _telegram_tool_error_reply
            else:
                reply = f"Here's what I found:\n\n{result_for_user}"
        else:
            reply = "I used a tool but couldn't format the result. Please try again."

    history.append({"role": "assistant", "content": reply})
    trim_telegram_history(history)

    # Extract and store memories if memory system is available and auto-extract is enabled
    if MEMORY_AVAILABLE and memory_manager:
        auto_extract = os.getenv("MEMORY_AUTO_EXTRACT", "true").lower() == "true"
        if auto_extract:
            # Extract memories from the latest user+assistant turns without blocking Telegram response delivery.
            recent_messages = history[-4:] if len(history) >= 4 else history
            memory_messages: List[Dict[str, str]] = []
            for msg in recent_messages:
                if not isinstance(msg, dict):
                    continue
                role = str(msg.get("role") or "").strip()
                content = str(msg.get("content") or "")
                if not role or not content:
                    continue
                memory_messages.append({"role": role, "content": content})
            if memory_messages:
                asyncio.create_task(_extract_memories_from_recent_messages_async(memory_messages))

    usage = data.get("usage") if isinstance(data, dict) else None

    await _finish_status_session(
        conversation_id=conversation_id,
        request_id=request_id,
        final_state="Done: response delivered",
        phase="done",
    )

    return TelegramChatResponse(
        reply=reply,
        conversation_id=conversation_id,
        usage=usage,
    )


@app.delete("/v1/telegram/chat/{conversation_id}")
async def telegram_clear_conversation(request: Request, conversation_id: str):
    """Clear cached Telegram conversation history for a user."""
    _validate_telegram_secret(request)

    removed = telegram_conversations.pop(conversation_id, None) is not None
    return {"conversation_id": conversation_id, "cleared": removed}

# ============================================================================
# STATUS EVENTS ENDPOINTS
# ============================================================================

@app.post("/v1/status/start")
async def status_start(request: StatusStartRequest):
    """Start a status session for progress updates."""
    state = (request.state or "On it. I'm working on that now.").strip()
    await _start_status_session(
        conversation_id=request.conversation_id,
        request_id=request.request_id,
        channel=request.channel or "web",
        initial_state=state,
        phase="start",
    )
    return {"ok": True}


@app.post("/v1/status/update")
async def status_update(request: StatusUpdateRequest):
    """Update the current status state for a session."""
    event = await _update_status_session(
        conversation_id=request.conversation_id,
        request_id=request.request_id,
        state=request.state,
        phase=request.phase,
    )
    return {"ok": True, "event": event}


@app.post("/v1/status/finish")
async def status_finish(request: StatusFinishRequest):
    """Finish a status session and stop heartbeat updates."""
    event = await _finish_status_session(
        conversation_id=request.conversation_id,
        request_id=request.request_id,
        final_state=request.final_state,
        phase=request.phase or "done",
    )
    return {"ok": True, "event": event}


@app.get("/v1/status/latest")
async def status_latest(conversation_id: str, request_id: Optional[str] = None):
    """Get the latest status event for a conversation (optionally by request_id)."""
    event = _get_latest_status_event(conversation_id, request_id=request_id)
    if not event:
        return {"found": False}
    key = (event.get("conversation_id"), event.get("request_id"))
    active = key in status_sessions
    return {"found": True, "active": active, "event": event}


@app.get("/v1/status/events")
async def status_events(conversation_id: str, request_id: str, since_seq: int = 0):
    """Get status events since a sequence number for a given request."""
    events = _get_status_events_since(conversation_id, request_id, since_seq)
    latest = status_latest_index.get((conversation_id, request_id))
    latest_seq = latest.get("seq") if latest else since_seq
    active = (conversation_id, request_id) in status_sessions
    return {"events": events, "latest_seq": latest_seq, "active": active}

# ============================================================================
# MEMORY SYSTEM ENDPOINTS
# ============================================================================

@app.post("/v1/memory/store", response_model=MemoryResponse)
async def store_memory(request: MemoryStoreRequest):
    """Store a memory explicitly."""
    if not MEMORY_AVAILABLE or not memory_manager:
        error_detail = "Memory system is not available."
        if not MEMORY_AVAILABLE:
            error_detail += f" Import failed: {memory_import_error}. Check /v1/memory/status for details."
        elif not memory_manager:
            error_detail += " Memory manager initialization failed. Check server logs and /v1/memory/status endpoint."
        raise HTTPException(
            status_code=503,
            detail=error_detail
        )
    
    try:
        # Store the memory
        memory_id = await memory_manager.store_memory(
            text=request.text,
            category=request.category,
            source=request.source or "explicit",
            metadata=request.metadata,
        )
        
        return MemoryResponse(
            success=True,
            message=f"Memory stored successfully",
            data={"memory_id": memory_id}
        )
    except Exception as e:
        print(f"Error storing memory: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to store memory: {str(e)}")

@app.post("/v1/memory/search", response_model=MemoryResponse)
async def search_memories(request: MemorySearchRequest):
    """Search memories by query."""
    if not MEMORY_AVAILABLE or not memory_manager:
        raise HTTPException(
            status_code=503,
            detail="Memory system is not available. Check MEMORY_ENABLED setting."
        )
    
    try:
        # Search for relevant memories
        results = await memory_manager.search_memories(
            query=request.query,
            limit=request.limit,
            similarity_threshold=request.similarity_threshold,
            category=request.category,
        )
        
        return MemoryResponse(
            success=True,
            message=f"Found {len(results)} relevant memories",
            data={"memories": results, "count": len(results)}
        )
    except Exception as e:
        print(f"Error searching memories: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to search memories: {str(e)}")

@app.get("/v1/memory/list", response_model=MemoryResponse)
async def list_memories(limit: Optional[int] = None):
    """List recent memories."""
    if not MEMORY_AVAILABLE or not memory_manager:
        raise HTTPException(
            status_code=503,
            detail="Memory system is not available. Check MEMORY_ENABLED setting."
        )
    
    try:
        # List memories
        memories = memory_manager.list_memories(limit=limit)
        
        return MemoryResponse(
            success=True,
            message=f"Retrieved {len(memories)} memories",
            data={"memories": memories, "count": len(memories), "total": memory_manager.count()}
        )
    except Exception as e:
        print(f"Error listing memories: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to list memories: {str(e)}")

@app.get("/v1/memory/{memory_id}", response_model=MemoryResponse)
async def get_memory(memory_id: str):
    """Get a specific memory by ID."""
    if not MEMORY_AVAILABLE or not memory_manager:
        raise HTTPException(
            status_code=503,
            detail="Memory system is not available. Check MEMORY_ENABLED setting."
        )
    
    try:
        # Get memory by ID
        memory = memory_manager.get_memory(memory_id)
        
        if not memory:
            raise HTTPException(status_code=404, detail=f"Memory {memory_id} not found")
        
        return MemoryResponse(
            success=True,
            message="Memory retrieved successfully",
            data={"memory": memory}
        )
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error getting memory: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to get memory: {str(e)}")

@app.post("/v1/memory/extract", response_model=MemoryResponse)
async def extract_memories(request: MemoryExtractRequest):
    """Extract and store memories from a conversation."""
    if not MEMORY_AVAILABLE or not memory_manager:
        error_detail = "Memory system is not available."
        if not MEMORY_AVAILABLE:
            error_detail += f" Import failed: {memory_import_error}. Check /v1/memory/status for details."
        elif not memory_manager:
            error_detail += " Memory manager initialization failed. Check server logs and /v1/memory/status endpoint."
        raise HTTPException(
            status_code=503,
            detail=error_detail
        )
    
    try:
        # Check if auto-extract is enabled
        auto_extract = os.getenv("MEMORY_AUTO_EXTRACT", "true").lower() == "true"
        if not auto_extract:
            return MemoryResponse(
                success=False,
                message="Automatic memory extraction is disabled via MEMORY_AUTO_EXTRACT=false",
                data={"extracted": 0}
            )
        
        # Extract memories from conversation
        max_memories = request.max_memories or 1
        memory_ids = await memory_manager.extract_memories_from_conversation(
            messages=request.messages,
            max_memories=max_memories,
        )
        
        return MemoryResponse(
            success=True,
            message=f"Extracted and stored {len(memory_ids)} memories",
            data={"extracted": len(memory_ids), "memory_ids": memory_ids}
        )
    except Exception as e:
        print(f"Error extracting memories: {e}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to extract memories: {str(e)}")

@app.get("/v1/memory/learning/events", response_model=MemoryResponse)
async def memory_learning_events(limit: int = 50, outcome: Optional[str] = None):
    """List recent task-learning events (captured from task execution outcomes)."""
    if not MEMORY_AVAILABLE or not memory_manager:
        raise HTTPException(
            status_code=503,
            detail="Memory system is not available. Check MEMORY_ENABLED setting.",
        )
    if not hasattr(memory_manager, "list_task_learning_events"):
        return MemoryResponse(
            success=False,
            message="Task learning events are not supported by this memory manager.",
            data={"events": [], "count": 0},
        )
    safe_limit = max(1, min(200, int(limit or 50)))
    try:
        events = memory_manager.list_task_learning_events(limit=safe_limit, outcome=outcome)
        return MemoryResponse(
            success=True,
            message=f"Retrieved {len(events)} task learning events",
            data={"events": events, "count": len(events)},
        )
    except Exception as e:
        print(f"Error listing task learning events: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to list task learning events: {str(e)}")

@app.post("/v1/memory/learning/context", response_model=MemoryResponse)
async def memory_learning_context(request: MemoryLearningContextRequest):
    """Get experience-based learning context and guidance for a task description."""
    if not MEMORY_AVAILABLE or not memory_manager:
        raise HTTPException(
            status_code=503,
            detail="Memory system is not available. Check MEMORY_ENABLED setting.",
        )
    if not hasattr(memory_manager, "get_task_learning_context"):
        return MemoryResponse(
            success=False,
            message="Task learning context is not supported by this memory manager.",
            data={"context": {}, "guidance": ""},
        )
    try:
        context = await memory_manager.get_task_learning_context(
            task_description=request.task_description,
            limit=request.limit,
            similarity_threshold=request.similarity_threshold,
        )
        guidance = ""
        if hasattr(memory_manager, "build_task_execution_guidance"):
            guidance = await memory_manager.build_task_execution_guidance(
                task_description=request.task_description,
                limit=request.limit or 6,
            )
        return MemoryResponse(
            success=True,
            message="Retrieved task learning context",
            data={"context": context, "guidance": guidance},
        )
    except Exception as e:
        print(f"Error retrieving task learning context: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to retrieve task learning context: {str(e)}")

@app.get("/v1/memory/status")
async def memory_status():
    """Get the status of the memory system."""
    status = {
        "available": MEMORY_AVAILABLE,
        "enabled": os.getenv("MEMORY_ENABLED", "true").lower() == "true",
        "initialized": memory_manager is not None,
        "memory_count": memory_manager.count() if memory_manager else 0,
    }
    if not MEMORY_AVAILABLE:
        status["error"] = f"Memory module not available (import failed: {memory_import_error})"
        if memory_import_error and "numpy" in memory_import_error.lower():
            status["fix"] = "Install numpy with: pip install numpy"
    elif not status["enabled"]:
        status["error"] = "Memory system disabled via MEMORY_ENABLED=false"
    elif not status["initialized"]:
        status["error"] = "Memory manager failed to initialize (check server logs)"
    return status

@app.delete("/v1/memory/{memory_id}", response_model=MemoryResponse)
async def delete_memory(memory_id: str):
    """Delete a memory by ID."""
    if not MEMORY_AVAILABLE or not memory_manager:
        error_detail = "Memory system is not available."
        if not MEMORY_AVAILABLE:
            error_detail += " Memory module import failed."
        elif not memory_manager:
            error_detail += " Memory manager initialization failed. Check server logs and /v1/memory/status endpoint."
        raise HTTPException(
            status_code=503,
            detail=error_detail
        )
    
    try:
        # Delete memory
        deleted = memory_manager.delete_memory(memory_id)
        
        if not deleted:
            raise HTTPException(status_code=404, detail=f"Memory {memory_id} not found")
        
        return MemoryResponse(
            success=True,
            message=f"Memory {memory_id} deleted successfully",
            data={"memory_id": memory_id}
        )
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error deleting memory: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to delete memory: {str(e)}")

# ============================================================================
# END MEMORY SYSTEM ENDPOINTS
# ============================================================================

# ============================================================================
# PHILOSOPHER MODE HELPER FUNCTIONS
# ============================================================================

async def get_all_available_tools() -> List[Dict]:
    """Get all available tools from all connected MCP servers and built-in proxy tools."""
    print(f"[PHILOSOPHER] get_all_available_tools called - MCP_AVAILABLE: {MCP_AVAILABLE}")
    
    all_tools = []
    
    # Add built-in proxy tools (web search, web scraper, news API)
    # These are always available if the server is running
    
    # 1. Web Search Tool
    all_tools.append({
        "name": "web_search",
        "description": "Search the web using Brave Search API (with DuckDuckGo fallback). Returns top search results with URLs, titles, and snippets.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "The search query to execute (e.g., 'latest AI developments 2024')",
                }
            },
            "required": ["query"]
        },
        "server_id": "proxy_server"
    })
    print("[PHILOSOPHER] Added web_search tool")
    
    # 2. Web Scraper/Fetcher Tool
    all_tools.append({
        "name": "web_scraper",
        "description": "Fetch and scrape readable content from a web URL. Supports JavaScript-rendered pages (Playwright/Selenium) for dynamic sites.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "url": {
                    "type": "string",
                    "description": "The URL to fetch and scrape (must include http:// or https://)",
                },
                "crawl": {
                    "type": "boolean",
                    "description": "Whether to crawl linked pages on the same domain",
                    "default": True,
                },
                "max_pages": {
                    "type": "integer",
                    "description": "Maximum pages to scrape including the start page",
                    "default": 3,
                },
                "max_depth": {
                    "type": "integer",
                    "description": "Maximum link depth from the start page",
                    "default": 1,
                },
                "render_js": {
                    "type": "boolean",
                    "description": "Render JavaScript before extraction (use for dynamic sites).",
                    "default": False,
                },
                "render_engine": {
                    "type": "string",
                    "description": "Renderer engine: auto, playwright, or selenium.",
                    "default": "auto",
                },
                "wait_for_selector": {
                    "type": "string",
                    "description": "Optional CSS selector to wait for before extracting content.",
                },
                "js_wait_ms": {
                    "type": "integer",
                    "description": "Extra milliseconds to wait after page load for dynamic content.",
                    "default": 2200,
                },
            },
            "required": ["url"]
        },
        "server_id": "proxy_server"
    })
    print("[PHILOSOPHER] Added web_scraper tool")
    
    # 3. News API Tool (only if API key is configured)
    news_api_key = os.getenv('NEWS_API_KEY')
    if news_api_key:
        all_tools.append({
            "name": "news_search",
            "description": "Search for recent news articles using News API. Returns articles with titles, URLs, descriptions, and publication dates.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "The search query for news articles (e.g., 'artificial intelligence')",
                    }
                },
                "required": ["query"]
            },
            "server_id": "proxy_server"
        })
        print("[PHILOSOPHER] Added news_search tool")
    else:
        print("[PHILOSOPHER] NEWS_API_KEY not configured, skipping news_search tool")

    all_tools.append({
        "name": "weather_info",
        "description": "Get weather information from Open-Meteo (api.open-meteo.com). Supports explicit location or memory-based location fallback.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "location": {"type": "string", "description": "City/suburb/postcode (optional if available in memory)"},
                "detail": {"type": "string", "description": "summary, current, or forecast", "default": "summary"}
            }
        },
        "server_id": "proxy_server"
    })
    print("[PHILOSOPHER] Added weather_info tool")
    
    # 4. Native delete tool (filesystem list/read/write/search are skill-backed)
    if FILE_OPS_AVAILABLE:
        all_tools.append({
            "name": "delete_file",
            "description": "Delete a file from the scratch workspace. Use filename only (e.g. notes.txt).",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "Name of the file to delete",
                    }
                },
                "required": ["filename"]
            },
            "server_id": "proxy_server"
        })
        print("[PHILOSOPHER] Added delete_file tool; filesystem file tools are skill-backed")
    else:
        print("[PHILOSOPHER] File ops not available, skipping delete_file tool")
    
    # 5. runWorkflow (AutoGen team - code generation and automation)
    if AUTOGEN_AVAILABLE:
        all_tools.append({
            "name": "runWorkflow",
            "description": "Execute workflows for code generation and automation tasks using an AutoGen team. Use for tasks like building apps, generating code, or multi-step automation.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "contentPrompt": {
                        "type": "string",
                        "description": "The task to execute (e.g. 'build a javascript todo app', 'write a Python script that parses CSV')",
                    }
                },
                "required": ["contentPrompt"]
            },
            "server_id": "proxy_server"
        })
        print("[PHILOSOPHER] Added runWorkflow tool")
    else:
        print("[PHILOSOPHER] AutoGen not available, skipping runWorkflow tool")

    # 5b. runCodexCli (Codex CLI for CATBot code changes and tooling)
    if CODEX_ENABLED:
        all_tools.append({
            "name": "runCodexCli",
            "description": "Run Codex CLI in non-interactive mode to make CATBot code changes or add new capabilities. Always provide a clear prompt describing the change. Output is written to a scratch summary file.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "prompt": {
                        "type": "string",
                        "description": "Task instructions for Codex (e.g. 'Add a new /v1/proxy/codex tool in the proxy server and update docs')",
                    },
                },
                "required": ["prompt"]
            },
            "server_id": "proxy_server"
        })
        print("[PHILOSOPHER] Added runCodexCli tool")
    else:
        print("[PHILOSOPHER] Codex CLI disabled, skipping runCodexCli tool")

    # 5c. Skill framework tools (manifest-driven tools from src/skills)
    skill_tools = _get_skill_tools_mcp_schema()
    if skill_tools:
        for tool in skill_tools:
            name = str(tool.get("name") or "").strip()
            if not name:
                continue
            input_schema = tool.get("inputSchema")
            if not isinstance(input_schema, dict):
                input_schema = {"type": "object", "properties": {}}
            all_tools.append(
                {
                    "name": name,
                    "description": str(tool.get("description") or "").strip(),
                    "inputSchema": input_schema,
                    "server_id": "skill_framework",
                }
            )
        print(f"[PHILOSOPHER] Added {len(skill_tools)} skill framework tools")
    else:
        print("[PHILOSOPHER] No skill framework tools available")
    
    # 6. run_browser_agent (standalone browser-use HTTP server; not tied to mcp_servers)
    # Browser-use runs as a separate HTTP server (MCP_BROWSER_USE_HTTP_URL); add tool when URL is configured
    run_browser_agent_tool = {
        "name": "run_browser_agent",
        "description": "Control a web browser using natural language commands. Executes browser automation tasks and returns results.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "instruction": {
                    "type": "string",
                    "description": "Natural language instruction for browser automation (e.g., 'Navigate to Google and search for cats')",
                },
                "max_steps": {
                    "type": "integer",
                    "description": "Maximum number of steps the agent should take",
                    "default": 10
                },
                "use_vision": {
                    "type": "boolean",
                    "description": "Whether to use vision for understanding page content",
                    "default": True
                }
            },
            "required": ["instruction"]
        },
        "server_id": "proxy_server"
    }
    if MCP_BROWSER_USE_HTTP_URL:
        all_tools.append(run_browser_agent_tool)
        print(f"[PHILOSOPHER] Added run_browser_agent (standalone browser-use HTTP at {MCP_BROWSER_USE_HTTP_URL})")
        # Deep research (same browser-use ecosystem; proxy forwards to MCP browser server on port 5001)
        all_tools.append({
            "name": "run_deep_research",
            "description": "Performs comprehensive multi-step web research on a topic using multiple browser agents. Gathers information from many sources and returns a detailed research report with citations and findings. Use for in-depth research or comprehensive analysis requests.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "research_task": {
                        "type": "string",
                        "description": "The research topic or question to investigate (e.g. 'What are the latest developments in quantum computing?', 'Compare the best electric vehicles available in 2024')",
                    },
                    "researchTask": {
                        "type": "string",
                        "description": "Same as research_task; the research topic or question.",
                    },
                    "max_parallel_browsers": {
                        "type": "integer",
                        "description": "Optional: maximum number of parallel browser instances (default 3, max 5)",
                        "default": 3,
                    },
                    "maxParallelBrowsers": {
                        "type": "integer",
                        "description": "Optional: same as max_parallel_browsers",
                    },
                },
                "required": [],
            },
            "server_id": "proxy_server"
        })
        print("[PHILOSOPHER] Added run_deep_research (deep research agent)")
        all_tools.append({
            "name": "health_check",
            "description": "Get browser-use server health and running background task status. Use when the user asks for update/progress/state of browser automation or deep research (e.g., still running vs completed).",
            "inputSchema": {
                "type": "object",
                "properties": {},
                "required": [],
            },
            "server_id": "proxy_server"
        })
        print("[PHILOSOPHER] Added health_check (browser-use status)")
    else:
        print("[PHILOSOPHER] MCP_BROWSER_USE_HTTP_URL not set, skipping run_browser_agent")
    
    # Add MCP tools if MCP is available
    if MCP_AVAILABLE:
        # Debug: Check what servers are available
        print(f"[PHILOSOPHER] Checking MCP tools - mcp_clients: {len(mcp_clients)} clients, mcp_servers: {len(mcp_servers)} servers")
        print(f"[PHILOSOPHER] Connected client IDs: {list(mcp_clients.keys())}")
        print(f"[PHILOSOPHER] Server IDs in mcp_servers: {list(mcp_servers.keys())}")
        
        # Optional: also add run_browser_agent from mcp_servers if user added browser-use as a server and connected it,
        # and we did not already add it from MCP_BROWSER_USE_HTTP_URL (avoid duplicate)
        if not MCP_BROWSER_USE_HTTP_URL:
            for server_id, server in mcp_servers.items():
                server_status = server.get("status", "disconnected")
                if server_status == "connected" and _is_browser_use_server(server):
                    print(f"[PHILOSOPHER] Found connected browser-use server in mcp_servers: {server_id}")
                    all_tools.append({
                        **run_browser_agent_tool,
                        "server_id": server_id,
                    })
                    break  # only one browser-use
        
        # Get tools from each connected MCP client (non-browser-use servers)
        for server_id, client in mcp_clients.items():
            try:
                print(f"[PHILOSOPHER] Processing MCP server: {server_id}")
                # Check if this is the browser-use server (shouldn't be, but check anyway)
                server = mcp_servers.get(server_id)
                print(f"[PHILOSOPHER] Server config for {server_id}: {server}")
                
                if server and _is_browser_use_server(server):
                    # Skip - browser-use is added from MCP_BROWSER_USE_HTTP_URL or from mcp_servers above
                    print(f"[PHILOSOPHER] Skipping browser-use server {server_id} (already handled)")
                    continue
                else:
                    # Get tools from MCP server
                    print(f"[PHILOSOPHER] Requesting tools/list from MCP server {server_id}")
                    result = await client.request(method="tools/list", params={})
                    print(f"[PHILOSOPHER] tools/list response from {server_id}: {result}")
                    if result and "tools" in result:
                        print(f"[PHILOSOPHER] Found {len(result['tools'])} tools from server {server_id}")
                        for tool in result["tools"]:
                            tool["server_id"] = server_id
                            all_tools.append(tool)
                    else:
                        print(f"[PHILOSOPHER] No tools in response from server {server_id}")
            except Exception as e:
                print(f"[PHILOSOPHER] Error getting tools from server {server_id}: {e}")
                import traceback
                print(traceback.format_exc())
                continue
    else:
        print("[PHILOSOPHER] MCP not available, skipping MCP tools")
    
    print(f"[PHILOSOPHER] Total tools collected: {len(all_tools)}")
    return all_tools

async def execute_tool_for_philosopher(tool_name: str, parameters: Dict) -> str:
    """Execute a tool for philosopher mode. Returns result as string."""
    _log_tool_invocation("philosopher", tool_name, parameters)
    
    # Handle built-in proxy server tools
    if tool_name == "web_search":
        try:
            query = parameters.get("query", "")
            if not query:
                return "Error: 'query' parameter is required for web_search"
            
            # Call the proxy search endpoint
            result = await proxy_search(query)
            
            # Format results as readable text
            if result and "results" in result:
                results = result["results"]
                if not results:
                    return f"No search results found for query: {query}"
                
                formatted_results = []
                for i, item in enumerate(results, 1):
                    formatted_results.append(
                        f"{i}. {item.get('title', 'No title')}\n"
                        f"   URL: {item.get('url', 'No URL')}\n"
                        f"   {item.get('snippet', 'No description')}"
                    )
                
                return f"Search results for '{query}':\n\n" + "\n\n".join(formatted_results)
            else:
                return f"Search returned no results for query: {query}"
        except HTTPException as e:
            return f"Error executing web_search: {e.detail}"
        except Exception as e:
            return f"Error executing web_search: {str(e)}"
    
    elif tool_name == "web_scraper":
        try:
            url = parameters.get("url", "")
            if not url:
                return "Error: 'url' parameter is required for web_scraper"

            crawl = _coerce_bool(parameters.get("crawl", True), default=True)
            try:
                max_pages = int(parameters.get("max_pages", 3) or 3)
            except (TypeError, ValueError):
                max_pages = 3
            try:
                max_depth = int(parameters.get("max_depth", 1) or 1)
            except (TypeError, ValueError):
                max_depth = 1
            render_js = _coerce_bool(parameters.get("render_js", False))
            render_engine = parameters.get("render_engine", "auto")
            wait_for_selector = parameters.get("wait_for_selector")
            js_wait_ms = parameters.get("js_wait_ms", 2200)

            # Call the proxy fetch endpoint with extraction + optional crawl
            result = await _do_proxy_fetch(
                url,
                crawl=crawl,
                max_pages=max_pages,
                max_depth=max_depth,
                render_js=render_js,
                render_engine=render_engine,
                wait_for_selector=wait_for_selector,
                js_wait_ms=js_wait_ms,
            )

            if result and "content" in result:
                content = result["content"] or ""
                page_count = result.get("page_count", 1)
                # Truncate if too long (keep first 10000 characters)
                if len(content) > 10000:
                    return (
                        f"Scraped content from {url} (pages: {page_count}, truncated):\n\n"
                        f"{content[:10000]}...\n\n[Content truncated to 10000 characters]"
                    )
                return f"Scraped content from {url} (pages: {page_count}):\n\n{content}"
            else:
                return f"No content retrieved from URL: {url}"
        except HTTPException as e:
            return f"Error executing web_scraper: {e.detail}"
        except Exception as e:
            return f"Error executing web_scraper: {str(e)}"
    
    elif tool_name == "news_search":
        try:
            query = parameters.get("query", "")
            if not query:
                return "Error: 'query' parameter is required for news_search"
            
            # Call the proxy news search endpoint
            result = await proxy_news_search(query)
            
            # Format results as readable text
            if result and "articles" in result:
                articles = result["articles"]
                if not articles:
                    return f"No news articles found for query: {query}"
                
                formatted_articles = []
                for i, article in enumerate(articles[:10], 1):  # Limit to top 10
                    formatted_articles.append(
                        f"{i}. {article.get('title', 'No title')}\n"
                        f"   Source: {article.get('source', 'Unknown')}\n"
                        f"   Published: {article.get('publishedAt', 'Unknown date')}\n"
                        f"   URL: {article.get('url', 'No URL')}\n"
                        f"   {article.get('description', 'No description')}"
                    )
                
                total = result.get("totalResults", len(articles))
                return f"News articles for '{query}' (showing {len(formatted_articles)} of {total}):\n\n" + "\n\n".join(formatted_articles)
            else:
                return f"News search returned no articles for query: {query}"
        except HTTPException as e:
            return f"Error executing news_search: {e.detail}"
        except Exception as e:
            return f"Error executing news_search: {str(e)}"
    
    elif tool_name == "weather_info":
        try:
            result = await _do_proxy_weather(
                location=parameters.get("location"),
                detail=parameters.get("detail", "summary"),
                user_id=parameters.get("user_id"),
                memory_manager=memory_manager if MEMORY_AVAILABLE else None,
            )
            return json.dumps(result, ensure_ascii=False)
        except HTTPException as e:
            return f"Error executing weather_info: {e.detail}"
        except Exception as e:
            return f"Error executing weather_info: {str(e)}"

    # Handle remaining native file tool
    elif tool_name == "delete_file":
        try:
            filename = parameters.get("filename", "").strip()
            if not filename:
                return "Error: 'filename' is required for delete_file"
            result = await _delete_file_internal(filename)
            if not result.get("success"):
                return result.get("message", "Delete failed")
            return result.get("message", "File deleted.")
        except Exception as e:
            return f"Error executing delete_file: {str(e)}"
    
    elif tool_name == "runWorkflow":
        try:
            content_prompt = (parameters.get("contentPrompt") or parameters.get("content_prompt") or "").strip()
            if not content_prompt:
                return "Error: 'contentPrompt' is required for runWorkflow"
            if not AUTOGEN_AVAILABLE:
                return "Error: Workflow (AutoGen) is not available."
            result = await _do_autogen(content_prompt)
            msg = result.get("output") or result.get("response") or str(result)
            return msg
        except HTTPException as e:
            return f"Error executing runWorkflow: {e.detail}"
        except Exception as e:
            return f"Error executing runWorkflow: {str(e)}"

    elif tool_name == "runCodexCli":
        try:
            codex_prompt = (parameters.get("prompt") or "").strip()
            if not codex_prompt:
                return "Error: 'prompt' is required for runCodexCli"
            result = await _run_codex_cli(codex_prompt)
            summary_file = result.get("summaryFile")
            exit_code = result.get("exitCode")
            timed_out = result.get("timedOut")
            return (
                f"Codex CLI finished (exit_code={exit_code}, timed_out={timed_out}). "
                f"Summary file: {summary_file}"
            )
        except HTTPException as e:
            return f"Error executing runCodexCli: {e.detail}"
        except Exception as e:
            return f"Error executing runCodexCli: {str(e)}"
    
    # Handle run_browser_agent (standalone browser-use HTTP server; no mcp_servers entry required)
    elif tool_name == "run_browser_agent":
        if not MCP_BROWSER_USE_HTTP_URL:
            return "Error: Browser-use is not configured (MCP_BROWSER_USE_HTTP_URL not set)."
        try:
            instruction = (parameters.get("instruction") or parameters.get("task") or "").strip()
            if not instruction:
                return "Error: 'instruction' is required for run_browser_agent"
            result = await _browser_use_http_call_tool("run_browser_agent", parameters)
            content = result.get("content", [])
            if isinstance(content, list):
                texts = [item.get("text", str(item)) for item in content if isinstance(item, dict)]
                return "\n".join(texts) if texts else "Browser agent completed (no output)."
            return str(content)
        except Exception as e:
            return f"Error executing run_browser_agent: {BROWSER_USE_HTTP_UNAVAILABLE_MSG} {str(e)}"
    
    # Handle run_deep_research (proxies to MCP browser server /api/deep-research on port 5001)
    elif tool_name == "run_deep_research":
        research_task = (parameters.get("research_task") or parameters.get("researchTask") or "").strip()
        if not research_task:
            return "Error: 'research_task' or 'researchTask' is required for run_deep_research."
        max_parallel = parameters.get("max_parallel_browsers") or parameters.get("maxParallelBrowsers")
        if max_parallel is not None:
            try:
                max_parallel = min(5, max(1, int(max_parallel)))
            except (TypeError, ValueError):
                max_parallel = 3
        body = {"research_task": research_task}
        if max_parallel is not None:
            body["max_parallel_browsers"] = max_parallel
        try:
            result = await _do_deep_research(body)
            # Flask returns { success, result, error }; proxy returns response.json()
            report = result.get("result") or result.get("message") or result.get("output")
            if report:
                return str(report)
            if not result.get("success", True) and result.get("error"):
                return f"Deep research failed: {result.get('error')}"
            return str(result)
        except HTTPException as e:
            return f"Error executing run_deep_research: {e.detail}"
        except Exception as e:
            return f"Error executing run_deep_research: {str(e)}"

    # Handle health_check (calls browser-use MCP tool health_check for running-task visibility)
    elif tool_name in {"health_check", "run_health_check"}:
        try:
            result = await _do_browser_health_check(parameters)
            payload = result.get("result")
            if isinstance(payload, (dict, list)):
                return json.dumps(payload, ensure_ascii=False, default=str, indent=2)
            return str(result.get("message") or payload or result)
        except HTTPException as e:
            return f"Error executing health_check: {e.detail}"
        except Exception as e:
            return f"Error executing health_check: {str(e)}"

    # Handle skill framework tools (qualified names like skill.tool or unique aliases)
    elif (qualified_skill_tool := _resolve_skill_tool_qualified_name(tool_name)):
        try:
            result = await _execute_skill_framework_tool(
                tool_name=tool_name,
                arguments=parameters,
                conversation_id=str(parameters.get("conversation_id") or ""),
                user_id=str(parameters.get("user_id") or ""),
                metadata={"channel": "philosopher"},
            )
            message = str(result.get("message") or "").strip()
            data = result.get("data")
            if not result.get("success", False):
                return f"Error executing {qualified_skill_tool}: {message or 'unknown error'}"
            if qualified_skill_tool.startswith("filesystem."):
                formatted = _format_filesystem_skill_tool_output(qualified_skill_tool, result)
                if formatted:
                    return formatted
            return _format_generic_skill_tool_output(qualified_skill_tool, result)
        except Exception as e:
            return f"Error executing {qualified_skill_tool}: {str(e)}"
    
    # Handle MCP tools
    elif MCP_AVAILABLE:
        # Find which server has this tool
        all_tools = await get_all_available_tools()
        matching_tools = [
            tool for tool in all_tools
            if isinstance(tool, dict) and str(tool.get("name") or "").strip() == tool_name
        ]
        if not matching_tools:
            return f"Error: Tool '{tool_name}' not found on any connected server"

        requested_server_id = str(
            parameters.get("server_id")
            or parameters.get("_server_id")
            or ""
        ).strip()
        if requested_server_id:
            matching_tools = [
                tool for tool in matching_tools
                if str(tool.get("server_id") or "").strip() == requested_server_id
            ]
            if not matching_tools:
                return (
                    f"Error: Tool '{tool_name}' is not available on server '{requested_server_id}'."
                )

        unique_server_ids = sorted(
            {
                str(tool.get("server_id") or "").strip()
                for tool in matching_tools
                if str(tool.get("server_id") or "").strip()
            }
        )
        if len(unique_server_ids) > 1:
            return (
                f"Error: Tool '{tool_name}' is available on multiple servers ({', '.join(unique_server_ids)}). "
                "Specify 'server_id' in parameters."
            )

        server_id = unique_server_ids[0] if unique_server_ids else None
        if not server_id:
            return f"Error: Tool '{tool_name}' has no executable server mapping."
        
        # Execute the tool using the existing call_tool logic
        try:
            forwarded_parameters = dict(parameters or {})
            forwarded_parameters.pop("server_id", None)
            forwarded_parameters.pop("_server_id", None)
            request = ToolCallRequest(toolName=tool_name, parameters=forwarded_parameters)
            result = await call_tool(server_id, request)
            
            # Extract result content
            if result and "result" in result:
                result_data = result["result"]
                if "content" in result_data:
                    content = result_data["content"]
                    if isinstance(content, list):
                        # Extract text from content array
                        text_parts = []
                        for item in content:
                            if isinstance(item, dict) and item.get("type") == "text":
                                text_parts.append(item.get("text", ""))
                        return "\n".join(text_parts)
                    return str(content)
                return str(result_data)
            
            return str(result)
        except Exception as e:
            return f"Error executing tool {tool_name}: {str(e)}"
    else:
        return f"Error: Tool '{tool_name}' is not a built-in tool and MCP is not available"

# ============================================================================
# PHILOSOPHER MODE ENDPOINTS
# ============================================================================

@app.post("/v1/philosopher/start", response_model=PhilosopherResponse)
async def philosopher_start(request: PhilosopherStartRequest):
    """Enable philosopher mode for a conversation."""
    if not PHILOSOPHER_MODE_AVAILABLE or not PhilosopherMode:
        raise HTTPException(
            status_code=503,
            detail="Philosopher mode is not available. Check server logs for details."
        )
    
    # Check if philosopher mode is enabled
    philosopher_enabled = os.getenv("PHILOSOPHER_MODE_ENABLED", "true").lower() == "true"
    if not philosopher_enabled:
        raise HTTPException(
            status_code=503,
            detail="Philosopher mode is disabled via PHILOSOPHER_MODE_ENABLED=false"
        )
    
    # Get conversation ID
    conversation_id = request.conversation_id or request.user_id or "default"
    
    # Check if already active
    if philosopher_mode_active.get(conversation_id, False):
        return PhilosopherResponse(
            success=True,
            message="Philosopher mode is already active for this conversation",
            data={"conversation_id": conversation_id, "active": True}
        )
    
    try:
        # Get API configuration
        api_base = os.getenv("OPENAI_API_BASE", "https://api.openai.com/v1")
        model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
        api_key = _first_non_empty_env(preferred_api_key_env_names(api_base, model))
        if not api_key:
            raise HTTPException(
                status_code=503,
                detail=(
                    "No compatible API key is configured for philosopher mode. "
                    "Set OPENAI_API_KEY / MCP_LLM_OPENAI_API_KEY, or MINIMAX_API_KEY / "
                    "MCP_LLM_MINIMAX_API_KEY when using Minimax."
                ),
            )
        
        # Use OPENAI_MODEL directly (not TELEGRAM_DEFAULT_MODEL) since philosopher mode is not Telegram-specific
        
        # Create philosopher mode instance with tool support
        philosopher = PhilosopherMode(
            api_key=api_key,
            api_base=api_base,
            model=model,
            memory_manager=memory_manager if MEMORY_AVAILABLE else None,
            max_cycles=int(os.getenv("PHILOSOPHER_MAX_CYCLES", "10")),
            similarity_threshold=float(os.getenv("PHILOSOPHER_SIMILARITY_THRESHOLD", "0.3")),
            memory_limit=int(os.getenv("PHILOSOPHER_MEMORY_LIMIT", "10")),
            conversation_history_limit=int(os.getenv("PHILOSOPHER_CONVERSATION_HISTORY_LIMIT", "3")),
            tool_executor=execute_tool_for_philosopher,
            get_tools_func=get_all_available_tools,
            diversification_threshold=int(os.getenv("PHILOSOPHER_DIVERSIFICATION_THRESHOLD", "7")),
        )
        
        # Store instance and activate mode
        philosopher_mode_instances[conversation_id] = philosopher
        philosopher_mode_active[conversation_id] = True
        
        return PhilosopherResponse(
            success=True,
            message="Philosopher mode activated",
            data={"conversation_id": conversation_id, "active": True}
        )
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error starting philosopher mode: {e}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to start philosopher mode: {str(e)}")

@app.post("/v1/philosopher/stop", response_model=PhilosopherResponse)
async def philosopher_stop(request: PhilosopherStopRequest):
    """Disable philosopher mode for a conversation."""
    # Get conversation ID
    conversation_id = request.conversation_id or request.user_id or "default"
    
    # Check if active
    if not philosopher_mode_active.get(conversation_id, False):
        return PhilosopherResponse(
            success=True,
            message="Philosopher mode is not active for this conversation",
            data={"conversation_id": conversation_id, "active": False}
        )
    
    try:
        # Deactivate mode
        philosopher_mode_active[conversation_id] = False
        # Remove instance (optional - could keep for reuse)
        philosopher_mode_instances.pop(conversation_id, None)
        
        return PhilosopherResponse(
            success=True,
            message="Philosopher mode deactivated",
            data={"conversation_id": conversation_id, "active": False}
        )
    except Exception as e:
        print(f"Error stopping philosopher mode: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to stop philosopher mode: {str(e)}")

@app.get("/v1/philosopher/status")
async def philosopher_status(conversation_id: Optional[str] = None, user_id: Optional[str] = None):
    """Check if philosopher mode is active for a conversation."""
    # Get conversation ID
    conv_id = conversation_id or user_id or "default"
    
    is_active = philosopher_mode_active.get(conv_id, False)
    
    return {
        "active": is_active,
        "conversation_id": conv_id,
        "available": PHILOSOPHER_MODE_AVAILABLE,
        "enabled": os.getenv("PHILOSOPHER_MODE_ENABLED", "true").lower() == "true"
    }

@app.post("/v1/philosopher/contemplate", response_model=PhilosopherResponse)
async def philosopher_contemplate(request: PhilosopherContemplateRequest):
    """Execute a single contemplation cycle."""
    if not PHILOSOPHER_MODE_AVAILABLE or not PhilosopherMode:
        raise HTTPException(
            status_code=503,
            detail="Philosopher mode is not available. Check server logs for details."
        )
    
    # Get conversation ID
    conversation_id = request.conversation_id or request.user_id or "default"
    
    # Check if mode is active
    if not philosopher_mode_active.get(conversation_id, False):
        raise HTTPException(
            status_code=400,
            detail="Philosopher mode is not active for this conversation. Start it first with /v1/philosopher/start"
        )
    
    # Get philosopher instance
    philosopher = philosopher_mode_instances.get(conversation_id)
    if not philosopher:
        raise HTTPException(
            status_code=500,
            detail="Philosopher mode instance not found. Try restarting philosopher mode."
        )
    
    monitor_run_id = _monitor_run_start(
        "philosopher",
        "contemplate",
        input_text=request.question or "",
        metadata={
            "conversation_id": conversation_id,
            "current_step": 0,
            "total_steps": getattr(philosopher, "max_cycles", 0),
            "phase": "queued",
        },
    )

    async def _philosopher_monitor(event: str, payload: Dict[str, Any]) -> None:
        workflow_name = _truncate_monitor_text(
            payload.get("workflow_name") or payload.get("question") or request.question or "Philosopher contemplation",
            max_chars=220,
        )
        metadata: Dict[str, Any] = {
            "conversation_id": conversation_id,
            "workflow_name": workflow_name,
            "phase": payload.get("phase") or event,
        }
        for key in (
            "question",
            "current_step",
            "total_steps",
            "tool_iteration",
            "max_tool_iterations",
            "tool_call_count",
            "completed_steps",
        ):
            value = payload.get(key)
            if value is not None:
                if key == "question":
                    value = _truncate_monitor_text(value, max_chars=320)
                metadata[key] = value
        summary = payload.get("message") or f"Philosopher event: {event}"
        _monitor_run_update(monitor_run_id, summary=summary, metadata=metadata)
        if payload.get("message"):
            _monitor_run_note(monitor_run_id, payload["message"])

    previous_progress_callback = getattr(philosopher, "progress_callback", None)
    philosopher.progress_callback = _philosopher_monitor

    try:
        # Generate question if not provided
        if request.question:
            question = request.question
            await _philosopher_monitor(
                "workflow_selected",
                {
                    "question": question,
                    "workflow_name": question,
                    "phase": "preparation",
                    "message": "Using provided philosopher workflow question.",
                    "current_step": 0,
                    "total_steps": getattr(philosopher, "max_cycles", 0),
                },
            )
        else:
            question = await philosopher.generate_contemplation_question()
            if not question:
                _monitor_run_finish(
                    monitor_run_id,
                    status="error",
                    summary="Failed to generate philosopher workflow question.",
                    metadata={
                        "conversation_id": conversation_id,
                        "phase": "question_generation",
                    },
                )
                raise HTTPException(status_code=500, detail="Failed to generate contemplation question")
        
        # Execute contemplation
        result = await philosopher.contemplate_question(question)
        
        # Store contemplation in memory
        memory_id = await philosopher.store_contemplation(
            question=result["question"],
            conclusion=result["conclusion"],
            cycle_count=result["cycle_count"]
        )
        _monitor_run_finish(
            monitor_run_id,
            status="completed",
            summary=f"Completed philosopher workflow in {result['cycle_count']} step(s).",
            metadata={
                "conversation_id": conversation_id,
                "workflow_name": result["question"],
                "question": result["question"],
                "current_step": result["cycle_count"],
                "total_steps": getattr(philosopher, "max_cycles", result["cycle_count"]),
                "completed_steps": len(result.get("contemplation_steps") or []),
                "phase": "completed",
                "memory_id": memory_id,
            },
        )
        
        return PhilosopherResponse(
            success=True,
            message="Contemplation completed",
            data={
                "question": result["question"],
                "conclusion": result["conclusion"],
                "contemplation_steps": result["contemplation_steps"],
                "cycle_count": result["cycle_count"],
                "memory_id": memory_id,
            }
        )
    except HTTPException:
        if monitor_run_id in monitor_active_runs:
            _monitor_run_finish(
                monitor_run_id,
                status="error",
                summary="Philosopher workflow failed.",
                metadata={"conversation_id": conversation_id},
            )
        raise
    except Exception as e:
        print(f"Error during contemplation: {e}")
        import traceback
        print(traceback.format_exc())
        _monitor_run_finish(
            monitor_run_id,
            status="error",
            summary=f"Philosopher workflow failed: {str(e)}",
            metadata={
                "conversation_id": conversation_id,
                "phase": "error",
            },
        )
        raise HTTPException(status_code=500, detail=f"Failed to contemplate: {str(e)}")
    finally:
        philosopher.progress_callback = previous_progress_callback

# ============================================================================
# END PHILOSOPHER MODE ENDPOINTS
# ============================================================================

# Simple test endpoint to verify requests are reaching the server
@app.get("/test")
async def test_endpoint():
    """Simple test endpoint that should always work."""
    import sys
    print("🧪 TEST endpoint called", flush=True)
    sys.stdout.flush()
    return {"message": "test successful", "timestamp": time.time()}

# Monitoring dashboard (standalone)
@app.get("/monitor")
async def monitor_dashboard():
    """Serve the monitoring dashboard HTML."""
    dashboard_file = _PROJECT_ROOT / "docs" / "monitoring_dashboard.html"
    if not dashboard_file.exists():
        return HTMLResponse(content="<h1>Monitoring dashboard not found.</h1>", status_code=404)
    return HTMLResponse(content=dashboard_file.read_text(encoding="utf-8"), status_code=200)


@app.get("/monitor/summary")
async def monitor_summary():
    """Return high-level system status summary."""
    uptime_seconds = max(0.0, time.time() - PROXY_START_TIME)
    openai_key_present = bool(
        os.getenv("OPENAI_API_KEY")
        or os.getenv("MCP_LLM_OPENAI_API_KEY")
        or os.getenv("MINIMAX_API_KEY")
        or os.getenv("MCP_LLM_MINIMAX_API_KEY")
    )
    return {
        "time": time.time(),
        "uptime_seconds": uptime_seconds,
        "memory_available": MEMORY_AVAILABLE,
        "telegram_tools_enabled": TELEGRAM_TOOLS_ENABLED,
        "status_sessions_active": len(status_sessions),
        "openai_api_key_configured": openai_key_present,
        "status_events_file": str(STATUS_EVENTS_FILE),
        "autogen_active_runs": _get_monitor_runs_payload("autogen")["active_count"],
        "browser_use_active_runs": _get_monitor_runs_payload("browser_use")["active_count"],
        "philosopher_active_runs": _get_monitor_runs_payload("philosopher")["active_count"],
        "task_execution_active_runs": _get_monitor_runs_payload("task_execution")["active_count"],
        "browser_use_log_file_configured": bool(BROWSER_USE_LOG_FILE),
    }


@app.get("/monitor/status")
async def monitor_status(limit: int = 50):
    """Return recent status events for dashboard display."""
    limit = max(1, min(200, limit))
    return {"events": _get_recent_status_events(limit)}


@app.get("/monitor/logs")
async def monitor_logs(limit: int = 200):
    """Return the last N lines from the proxy log file if configured."""
    limit = max(1, min(1000, limit))
    log_path = PROXY_LOG_FILE
    if not log_path or not log_path.exists():
        return {"available": False, "lines": []}
    text = _tail_text_file(log_path, max_lines=limit)
    lines = text.splitlines()
    return {"available": True, "lines": lines[-limit:], "path": str(log_path)}


@app.get("/monitor/workflows")
async def monitor_workflows():
    """Return recent AutoGen, Browser-use, Philosopher, and task execution activity."""
    browser_health = dict(monitor_browser_health_snapshot)
    if _monitor_browser_health_is_stale(browser_health):
        try:
            await _do_browser_health_check({})
            browser_health = dict(monitor_browser_health_snapshot)
        except HTTPException:
            browser_health = dict(monitor_browser_health_snapshot)
    return {
        "autogen": _get_monitor_runs_payload("autogen"),
        "browser_use": {
            **_get_monitor_runs_payload("browser_use"),
            "health": browser_health,
            "log_file_configured": bool(BROWSER_USE_LOG_FILE),
        },
        "philosopher": _get_monitor_runs_payload("philosopher"),
        "task_execution": _get_monitor_runs_payload("task_execution"),
    }


@app.get("/monitor/workflows/log/{run_id}")
async def monitor_workflow_log(run_id: str):
    """Return the full scratch log for a recent workflow run when available."""
    run = _find_monitor_run(run_id)
    if not run:
        raise HTTPException(status_code=404, detail="Workflow run not found")
    log_file = run.get("log_file")
    log_path = _resolve_monitor_log_path(log_file)
    payload = _read_monitor_run_log(log_path)
    return {
        "run_id": run_id,
        "status": run.get("status"),
        "log_file": log_file,
        **payload,
    }


@app.get("/monitor/logs/browser-use")
async def monitor_browser_use_logs(limit: int = 200):
    """Return the last N lines from the browser-use log file when configured."""
    limit = max(1, min(1000, limit))
    log_path = Path(BROWSER_USE_LOG_FILE) if BROWSER_USE_LOG_FILE else None
    if not log_path or not log_path.exists():
        return {"available": False, "lines": []}
    text = _tail_text_file(log_path, max_lines=limit)
    lines = text.splitlines()
    return {"available": True, "lines": lines[-limit:], "path": str(log_path)}

# Health check endpoint
@app.get("/health")
async def health_check():
    """Health check endpoint."""
    # Add explicit logging to debug
    import sys
    print("ðŸ¥ Health check endpoint called", flush=True)
    sys.stdout.flush()
    try:
        result = {"status": "healthy", "timestamp": time.time()}
        print(f"ðŸ¥ Health check returning: {result}", flush=True)
        sys.stdout.flush()
        return result
    except Exception as e:
        import traceback
        print(f"âŒ Health check error: {e}", flush=True)
        print(traceback.format_exc(), flush=True)
        sys.stdout.flush()
        raise


@app.get("/v1/client-config")
async def client_config():
    """Expose safe client defaults from .env (no secrets)."""
    tts_model = _env_str("TTS_MODEL") or _env_str("TELEGRAM_TTS_MODEL")
    tts_voice = _env_str("TTS_VOICE") or _env_str("TELEGRAM_TTS_VOICE")
    return {
        "soulPrompt": _get_soul_prompt_text(),
        "ttsEndpoint": _env_str("TTS_ENDPOINT"),
        "ttsModel": tts_model,
        "ttsVoice": tts_voice,
    }

# Models list proxy endpoint to handle CORS and mixed content
@app.get("/v1/proxy/models")
async def proxy_models(request: Request, endpoint: Optional[str] = None):
    """
    Proxy models list requests to handle CORS and avoid mixed content issues.
    Routes requests to the OpenAI-compatible API endpoint specified in the request.
    """
    try:
        if not endpoint:
            endpoint = request.query_params.get('endpoint', '')

        if not endpoint:
            endpoint = os.getenv('OPENAI_API_BASE', 'http://localhost:1234/v1/models')
        else:
            if not endpoint.endswith('/models'):
                endpoint = endpoint.rstrip('/') + '/models'

        auth_header = request.headers.get('Authorization', '')
        headers = {}
        if auth_header:
            headers['Authorization'] = auth_header

        org_header = request.headers.get('OpenAI-Organization')
        if org_header:
            headers['OpenAI-Organization'] = org_header

        project_header = request.headers.get('OpenAI-Project')
        if project_header:
            headers['OpenAI-Project'] = project_header

        print(f"Proxying models list request to: {endpoint}")

        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(endpoint, headers=headers)

        print(f"Models list response status: {response.status_code}")

        if response.status_code != 200:
            print(f"LLM service returned error: {response.status_code}")
            print(f"   Response text: {response.text[:500]}")
            return JSONResponse(
                content=response.json() if response.headers.get('content-type', '').startswith('application/json') else {"error": response.text},
                status_code=response.status_code,
            )

        try:
            response_data = response.json()
            return JSONResponse(content=response_data, status_code=200)
        except Exception as json_error:
            print(f"Failed to parse JSON response: {json_error}")
            return JSONResponse(
                content={"error": "Invalid JSON response from LLM service"},
                status_code=500,
            )

    except httpx.ConnectError:
        print("Connection error: Could not connect to LLM service")
        raise HTTPException(
            status_code=503,
            detail="Could not connect to LLM service. Please check the endpoint configuration.",
        )
    except Exception as e:
        print(f"Models list proxy error: {e}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to proxy models list request: {str(e)}")

# Shared browser-agent logic for route and Telegram tool runner
async def _do_browser_agent(body: Dict[str, Any]) -> Dict[str, Any]:
    """Forward browser-agent request to MCP browser server. Returns response dict or raises HTTPException."""
    normalized_body = dict(body or {}) if isinstance(body, dict) else {}
    if "instruction" in normalized_body and "task" not in normalized_body:
        normalized_body["task"] = normalized_body.get("instruction")
    instruction_preview = str(
        normalized_body.get("task") or normalized_body.get("instruction") or normalized_body.get("url") or ""
    )
    monitor_run_id = _monitor_run_start("browser_use", "browser-agent", input_text=instruction_preview)
    mcp_browser_url = os.getenv('MCP_BROWSER_SERVER_URL', None)
    if not mcp_browser_url:
        try:
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            s.connect(("8.8.8.8", 80))
            local_ip = s.getsockname()[0]
            s.close()
            mcp_browser_url = f"http://{local_ip}:5001"
            print(f"   Detected local IP: {local_ip}, using {mcp_browser_url}")
        except Exception:
            mcp_browser_url = "http://127.0.0.1:5001"
            print(f"   Using default: {mcp_browser_url}")
    else:
        print(f"   Using configured MCP_BROWSER_SERVER_URL: {mcp_browser_url}")
    endpoint = f"{mcp_browser_url.rstrip('/')}/api/browser-agent"
    _monitor_run_note(monitor_run_id, f"Proxying browser-agent request to {endpoint}")
    print(f"ðŸŒ Proxying browser-agent request to: {endpoint}")
    health_endpoint = f"{mcp_browser_url.rstrip('/')}/api/health"
    health_check_passed = False
    try:
        async with httpx.AsyncClient(timeout=5.0) as health_client:
            health_response = await health_client.get(health_endpoint)
            if health_response.status_code == 200:
                health_check_passed = True
    except Exception as health_err:
        _monitor_run_note(monitor_run_id, f"Health check warning: {health_err}")
        print(f"   âš ï¸  MCP browser server health check failed: {health_err}")
        if not mcp_browser_url.startswith("http://127.0.0.1"):
            mcp_browser_url = "http://127.0.0.1:5001"
            endpoint = f"{mcp_browser_url.rstrip('/')}/api/browser-agent"
            try:
                async with httpx.AsyncClient(timeout=5.0) as hc:
                    hr = await hc.get(f"{mcp_browser_url.rstrip('/')}/api/health")
                    if hr.status_code == 200:
                        health_check_passed = True
            except Exception:
                pass
    if not health_check_passed:
        _monitor_run_note(monitor_run_id, "Browser server health check did not pass, continuing anyway.")
        print(f"   âš ï¸  Warning: Health check failed, but continuing with request")
    timeout = httpx.Timeout(connect=10.0, read=10800.0, write=10.0, pool=10.0)
    async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
        try:
            response = await client.post(
                endpoint,
                json=normalized_body,
                headers={'Content-Type': 'application/json'},
            )
        except httpx.ConnectError as conn_err:
            _monitor_run_finish(monitor_run_id, status="error", summary=f"Browser-agent connection error: {conn_err}")
            print(f"âŒ Connection error to MCP browser server: {conn_err}")
            raise HTTPException(
                status_code=503,
                detail="Could not connect to MCP browser server. Please ensure it's running on port 5001."
            )
        except httpx.ReadTimeout:
            _monitor_run_finish(monitor_run_id, status="error", summary="Browser-agent request timed out.")
            raise HTTPException(
                status_code=504,
                detail="Browser automation task timed out. Please try again or check the MCP browser server logs."
            )
        except httpx.TimeoutException:
            _monitor_run_finish(monitor_run_id, status="error", summary="Browser-agent request timed out.")
            raise HTTPException(
                status_code=504,
                detail="Browser automation task timed out. Please try again or check the MCP browser server logs."
            )
    print(f"✅ Browser-agent response status: {response.status_code}")
    if response.status_code != 200:
        error_content = response.json() if response.headers.get('content-type', '').startswith('application/json') else {"error": response.text}
        _monitor_run_finish(
            monitor_run_id,
            status="error",
            summary=error_content.get("error", str(error_content)),
            metadata={"http_status": response.status_code},
        )
        print(f"   Error response: {error_content}")
        raise HTTPException(status_code=response.status_code, detail=error_content.get("error", str(error_content)))
    response_json = response.json()
    if isinstance(response_json, dict) and not str(response_json.get("message") or "").strip():
        result_text = response_json.get("result")
        if isinstance(result_text, str) and result_text.strip():
            response_json["message"] = result_text.strip()
    _monitor_run_finish(
        monitor_run_id,
        status="completed",
        summary=_truncate_monitor_text(response_json.get("message") or response_json.get("result") or response_json),
        metadata={"http_status": response.status_code},
    )
    return response_json


@app.post("/v1/proxy/browser-agent")
async def proxy_browser_agent(request: Request):
    """Proxy browser automation requests to the MCP browser server."""
    body = await request.json()
    result = await _do_browser_agent(body)
    return JSONResponse(content=result, status_code=200)


# Shared deep-research logic for route and Telegram tool runner
def _normalize_deep_research_body(body: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Normalize deep-research payload aliases before forwarding to browser-use."""
    normalized = dict(body or {}) if isinstance(body, dict) else {}

    research_task = (
        normalized.get("research_task")
        or normalized.pop("researchTask", None)
        or normalized.get("topic")
    )
    if research_task is not None:
        research_text = str(research_task).strip()
        if research_text:
            normalized["research_task"] = research_text

    max_parallel = normalized.get("max_parallel_browsers")
    if max_parallel is None:
        max_parallel = normalized.pop("maxParallelBrowsers", None)
    if max_parallel is not None:
        normalized["max_parallel_browsers"] = max_parallel

    return normalized


async def _do_deep_research(body: Dict[str, Any]) -> Dict[str, Any]:
    """Forward deep-research request to MCP browser server. Returns response dict or raises HTTPException."""
    normalized_body = _normalize_deep_research_body(body)
    research_preview = str(normalized_body.get("research_task") or normalized_body.get("topic") or "")
    monitor_run_id = _monitor_run_start("browser_use", "deep-research", input_text=research_preview)
    mcp_browser_url = os.getenv('MCP_BROWSER_SERVER_URL', None)
    if not mcp_browser_url:
        try:
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            s.connect(("8.8.8.8", 80))
            local_ip = s.getsockname()[0]
            s.close()
            mcp_browser_url = f"http://{local_ip}:5001"
            print(f"   Detected local IP: {local_ip}, using {mcp_browser_url}")
        except Exception:
            mcp_browser_url = "http://127.0.0.1:5001"
            print(f"   Using default: {mcp_browser_url}")
    else:
        print(f"   Using configured MCP_BROWSER_SERVER_URL: {mcp_browser_url}")
    endpoint = f"{mcp_browser_url.rstrip('/')}/api/deep-research"
    _monitor_run_note(monitor_run_id, f"Proxying deep-research request to {endpoint}")
    print(f"🔬 Proxying deep-research request to: {endpoint}")
    timeout = httpx.Timeout(connect=10.0, read=10800.0, write=10.0, pool=10.0)
    async with httpx.AsyncClient(timeout=timeout) as client:
        try:
            response = await client.post(endpoint, json=normalized_body, headers={'Content-Type': 'application/json'})
        except httpx.ConnectError as conn_err:
            print(f"âŒ Connection error to MCP browser server: {conn_err}")
            raise HTTPException(
                status_code=503,
                detail="Could not connect to MCP browser server. Please ensure it's running on port 5001."
            )
        except httpx.ReadTimeout as timeout_err:
            _monitor_run_finish(monitor_run_id, status="error", summary=f"Deep research timed out: {timeout_err}")
            print(f"âŒ Read timeout from MCP browser server: {timeout_err}")
            raise HTTPException(
                status_code=504,
                detail="Deep research task timed out. Please try again or check the MCP browser server logs."
            )
    print(f"✅ Deep-research response status: {response.status_code}")
    if response.status_code != 200:
        error_content = response.json() if response.headers.get('content-type', '').startswith('application/json') else {"error": response.text}
        _monitor_run_finish(
            monitor_run_id,
            status="error",
            summary=error_content.get("error", str(error_content)),
            metadata={"http_status": response.status_code},
        )
        raise HTTPException(status_code=response.status_code, detail=error_content.get("error", str(error_content)))
    response_json = response.json()
    if isinstance(response_json, dict) and not str(response_json.get("message") or "").strip():
        for key in ("report", "result", "output"):
            candidate = response_json.get(key)
            if isinstance(candidate, str) and candidate.strip():
                response_json["message"] = candidate.strip()
                break
    _monitor_run_finish(
        monitor_run_id,
        status="completed",
        summary=_truncate_monitor_text(response_json.get("message") or response_json.get("report") or response_json),
        metadata={"http_status": response.status_code},
    )
    return response_json


async def _acquire_telegram_deep_research_slot(
    conversation_id: str,
    request_id: str,
) -> Tuple[bool, Optional[Dict[str, Any]]]:
    """Acquire per-conversation deep-research slot. Returns (acquired, active_holder_if_busy)."""
    now_ts = time.time()
    stale_before = now_ts - float(TELEGRAM_DEEP_RESEARCH_STALE_SECONDS)
    async with telegram_deep_research_lock:
        stale_keys = [
            cid
            for cid, meta in telegram_deep_research_active.items()
            if float(meta.get("started_at_ts") or 0.0) < stale_before
        ]
        for cid in stale_keys:
            telegram_deep_research_active.pop(cid, None)

        active = telegram_deep_research_active.get(conversation_id)
        if active and active.get("request_id") != request_id:
            return False, dict(active)

        telegram_deep_research_active[conversation_id] = {
            "request_id": request_id,
            "started_at_ts": now_ts,
            "started_at": datetime.now(timezone.utc).isoformat(),
        }
        return True, None


async def _release_telegram_deep_research_slot(conversation_id: str, request_id: str) -> None:
    """Release per-conversation deep-research slot if owned by request_id."""
    async with telegram_deep_research_lock:
        active = telegram_deep_research_active.get(conversation_id)
        if active and active.get("request_id") == request_id:
            telegram_deep_research_active.pop(conversation_id, None)


async def _do_deep_research_for_telegram(
    conversation_id: str,
    request_id: str,
    body: Dict[str, Any],
) -> Dict[str, Any]:
    """
    Run deep research with per-conversation concurrency guard for Telegram adhoc flows.
    Prevents overlapping long-running runs when users send repeated messages before prior completion.
    """
    acquired, active = await _acquire_telegram_deep_research_slot(conversation_id, request_id)
    if not acquired:
        active_request_id = str((active or {}).get("request_id") or "").strip()
        active_started = str((active or {}).get("started_at") or "").strip()
        suffix = []
        if active_started:
            suffix.append(f"started {active_started}")
        if active_request_id:
            suffix.append(f"request_id={active_request_id}")
        suffix_text = f" ({', '.join(suffix)})" if suffix else ""
        return {
            "success": False,
            "message": (
                "A deep research request is already running for this chat"
                f"{suffix_text}. Please wait for it to finish before starting another."
            ),
            "busy": True,
            "active_request_id": active_request_id or None,
        }

    try:
        return await _do_deep_research(body)
    except HTTPException as e:
        detail = e.detail if isinstance(e.detail, str) else str(e.detail)
        return {"success": False, "message": detail or "Deep research request failed."}
    except Exception as e:
        return {"success": False, "message": f"Deep research request failed: {e}"}
    finally:
        await _release_telegram_deep_research_slot(conversation_id, request_id)


async def _do_browser_health_check(body: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """Call browser-use MCP health_check and normalize to JSON for CATBot tool consumers."""
    if not MCP_BROWSER_USE_HTTP_URL:
        monitor_browser_health_snapshot.update(
            {
                "checked_at": datetime.now(timezone.utc).isoformat(),
                "ok": False,
                "message": "Browser-use is not configured (MCP_BROWSER_USE_HTTP_URL not set).",
                "result": None,
            }
        )
        raise HTTPException(status_code=503, detail="Browser-use is not configured (MCP_BROWSER_USE_HTTP_URL not set).")
    # Upstream health_check tool currently accepts no arguments.
    # Ignore caller-provided payload to avoid schema mismatch errors.
    payload: Dict[str, Any] = {}
    try:
        result = await _browser_use_http_call_tool("health_check", payload)
    except Exception as e:
        monitor_browser_health_snapshot.update(
            {
                "checked_at": datetime.now(timezone.utc).isoformat(),
                "ok": False,
                "message": f"{BROWSER_USE_HTTP_UNAVAILABLE_MSG} {str(e)}",
                "result": None,
            }
        )
        raise HTTPException(status_code=503, detail=f"{BROWSER_USE_HTTP_UNAVAILABLE_MSG} {str(e)}")

    content = result.get("content", [])
    text_parts: List[str] = []
    if isinstance(content, list):
        for item in content:
            if isinstance(item, dict):
                text = item.get("text")
                if text is not None:
                    text_parts.append(str(text))
            elif item is not None:
                text_parts.append(str(item))
    text_blob = "\n".join(part for part in text_parts if part).strip()

    parsed: Optional[Any] = None
    if text_blob:
        try:
            parsed = json.loads(text_blob)
        except Exception:
            parsed = None

    if isinstance(parsed, dict):
        status = str(parsed.get("status") or "unknown")
        running = parsed.get("running_tasks")
        uptime = parsed.get("uptime_seconds")
        message = f"Browser-use status: {status}. Running tasks: {running}. Uptime: {uptime}s."
        monitor_browser_health_snapshot.update(
            {
                "checked_at": datetime.now(timezone.utc).isoformat(),
                "ok": True,
                "message": message,
                "result": parsed,
            }
        )
        return {"success": True, "message": message, "result": parsed}

    fallback = text_blob or "Health check completed but returned no content."
    monitor_browser_health_snapshot.update(
        {
            "checked_at": datetime.now(timezone.utc).isoformat(),
            "ok": True,
            "message": fallback,
            "result": {"raw": fallback},
        }
    )
    return {"success": True, "message": fallback, "result": {"raw": fallback}}


@app.post("/v1/proxy/deep-research")
async def proxy_deep_research(request: Request):
    """Proxy deep research requests to the MCP browser server."""
    try:
        body = await request.json()
        result = await _do_deep_research(body)
        return JSONResponse(content=result, status_code=200)
    except HTTPException:
        raise
    except Exception as e:
        print(f"âŒ Deep-research proxy error: {e}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to proxy deep-research request: {str(e)}")


@app.post("/v1/proxy/browser-health")
async def proxy_browser_health(request: Request):
    """Proxy browser-use health/status requests to MCP health_check tool."""
    try:
        try:
            body = await request.json()
        except Exception:
            body = {}
        result = await _do_browser_health_check(body if isinstance(body, dict) else {})
        return JSONResponse(content=result, status_code=200)
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Browser-health proxy error: {e}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to proxy browser health request: {str(e)}")

# Chat completions proxy endpoint to handle CORS and mixed content
@app.post("/v1/proxy/chat/completions")
async def proxy_chat_completions(request: Request):
    """
    Proxy chat completions requests to handle CORS and avoid mixed content issues.
    Routes requests to the OpenAI-compatible API endpoint specified in the request.
    """
    try:
        # Get the request body
        body = await request.json()

        # Get the endpoint from query parameter or request body, or use default
        endpoint = request.query_params.get('endpoint', '')
        if not endpoint:
            endpoint = body.get('_endpoint', '')

        # If still no endpoint, use default from environment or localhost
        if not endpoint:
            endpoint = os.getenv('OPENAI_API_BASE', 'http://localhost:1234/v1/chat/completions')
        endpoint = _normalize_chat_endpoint(endpoint)

        # Remove internal endpoint parameter from body before forwarding
        body_clean = {k: v for k, v in body.items() if k != '_endpoint'}
        if not body_clean.get("model"):
            default_openai_model = (os.getenv("OPENAI_MODEL") or "").strip()
            if default_openai_model:
                body_clean["model"] = default_openai_model

        # Get Authorization header from the request (or server default key)
        auth_header = (request.headers.get('Authorization', '') or "").strip()
        if not auth_header:
            primary_api_key = _first_non_empty_env(
                preferred_api_key_env_names(endpoint, body_clean.get("model"))
            )
            if primary_api_key:
                auth_header = f"Bearer {primary_api_key}"

        # Build headers for the forwarded request
        headers = {
            'Content-Type': 'application/json'
        }
        if auth_header:
            headers['Authorization'] = auth_header

        # Add organization/project headers if present in original request
        org_header = request.headers.get('OpenAI-Organization')
        if org_header:
            headers['OpenAI-Organization'] = org_header

        project_header = request.headers.get('OpenAI-Project')
        if project_header:
            headers['OpenAI-Project'] = project_header

        print(f"Proxying chat completions request to: {endpoint}")
        print(f"   Model: {body_clean.get('model', 'unknown')}")

        max_tokens = _get_max_tokens_from_payload(body_clean)
        summarized = False
        if isinstance(body_clean.get("messages"), list):
            estimated = _estimate_total_tokens(body_clean["messages"], max_tokens)
            if estimated > get_max_token_limit():
                body_clean["messages"] = await _summarize_messages_for_budget_proxy(
                    body_clean["messages"],
                    endpoint=endpoint,
                    headers=headers,
                    model_name=body_clean.get("model", ""),
                    max_tokens=max_tokens,
                    large_model=LARGE_PAYLOAD_MODEL,
                    large_endpoint=LARGE_PAYLOAD_ENDPOINT,
                )
                summarized = True

        response: Optional[httpx.Response] = None
        primary_request_error: Optional[Exception] = None
        try:
            response = await _call_chat_completion(endpoint, headers, body_clean, timeout_seconds=120.0)
        except httpx.RequestError as exc:
            primary_request_error = exc

        if response is None:
            print(f"Primary chat request error: {primary_request_error}")
            fallback_response, fallback_error = await _attempt_mcp_chat_fallback(
                primary_headers=headers,
                payload=body_clean,
                timeout_seconds=120.0,
                source_label="proxy_chat_completions_request_error",
            )
            if fallback_response is None:
                raise HTTPException(
                    status_code=503,
                    detail=(
                        "Could not connect to LLM service. "
                        f"Primary error: {primary_request_error}. "
                        f"Fallback error: {fallback_error or 'not configured'}"
                    ),
                )
            response = fallback_response

        print(f"Chat completions response status: {response.status_code}")

        # Check if the response is successful
        if response.status_code != 200:
            print(f"LLM service returned error: {response.status_code}")
            print(f"   Response text: {response.text[:500]}")
            error_text = response.text or ""
            if is_context_limit_error(response.status_code, error_text):
                # Retry with summarization if we haven't yet.
                if not summarized and isinstance(body_clean.get("messages"), list):
                    body_clean["messages"] = await _summarize_messages_for_budget_proxy(
                        body_clean["messages"],
                        endpoint=endpoint,
                        headers=headers,
                        model_name=body_clean.get("model", ""),
                        max_tokens=max_tokens,
                        large_model=LARGE_PAYLOAD_MODEL,
                        large_endpoint=LARGE_PAYLOAD_ENDPOINT,
                    )
                    summarized = True
                    try:
                        response = await _call_chat_completion(endpoint, headers, body_clean, timeout_seconds=120.0)
                    except httpx.RequestError:
                        fallback_response, _ = await _attempt_mcp_chat_fallback(
                            primary_headers=headers,
                            payload=body_clean,
                            timeout_seconds=120.0,
                            source_label="proxy_chat_completions_context_retry_error",
                        )
                        if fallback_response is None:
                            raise
                        response = fallback_response
                # Retry with large payload model if configured
                if response.status_code != 200 and LARGE_PAYLOAD_MODEL:
                    body_clean["model"] = LARGE_PAYLOAD_MODEL
                    large_endpoint = _normalize_chat_endpoint(LARGE_PAYLOAD_ENDPOINT or endpoint)
                    try:
                        response = await _call_chat_completion(large_endpoint, headers, body_clean, timeout_seconds=120.0)
                    except httpx.RequestError:
                        fallback_response, _ = await _attempt_mcp_chat_fallback(
                            primary_headers=headers,
                            payload=body_clean,
                            timeout_seconds=120.0,
                            source_label="proxy_chat_completions_large_retry_error",
                        )
                        if fallback_response is None:
                            raise
                        response = fallback_response
                if response.status_code == 200:
                    try:
                        response_data = response.json()
                        return JSONResponse(content=response_data, status_code=200)
                    except Exception as json_error:
                        print(f"Failed to parse JSON response after retry: {json_error}")
                        return JSONResponse(content={"error": "Invalid JSON response from LLM service"}, status_code=500)

            fallback_response, _fallback_error = await _attempt_mcp_chat_fallback(
                primary_headers=headers,
                payload=body_clean,
                timeout_seconds=120.0,
                source_label="proxy_chat_completions_non_200",
            )
            if fallback_response is not None:
                response = fallback_response

            return JSONResponse(
                content=response.json() if response.headers.get('content-type', '').startswith('application/json') else {"error": response.text},
                status_code=response.status_code
            )

        # Return the JSON response
        try:
            response_data = response.json()
            return JSONResponse(content=response_data, status_code=200)
        except Exception as json_error:
            print(f"Failed to parse JSON response: {json_error}")
            return JSONResponse(
                content={"error": "Invalid JSON response from LLM service"},
                status_code=500
            )

    except HTTPException:
        raise
    except httpx.ConnectError as e:
        print(f"Connection error: Could not connect to LLM service")
        raise HTTPException(
            status_code=503,
            detail=f"Could not connect to LLM service. Please check the endpoint configuration."
        )
    except Exception as e:
        print(f"Chat completions proxy error: {e}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to proxy chat completions request: {str(e)}")

# OPTIONS handler for Whisper endpoint to handle CORS preflight
@app.options("/v1/audio/transcriptions")
async def proxy_whisper_options(request: Request):
    """Handle CORS preflight requests for Whisper endpoint."""
    return JSONResponse(
        content={},
        status_code=200,
        headers={
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "POST, OPTIONS",
            "Access-Control-Allow-Headers": "*",
            "Access-Control-Max-Age": "3600",
        }
    )

# Whisper proxy endpoint to handle CORS; compatible with OpenAI whisper-1 (passes key when needed)
@app.post("/v1/audio/transcriptions")
async def proxy_whisper(request: Request):
    """Proxy Whisper transcription requests to handle CORS. Client Authorization is for proxy auth only; upstream uses WHISPER_API_KEY."""
    try:
        # Get the form data from the request
        form_data = await request.form()
        
        # Get the Whisper endpoint (defaulting to localhost:8001)
        whisper_endpoint = os.getenv('WHISPER_ENDPOINT', 'http://localhost:8001/v1/audio/transcriptions')
        
        print(f"ðŸ“ Proxying Whisper request to: {whisper_endpoint}")
        
        # Outgoing auth to Whisper/OpenAI uses only WHISPER_API_KEY (client Authorization is for proxy auth, not forwarded)
        forward_headers = {}
        whisper_api_key = os.getenv("WHISPER_API_KEY")
        if whisper_api_key:
            forward_headers["Authorization"] = f"Bearer {whisper_api_key}"
        
        # Prepare the files and data for forwarding
        files = {}
        data = {}
        
        for key, value in form_data.items():
            if hasattr(value, 'read'):  # This is a file
                # Read the file content
                file_content = await value.read()
                files[key] = (value.filename, file_content, value.content_type)
                print(f"  📎 File: {value.filename} ({len(file_content)} bytes)")
            else:  # This is a regular form field
                data[key] = value
                print(f"  📄 Field: {key} = {value}")
        
        # Forward the request to the Whisper service (server key only; client auth stays between client and proxy)
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.post(
                whisper_endpoint,
                files=files,
                data=data,
                headers=forward_headers
            )
        
        print(f"✅ Whisper response status: {response.status_code}")
        print(f"📄 Response content type: {response.headers.get('content-type', 'unknown')}")
        
        # Check if the response is successful
        if response.status_code != 200:
            print(f"âŒ Whisper service returned error: {response.status_code}")
            print(f"   Response text: {response.text}")
            return JSONResponse(
                content={"error": f"Whisper service error: {response.text}"},
                status_code=response.status_code
            )
        
        # Try to parse the JSON response
        try:
            response_data = response.json()
            print(f"✅ Parsed JSON response: {response_data}")
            return JSONResponse(content=response_data, status_code=200)
        except Exception as json_error:
            print(f"âŒ Failed to parse JSON response: {json_error}")
            print(f"   Raw response text: {response.text[:200]}")
            # Return the raw text if JSON parsing fails
            return JSONResponse(
                content={"text": response.text},
                status_code=200
            )
    
    except httpx.ConnectError as e:
        print(f"âŒ Connection error: Could not connect to Whisper service at {whisper_endpoint}")
        print(f"   Make sure the Whisper service is running on port 8001")
        raise HTTPException(
            status_code=503,
            detail=f"Could not connect to Whisper service. Make sure it's running on port 8001."
        )
    except Exception as e:
        print(f"âŒ Whisper proxy error: {e}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to proxy Whisper request: {str(e)}")

def _build_wav_header(sample_rate: int, channels: int, bits_per_sample: int, pcm_bytes_len: int) -> bytes:
    channels = max(1, int(channels))
    sample_rate = max(8000, int(sample_rate))
    bits_per_sample = max(8, int(bits_per_sample))
    block_align = channels * (bits_per_sample // 8)
    byte_rate = sample_rate * block_align
    return (
        b"RIFF"
        + struct.pack("<I", 36 + pcm_bytes_len)
        + b"WAVE"
        + b"fmt "
        + struct.pack("<I", 16)
        + struct.pack("<H", 1)
        + struct.pack("<H", channels)
        + struct.pack("<I", sample_rate)
        + struct.pack("<I", byte_rate)
        + struct.pack("<H", block_align)
        + struct.pack("<H", bits_per_sample)
        + b"data"
        + struct.pack("<I", pcm_bytes_len)
    )


def _float_audio_to_pcm16_bytes(audio: Any) -> bytes:
    try:
        import numpy as np  # type: ignore

        arr = np.asarray(audio).reshape(-1)
        if arr.size == 0:
            return b""
        if np.issubdtype(arr.dtype, np.integer):
            return arr.astype(np.int16, copy=False).tobytes()

        arr = arr.astype(np.float32, copy=False)
        peak = float(np.max(np.abs(arr))) if arr.size else 0.0
        if peak > 1.5:
            arr = np.clip(arr, -32768.0, 32767.0)
            return arr.astype(np.int16).tobytes()

        arr = np.clip(arr, -1.0, 1.0)
        return (arr * 32767.0).astype(np.int16).tobytes()
    except Exception:
        pass

    out = bytearray()
    try:
        for value in (audio or []):
            f = float(value)
            if f > 1.0:
                f = 1.0
            elif f < -1.0:
                f = -1.0
            out.extend(struct.pack("<h", int(f * 32767.0)))
    except Exception as exc:
        raise RuntimeError(f"Could not convert generated audio to PCM16: {exc}") from exc
    return bytes(out)


def _embedded_tts_endpoint_enabled() -> bool:
    return EMBEDDED_KITTEN_TTS_ENABLED or EMBEDDED_POCKET_TTS_ENABLED


def _is_embedded_pocket_model(model_name: Optional[str]) -> bool:
    raw = (model_name or "").strip().lower()
    if not raw:
        return False
    normalized = raw.replace("_", "-")
    if normalized in EMBEDDED_POCKET_MODEL_ALIASES:
        return True
    if normalized.startswith("pocket-tts"):
        return True
    return "kyutai/pocket-tts" in normalized


def _resolve_embedded_tts_backend(model_name: Optional[str] = None) -> str:
    requested = (model_name or "").strip()
    if requested:
        if _is_embedded_pocket_model(requested):
            if not EMBEDDED_POCKET_TTS_ENABLED:
                raise HTTPException(
                    status_code=503,
                    detail="Embedded Pocket TTS is disabled. Set EMBEDDED_POCKET_TTS_ENABLED=true.",
                )
            return "pocket"
        if not EMBEDDED_KITTEN_TTS_ENABLED:
            if EMBEDDED_POCKET_TTS_ENABLED:
                raise HTTPException(
                    status_code=400,
                    detail=(
                        f"Unsupported embedded TTS model '{requested}'. "
                        f"Pocket aliases: {', '.join(sorted(EMBEDDED_POCKET_MODEL_ALIASES))}"
                    ),
                )
            raise HTTPException(
                status_code=503,
                detail="Embedded Kitten TTS is disabled. Set EMBEDDED_KITTEN_TTS_ENABLED=true.",
            )
        return "kitten"

    if EMBEDDED_KITTEN_TTS_ENABLED:
        return "kitten"
    if EMBEDDED_POCKET_TTS_ENABLED:
        return "pocket"

    raise HTTPException(
        status_code=404,
        detail="Embedded TTS endpoint is disabled. Enable EMBEDDED_KITTEN_TTS_ENABLED or EMBEDDED_POCKET_TTS_ENABLED.",
    )


def _normalize_embedded_pocket_model_name(model_name: Optional[str]) -> str:
    raw = (model_name or EMBEDDED_POCKET_MODEL or "").strip()
    return raw or "pocket-tts-realtime"


def _normalize_embedded_pocket_voice_key(prompt: str) -> str:
    value = (prompt or "").strip()
    lowered = value.lower()
    if lowered in {voice.lower() for voice in EMBEDDED_POCKET_VOICES}:
        return lowered
    if lowered in EMBEDDED_POCKET_COMPAT_VOICE_ALIASES:
        return lowered
    return value


def _resolve_embedded_pocket_voice(requested_voice: str) -> str:
    requested = (requested_voice or "").strip()
    if not requested:
        return EMBEDDED_POCKET_DEFAULT_VOICE

    available_lookup = {voice.lower(): voice for voice in EMBEDDED_POCKET_VOICES}
    if requested.lower() in available_lookup:
        return available_lookup[requested.lower()]

    alias_target = EMBEDDED_POCKET_COMPAT_VOICE_ALIASES.get(requested.lower())
    if alias_target and alias_target.lower() in available_lookup:
        return available_lookup[alias_target.lower()]

    if requested.startswith("hf://") or requested.lower().endswith((".wav", ".mp3", ".flac", ".safetensors")):
        return requested

    available_preview = ", ".join(EMBEDDED_POCKET_VOICES)
    raise HTTPException(
        status_code=400,
        detail=f"Unsupported voice '{requested}'. Available voices: {available_preview}",
    )


def _load_embedded_pocket_model_sync() -> Any:
    if not EMBEDDED_POCKET_IMPORT_AVAILABLE or EmbeddedPocketTTSModel is None:
        raise RuntimeError("Embedded Pocket TTS import unavailable.")

    load_model = getattr(EmbeddedPocketTTSModel, "load_model", None)
    if callable(load_model):
        return load_model()
    return EmbeddedPocketTTSModel()


async def _get_embedded_pocket_model() -> Any:
    global _embedded_pocket_model_instance

    if not EMBEDDED_POCKET_TTS_ENABLED:
        raise HTTPException(
            status_code=503,
            detail="Embedded Pocket TTS is disabled. Set EMBEDDED_POCKET_TTS_ENABLED=true.",
        )

    if not EMBEDDED_POCKET_IMPORT_AVAILABLE or EmbeddedPocketTTSModel is None:
        raise HTTPException(
            status_code=503,
            detail="Embedded Pocket TTS is unavailable. Install pocket-tts support dependencies.",
        )

    async with _embedded_pocket_model_lock:
        if _embedded_pocket_model_instance is None:
            try:
                print("Loading embedded Pocket TTS model")
                _embedded_pocket_model_instance = await asyncio.to_thread(_load_embedded_pocket_model_sync)
            except Exception as exc:
                raise HTTPException(
                    status_code=503,
                    detail=f"Failed to load embedded Pocket TTS model '{_normalize_embedded_pocket_model_name(None)}': {exc}",
                ) from exc
            print("Embedded Pocket TTS model loaded")

    return _embedded_pocket_model_instance


async def _get_embedded_pocket_voice_state(voice_prompt: str) -> Any:
    model = await _get_embedded_pocket_model()
    cache_key = _normalize_embedded_pocket_voice_key(voice_prompt)

    async with _embedded_pocket_voice_states_lock:
        if cache_key in _embedded_pocket_voice_states:
            return _embedded_pocket_voice_states[cache_key]

        try:
            state = await asyncio.to_thread(model.get_state_for_audio_prompt, voice_prompt)
        except Exception as exc:
            raise HTTPException(
                status_code=400,
                detail=f"Embedded Pocket TTS could not load voice '{voice_prompt}': {exc}",
            ) from exc
        _embedded_pocket_voice_states[cache_key] = state
        return state


async def _get_embedded_pocket_sample_rate() -> int:
    model = await _get_embedded_pocket_model()
    try:
        return max(8000, int(getattr(model, "sample_rate", EMBEDDED_KITTEN_SAMPLE_RATE)))
    except Exception:
        return EMBEDDED_KITTEN_SAMPLE_RATE


async def _iter_embedded_pocket_pcm_chunks(
    text: str,
    voice: str,
    model_name: Optional[str] = None,
    speed: Optional[float] = None,
    sample_rate: Optional[int] = None,
):
    del model_name, speed, sample_rate
    normalized_text = (text or "").strip()
    if not normalized_text:
        raise HTTPException(status_code=400, detail="Input text is required.")

    model = await _get_embedded_pocket_model()
    resolved_voice = _resolve_embedded_pocket_voice(voice)
    voice_state = await _get_embedded_pocket_voice_state(resolved_voice)

    try:
        stream_factory = getattr(model, "generate_audio_stream", None)
        has_streaming = callable(stream_factory)
    except Exception:
        has_streaming = False

    if has_streaming:
        loop = asyncio.get_running_loop()
        queue: asyncio.Queue[Any] = asyncio.Queue()
        sentinel = object()

        def _stream_worker() -> None:
            try:
                for generated_chunk in stream_factory(voice_state, normalized_text):
                    chunk_pcm = _float_audio_to_pcm16_bytes(generated_chunk)
                    if chunk_pcm:
                        loop.call_soon_threadsafe(queue.put_nowait, chunk_pcm)
            except Exception as exc:
                loop.call_soon_threadsafe(queue.put_nowait, exc)
            finally:
                loop.call_soon_threadsafe(queue.put_nowait, sentinel)

        worker_task = asyncio.create_task(asyncio.to_thread(_stream_worker))
        chunk_index = 0
        try:
            while True:
                item = await queue.get()
                if item is sentinel:
                    break
                if isinstance(item, Exception):
                    raise HTTPException(status_code=400, detail=f"Embedded Pocket TTS streaming failed: {item}") from item
                chunk_index += 1
                yield item
        finally:
            await worker_task

        if chunk_index == 0:
            raise RuntimeError("Embedded Pocket TTS generated empty audio.")
        return
    try:
        generated = await asyncio.to_thread(model.generate_audio, voice_state, normalized_text)
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Embedded Pocket TTS rejected request: {exc}") from exc

    chunk_pcm = _float_audio_to_pcm16_bytes(generated)
    if not chunk_pcm:
        raise RuntimeError("Embedded Pocket TTS generated empty audio.")
    yield chunk_pcm


async def _generate_embedded_pocket_pcm(
    text: str,
    voice: str,
    model_name: Optional[str] = None,
    speed: Optional[float] = None,
    sample_rate: Optional[int] = None,
) -> bytes:
    pcm_parts: List[bytes] = []
    async for piece in _iter_embedded_pocket_pcm_chunks(
        text=text,
        voice=voice,
        model_name=model_name,
        speed=speed,
        sample_rate=sample_rate,
    ):
        pcm_parts.append(piece)
    pcm_bytes = b"".join(pcm_parts)
    if not pcm_bytes:
        raise RuntimeError("Embedded Pocket TTS generated empty audio.")
    return pcm_bytes


def _normalize_embedded_kitten_repo_id(model_name: Optional[str]) -> str:
    raw = (model_name or EMBEDDED_KITTEN_MODEL or "").strip()
    if not raw:
        raw = "KittenML/kitten-tts-nano-0.2"
    if "/" not in raw:
        return f"KittenML/{raw}"
    return raw


def _get_embedded_kitten_available_voices(model: Any) -> List[str]:
    available = getattr(model, "available_voices", None)
    if not available:
        return []
    if isinstance(available, (list, tuple, set)):
        return [str(v).strip() for v in available if str(v).strip()]
    return []


def _fetch_embedded_kitten_config(repo_id: str) -> Dict[str, Any]:
    if _hf_hub_download is None:
        return {}
    try:
        config_path = _hf_hub_download(repo_id=repo_id, filename="config.json")
        with open(config_path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _extract_embedded_kitten_voice_aliases(config: Dict[str, Any]) -> Dict[str, str]:
    aliases: Dict[str, str] = {}
    raw_aliases = config.get("voice_aliases")
    if isinstance(raw_aliases, dict):
        for alias, target in raw_aliases.items():
            alias_key = str(alias).strip().lower()
            target_value = str(target).strip()
            if alias_key and target_value:
                aliases[alias_key] = target_value
    return aliases


def _coerce_embedded_kitten_style_embeddings(model: Any) -> bool:
    """Normalize legacy voice embeddings to match ONNX style tensor shape."""
    session = getattr(model, "session", None)
    voices_obj = getattr(model, "voices", None)
    if session is None or voices_obj is None:
        return False

    try:
        style_input = next((inp for inp in session.get_inputs() if inp.name == "style"), None)
    except Exception:
        return False
    if style_input is None:
        return False

    shape = list(getattr(style_input, "shape", []) or [])
    expected_batch = shape[0] if len(shape) >= 1 else None
    expected_width = shape[1] if len(shape) >= 2 else None

    keys: List[str] = []
    if isinstance(voices_obj, dict):
        keys = [str(k) for k in voices_obj.keys()]
    elif hasattr(voices_obj, "files"):
        keys = [str(k) for k in getattr(voices_obj, "files", [])]
    if not keys:
        return False

    patched: Dict[str, Any] = {}
    changed = False
    for key in keys:
        try:
            style_embedding = voices_obj[key]
        except Exception:
            continue

        embedding_shape = tuple(getattr(style_embedding, "shape", ()) or ())
        normalized = style_embedding

        if expected_batch == 1:
            if len(embedding_shape) == 1:
                if expected_width is None or embedding_shape[0] == expected_width:
                    try:
                        normalized = style_embedding.reshape((1, embedding_shape[0]))
                        changed = True
                    except Exception:
                        normalized = style_embedding
            elif len(embedding_shape) >= 2 and embedding_shape[0] != 1:
                if expected_width is None or embedding_shape[-1] == expected_width:
                    normalized = style_embedding[:1]
                    changed = True

        patched[key] = normalized

    if not changed:
        return False

    model.voices = patched
    if hasattr(model, "available_voices"):
        available = [k for k in keys if k in patched]
        if available:
            model.available_voices = available
    print("Normalized embedded KittenTTS voice embeddings for ONNX style input compatibility")
    return True


def _load_embedded_kitten_model_sync(repo_id: str) -> Tuple[Any, Dict[str, str]]:
    if not EMBEDDED_KITTEN_IMPORT_AVAILABLE or EmbeddedKittenTTS is None:
        raise RuntimeError("Embedded Kitten TTS import unavailable.")

    config = _fetch_embedded_kitten_config(repo_id)
    alias_map = _extract_embedded_kitten_voice_aliases(config)

    try:
        model = EmbeddedKittenTTS(repo_id)
        _coerce_embedded_kitten_style_embeddings(model)
        return model, alias_map
    except ValueError as exc:
        if "Unsupported model type" not in str(exc):
            raise
        if not EMBEDDED_KITTEN_FALLBACK_LOADER_AVAILABLE or _hf_hub_download is None or _EmbeddedKittenOnnxModel is None:
            raise RuntimeError(
                f"Model {repo_id} is unsupported by installed kittentts and fallback loader is unavailable."
            ) from exc

        config = config or _fetch_embedded_kitten_config(repo_id)
        model_type = str(config.get("type", "")).strip().upper()
        model_file = str(config.get("model_file", "")).strip()
        voices_file = str(config.get("voices", "")).strip()
        if model_type not in {"ONNX1", "ONNX2"} or not model_file or not voices_file:
            raise RuntimeError(
                f"Model {repo_id} uses unsupported type '{model_type or 'unknown'}'."
            ) from exc

        model_path = _hf_hub_download(repo_id=repo_id, filename=model_file)
        voices_path = _hf_hub_download(repo_id=repo_id, filename=voices_file)
        alias_map = alias_map or _extract_embedded_kitten_voice_aliases(config)
        model = _EmbeddedKittenOnnxModel(model_path=model_path, voices_path=voices_path)
        _coerce_embedded_kitten_style_embeddings(model)
        return model, alias_map


def _resolve_embedded_kitten_voice(
    requested_voice: str,
    available_voices: List[str],
    alias_map: Dict[str, str],
) -> str:
    requested = (requested_voice or "").strip()
    if not requested:
        return available_voices[0] if available_voices else EMBEDDED_KITTEN_DEFAULT_VOICE

    available_lookup = {voice.lower(): voice for voice in available_voices}
    if requested in available_voices:
        return requested
    if requested.lower() in available_lookup:
        return available_lookup[requested.lower()]

    alias_target = alias_map.get(requested.lower()) or EMBEDDED_KITTEN_COMPAT_VOICE_ALIASES.get(requested.lower())
    if alias_target:
        mapped = available_lookup.get(alias_target.lower())
        if mapped:
            return mapped
        if alias_target in available_voices:
            return alias_target

    if available_voices:
        available_preview = ", ".join(available_voices)
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported voice '{requested}'. Available voices: {available_preview}",
        )
    return requested


async def _get_embedded_kitten_model(model_name: Optional[str] = None) -> Tuple[Any, Dict[str, str]]:
    global _embedded_kitten_model_instance, _embedded_kitten_model_repo_id, _embedded_kitten_voice_aliases

    if not EMBEDDED_KITTEN_TTS_ENABLED:
        raise HTTPException(
            status_code=503,
            detail="Embedded Kitten TTS is disabled. Set EMBEDDED_KITTEN_TTS_ENABLED=true.",
        )

    if not EMBEDDED_KITTEN_IMPORT_AVAILABLE or EmbeddedKittenTTS is None:
        raise HTTPException(
            status_code=503,
            detail="Embedded Kitten TTS is unavailable. Install kitten-tts support dependencies.",
        )

    target_repo_id = _normalize_embedded_kitten_repo_id(model_name)

    async with _embedded_kitten_model_lock:
        if _embedded_kitten_model_instance is None or _embedded_kitten_model_repo_id != target_repo_id:
            try:
                print(f"Loading embedded KittenTTS model: {target_repo_id}")
                model, alias_map = await asyncio.to_thread(_load_embedded_kitten_model_sync, target_repo_id)
            except Exception as exc:
                raise HTTPException(
                    status_code=503,
                    detail=f"Failed to load embedded Kitten TTS model '{target_repo_id}': {exc}",
                ) from exc

            _embedded_kitten_model_instance = model
            _embedded_kitten_model_repo_id = target_repo_id
            _embedded_kitten_voice_aliases = alias_map
            print("Embedded KittenTTS model loaded")

    return _embedded_kitten_model_instance, _embedded_kitten_voice_aliases


def _split_embedded_kitten_text_chunks(raw_text: str, max_chars: int) -> List[str]:
    collapsed = re.sub(r"\s+", " ", (raw_text or "").strip())
    if not collapsed:
        return []
    if len(collapsed) <= max_chars:
        return [collapsed]

    chunks: List[str] = []
    current = ""

    def append_piece(piece: str) -> None:
        nonlocal current
        piece = piece.strip()
        if not piece:
            return
        if not current:
            current = piece
            return
        candidate = f"{current} {piece}"
        if len(candidate) <= max_chars:
            current = candidate
            return
        chunks.append(current)
        current = piece

    segments = re.split(r"(?<=[.!?])\s+", collapsed)
    for segment in segments:
        segment = segment.strip()
        if not segment:
            continue
        if len(segment) <= max_chars:
            append_piece(segment)
            continue

        words = segment.split(" ")
        word_buffer = ""
        for word in words:
            word = word.strip()
            if not word:
                continue
            candidate = f"{word_buffer} {word}".strip() if word_buffer else word
            if len(candidate) <= max_chars:
                word_buffer = candidate
                continue
            if word_buffer:
                append_piece(word_buffer)
                word_buffer = ""

            if len(word) <= max_chars:
                word_buffer = word
                continue

            for i in range(0, len(word), max_chars):
                token = word[i:i + max_chars].strip()
                if token:
                    append_piece(token)

        if word_buffer:
            append_piece(word_buffer)

    if current:
        chunks.append(current)
    return [chunk for chunk in chunks if chunk]


def _build_embedded_kitten_pcm16_silence_bytes(target_sample_rate: int, milliseconds: int) -> bytes:
    if milliseconds <= 0:
        return b""
    frame_count = max(0, int((target_sample_rate * milliseconds) / 1000))
    if frame_count <= 0:
        return b""
    return b"\x00\x00" * frame_count


async def _iter_embedded_kitten_pcm_chunks(
    text: str,
    voice: str,
    model_name: Optional[str] = None,
    speed: Optional[float] = None,
    sample_rate: Optional[int] = None,
):
    model, alias_map = await _get_embedded_kitten_model(model_name=model_name)
    resolved_voice = _resolve_embedded_kitten_voice(
        requested_voice=voice,
        available_voices=_get_embedded_kitten_available_voices(model),
        alias_map=alias_map,
    )

    chunks = _split_embedded_kitten_text_chunks(text, EMBEDDED_KITTEN_MAX_INPUT_CHARS)
    if not chunks:
        raise HTTPException(status_code=400, detail="Input text is required.")
    if len(chunks) > 1:
        print(
            f"[TTS] Chunking embedded Kitten request into {len(chunks)} segments "
            f"(max {EMBEDDED_KITTEN_MAX_INPUT_CHARS} chars each)"
        )

    target_sample_rate = max(8000, int(sample_rate or EMBEDDED_KITTEN_SAMPLE_RATE))
    inter_chunk_silence = _build_embedded_kitten_pcm16_silence_bytes(
        target_sample_rate, EMBEDDED_KITTEN_CHUNK_SILENCE_MS
    )
    speed_supported = speed is not None

    for index, chunk in enumerate(chunks):
        kwargs: Dict[str, Any] = {"voice": resolved_voice}
        if speed_supported and speed is not None:
            kwargs["speed"] = speed

        try:
            generated = await asyncio.to_thread(model.generate, chunk, **kwargs)
        except TypeError:
            if "speed" in kwargs:
                speed_supported = False
                kwargs.pop("speed", None)
                generated = await asyncio.to_thread(model.generate, chunk, **kwargs)
            else:
                raise
        except ValueError as exc:
            detail = (
                f"Embedded Kitten TTS rejected request on chunk {index + 1}/{len(chunks)}: {exc}"
                if len(chunks) > 1
                else f"Embedded Kitten TTS rejected request: {exc}"
            )
            raise HTTPException(status_code=400, detail=detail) from exc

        chunk_pcm = _float_audio_to_pcm16_bytes(generated)
        if not chunk_pcm:
            raise RuntimeError(f"Embedded Kitten TTS generated empty audio for chunk {index + 1}/{len(chunks)}.")
        yield chunk_pcm
        if inter_chunk_silence and index < len(chunks) - 1:
            yield inter_chunk_silence


async def _generate_embedded_kitten_pcm(
    text: str,
    voice: str,
    model_name: Optional[str] = None,
    speed: Optional[float] = None,
    sample_rate: Optional[int] = None,
) -> bytes:
    pcm_parts: List[bytes] = []
    async for piece in _iter_embedded_kitten_pcm_chunks(
        text=text,
        voice=voice,
        model_name=model_name,
        speed=speed,
        sample_rate=sample_rate,
    ):
        pcm_parts.append(piece)
    pcm_bytes = b"".join(pcm_parts)
    if not pcm_bytes:
        raise RuntimeError("Embedded Kitten TTS generated empty audio.")
    return pcm_bytes


def _list_embedded_tts_voices(model_name: Optional[str] = None) -> List[Dict[str, str]]:
    backend = _resolve_embedded_tts_backend(model_name)
    if backend == "pocket":
        voices = EMBEDDED_POCKET_VOICES or [EMBEDDED_POCKET_DEFAULT_VOICE]
    else:
        voices = EMBEDDED_KITTEN_VOICES or [EMBEDDED_KITTEN_DEFAULT_VOICE]

    return [
        {"id": voice_id, "object": "voice", "name": voice_id}
        for voice_id in voices
    ]


async def _get_embedded_tts_sample_rate(model_name: Optional[str], requested_sample_rate: Optional[int]) -> int:
    backend = _resolve_embedded_tts_backend(model_name)
    if backend == "pocket":
        return await _get_embedded_pocket_sample_rate()
    return max(8000, int(requested_sample_rate or EMBEDDED_KITTEN_SAMPLE_RATE))


async def _iter_embedded_tts_pcm_chunks(
    text: str,
    voice: str,
    model_name: Optional[str] = None,
    speed: Optional[float] = None,
    sample_rate: Optional[int] = None,
):
    backend = _resolve_embedded_tts_backend(model_name)
    if backend == "pocket":
        async for piece in _iter_embedded_pocket_pcm_chunks(
            text=text,
            voice=voice,
            model_name=model_name,
            speed=speed,
            sample_rate=sample_rate,
        ):
            yield piece
        return

    async for piece in _iter_embedded_kitten_pcm_chunks(
        text=text,
        voice=voice,
        model_name=model_name,
        speed=speed,
        sample_rate=sample_rate,
    ):
        yield piece


async def _generate_embedded_tts_pcm(
    text: str,
    voice: str,
    model_name: Optional[str] = None,
    speed: Optional[float] = None,
    sample_rate: Optional[int] = None,
) -> bytes:
    backend = _resolve_embedded_tts_backend(model_name)
    if backend == "pocket":
        return await _generate_embedded_pocket_pcm(
            text=text,
            voice=voice,
            model_name=model_name,
            speed=speed,
            sample_rate=sample_rate,
        )

    return await _generate_embedded_kitten_pcm(
        text=text,
        voice=voice,
        model_name=model_name,
        speed=speed,
        sample_rate=sample_rate,
    )

@app.get("/v1/audio/voices")
async def embedded_audio_voices(model: Optional[str] = None):
    """OpenAI-compatible voices endpoint served directly by proxy_server for embedded TTS backends."""
    if not _embedded_tts_endpoint_enabled():
        raise HTTPException(status_code=404, detail="Embedded TTS endpoint is disabled.")
    return JSONResponse(
        content={
            "object": "list",
            "data": _list_embedded_tts_voices(model_name=model),
        },
        status_code=200,
    )


@app.post("/v1/audio/speech")
async def embedded_audio_speech(payload: EmbeddedTtsSpeechRequest):
    """OpenAI-compatible TTS speech endpoint served directly by proxy_server for embedded TTS backends."""
    if not _embedded_tts_endpoint_enabled():
        raise HTTPException(status_code=404, detail="Embedded TTS endpoint is disabled.")

    text = (payload.input or "").strip()
    if not text:
        raise HTTPException(status_code=400, detail="Input text is required.")

    backend = _resolve_embedded_tts_backend(payload.model)
    if backend == "pocket":
        requested_model = (payload.model or EMBEDDED_POCKET_MODEL or "").strip() or EMBEDDED_POCKET_MODEL
        requested_voice = (payload.voice or EMBEDDED_POCKET_DEFAULT_VOICE or "").strip() or EMBEDDED_POCKET_DEFAULT_VOICE
    else:
        requested_model = (payload.model or EMBEDDED_KITTEN_MODEL or "").strip() or EMBEDDED_KITTEN_MODEL
        requested_voice = (payload.voice or EMBEDDED_KITTEN_DEFAULT_VOICE or "").strip() or EMBEDDED_KITTEN_DEFAULT_VOICE
    sample_rate = await _get_embedded_tts_sample_rate(payload.model, payload.sample_rate)
    channels = max(1, int(payload.channels or 1))
    response_format = (payload.response_format or "wav").strip().lower()
    stream_mode = bool(payload.stream)
    stream_chunk_bytes = EMBEDDED_POCKET_STREAM_CHUNK_BYTES if backend == "pocket" else EMBEDDED_KITTEN_STREAM_CHUNK_BYTES

    headers = {
        "X-Audio-Sample-Rate": str(sample_rate),
        "X-Audio-Channels": str(channels),
        "Cache-Control": "no-store",
    }

    if response_format in {"pcm", "l16", "s16le"}:
        media_type = "audio/pcm"
        if stream_mode:
            async def pcm_stream():
                async for generated_chunk in _iter_embedded_tts_pcm_chunks(
                    text=text,
                    voice=requested_voice,
                    model_name=requested_model,
                    speed=payload.speed,
                    sample_rate=sample_rate,
                ):
                    for i in range(0, len(generated_chunk), stream_chunk_bytes):
                        yield generated_chunk[i:i + stream_chunk_bytes]
                        await asyncio.sleep(0)

            return StreamingResponse(pcm_stream(), media_type=media_type, headers=headers)
        pcm_bytes = await _generate_embedded_tts_pcm(
            text=text,
            voice=requested_voice,
            model_name=requested_model,
            speed=payload.speed,
            sample_rate=sample_rate,
        )
        return Response(content=pcm_bytes, media_type=media_type, headers=headers, status_code=200)

    pcm_bytes = await _generate_embedded_tts_pcm(
        text=text,
        voice=requested_voice,
        model_name=requested_model,
        speed=payload.speed,
        sample_rate=sample_rate,
    )
    wav_bytes = _build_wav_header(sample_rate=sample_rate, channels=channels, bits_per_sample=16, pcm_bytes_len=len(pcm_bytes)) + pcm_bytes
    if stream_mode:
        async def wav_stream():
            yield wav_bytes[:44]
            await asyncio.sleep(0)
            payload_bytes = wav_bytes[44:]
            for i in range(0, len(payload_bytes), stream_chunk_bytes):
                yield payload_bytes[i:i + stream_chunk_bytes]
                await asyncio.sleep(0)

        return StreamingResponse(wav_stream(), media_type="audio/wav", headers=headers)
    return Response(content=wav_bytes, media_type="audio/wav", headers=headers, status_code=200)

# TTS voices proxy endpoint to handle CORS
def _should_skip_tts_tls_verify(base_url: str) -> bool:
    """Return True for local HTTPS TTS endpoints where self-signed certs are common."""
    try:
        parsed = urlparse(base_url)
        host = (parsed.hostname or "").lower()
        if parsed.scheme != "https":
            return False
        if host in {"localhost", "127.0.0.1", "::1"}:
            return True
        return host.endswith(".local")
    except Exception:
        return False


@app.get("/v1/proxy/tts/voices")
async def proxy_tts_voices(endpoint: str, model: Optional[str] = None):
    """Proxy TTS voices requests to handle CORS. Tries /voices first, then /v1/audio/voices."""
    if not endpoint:
        raise HTTPException(status_code=400, detail="Endpoint parameter is required")
    
    try:
        # Normalize the endpoint URL (remove trailing slash)
        base_endpoint = endpoint.rstrip('/')
        
        # Extract base URL (origin: protocol + host + port) to avoid path duplication
        try:
            # Parse the endpoint URL to extract the origin (protocol + host + port)
            if not base_endpoint.startswith('http://') and not base_endpoint.startswith('https://'):
                # If no protocol, assume http://
                base_endpoint = f"http://{base_endpoint}"
            
            parsed_url = urlparse(base_endpoint)
            # Reconstruct base URL with scheme, hostname, and port (if present)
            base_url = f"{parsed_url.scheme}://{parsed_url.netloc}"
        except Exception as e:
            # Fallback to simple string replacement if URL parsing fails
            # Extract origin manually using regex
            match = re.match(r'^(https?://[^/]+)', base_endpoint)
            if match:
                base_url = match.group(1)
            else:
                # Last resort: remove /v1 and any path
                base_url = base_endpoint.split('/')[0] if '/' in base_endpoint else base_endpoint
                if not base_url.startswith('http'):
                    base_url = f"http://{base_url}"
        
        skip_tls_verify = _should_skip_tts_tls_verify(base_url)
        if skip_tls_verify:
            print(f"[WARN] TTS voices proxy: TLS verification disabled for local endpoint: {base_url}")

        model_query_suffix = f"?{httpx.QueryParams({'model': model})}" if model else ""

        # Try /voices first (Chatterbox style)
        voices_url_primary = f"{base_url}/voices{model_query_suffix}"
        print(f"🎤 Trying primary TTS voices endpoint: {voices_url_primary}")
        
        response = None
        response_data = None
        
        async with httpx.AsyncClient(timeout=10.0, verify=(not skip_tls_verify)) as client:
            try:
                # Try the primary endpoint first
                response = await client.get(
                    voices_url_primary,
                    headers={
                        'Content-Type': 'application/json',
                        'Accept': 'application/json'
                    }
                )
                
                print(f"✅ Primary TTS voices response status: {response.status_code}")
                
                # If primary endpoint succeeds, use it
                if response.status_code == 200:
                    try:
                        response_data = response.json()
                        print(f"✅ Parsed TTS voices JSON response from primary endpoint")
                        return JSONResponse(content=response_data, status_code=200)
                    except Exception as json_error:
                        print(f"âŒ Failed to parse JSON response: {json_error}")
                        print(f"   Raw response text: {response.text[:200]}")
                        # Return the raw text if JSON parsing fails
                        return JSONResponse(
                            content={"text": response.text},
                            status_code=200
                        )
            except (httpx.ConnectError, httpx.HTTPStatusError) as e:
                # Primary endpoint failed, try fallback
                print(f"âš ï¸ Primary endpoint failed: {e}")
                response = None
            
            # If primary failed, try /v1/audio/voices (OpenAI-compatible style)
            if not response or response.status_code != 200:
                voices_url_fallback = f"{base_url}/v1/audio/voices{model_query_suffix}"
                print(f"🎤 Trying fallback TTS voices endpoint: {voices_url_fallback}")
                
                try:
                    response = await client.get(
                        voices_url_fallback,
                        headers={
                            'Content-Type': 'application/json',
                            'Accept': 'application/json'
                        }
                    )
                    
                    print(f"✅ Fallback TTS voices response status: {response.status_code}")
                    
                    # Check if the fallback response is successful
                    if response.status_code == 200:
                        try:
                            response_data = response.json()
                            print(f"✅ Parsed TTS voices JSON response from fallback endpoint")
                            return JSONResponse(content=response_data, status_code=200)
                        except Exception as json_error:
                            print(f"âŒ Failed to parse JSON response: {json_error}")
                            print(f"   Raw response text: {response.text[:200]}")
                            # Return the raw text if JSON parsing fails
                            return JSONResponse(
                                content={"text": response.text},
                                status_code=200
                            )
                    else:
                        # Fallback also failed
                        print(f"âŒ Fallback TTS service returned error: {response.status_code}")
                        print(f"   Response text: {response.text[:200]}")
                        raise HTTPException(
                            status_code=response.status_code,
                            detail=f"TTS service error: {response.text[:200]}"
                        )
                except httpx.ConnectError as e:
                    print(f"âŒ Connection error: Could not connect to TTS service at {voices_url_fallback}")
                    raise HTTPException(
                        status_code=503,
                        detail=f"Could not connect to TTS service. Tried {voices_url_primary} and {voices_url_fallback}"
                    )
                except httpx.HTTPStatusError as e:
                    print(f"âŒ HTTP error from fallback TTS service: {e.response.status_code}")
                    raise HTTPException(
                        status_code=e.response.status_code,
                        detail=f"TTS service returned error: {str(e)}"
                    )
        
        # If we get here, both attempts failed
        if response and response.status_code != 200:
            raise HTTPException(
                status_code=response.status_code,
                detail=f"TTS service error: {response.text[:200]}"
            )
        else:
            raise HTTPException(
                status_code=503,
                detail=f"Could not connect to TTS service. Tried {voices_url_primary} and {base_url}/v1/audio/voices"
            )
    
    except HTTPException:
        # Re-raise HTTP exceptions as-is
        raise
    except httpx.ConnectError as e:
        print(f"âŒ Connection error: Could not connect to TTS service at {endpoint}")
        raise HTTPException(
            status_code=503,
            detail=f"Could not connect to TTS service at {endpoint}"
        )
    except httpx.HTTPStatusError as e:
        print(f"âŒ HTTP error from TTS service: {e.response.status_code}")
        raise HTTPException(
            status_code=e.response.status_code,
            detail=f"TTS service returned error: {str(e)}"
        )
    except Exception as e:
        print(f"âŒ TTS voices proxy error: {e}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to proxy TTS voices request: {str(e)}")

# TTS speech proxy endpoint to handle CORS and streaming
@app.post("/v1/proxy/tts/speech")
async def proxy_tts_speech(request: Request, endpoint: Optional[str] = None):
    """Proxy TTS speech requests to handle CORS. Supports streaming responses."""
    try:
        buffer_response = request.query_params.get("buffer", "").lower() in {"1", "true", "yes", "on"}
        # Get endpoint from query parameter or try to extract from request
        if not endpoint:
            # Try to get from query params
            endpoint = request.query_params.get('endpoint', '')
        
        # If still no endpoint, try to get from TTS_ENDPOINT env var or use default
        if not endpoint:
            tts_endpoint = os.getenv('TTS_ENDPOINT', 'http://localhost:4123/v1')
            endpoint = tts_endpoint.rstrip('/')
        
        # Normalize the endpoint URL (remove trailing slash)
        base_endpoint = endpoint.rstrip('/')
        
        # Extract base URL (origin: protocol + host + port) to avoid path duplication
        try:
            # Parse the endpoint URL to extract the origin (protocol + host + port)
            if not base_endpoint.startswith('http://') and not base_endpoint.startswith('https://'):
                # If no protocol, assume http://
                base_endpoint = f"http://{base_endpoint}"
            
            parsed_url = urlparse(base_endpoint)
            # Reconstruct base URL with scheme, hostname, and port (if present)
            base_url = f"{parsed_url.scheme}://{parsed_url.netloc}"
        except Exception as e:
            # Fallback to simple string replacement if URL parsing fails
            # Extract origin manually using regex
            match = re.match(r'^(https?://[^/]+)', base_endpoint)
            if match:
                base_url = match.group(1)
            else:
                # Last resort: remove /v1 and any path
                base_url = base_endpoint.split('/')[0] if '/' in base_endpoint else base_endpoint
                if not base_url.startswith('http'):
                    base_url = f"http://{base_url}"
        
        # Construct the speech endpoint URL
        speech_url = f"{base_url}/v1/audio/speech"
        
        print(f"🎤 Proxying TTS speech request to: {speech_url}")
        
        # Get the request body
        try:
            request_body = await request.json()
        except Exception:
            request_body = {}
        
        # Get headers from the original request
        # Forward Accept header to support both SSE (Chatterbox) and binary audio (VibeVoice/OpenAI-compatible)
        # If client requests SSE, forward it; otherwise let TTS service decide (defaults to binary audio)
        accept_header = request.headers.get('Accept', '')
        forward_headers = {
            'Content-Type': 'application/json',
        }
        
        # Forward Accept header if present (allows Chatterbox to return SSE, VibeVoice will ignore and return binary)
        if accept_header:
            forward_headers['Accept'] = accept_header
        
        # Forward Authorization header if present
        auth_header = request.headers.get('Authorization')
        if auth_header:
            forward_headers['Authorization'] = auth_header

        skip_tls_verify = _should_skip_tts_tls_verify(base_url)
        if skip_tls_verify:
            print(f"[WARN] TTS speech proxy: TLS verification disabled for local endpoint: {base_url}")

        if buffer_response:
            async with httpx.AsyncClient(timeout=TTS_PROXY_TIMEOUT_SECONDS, verify=(not skip_tls_verify)) as client:
                response = await client.post(
                    speech_url,
                    json=request_body,
                    headers=forward_headers,
                )
            content_type = response.headers.get('content-type', 'audio/mpeg')
            # Safari/iOS can fail to decode blobs labeled as audio/mp3; normalize to audio/mpeg.
            if isinstance(content_type, str) and content_type.lower().startswith('audio/mp3'):
                content_type = 'audio/mpeg'
            if response.status_code != 200:
                print(f"âŒ TTS service returned error (buffered): {response.status_code}")
                return Response(
                    content=response.content,
                    media_type=content_type,
                    status_code=response.status_code,
                )
            return Response(
                content=response.content,
                media_type=content_type,
                status_code=200,
            )
        
        # Open upstream stream first so we can forward the *actual* content-type.
        # This avoids mislabeling MP3 as SSE, which breaks browser playback/parsing.
        client = httpx.AsyncClient(timeout=TTS_PROXY_TIMEOUT_SECONDS, verify=(not skip_tls_verify))
        try:
            upstream_response = await client.send(
                client.build_request(
                    "POST",
                    speech_url,
                    json=request_body,
                    headers=forward_headers,
                ),
                stream=True,
            )
        except Exception:
            await client.aclose()
            raise

        print(f"✅ TTS speech response status: {upstream_response.status_code}")
        print(f"📤 TTS request body: {json.dumps(request_body, indent=2)[:500]}")
        print(f"📤 TTS request headers: {forward_headers}")
        upstream_content_type = upstream_response.headers.get("content-type", "audio/mpeg")
        # Normalize non-standard MP3 MIME type for better iOS Safari compatibility.
        if isinstance(upstream_content_type, str) and upstream_content_type.lower().startswith("audio/mp3"):
            upstream_content_type = "audio/mpeg"
        print(f"📦 TTS response content-type: {upstream_content_type}")

        if upstream_response.status_code != 200:
            error_body = await upstream_response.aread()
            print(f"âŒ TTS service returned error: {upstream_response.status_code}")
            print(f"   Response text: {error_body[:200]}")
            await upstream_response.aclose()
            await client.aclose()
            return Response(
                content=error_body,
                status_code=upstream_response.status_code,
                headers={"Content-Type": upstream_content_type},
            )

        async def stream_upstream():
            try:
                async for chunk in upstream_response.aiter_bytes():
                    if chunk:
                        yield chunk
            finally:
                await upstream_response.aclose()
                await client.aclose()

        return StreamingResponse(
            stream_upstream(),
            status_code=200,
            headers={"Content-Type": upstream_content_type},
        )
    
    except Exception as e:
        print(f"âŒ TTS speech proxy error: {e}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to proxy TTS speech request: {str(e)}")

# ============================================================================
# FILE OPERATIONS ENDPOINTS
# ============================================================================

def _normalize_scratch_relative_input(raw_input: str) -> str:
    """Normalize user-provided scratch-relative paths and optional 'scratch/' prefix."""
    normalized = str(raw_input or "").strip()
    if not normalized:
        return normalized
    normalized = normalized.replace("\\", "/")
    lowered = normalized.lower()
    if lowered == "scratch":
        return "."
    if lowered.startswith("scratch/"):
        trimmed = normalized[len("scratch/") :]
        return trimmed or "."
    return normalized


def resolve_scratch_path(filename: str, allowed_extensions: Optional[Set[str]] = None) -> Path:
    """
    Resolve a user-supplied filename to a path under SCRATCH_DIR.
    Rejects absolute paths, traversal (..), and disallowed extensions.
    Returns the canonical path for safe I/O. Raises HTTPException 400 on invalid input.
    """
    normalized_name = _normalize_scratch_relative_input(filename)

    # Reject empty or whitespace-only filename
    if not normalized_name or not normalized_name.strip():
        raise HTTPException(status_code=400, detail="Invalid filename")

    # Reject absolute paths (Unix / or Windows drive/root)
    if (
        os.path.isabs(normalized_name)
        or normalized_name.startswith("/")
        or bool(re.match(r"^[A-Za-z]:/", normalized_name))
    ):
        raise HTTPException(status_code=400, detail="Invalid filename")

    # Reject path traversal components
    parts = PurePosixPath(normalized_name).parts
    if ".." in parts:
        raise HTTPException(status_code=400, detail="Invalid filename")

    # Build candidate path and resolve to canonical form
    candidate = SCRATCH_DIR / Path(normalized_name)
    try:
        resolved = candidate.resolve()
    except (OSError, RuntimeError) as e:
        raise HTTPException(status_code=400, detail="Invalid filename") from e
    # Enforce containment under SCRATCH_DIR (Python 3.9+ relative_to)
    root = SCRATCH_DIR.resolve()
    try:
        resolved.relative_to(root)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid filename")
    # Optional: enforce extension allowlist
    if allowed_extensions is not None:
        suffix = resolved.suffix.lower()
        if suffix not in allowed_extensions:
            raise HTTPException(status_code=400, detail="Invalid filename")
    return resolved


def _sanitize_attachment_component(value: str, *, fallback: str) -> str:
    sanitized = re.sub(r"[^A-Za-z0-9._-]+", "-", str(value or "").strip())
    sanitized = sanitized.strip(".-_")
    return sanitized or fallback


def _sanitize_attachment_filename(filename: str) -> str:
    raw_name = Path(str(filename or "").strip() or "attachment.bin").name
    ext = Path(raw_name).suffix.lower()
    if ext not in ATTACHMENT_ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=(
                "Unsupported attachment type. Allowed: "
                + ", ".join(sorted(ATTACHMENT_ALLOWED_EXTENSIONS))
            ),
        )
    stem = _sanitize_attachment_component(Path(raw_name).stem, fallback="attachment")
    return f"{stem}{ext}"


def _build_attachment_relative_path(
    *,
    source: str,
    conversation_id: str,
    filename: str,
    index: int,
) -> str:
    safe_source = _sanitize_attachment_component(source, fallback="web")
    safe_conversation = _sanitize_attachment_component(conversation_id, fallback="default")
    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d-%H%M%S")
    return f"attachments/{safe_source}/{safe_conversation}/{timestamp}-{index:02d}-{filename}"


def _store_attachment_bytes(
    *,
    content: bytes,
    filename: str,
    mime_type: Optional[str],
    conversation_id: str,
    source: str,
    index: int,
) -> Dict[str, Any]:
    if not FILE_OPS_AVAILABLE:
        raise HTTPException(
            status_code=503,
            detail="File operations not available. Install: pip install python-docx openpyxl PyPDF2 reportlab Pillow",
        )

    size_bytes = len(content)
    if size_bytes <= 0:
        raise HTTPException(status_code=400, detail="Attachment is empty.")
    if size_bytes > FILE_OPS_MAX_SIZE_BYTES:
        raise HTTPException(
            status_code=400,
            detail=f"Attachment '{filename}' is too large. Limit is {FILE_OPS_MAX_SIZE_BYTES} bytes.",
        )

    safe_filename = _sanitize_attachment_filename(filename)
    relative_path = _build_attachment_relative_path(
        source=source,
        conversation_id=conversation_id,
        filename=safe_filename,
        index=index,
    )
    filepath = resolve_scratch_path(relative_path, ATTACHMENT_ALLOWED_EXTENSIONS)
    filepath.parent.mkdir(parents=True, exist_ok=True)
    with filepath.open("wb") as handle:
        handle.write(content)

    resolved_mime = (mime_type or "").strip() or mimetypes.guess_type(filepath.name)[0] or "application/octet-stream"
    return {
        "filename": filepath.name,
        "relative_path": filepath.relative_to(SCRATCH_DIR.resolve()).as_posix(),
        "original_filename": Path(str(filename or "")).name or filepath.name,
        "mime_type": resolved_mime,
        "size_bytes": size_bytes,
    }


def _store_json_attachments(
    attachments: Optional[List[ChatAttachment]],
    *,
    conversation_id: str,
    source: str,
) -> List[Dict[str, Any]]:
    items = attachments or []
    if not items:
        return []
    if len(items) > ATTACHMENT_MAX_FILES_PER_REQUEST:
        raise HTTPException(
            status_code=400,
            detail=f"Too many attachments. Limit is {ATTACHMENT_MAX_FILES_PER_REQUEST} files per request.",
        )

    stored: List[Dict[str, Any]] = []
    for index, attachment in enumerate(items, start=1):
        try:
            content = base64.b64decode(str(attachment.content_base64 or ""), validate=True)
        except (ValueError, binascii.Error):
            raise HTTPException(status_code=400, detail=f"Attachment '{attachment.filename}' is not valid base64.")
        stored.append(
            _store_attachment_bytes(
                content=content,
                filename=attachment.filename,
                mime_type=attachment.mime_type,
                conversation_id=conversation_id,
                source=source,
                index=index,
            )
        )
    return stored


def _build_attachment_manifest(attachments: List[Dict[str, Any]]) -> str:
    if not attachments:
        return ""
    lines = ["Attached files saved in scratch:"]
    for item in attachments:
        rel_path = str(item.get("relative_path") or "")
        original = str(item.get("original_filename") or item.get("filename") or "")
        mime_type = str(item.get("mime_type") or "application/octet-stream")
        size_bytes = int(item.get("size_bytes") or 0)
        lines.append(
            f"- {rel_path} (original: {original}, type: {mime_type}, size: {size_bytes} bytes)"
        )
    lines.append("Use filesystem.read_text with these scratch-relative filenames to inspect attachments before answering.")
    lines.append("For attached PDFs, DOCX, XLSX, text files, Markdown files, and images, inspect the attachment first instead of guessing.")
    lines.append("Do not use pdfToPowerPoint unless the user explicitly asks to convert a PDF or Markdown document into a PowerPoint or slide deck.")
    return "\n".join(lines)


def _augment_message_with_attachments(message_text: str, attachments: List[Dict[str, Any]]) -> str:
    manifest = _build_attachment_manifest(attachments)
    base_message = str(message_text or "").strip()
    if not base_message and attachments:
        base_message = "Please review the attached file(s)."
    if not manifest:
        return base_message
    return f"{base_message}\n\n{manifest}"


def _is_vision_image_attachment(*, mime_type: Optional[str], filename: str) -> bool:
    resolved = str(mime_type or "").strip().lower()
    if not resolved:
        resolved = (mimetypes.guess_type(filename)[0] or "").strip().lower()
    return resolved in {"image/png", "image/jpeg", "image/jpg"}


def _build_attachment_vision_parts(
    attachments: Optional[List[ChatAttachment]],
) -> List[Dict[str, Any]]:
    parts: List[Dict[str, Any]] = []
    for attachment in attachments or []:
        filename = Path(str(attachment.filename or "")).name or "attachment"
        mime_type = str(attachment.mime_type or "").strip() or mimetypes.guess_type(filename)[0] or "application/octet-stream"
        if not _is_vision_image_attachment(mime_type=mime_type, filename=filename):
            continue
        content_base64 = str(attachment.content_base64 or "").strip()
        if not content_base64:
            continue
        parts.append(
            {
                "type": "image_url",
                "image_url": {
                    "url": f"data:{mime_type};base64,{content_base64}",
                    "detail": "auto",
                },
            }
        )
    return parts


def _attach_vision_parts_to_latest_user_message(
    messages: List[Dict[str, Any]],
    attachments: Optional[List[ChatAttachment]],
) -> List[Dict[str, Any]]:
    vision_parts = _build_attachment_vision_parts(attachments)
    if not vision_parts:
        return messages
    for index in range(len(messages) - 1, -1, -1):
        message = messages[index]
        if str(message.get("role") or "") != "user":
            continue
        current_text = coerce_message_text(message.get("content") or "").strip()
        content_parts: List[Dict[str, Any]] = []
        if current_text:
            content_parts.append({"type": "text", "text": current_text})
        content_parts.extend(vision_parts)
        messages[index] = {**message, "content": content_parts}
        break
    return messages


def write_text_file(filepath: Path, content: str) -> None:
    """Write content to a plain text file"""
    # Write the content with UTF-8 encoding
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(content)

def write_docx_file(filepath: Path, content: str) -> None:
    """Write content to a Word document"""
    # Create a new document
    doc = Document()
    
    # Split content into paragraphs and add each to the document
    for paragraph in content.split('\n'):
        doc.add_paragraph(paragraph)
    
    # Save the document
    doc.save(filepath)

def write_xlsx_file(filepath: Path, content: str) -> None:
    """Write content to an Excel file"""
    # Create a new workbook
    wb = openpyxl.Workbook()
    # Get the active sheet
    ws = wb.active
    ws.title = "Sheet1"
    
    # Split content into rows
    rows = content.split('\n')
    
    # Write each row to the Excel file
    for row_idx, row_content in enumerate(rows, start=1):
        # Split row by tabs or commas
        if '\t' in row_content:
            cells = row_content.split('\t')
        else:
            cells = row_content.split(',')
        
        # Write each cell
        for col_idx, cell_content in enumerate(cells, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=cell_content.strip())
            
            # Apply formatting to the first row (header)
            if row_idx == 1:
                cell.font = Font(bold=True)
                cell.alignment = Alignment(horizontal='center')
    
    # Auto-adjust column widths
    for column in ws.columns:
        max_length = 0
        column_letter = column[0].column_letter
        for cell in column:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        adjusted_width = min(max_length + 2, 50)  # Cap at 50 characters
        ws.column_dimensions[column_letter].width = adjusted_width
    
    # Save the workbook
    wb.save(filepath)

def write_pdf_file(filepath: Path, content: str) -> None:
    """Write content to a PDF file using reportlab"""
    try:
        # Import reportlab for PDF creation
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas as pdf_canvas
        from reportlab.lib.units import inch
        
        # Create a canvas for PDF generation
        c = pdf_canvas.Canvas(str(filepath), pagesize=letter)
        width, height = letter
        
        # Set up text formatting
        text_object = c.beginText(1 * inch, height - 1 * inch)
        text_object.setFont("Helvetica", 12)
        
        # Split content into lines and add to PDF
        lines = content.split('\n')
        for line in lines:
            # Handle long lines by wrapping
            if len(line) > 80:
                words = line.split(' ')
                current_line = ''
                for word in words:
                    if len(current_line + word) < 80:
                        current_line += word + ' '
                    else:
                        text_object.textLine(current_line)
                        current_line = word + ' '
                if current_line:
                    text_object.textLine(current_line)
            else:
                text_object.textLine(line)
            
            # Add new page if needed (simple check)
            if text_object.getY() < 1 * inch:
                c.drawText(text_object)
                c.showPage()
                text_object = c.beginText(1 * inch, height - 1 * inch)
                text_object.setFont("Helvetica", 12)
        
        # Draw the text and save the PDF
        c.drawText(text_object)
        c.save()
        
    except ImportError:
        # If reportlab is not available, raise an error
        raise HTTPException(
            status_code=500,
            detail="PDF writing requires reportlab library. Install with: pip install reportlab"
        )


def _truncate_text_middle(text: str, max_chars: Optional[int]) -> tuple[str, bool]:
    if max_chars is None or max_chars < 1 or len(text) <= max_chars:
        return text, False
    if max_chars < 80:
        return text[:max_chars], True
    head = max_chars // 2
    tail = max_chars - head
    removed = len(text) - max_chars
    return (
        f"{text[:head]}\n\n...[truncated {removed} chars]...\n\n{text[-tail:]}",
        True,
    )


def _slice_text_for_read(
    content: str,
    *,
    start_line: Optional[int] = None,
    end_line: Optional[int] = None,
    max_chars: Optional[int] = None,
    include_line_numbers: bool = False,
) -> Dict[str, Any]:
    lines = content.splitlines()
    total_lines = len(lines)
    excerpt_start_line = 1 if total_lines > 0 else 0
    excerpt_end_line = total_lines
    line_filtered = False
    rendered_lines = lines

    if start_line is not None or end_line is not None:
        start_index = max(0, (start_line or 1) - 1)
        resolved_end_line = end_line
        if resolved_end_line is not None and start_line is not None and resolved_end_line < start_line:
            resolved_end_line = start_line
        end_index = total_lines if resolved_end_line is None else min(total_lines, resolved_end_line)
        rendered_lines = lines[start_index:end_index]
        excerpt_start_line = start_index + 1 if rendered_lines else 0
        excerpt_end_line = start_index + len(rendered_lines) if rendered_lines else 0
        line_filtered = True

    if include_line_numbers and rendered_lines:
        base_line = excerpt_start_line
        rendered_text = "\n".join(
            f"{base_line + index}: {line}" for index, line in enumerate(rendered_lines)
        )
    else:
        rendered_text = "\n".join(rendered_lines)

    truncated_text, truncated = _truncate_text_middle(rendered_text, max_chars)
    return {
        "content": truncated_text,
        "total_lines": total_lines,
        "excerpt_start_line": excerpt_start_line,
        "excerpt_end_line": excerpt_end_line,
        "line_filtered": line_filtered,
        "truncated": truncated,
    }


def _extract_query_match(
    content: str,
    query: str,
    *,
    case_sensitive: bool = False,
) -> Optional[Dict[str, Any]]:
    if not query:
        return None
    haystack = content if case_sensitive else content.lower()
    needle = query if case_sensitive else query.lower()
    index = haystack.find(needle)
    if index < 0:
        return None
    line_number = content.count("\n", 0, index) + 1
    start = max(0, index - 90)
    end = min(len(content), index + len(query) + 140)
    excerpt = content[start:end].replace("\r", " ").replace("\n", " ").strip()
    if start > 0:
        excerpt = "..." + excerpt
    if end < len(content):
        excerpt = excerpt + "..."
    return {
        "line_number": line_number,
        "excerpt": excerpt,
        "match_count": haystack.count(needle),
    }


def _read_arg_text(arguments: Dict[str, Any], *keys: str, default: str = "") -> str:
    for key in keys:
        value = arguments.get(key)
        if value is None:
            continue
        text = str(value).strip()
        if text:
            return text
    return default


# Internal file ops for Telegram tool runner (no auth; same security as routes)
async def _read_file_internal(
    filename: str,
    *,
    max_chars: Optional[int] = None,
    start_line: Optional[int] = None,
    end_line: Optional[int] = None,
    include_line_numbers: bool = False,
) -> Dict[str, Any]:
    """Read file from scratch dir. Returns dict with success, message, data (content/type). Used by Telegram tools only."""
    if not FILE_OPS_AVAILABLE:
        return {"success": False, "message": "File operations not available."}
    try:
        filepath = resolve_scratch_path(filename, READ_ALLOWED_EXTENSIONS)
        if not filepath.exists():
            return {"success": False, "message": f"File not found: {filename}"}
        if filepath.stat().st_size > FILE_OPS_MAX_SIZE_BYTES:
            return {"success": False, "message": "File too large"}
        ext = filepath.suffix.lower()
        if ext in TEXT_FILE_EXTENSIONS:
            raw_content = read_text_file(filepath)
            excerpt = _slice_text_for_read(
                raw_content,
                start_line=start_line,
                end_line=end_line,
                max_chars=max_chars,
                include_line_numbers=include_line_numbers,
            )
            return {
                "success": True,
                "message": f"Read {filename}",
                "data": {
                    "content": excerpt["content"],
                    "type": "text",
                    "size_bytes": filepath.stat().st_size,
                    "total_lines": excerpt["total_lines"],
                    "excerpt_start_line": excerpt["excerpt_start_line"],
                    "excerpt_end_line": excerpt["excerpt_end_line"],
                    "truncated": excerpt["truncated"],
                    "line_filtered": excerpt["line_filtered"],
                },
            }
        if ext == '.docx':
            raw_content = read_docx_file(filepath)
            excerpt = _slice_text_for_read(
                raw_content,
                start_line=start_line,
                end_line=end_line,
                max_chars=max_chars,
                include_line_numbers=include_line_numbers,
            )
            return {
                "success": True,
                "message": f"Read {filename}",
                "data": {
                    "content": excerpt["content"],
                    "type": "text",
                    "size_bytes": filepath.stat().st_size,
                    "total_lines": excerpt["total_lines"],
                    "excerpt_start_line": excerpt["excerpt_start_line"],
                    "excerpt_end_line": excerpt["excerpt_end_line"],
                    "truncated": excerpt["truncated"],
                    "line_filtered": excerpt["line_filtered"],
                },
            }
        if ext in ['.xlsx', '.xls']:
            raw_content = read_xlsx_file(filepath)
            excerpt = _slice_text_for_read(
                raw_content,
                start_line=start_line,
                end_line=end_line,
                max_chars=max_chars,
                include_line_numbers=include_line_numbers,
            )
            return {
                "success": True,
                "message": f"Read {filename}",
                "data": {
                    "content": excerpt["content"],
                    "type": "text",
                    "size_bytes": filepath.stat().st_size,
                    "total_lines": excerpt["total_lines"],
                    "excerpt_start_line": excerpt["excerpt_start_line"],
                    "excerpt_end_line": excerpt["excerpt_end_line"],
                    "truncated": excerpt["truncated"],
                    "line_filtered": excerpt["line_filtered"],
                },
            }
        if ext == '.pdf':
            raw_content = read_pdf_file(filepath)
            excerpt = _slice_text_for_read(
                raw_content,
                start_line=start_line,
                end_line=end_line,
                max_chars=max_chars,
                include_line_numbers=include_line_numbers,
            )
            return {
                "success": True,
                "message": f"Read {filename}",
                "data": {
                    "content": excerpt["content"],
                    "type": "text",
                    "size_bytes": filepath.stat().st_size,
                    "total_lines": excerpt["total_lines"],
                    "excerpt_start_line": excerpt["excerpt_start_line"],
                    "excerpt_end_line": excerpt["excerpt_end_line"],
                    "truncated": excerpt["truncated"],
                    "line_filtered": excerpt["line_filtered"],
                },
            }
        if ext in ['.png', '.jpg', '.jpeg']:
            image_data = read_png_file(filepath)
            return {"success": True, "message": f"Read {filename}", "data": {"content": image_data.get("description", ""), "type": "image", "image_data": image_data}}
        return {"success": False, "message": f"Unsupported file type: {ext}"}
    except HTTPException as e:
        return {"success": False, "message": e.detail or "Invalid filename"}
    except Exception as e:
        return {"success": False, "message": str(e)}


async def _write_file_internal(
    filename: str,
    content: str,
    format: str = "txt",
    *,
    append: bool = False,
) -> Dict[str, Any]:
    """Write file to scratch dir. Returns dict with success, message. Used by Telegram tools only."""
    if not FILE_OPS_AVAILABLE:
        return {"success": False, "message": "File operations not available."}
    try:
        content_bytes = content.encode("utf-8")
        if len(content_bytes) > FILE_OPS_MAX_SIZE_BYTES:
            return {"success": False, "message": "Content too large"}
        logical_name = (filename or "").strip()
        if not logical_name:
            return {"success": False, "message": "Filename is required"}
        if not Path(logical_name).suffix:
            logical_name = f"{logical_name}.{format.lower()}"
        filepath = resolve_scratch_path(logical_name, WRITE_ALLOWED_EXTENSIONS)
        ext = filepath.suffix.lower()
        if ext in TEXT_FILE_EXTENSIONS:
            filepath.parent.mkdir(parents=True, exist_ok=True)
            mode = "a" if append else "w"
            with filepath.open(mode, encoding="utf-8") as handle:
                handle.write(content)
        elif ext == '.docx':
            if append:
                return {"success": False, "message": "append is only supported for text files"}
            filepath.parent.mkdir(parents=True, exist_ok=True)
            write_docx_file(filepath, content)
        elif ext in ['.xlsx', '.xls']:
            if append:
                return {"success": False, "message": "append is only supported for text files"}
            filepath.parent.mkdir(parents=True, exist_ok=True)
            write_xlsx_file(filepath, content)
        elif ext == '.pdf':
            if append:
                return {"success": False, "message": "append is only supported for text files"}
            filepath.parent.mkdir(parents=True, exist_ok=True)
            write_pdf_file(filepath, content)
        else:
            return {"success": False, "message": f"Unsupported file type for writing: {ext}"}
        action = "Appended to" if append else "Wrote"
        return {
            "success": True,
            "message": f"{action} {filepath.name}",
            "data": {"filepath": str(filepath), "size": filepath.stat().st_size, "appended": append},
        }
    except HTTPException as e:
        return {"success": False, "message": e.detail or "Invalid filename"}
    except Exception as e:
        return {"success": False, "message": str(e)}


async def _list_files_internal(
    path: str = "",
    recursive: bool = False,
    offset: int = 0,
    max_entries: Optional[int] = None,
) -> Dict[str, Any]:
    """List files in scratch dir (optionally scoped to a subdirectory)."""
    try:
        requested_path = str(path or "").strip()
        recursive_mode = _coerce_bool(recursive, default=False)
        page_offset = _coerce_bounded_int(offset, default=0, minimum=0)
        page_size = None
        if max_entries is not None:
            page_size = _coerce_bounded_int(max_entries, default=100, minimum=1, maximum=500)
        target_dir = SCRATCH_DIR.resolve()
        if requested_path:
            target_dir = resolve_scratch_path(requested_path)
        if not target_dir.exists():
            return {"success": False, "message": f"Path not found: {requested_path or '.'}", "files": []}
        if not target_dir.is_dir():
            return {"success": False, "message": f"Path is not a directory: {requested_path}", "files": []}

        files = []
        skipped_count = 0
        scratch_root = SCRATCH_DIR.resolve()
        pending_dirs = [target_dir]
        seen_dirs: Set[Path] = set()

        while pending_dirs:
            current_dir = pending_dirs.pop()
            try:
                current_resolved = current_dir.resolve()
                current_resolved.relative_to(scratch_root)
            except (OSError, RuntimeError, ValueError):
                skipped_count += 1
                continue
            if current_resolved in seen_dirs:
                continue
            seen_dirs.add(current_resolved)

            try:
                entries = list(current_dir.iterdir())
            except OSError:
                skipped_count += 1
                continue

            for entry in entries:
                try:
                    if entry.is_symlink():
                        skipped_count += 1
                        continue

                    entry_resolved = entry.resolve()
                    entry_resolved.relative_to(scratch_root)
                except (OSError, RuntimeError, ValueError):
                    skipped_count += 1
                    continue

                try:
                    if not entry.is_dir() and not entry.is_file():
                        continue
                    entry_stat = entry.stat()
                    is_dir = entry.is_dir()
                    if is_dir and recursive_mode:
                        pending_dirs.append(entry)
                except OSError:
                    skipped_count += 1
                    continue

                rel_path = entry.relative_to(scratch_root).as_posix()
                files.append({
                    "name": rel_path,
                    "relative_path": rel_path,
                    "size": None if is_dir else entry_stat.st_size,
                    "modified": entry_stat.st_mtime,
                    "extension": "" if is_dir else entry.suffix,
                    "type": "directory" if is_dir else "file",
                })

            if not recursive_mode:
                break

        files.sort(
            key=lambda x: (
                x.get("type") != "directory",
                str(x.get("relative_path", x.get("name", ""))).lower(),
            )
        )
        total_count = len(files)
        directory_count = sum(1 for item in files if item.get("type") == "directory")
        file_count = total_count - directory_count
        paged_files = files
        if page_size is not None:
            paged_files = files[page_offset:page_offset + page_size]
        returned_count = len(paged_files)
        remaining_count = max(0, total_count - (page_offset + returned_count))
        result = {
            "success": True,
            "files": paged_files,
            "count": total_count,
            "directory_count": directory_count,
            "file_count": file_count,
            "total_count": total_count,
            "returned_count": returned_count,
            "remaining_count": remaining_count,
            "scratch_dir": str(SCRATCH_DIR),
            "path": requested_path or ".",
            "recursive": recursive_mode,
            "offset": page_offset,
            "max_entries": page_size,
            "has_more": remaining_count > 0,
            "next_offset": (page_offset + returned_count) if remaining_count > 0 else None,
            "skipped_count": skipped_count,
        }
        if skipped_count > 0:
            result["message"] = (
                f"Listed {returned_count} of {total_count} files. "
                f"Skipped {skipped_count} inaccessible or unsafe entries."
            )
        return result
    except Exception as e:
        return {"success": False, "message": str(e), "files": []}


async def _search_files_internal(
    query: str,
    *,
    path: str = "",
    recursive: bool = True,
    offset: int = 0,
    max_results: Optional[int] = None,
    case_sensitive: bool = False,
    filename_only: bool = False,
) -> Dict[str, Any]:
    """Search scratch files by filename and text content."""
    try:
        search_query = str(query or "").strip()
        if not search_query:
            return {"success": False, "message": "query is required", "matches": []}
        requested_path = str(path or "").strip()
        recursive_mode = _coerce_bool(recursive, default=True)
        case_sensitive_mode = _coerce_bool(case_sensitive, default=False)
        filename_only_mode = _coerce_bool(filename_only, default=False)
        page_offset = _coerce_bounded_int(offset, default=0, minimum=0)
        page_size = _coerce_bounded_int(
            max_results,
            default=20,
            minimum=1,
            maximum=100,
        )
        target_dir = SCRATCH_DIR.resolve()
        if requested_path:
            target_dir = resolve_scratch_path(requested_path)
        if not target_dir.exists():
            return {"success": False, "message": f"Path not found: {requested_path or '.'}", "matches": []}
        if not target_dir.is_dir():
            return {"success": False, "message": f"Path is not a directory: {requested_path}", "matches": []}

        scratch_root = SCRATCH_DIR.resolve()
        pending_dirs = [target_dir]
        seen_dirs: Set[Path] = set()
        matches: List[Dict[str, Any]] = []
        searched_file_count = 0
        skipped_count = 0

        while pending_dirs:
            current_dir = pending_dirs.pop()
            try:
                current_resolved = current_dir.resolve()
                current_resolved.relative_to(scratch_root)
            except (OSError, RuntimeError, ValueError):
                skipped_count += 1
                continue
            if current_resolved in seen_dirs:
                continue
            seen_dirs.add(current_resolved)

            try:
                entries = sorted(current_dir.iterdir(), key=lambda item: item.name.lower())
            except OSError:
                skipped_count += 1
                continue

            for entry in entries:
                try:
                    if entry.is_symlink():
                        skipped_count += 1
                        continue
                    entry_resolved = entry.resolve()
                    entry_resolved.relative_to(scratch_root)
                except (OSError, RuntimeError, ValueError):
                    skipped_count += 1
                    continue

                try:
                    if entry.is_dir():
                        if recursive_mode:
                            pending_dirs.append(entry)
                        continue
                    if not entry.is_file():
                        continue
                except OSError:
                    skipped_count += 1
                    continue

                searched_file_count += 1
                rel_path = entry.relative_to(scratch_root).as_posix()
                rel_haystack = rel_path if case_sensitive_mode else rel_path.lower()
                query_haystack = search_query if case_sensitive_mode else search_query.lower()
                filename_match = query_haystack in rel_haystack
                content_match = None

                if not filename_only_mode and entry.suffix.lower() in SEARCHABLE_TEXT_EXTENSIONS:
                    try:
                        if entry.stat().st_size <= SEARCH_FILE_MAX_SIZE_BYTES:
                            text_content, _ = read_supported_file_text(entry, TEXT_FILE_EXTENSIONS)
                            content_match = _extract_query_match(
                                text_content,
                                search_query,
                                case_sensitive=case_sensitive_mode,
                            )
                    except Exception:
                        skipped_count += 1
                        continue

                if not filename_match and not content_match:
                    continue

                matches.append(
                    {
                        "name": rel_path,
                        "relative_path": rel_path,
                        "size": entry.stat().st_size,
                        "extension": entry.suffix,
                        "type": "file",
                        "match_types": [
                            label
                            for label, matched in (("filename", filename_match), ("content", bool(content_match)))
                            if matched
                        ],
                        "line_number": content_match.get("line_number") if content_match else None,
                        "excerpt": content_match.get("excerpt", "") if content_match else "",
                        "match_count": content_match.get("match_count", 1) if content_match else 1,
                    }
                )

            if not recursive_mode:
                break

        matches.sort(
            key=lambda item: (
                "filename" not in item.get("match_types", []),
                -int(item.get("match_count", 0) or 0),
                str(item.get("relative_path", "")).lower(),
            )
        )
        total_matches = len(matches)
        paged_matches = matches[page_offset:page_offset + page_size]
        returned_count = len(paged_matches)
        remaining_count = max(0, total_matches - (page_offset + returned_count))
        return {
            "success": True,
            "matches": paged_matches,
            "query": search_query,
            "total_matches": total_matches,
            "returned_count": returned_count,
            "remaining_count": remaining_count,
            "offset": page_offset,
            "max_results": page_size,
            "has_more": remaining_count > 0,
            "next_offset": (page_offset + returned_count) if remaining_count > 0 else None,
            "searched_file_count": searched_file_count,
            "path": requested_path or ".",
            "recursive": recursive_mode,
            "filename_only": filename_only_mode,
            "case_sensitive": case_sensitive_mode,
            "skipped_count": skipped_count,
            "scratch_dir": str(SCRATCH_DIR),
        }
    except Exception as e:
        return {"success": False, "message": str(e), "matches": []}


def _get_list_files_tool_max_entries() -> int:
    """Return max entries rendered in list-files tool replies."""
    raw = (os.getenv("LIST_FILES_TOOL_MAX_ENTRIES", "60") or "60").strip()
    try:
        parsed = int(raw)
    except ValueError:
        parsed = 60
    return max(1, parsed)


def _format_list_files_for_tool_output(
    files: List[Dict[str, Any]],
    include_sizes: bool = False,
    *,
    total_count: Optional[Any] = None,
    offset: Optional[Any] = None,
    has_more: Optional[Any] = None,
    next_offset: Optional[Any] = None,
    limit: Optional[int] = None,
) -> str:
    """Format scratch files for LLM-facing tool output with explicit pagination metadata."""
    if not files:
        return "Scratch workspace is empty."
    render_limit = _coerce_bounded_int(
        limit,
        default=_get_list_files_tool_max_entries(),
        minimum=1,
        maximum=500,
    )
    shown = files[:render_limit]
    lines = []
    for item in shown:
        name = str(item.get("name", "?"))
        is_dir = str(item.get("type", "")).lower() == "directory"
        if is_dir:
            lines.append(f"{name}/ [dir]")
            continue
        if include_sizes:
            lines.append(f"{name} ({item.get('size', 0)} bytes)")
        else:
            lines.append(name)
    parsed_total = _coerce_bounded_int(total_count, default=len(files), minimum=0)
    parsed_offset = _coerce_bounded_int(offset, default=0, minimum=0)
    inferred_remaining = max(0, parsed_total - (parsed_offset + len(shown)))
    parsed_has_more = bool(has_more) if has_more is not None else inferred_remaining > 0
    header = "Files in scratch workspace:"
    if parsed_total > len(shown) or parsed_offset > 0:
        header += f" (showing {parsed_offset + 1}-{parsed_offset + len(shown)} of {parsed_total})"
    if len(files) > len(shown):
        lines.append(f"... and {len(files) - len(shown)} more files in this page.")
    if parsed_has_more:
        continuation_offset = _coerce_bounded_int(
            next_offset,
            default=parsed_offset + len(shown),
            minimum=0,
        )
        lines.append(f"... more files available. Continue with offset={continuation_offset}.")
    return header + "\n" + "\n".join(lines)


def _format_search_files_for_tool_output(
    matches: List[Dict[str, Any]],
    *,
    query: str,
    total_matches: Optional[Any] = None,
    offset: Optional[Any] = None,
    has_more: Optional[Any] = None,
    next_offset: Optional[Any] = None,
    read_tool_name: str = "read_file",
) -> str:
    if not matches:
        return f'No matching files found for "{query}".'
    parsed_total = _coerce_bounded_int(total_matches, default=len(matches), minimum=0)
    parsed_offset = _coerce_bounded_int(offset, default=0, minimum=0)
    lines = [f'Search results for "{query}":']
    if parsed_total > len(matches) or parsed_offset > 0:
        lines[0] += f" (showing {parsed_offset + 1}-{parsed_offset + len(matches)} of {parsed_total})"
    for index, item in enumerate(matches, start=parsed_offset + 1):
        rel_path = str(item.get("relative_path") or item.get("name") or "?")
        match_types = ",".join(item.get("match_types") or []) or "unknown"
        line_number = item.get("line_number")
        excerpt = str(item.get("excerpt") or "").strip()
        suffix = f" line {line_number}" if isinstance(line_number, int) and line_number > 0 else ""
        if excerpt:
            lines.append(f"{index}. {rel_path} [{match_types}]{suffix}: {excerpt}")
        else:
            lines.append(f"{index}. {rel_path} [{match_types}]")
    parsed_has_more = bool(has_more) if has_more is not None else False
    if parsed_has_more:
        continuation_offset = _coerce_bounded_int(
            next_offset,
            default=parsed_offset + len(matches),
            minimum=0,
        )
        lines.append(f"More results available. Continue with offset={continuation_offset}.")
    lines.append(f"Use {read_tool_name} with filename/path and optional start_line/end_line for more context.")
    return "\n".join(lines)


async def _delete_file_internal(filename: str) -> Dict[str, Any]:
    """Delete file from scratch dir. Returns dict with success, message. Used by philosopher and other server-side callers."""
    if not FILE_OPS_AVAILABLE:
        return {"success": False, "message": "File operations not available."}
    try:
        filepath = resolve_scratch_path(filename, READ_ALLOWED_EXTENSIONS)
        if not filepath.exists():
            return {"success": False, "message": f"File not found: {filename}"}
        filepath.unlink()
        return {"success": True, "message": f"Deleted {filename}"}
    except HTTPException as e:
        return {"success": False, "message": e.detail or "Invalid filename"}
    except Exception as e:
        return {"success": False, "message": str(e)}


@app.post("/v1/files/read", response_model=FileResponse)
async def read_file(
    request: ReadFileRequest,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """
    Read a file from the scratch directory
    Supports: txt, md, csv, docx, xlsx, xls, pdf, png, jpg, jpeg, py, js, html
    """
    if not FILE_OPS_AVAILABLE:
        raise HTTPException(
            status_code=503,
            detail="File operations not available. Install: pip install python-docx openpyxl PyPDF2 reportlab Pillow"
        )
    resolve_scratch_path(request.filename, READ_ALLOWED_EXTENSIONS)
    result = await _read_file_internal(
        request.filename,
        max_chars=request.max_chars,
        start_line=request.start_line,
        end_line=request.end_line,
        include_line_numbers=request.include_line_numbers,
    )
    return FileResponse(
        success=bool(result.get("success")),
        message=str(result.get("message", "")),
        data=result.get("data"),
    )


@app.get("/v1/files/content")
async def read_file_content(
    path: str = Query(..., min_length=1),
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Return raw file bytes from the scratch directory for browser-side consumers."""
    filepath = resolve_scratch_path(path, READ_ALLOWED_EXTENSIONS)
    if not filepath.exists() or not filepath.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    try:
        size_bytes = filepath.stat().st_size
    except OSError as exc:
        raise HTTPException(status_code=404, detail="File not found") from exc
    if size_bytes > FILE_OPS_MAX_SIZE_BYTES:
        raise HTTPException(
            status_code=400,
            detail=f"File is too large. Limit is {FILE_OPS_MAX_SIZE_BYTES} bytes.",
        )
    mime_type = mimetypes.guess_type(filepath.name)[0] or "application/octet-stream"
    try:
        content = filepath.read_bytes()
    except OSError as exc:
        raise HTTPException(status_code=500, detail="Failed to read file") from exc
    return Response(
        content=content,
        media_type=mime_type,
        headers={"Content-Disposition": f'inline; filename="{filepath.name}"'},
    )

@app.post("/v1/files/write", response_model=FileResponse)
async def write_file(
    request: WriteFileRequest,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """
    Write content to a file in the scratch directory
    Supports: txt, md, csv, docx, xlsx, xls, pdf, py, js, html
    """
    if not FILE_OPS_AVAILABLE:
        raise HTTPException(
            status_code=503,
            detail="File operations not available. Install: pip install python-docx openpyxl PyPDF2 reportlab Pillow"
        )
    logical_name = (request.filename or "").strip()
    if logical_name and not Path(logical_name).suffix:
        logical_name = f"{logical_name}.{(request.format or 'txt').lower()}"
    resolve_scratch_path(logical_name or request.filename, WRITE_ALLOWED_EXTENSIONS)
    result = await _write_file_internal(
        request.filename,
        request.content,
        format=request.format or "txt",
        append=bool(request.append),
    )
    return FileResponse(
        success=bool(result.get("success")),
        message=str(result.get("message", "")),
        data=result.get("data"),
    )


@app.post("/v1/files/attachments")
async def upload_attachments(
    files: List[UploadFile] = File(...),
    conversation_id: Optional[str] = Form(default=None),
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Upload one or more user attachments into scratch for web chat turns."""
    upload_files = files or []
    if not upload_files:
        raise HTTPException(status_code=400, detail="At least one file is required.")
    if len(upload_files) > ATTACHMENT_MAX_FILES_PER_REQUEST:
        raise HTTPException(
            status_code=400,
            detail=f"Too many attachments. Limit is {ATTACHMENT_MAX_FILES_PER_REQUEST} files per request.",
        )

    effective_conversation_id = str(conversation_id or current_user.get("sub") or "default").strip() or "default"
    stored: List[Dict[str, Any]] = []
    for index, upload in enumerate(upload_files, start=1):
        try:
            file_bytes = await upload.read()
            stored.append(
                _store_attachment_bytes(
                    content=file_bytes,
                    filename=upload.filename or f"attachment-{index}.bin",
                    mime_type=upload.content_type,
                    conversation_id=effective_conversation_id,
                    source="web",
                    index=index,
                )
            )
        finally:
            with suppress(Exception):
                await upload.close()

    return {
        "success": True,
        "message": f"Uploaded {len(stored)} attachment(s).",
        "attachments": stored,
    }


@app.get("/v1/files/list")
async def list_files(
    path: Optional[str] = None,
    recursive: bool = False,
    offset: int = Query(default=0, ge=0),
    max_entries: Optional[int] = Query(default=None, ge=1, le=500),
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """List files in the scratch directory (optionally scoped by subdirectory and recursion)."""
    try:
        result = await _list_files_internal(
            path=(path or ""),
            recursive=recursive,
            offset=offset,
            max_entries=max_entries,
        )
        if not result.get("success"):
            return {
                "success": False,
                "message": result.get("message", "Error listing files"),
                "files": [],
            }
        return result
    
    except Exception as e:
        # Handle any errors during directory listing
        return {
            'success': False,
            'message': f"Error listing files: {str(e)}"
        }


@app.get("/v1/files/search")
async def search_files(
    query: str,
    path: Optional[str] = None,
    recursive: bool = True,
    offset: int = Query(default=0, ge=0),
    max_results: Optional[int] = Query(default=20, ge=1, le=100),
    case_sensitive: bool = False,
    filename_only: bool = False,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Search scratch files by filename and text content."""
    try:
        result = await _search_files_internal(
            query,
            path=(path or ""),
            recursive=recursive,
            offset=offset,
            max_results=max_results,
            case_sensitive=case_sensitive,
            filename_only=filename_only,
        )
        if not result.get("success"):
            return {
                "success": False,
                "message": result.get("message", "Error searching files"),
                "matches": [],
            }
        return result
    except Exception as e:
        return {
            "success": False,
            "message": f"Error searching files: {str(e)}",
            "matches": [],
        }

@app.delete("/v1/files/delete/{filename}")
async def delete_file(
    filename: str,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Delete a file from the scratch directory"""
    # Resolve path with containment and extension checks (blocks path traversal)
    filepath = resolve_scratch_path(filename, READ_ALLOWED_EXTENSIONS)
    try:
        # Check if file exists
        if not filepath.exists():
            return FileResponse(
                success=False,
                message=f"File not found: {filename}"
            )
        
        # Delete the file
        filepath.unlink()
        
        return FileResponse(
            success=True,
            message=f"Successfully deleted {filename}"
        )
    
    except Exception as e:
        # Handle any errors during file deletion
        return FileResponse(
            success=False,
            message=f"Error deleting file: {str(e)}"
        )


# ============================================================================
# COMPANIONS API (CATBot tool settings saved as static config files)
# ============================================================================

# Safe filename: only alphanumeric, hyphen, underscore (no path traversal, .json only)
_COMPANION_ID_RE = re.compile(r"^[a-zA-Z0-9_-]+$")


def resolve_companion_path(companion_id: str) -> Path:
    """
    Resolve a companion id to a path under COMPANIONS_DIR.
    Only allows safe ids (alphanumeric, hyphen, underscore); enforces .json.
    Raises HTTPException 400 on invalid input.
    """
    if not companion_id or not companion_id.strip():
        raise HTTPException(status_code=400, detail="Invalid companion id")
    if not _COMPANION_ID_RE.match(companion_id):
        raise HTTPException(status_code=400, detail="Invalid companion id: only letters, numbers, hyphen, underscore allowed")
    # Build path: COMPANIONS_DIR / {id}.json
    candidate = COMPANIONS_DIR / f"{companion_id}.json"
    try:
        resolved = candidate.resolve()
    except (OSError, RuntimeError) as e:
        raise HTTPException(status_code=400, detail="Invalid companion id") from e
    # Enforce containment under COMPANIONS_DIR
    root = COMPANIONS_DIR.resolve()
    try:
        resolved.relative_to(root)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid companion id")
    if resolved.suffix.lower() != ".json":
        raise HTTPException(status_code=400, detail="Invalid companion id")
    return resolved


@app.get("/v1/companions", response_model=List[CompanionResponse])
async def list_companions(
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """List all saved companions (id and name only for list view)."""
    try:
        result = []
        for path in COMPANIONS_DIR.glob("*.json"):
            try:
                data = json.loads(path.read_text(encoding="utf-8"))
                cid = data.get("id") or path.stem
                name = data.get("name") or path.stem
                result.append(CompanionResponse(id=cid, name=name, settings=None))
            except (json.JSONDecodeError, OSError):
                continue
        result.sort(key=lambda c: c.name.lower())
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/v1/companions/{companion_id}", response_model=CompanionResponse)
async def get_companion(
    companion_id: str,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Get one companion by id (full record including settings)."""
    filepath = resolve_companion_path(companion_id)
    if not filepath.exists():
        raise HTTPException(status_code=404, detail="Companion not found")
    try:
        data = json.loads(filepath.read_text(encoding="utf-8"))
        return CompanionResponse(
            id=data.get("id", companion_id),
            name=data.get("name", companion_id),
            settings=data.get("settings"),
        )
    except (json.JSONDecodeError, OSError) as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/v1/companions", response_model=CompanionResponse)
async def create_companion(
    request: CompanionCreateRequest,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Create a new companion; server generates id (UUID), writes config/companions/{id}.json."""
    import uuid
    name = (request.name or "").strip()
    if not name:
        raise HTTPException(status_code=400, detail="Companion name is required")
    # Generate a URL-safe id (use uuid4 hex for simplicity and uniqueness)
    cid = uuid.uuid4().hex
    filepath = resolve_companion_path(cid)
    record = {"id": cid, "name": name, "settings": request.settings or {}}
    try:
        filepath.write_text(json.dumps(record, indent=2), encoding="utf-8")
        return CompanionResponse(id=cid, name=name, settings=record["settings"])
    except OSError as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/v1/companions/{companion_id}")
async def delete_companion(
    companion_id: str,
    current_user: Dict[str, Any] = Depends(get_current_user),
):
    """Delete a companion by id."""
    filepath = resolve_companion_path(companion_id)
    if not filepath.exists():
        raise HTTPException(status_code=404, detail="Companion not found")
    try:
        filepath.unlink()
        return {"success": True, "message": f"Companion {companion_id} deleted"}
    except OSError as e:
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================================
# END COMPANIONS API
# ============================================================================

# ============================================================================
# MODEL AVATAR SCAN (discover *.model3.json and *.vrm under model_avatar/)
# ============================================================================

def scan_model_avatar_dir() -> Dict[str, Any]:
    """
    Recursively scan model_avatar directory for Live2D (.model3.json) and VRM (.vrm) files.
    Returns paths relative to project root with ./ prefix and forward slashes for UI consistency.
    """
    model_avatar_dir = _PROJECT_ROOT / "model_avatar"
    result = {"live2d": [], "vrm": []}
    if not model_avatar_dir.is_dir():
        result["message"] = "model_avatar directory not found"
        return result
    try:
        # Collect Live2D model paths (one per file, any depth)
        for p in model_avatar_dir.rglob("*.model3.json"):
            if p.is_file():
                rel = p.relative_to(_PROJECT_ROOT)
                path_str = "./" + str(rel).replace("\\", "/")
                result["live2d"].append(path_str)
        result["live2d"].sort()
        # Collect VRM model paths (one per file, any depth)
        for p in model_avatar_dir.rglob("*.vrm"):
            if p.is_file():
                rel = p.relative_to(_PROJECT_ROOT)
                path_str = "./" + str(rel).replace("\\", "/")
                result["vrm"].append(path_str)
        result["vrm"].sort()
    except Exception as e:
        result["message"] = str(e)
    return result


@app.get("/v1/model-avatar/scan")
async def model_avatar_scan(
    current_user: Dict[str, Any] = Depends(get_current_user),
) -> Dict[str, Any]:
    """Return lists of discovered Live2D and VRM model paths under model_avatar/ for tools settings."""
    return scan_model_avatar_dir()


# ============================================================================
# GOOGLE DRIVE UPLOAD ENDPOINT
# ============================================================================

# Internal drive upload for Telegram tool runner (no auth; same path/credential checks)
async def _upload_drive_internal(file_path: str, file_name: Optional[str] = None) -> Dict[str, Any]:
    """Upload a file from scratch dir to Google Drive. Returns dict with success, message, fileId, etc. Used by Telegram tools only."""
    try:
        if not file_path or not str(file_path).strip():
            return {"success": False, "message": "filePath is required"}
        file_path_obj = resolve_scratch_path(str(file_path).strip(), DRIVE_UPLOAD_EXTENSIONS)
        if not file_path_obj.exists():
            return {"success": False, "message": f"File not found: {file_path}"}
        folder_id = os.getenv('GOOGLE_DRIVE_FOLDER_ID')
        if not folder_id:
            return {"success": False, "message": "GOOGLE_DRIVE_FOLDER_ID is not set"}
        project_id = os.getenv('GOOGLE_DRIVE_PROJECT_ID')
        private_key_id = os.getenv('GOOGLE_DRIVE_PRIVATE_KEY_ID')
        private_key = os.getenv('GOOGLE_DRIVE_PRIVATE_KEY')
        client_email = os.getenv('GOOGLE_DRIVE_CLIENT_EMAIL')
        if not all([project_id, private_key_id, private_key, client_email]):
            missing = [k for k, v in {
                'GOOGLE_DRIVE_PROJECT_ID': project_id,
                'GOOGLE_DRIVE_PRIVATE_KEY_ID': private_key_id,
                'GOOGLE_DRIVE_PRIVATE_KEY': private_key,
                'GOOGLE_DRIVE_CLIENT_EMAIL': client_email
            }.items() if not v]
            return {"success": False, "message": f"Missing Google Drive credentials: {', '.join(missing)}"}
        private_key_formatted = private_key.replace('\\n', '\n')
        credentials_dict = {
            "type": "service_account",
            "project_id": project_id,
            "private_key_id": private_key_id,
            "private_key": private_key_formatted,
            "client_email": client_email,
            "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            "token_uri": "https://oauth2.googleapis.com/token",
            "auth_provider_x509_cert_url": "https://www.googleapis.com/oauth2/v1/certs",
            "client_x509_cert_url": f"https://www.googleapis.com/robot/v1/metadata/x509/{client_email}"
        }
        try:
            from google.oauth2 import service_account
            from googleapiclient.discovery import build
            from googleapiclient.http import MediaFileUpload
        except ImportError:
            return {"success": False, "message": "Google Drive API libraries not available. Install: pip install google-auth google-api-python-client"}
        credentials = service_account.Credentials.from_service_account_info(
            credentials_dict,
            scopes=['https://www.googleapis.com/auth/drive.file']
        )
        drive_service = build('drive', 'v3', credentials=credentials)
        upload_file_name = file_name if file_name else file_path_obj.name
        file_metadata = {'name': upload_file_name, 'parents': [folder_id]}
        media = MediaFileUpload(str(file_path_obj), resumable=True)
        file = drive_service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id, name, webViewLink'
        ).execute()
        return {
            'success': True,
            'fileId': file.get('id'),
            'fileName': file.get('name'),
            'webViewLink': file.get('webViewLink'),
            'message': f"File successfully uploaded to Google Drive with ID: {file.get('id')}"
        }
    except HTTPException as e:
        return {"success": False, "message": e.detail or "Upload failed"}
    except Exception as e:
        return {"success": False, "message": str(e)}


@app.post("/v1/proxy/upload-to-drive")
async def upload_to_drive(request: Request, current_user: Dict[str, Any] = Depends(get_current_user)):
    """
    Upload a file from the scratch directory to Google Drive using service account credentials from .env file.
    filePath must be a filename relative to the scratch directory (e.g. report.docx). Path traversal and absolute paths are rejected.
    """
    folder_id = None
    file_path_obj = None
    try:
        form_data = await request.form()
        file_path = form_data.get('filePath')
        if not file_path or not str(file_path).strip():
            raise HTTPException(status_code=400, detail="filePath is required")
        file_path_obj = resolve_scratch_path(str(file_path).strip(), DRIVE_UPLOAD_EXTENSIONS)
        if not file_path_obj.exists():
            raise HTTPException(status_code=404, detail=f"File not found: {file_path}")
        file_name = form_data.get('fileName')
        folder_id = form_data.get('folderId') or os.getenv('GOOGLE_DRIVE_FOLDER_ID')
        if not folder_id:
            raise HTTPException(status_code=400, detail="folderId is required (provide in request or set GOOGLE_DRIVE_FOLDER_ID in .env)")
        
        # Read Google Drive credentials from environment variables (loaded from .env file)
        project_id = os.getenv('GOOGLE_DRIVE_PROJECT_ID')
        private_key_id = os.getenv('GOOGLE_DRIVE_PRIVATE_KEY_ID')
        private_key = os.getenv('GOOGLE_DRIVE_PRIVATE_KEY')
        client_email = os.getenv('GOOGLE_DRIVE_CLIENT_EMAIL')
        
        # Validate that all required credentials are present
        if not all([project_id, private_key_id, private_key, client_email]):
            missing = [k for k, v in {
                'GOOGLE_DRIVE_PROJECT_ID': project_id,
                'GOOGLE_DRIVE_PRIVATE_KEY_ID': private_key_id,
                'GOOGLE_DRIVE_PRIVATE_KEY': private_key,
                'GOOGLE_DRIVE_CLIENT_EMAIL': client_email
            }.items() if not v]
            raise HTTPException(
                status_code=500,
                detail=f"Missing Google Drive credentials in .env file: {', '.join(missing)}"
            )
        
        # Construct credentials object from environment variables
        # Replace \\n with actual newlines in private key
        private_key_formatted = private_key.replace('\\n', '\n')
        
        credentials_dict = {
            "type": "service_account",
            "project_id": project_id,
            "private_key_id": private_key_id,
            "private_key": private_key_formatted,
            "client_email": client_email,
            "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            "token_uri": "https://oauth2.googleapis.com/token",
            "auth_provider_x509_cert_url": "https://www.googleapis.com/oauth2/v1/certs",
            "client_x509_cert_url": f"https://www.googleapis.com/robot/v1/metadata/x509/{client_email}"
        }
        
        # Try to import Google Drive API libraries
        try:
            from google.oauth2 import service_account
            from googleapiclient.discovery import build
            from googleapiclient.http import MediaFileUpload
            GOOGLE_DRIVE_AVAILABLE = True
        except ImportError:
            GOOGLE_DRIVE_AVAILABLE = False
            raise HTTPException(
                status_code=503,
                detail="Google Drive API libraries not available. Install with: pip install google-auth google-auth-oauthlib google-auth-httplib2 google-api-python-client"
            )
        
        # Authenticate with Google Drive using service account credentials
        credentials = service_account.Credentials.from_service_account_info(
            credentials_dict,
            scopes=['https://www.googleapis.com/auth/drive.file']
        )
        
        # Build the Drive service
        drive_service = build('drive', 'v3', credentials=credentials)
        
        # Determine file name to use
        upload_file_name = file_name if file_name else file_path_obj.name
        
        # Prepare file metadata
        file_metadata = {
            'name': upload_file_name,
            'parents': [folder_id] if folder_id else []
        }
        
        # Upload file to Google Drive
        media = MediaFileUpload(
            str(file_path_obj),
            resumable=True
        )
        
        # Perform the upload
        file = drive_service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id, name, webViewLink'
        ).execute()
        
        # Audit log: user, filename, folder_id, success, file_id
        user_sub = current_user.get("sub") or "unknown"
        print(f"[AUDIT] upload-to-drive user={user_sub} filename={file_path_obj.name} folder_id={folder_id} success=true file_id={file.get('id')}")
        # Return success response with file ID
        return {
            'success': True,
            'fileId': file.get('id'),
            'fileName': file.get('name'),
            'webViewLink': file.get('webViewLink'),
            'message': f"File successfully uploaded to Google Drive with ID: {file.get('id')}"
        }
        
    except HTTPException as exc:
        # Audit log on auth/path/not-found/credential errors
        user_sub = current_user.get("sub") if current_user else "unknown"
        filename_log = file_path_obj.name if file_path_obj else "n/a"
        print(f"[AUDIT] upload-to-drive user={user_sub} filename={filename_log} folder_id={folder_id or 'n/a'} success=false status={exc.status_code} detail={exc.detail}")
        # Re-raise HTTP exceptions
        raise
    except Exception as e:
        # Audit log on unexpected errors
        user_sub = current_user.get("sub") if current_user else "unknown"
        print(f"[AUDIT] upload-to-drive user={user_sub} folder_id={folder_id or 'n/a'} success=false error={str(e)}")
        # Handle any other errors
        print(f"âŒ Google Drive upload error: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to upload file to Google Drive: {str(e)}"
        )

# ============================================================================
# END GOOGLE DRIVE UPLOAD ENDPOINT
# ============================================================================

# ============================================================================
# SSL CERTIFICATE UTILITIES
# ============================================================================

def get_local_ip() -> Optional[str]:
    """
    Get the local IP address of this machine.
    Returns the IP address or None if unable to determine.
    """
    try:
        # Connect to a remote address to determine local IP (doesn't actually send data)
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(("8.8.8.8", 80))
        ip = s.getsockname()[0]
        s.close()
        return ip
    except Exception:
        return None

def find_mkcert_certificates() -> Tuple[Optional[str], Optional[str]]:
    """
    Find mkcert-generated certificates in project root or certs/ directory.
    Returns a tuple of (cert_file, key_file) or (None, None) if not found.
    mkcert creates files like: <hostname>+2.pem, <hostname>+2-key.pem, etc.
    Hostname comes from HTTPS_CERT_HOSTNAME in .env.
    """
    # Search in certs/ first, then project root (glob uses sanitized hostname)
    search_dirs = [_PROJECT_ROOT / "certs", _PROJECT_ROOT]
    for search_dir in search_dirs:
        if not search_dir.exists():
            continue
        cert_files = list(search_dir.glob(f"{_SSL_CERT_HOSTNAME_GLOB}*.pem"))
        # Filter out key files and find matching pairs
        cert_key_pairs = []
        for cert_path in cert_files:
            if "-key" in cert_path.name:
                continue
            key_path = cert_path.parent / (cert_path.stem + "-key.pem")
            if key_path.exists():
                cert_key_pairs.append((str(cert_path), str(key_path), cert_path.stat().st_mtime))
        if cert_key_pairs:
            cert_key_pairs.sort(key=lambda x: x[2], reverse=True)
            return cert_key_pairs[0][0], cert_key_pairs[0][1]
    return None, None

def get_ssl_certificates() -> Tuple[Optional[str], Optional[str]]:
    """
    Get SSL certificate files for HTTPS server.
    First tries to find mkcert certificates in certs/ or project root.
    Returns a tuple of (cert_file, key_file) or (None, None) if not found.
    """
    # Try to find mkcert certificates first
    cert_file, key_file = find_mkcert_certificates()
    if cert_file and key_file and os.path.exists(cert_file) and os.path.exists(key_file):
        print(f"[SSL] Found mkcert certificate: {cert_file}")
        return cert_file, key_file
    
    # Fall back to default certificate file names in certs/ or project root (hostname from .env)
    for base in [_PROJECT_ROOT / "certs", _PROJECT_ROOT]:
        default_cert = base / f"{_SSL_CERT_HOSTNAME}+2.pem"
        default_key = base / f"{_SSL_CERT_HOSTNAME}+2-key.pem"
        if default_cert.exists() and default_key.exists():
            print(f"[SSL] Using default certificate: {default_cert}")
            return str(default_cert), str(default_key)
    
    # Return None if no certificates found
    print("[SSL] No SSL certificates found. Server will run without HTTPS.")
    return None, None

# ============================================================================
# END SSL CERTIFICATE UTILITIES
# ============================================================================

if __name__ == "__main__":
    # Start the server
    print("[START] Starting CATBot Proxy Server with File Operations...")
    print(f"[INFO] Scratch directory: {SCRATCH_DIR}")
    
    # Get SSL certificates for HTTPS
    cert_file, key_file = get_ssl_certificates()
    
    # Configure uvicorn with SSL if certificates are available
    if cert_file and key_file:
        print(f"[SSL] Starting HTTPS server on port 8002")
        print(f"[SSL] Certificate: {cert_file}")
        print(f"[SSL] Key: {key_file}")
        uvicorn.run(
            "src.servers.proxy_server:app",
            host="0.0.0.0",
            port=8002,
            reload=PROXY_RELOAD,
            log_level="info",
            ssl_keyfile=key_file,
            ssl_certfile=cert_file
        )
    else:
        print("[WARN] Starting HTTP server (no SSL certificates found)")
        print("[INFO] To enable HTTPS, ensure mkcert certificate files are in certs/ directory")
        uvicorn.run(
            "src.servers.proxy_server:app",
            host="0.0.0.0",
            port=8002,
            reload=PROXY_RELOAD,
            log_level="info"
        )
